"""
导出路由
"""

import sys
import base64
import asyncio
import tempfile
import shutil
from dataclasses import replace
from pathlib import Path
from fastapi import APIRouter, HTTPException
from fastapi.responses import FileResponse
from starlette.background import BackgroundTask

# 添加项目根目录到 Python 路径
project_root = Path(__file__).parent.parent.parent
sys.path.insert(0, str(project_root))

from src.exporter import PDFExporter
from src.generative_editable_composer import compose_deck_from_manifests
from src.generative_editable_config import (
    GenerativeEditableConfig,
    GenerativeEditableConfigError,
    ProviderConfig,
    load_generative_editable_config,
)
from src.generative_editable_manifest import read_deck_manifest, read_page_manifest, write_manifest
from src.generative_editable_pipeline import (
    GenerativeEditableFallbackError,
    GenerativeEditablePipelineDependencies,
    GenerativeEditableSlideInput,
    GenerativeEditableValidationError,
    finalize_validated_export,
)
from src.generative_editable_vlm_reconstruction import (
    OpenAIChatVLMPageAnalysisProvider,
    VLMEditablePipelineDependencies,
    run_vlm_editable_pptx_pipeline,
)
from src.generative_editable_providers import (
    FakeImageEditProvider,
    FakeImageGenerationProvider,
    FakeOCRProvider,
    LocalTesseractOCRProvider,
    OpenAIChatImageEditProvider,
    OpenAIChatImageGenerationProvider,
    OpenAIChatOCRProvider,
    ProviderError,
    ProviderTimeoutError,
)
from src.generative_editable_preview_validator import ValidationIssue, ValidationReport
from src.model_profiles import load_default_profiles
from ..models import ExportRequest

router = APIRouter(prefix="/api", tags=["export"])


@router.post("/export")
async def export_presentation(request: ExportRequest):
    """
    导出演示文稿

    Args:
        request: 导出请求，包含所有幻灯片和格式

    Returns:
        FileResponse: 导出的文件
    """
    try:
        temp_dir = Path(tempfile.mkdtemp())
        image_paths = []

        # 解码并保存所有图片
        for idx, slide in enumerate(request.slides):
            image_data = base64.b64decode(slide.image_base64)
            image_path = temp_dir / f"slide_{idx + 1}.png"

            with open(image_path, "wb") as f:
                f.write(image_data)

            image_paths.append(str(image_path))

        # 根据格式导出
        if request.format == "pdf":
            output_path = temp_dir / "presentation.pdf"
            exporter = PDFExporter()
            await asyncio.to_thread(exporter.export, image_paths, str(output_path))

            return FileResponse(
                path=str(output_path),
                media_type="application/pdf",
                filename="presentation.pdf",
                background=BackgroundTask(shutil.rmtree, temp_dir, ignore_errors=True),
            )

        elif request.format == "pptx":
            output_path = temp_dir / "presentation.pptx"
            await asyncio.to_thread(
                _export_pptx, image_paths, str(output_path), aspect_ratio=request.aspect_ratio
            )

            return FileResponse(
                path=str(output_path),
                media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                filename="presentation.pptx",
                background=BackgroundTask(shutil.rmtree, temp_dir, ignore_errors=True),
            )

        elif request.format == "generative_editable_pptx":
            output_path = temp_dir / "presentation.generative-editable.pptx"
            export_result = await asyncio.to_thread(
                _export_generative_editable_pptx,
                image_paths,
                str(output_path),
                aspect_ratio=request.aspect_ratio,
                slides=request.slides,
                editable_options=getattr(request, "editable_options", None),
                slide_order=getattr(request, "slide_order", None),
            )
            response_path = getattr(export_result, "output_path", str(output_path))

            return FileResponse(
                path=str(response_path),
                media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                filename="presentation.generative-editable.pptx",
                headers=_generative_editable_response_headers(export_result),
                background=BackgroundTask(shutil.rmtree, temp_dir, ignore_errors=True),
            )

        else:
            raise HTTPException(status_code=400, detail=f"不支持的导出格式: {request.format}")

    except HTTPException:
        if "temp_dir" in locals():
            shutil.rmtree(temp_dir, ignore_errors=True)
        raise

    except Exception as e:
        if "temp_dir" in locals():
            shutil.rmtree(temp_dir, ignore_errors=True)
        raise _export_http_exception(e)


def _export_pptx(image_paths: list, output_path: str, aspect_ratio: str = "16:9"):
    """
    导出为 PPTX 格式

    Args:
        image_paths: 图片路径列表
        output_path: 输出路径
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        raise Exception("需要安装 python-pptx: pip install python-pptx")

    # 创建演示文稿
    prs = Presentation()

    # 设置幻灯片尺寸
    if aspect_ratio == "4:3":
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
    else:
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)

    # 添加每一页
    for image_path in image_paths:
        # 使用空白布局
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)

        # 添加图片，填充整个幻灯片
        slide.shapes.add_picture(image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    # 保存
    prs.save(output_path)


def _export_generative_editable_pptx(
    image_paths: list,
    output_path: str,
    *,
    aspect_ratio: str = "16:9",
    slides: list | None = None,
    editable_options: object | None = None,
    slide_order: list[str] | None = None,
):
    """Export a generative editable PPTX through the reconstruction pipeline."""
    config = load_generative_editable_config()
    fallback_policy = _editable_fallback_policy(editable_options)
    artifact_root = Path(output_path).with_suffix(".artifacts")
    job_id = "export"
    slide_inputs = _generative_editable_slide_inputs(
        image_paths=image_paths,
        slides=slides,
        slide_order=slide_order,
    )
    ordered_image_paths = [slide.image_path for slide in slide_inputs]
    if config.reconstruction_mode != "vlm_first":
        raise GenerativeEditableConfigError(
            "VLM-first reconstruction is required for generative_editable_pptx export"
        )
    fallback_factory = None
    if fallback_policy == "raster_pptx":

        def fallback_factory():
            return _export_raster_pptx_fallback(
                image_paths=ordered_image_paths,
                output_path=output_path,
                aspect_ratio=aspect_ratio,
            )
    elif fallback_policy == "text_editable_background":

        def fallback_factory():
            return _export_text_editable_background_fallback(
                artifact_root=artifact_root,
                job_id=job_id,
                output_path=output_path,
            )

    try:
        try:
            return run_vlm_editable_pptx_pipeline(
                slides=slide_inputs,
                output_path=output_path,
                artifact_root=str(artifact_root),
                job_id=job_id,
                dependencies=_build_vlm_editable_pipeline_dependencies(config),
                aspect_ratio=aspect_ratio,
                cleanup_artifacts=fallback_policy == "fail",
            )
        except GenerativeEditableValidationError as exc:
            return finalize_validated_export(
                validation_report=exc.validation_report,
                output_path=output_path,
                fallback_policy=fallback_policy,
                fallback_output_factory=fallback_factory,
            )
        except (ProviderError, ProviderTimeoutError) as exc:
            if fallback_policy == "fail":
                raise
            return finalize_validated_export(
                validation_report=_provider_failure_validation_report(exc, len(slide_inputs)),
                output_path=output_path,
                fallback_policy=fallback_policy,
                fallback_output_factory=fallback_factory,
            )
    except Exception:
        Path(output_path).unlink(missing_ok=True)
        raise


def _provider_failure_validation_report(
    error: ProviderError,
    slide_count: int,
) -> ValidationReport:
    details: dict[str, object] = {
        "provider_role": error.provider_role,
        "operation": error.operation,
        "retryable": bool(error.retryable),
    }
    if error.status_code is not None:
        details["status_code"] = error.status_code
    if error.provider_error_code:
        details["provider_error_code"] = error.provider_error_code
    if isinstance(error, ProviderTimeoutError):
        details["timeout_seconds"] = error.timeout_seconds
    return ValidationReport(
        status="failed",
        checked_pages=slide_count,
        issues=[
            ValidationIssue(
                code="provider_timeout"
                if isinstance(error, ProviderTimeoutError)
                else "provider_failure",
                message=str(error) or error.__class__.__name__,
                details=details,
            )
        ],
    )


def _editable_fallback_policy(editable_options: object | None) -> str:
    if editable_options is None:
        return "fail"
    if isinstance(editable_options, dict):
        return str(editable_options.get("fallback_policy") or "fail")
    return str(getattr(editable_options, "fallback_policy", "fail") or "fail")


def _export_raster_pptx_fallback(*, image_paths: list, output_path: str, aspect_ratio: str) -> str:
    _export_pptx(image_paths, output_path, aspect_ratio=aspect_ratio)
    return output_path


def _text_editable_background_fallback_error(reason: str) -> GenerativeEditableFallbackError:
    return GenerativeEditableFallbackError(
        validation_report=ValidationReport(
            status="failed",
            checked_pages=0,
            issues=[
                ValidationIssue(
                    code="fallback_precondition_failed",
                    message=reason,
                )
            ],
        ),
        fallback_policy="text_editable_background",
        fallback_failure_reason=reason,
    )


def _export_text_editable_background_fallback(
    *, artifact_root: Path, job_id: str, output_path: str
) -> str:
    job_dir = artifact_root / job_id
    deck_path = job_dir / "deck.json"
    if not deck_path.is_file():
        raise _text_editable_background_fallback_error(
            "text_clean_background fallback failed: deck manifest is missing because the pipeline failed before creating artifacts"
        )
    deck = read_deck_manifest(deck_path)
    fallback_dir = job_dir / "pages" / "text-editable-background-fallback"
    fallback_dir.mkdir(parents=True, exist_ok=True)
    page_manifest_paths = []
    for page_ref in deck.page_manifest_paths:
        page = read_page_manifest(job_dir / page_ref)
        background = (
            page.text_clean_background or page.base_clean_background or page.chosen_background
        )
        if not background:
            raise _text_editable_background_fallback_error(
                "text_clean_background fallback requires a cleaned background artifact"
            )
        fallback_page = replace(
            page,
            chosen_background=background,
            native_shapes=[],
            bitmap_assets=[],
            asset_sheets=[],
            repair_attempts=[],
            validation_status="pending",
        )
        fallback_page_path = fallback_dir / Path(page_ref).name
        write_manifest(fallback_page_path, fallback_page)
        page_manifest_paths.append(
            f"pages/text-editable-background-fallback/{fallback_page_path.name}"
        )
    fallback_deck = replace(
        deck,
        page_manifest_paths=page_manifest_paths,
        fallback_policy="text_editable_background",
        validation_status="pending",
    )
    fallback_deck_path = fallback_dir / "deck.json"
    write_manifest(fallback_deck_path, fallback_deck)
    compose_deck_from_manifests(fallback_deck_path, job_dir, output_path)
    return output_path


def _export_http_exception(error: Exception) -> HTTPException:
    if isinstance(error, GenerativeEditableConfigError):
        return HTTPException(status_code=400, detail=f"导出失败: {error}")
    if isinstance(error, ProviderTimeoutError):
        return HTTPException(status_code=504, detail=f"导出失败: {error}")
    if isinstance(error, ProviderError):
        return HTTPException(status_code=502, detail=f"导出失败: {error}")
    if isinstance(error, GenerativeEditableValidationError):
        return HTTPException(status_code=422, detail=f"导出失败: {error}")
    if isinstance(error, GenerativeEditableFallbackError):
        return HTTPException(status_code=502, detail=f"导出失败: {error}")
    return HTTPException(status_code=500, detail=f"导出失败: {error}")


def _build_generative_editable_pipeline_dependencies() -> GenerativeEditablePipelineDependencies:
    config = load_generative_editable_config()
    return GenerativeEditablePipelineDependencies(
        ocr_provider=_ocr_provider_for_config(
            config.ocr, timeout_seconds=config.timeouts.provider_call
        ),
        image_edit_provider=_image_edit_provider_for_config(config.clean_base_model),
        asset_sheet_image_edit_provider=_image_edit_provider_for_config(config.asset_sheet_model),
        repair_image_edit_provider=_image_edit_provider_for_config(config.repair_model),
        image_generation_provider=_image_generation_provider_for_config(config.generation_model),
        preview_similarity_threshold=config.quality.preview_similarity_threshold,
        require_preview_validation=config.quality.require_preview_validation,
        max_repair_attempts=config.quality.max_repair_attempts,
        provider_timeout_seconds=config.timeouts.provider_call,
        provider_max_attempts=config.retries.provider_max_attempts,
        provider_retry_backoff_seconds=config.retries.backoff_seconds,
        use_aippt_metadata_first=config.use_aippt_metadata_first,
        ocr_min_confidence=config.ocr_min_confidence,
    )


def _build_vlm_editable_pipeline_dependencies(
    config: GenerativeEditableConfig,
) -> VLMEditablePipelineDependencies:
    profiles = load_default_profiles()
    if not profiles or not profiles.vlm:
        raise GenerativeEditableConfigError(
            "Missing VLM model profile for vlm_first reconstruction"
        )
    vlm_config = ProviderConfig(
        role="VLM",
        model=profiles.vlm.model,
        base_url=profiles.vlm.base_url,
        api_key=profiles.vlm.api_key,
        adapter=profiles.vlm.adapter,
    )
    return VLMEditablePipelineDependencies(
        vlm_provider=OpenAIChatVLMPageAnalysisProvider(vlm_config),
        image_edit_provider=_image_edit_provider_for_config(config.clean_base_model),
        asset_sheet_image_edit_provider=_image_edit_provider_for_config(config.asset_sheet_model),
        ocr_provider=_ocr_provider_for_config(
            config.ocr, timeout_seconds=config.timeouts.provider_call
        ),
        provider_timeout_seconds=config.timeouts.provider_call,
        page_timeout_seconds=config.timeouts.page,
        provider_max_attempts=config.retries.provider_max_attempts,
        provider_retry_backoff_seconds=config.retries.backoff_seconds,
        ocr_min_confidence=config.ocr_min_confidence,
        preview_similarity_threshold=config.quality.preview_similarity_threshold,
        require_preview_validation=config.quality.require_preview_validation,
    )


def _build_fake_generative_editable_pipeline_dependencies(
    *,
    validation_status: str = "passed",
    validation_code: str = "preview_similarity_failed",
) -> GenerativeEditablePipelineDependencies:
    from src.generative_editable_preview_validator import ValidationIssue, ValidationReport

    config = load_generative_editable_config(use_fake=True)

    def structure_validator(**kwargs):
        return ValidationReport(status="passed", checked_pages=1, issues=[])

    def preview_validator(**kwargs):
        if validation_status == "failed":
            return ValidationReport(
                status="failed",
                checked_pages=1,
                issues=[
                    ValidationIssue(
                        code=validation_code,
                        message="forced preview failure",
                        slide_id=kwargs["slide_id"],
                    )
                ],
            )
        return ValidationReport(status="passed", checked_pages=1, issues=[])

    return GenerativeEditablePipelineDependencies(
        ocr_provider=FakeOCRProvider(config.ocr),
        image_edit_provider=FakeImageEditProvider(config.clean_base_model),
        asset_sheet_image_edit_provider=FakeImageEditProvider(config.asset_sheet_model),
        repair_image_edit_provider=FakeImageEditProvider(config.repair_model),
        image_generation_provider=FakeImageGenerationProvider(config.generation_model),
        preview_similarity_threshold=config.quality.preview_similarity_threshold,
        require_preview_validation=config.quality.require_preview_validation,
        max_repair_attempts=config.quality.max_repair_attempts,
        provider_timeout_seconds=config.timeouts.provider_call,
        provider_max_attempts=config.retries.provider_max_attempts,
        provider_retry_backoff_seconds=config.retries.backoff_seconds,
        use_aippt_metadata_first=config.use_aippt_metadata_first,
        ocr_min_confidence=config.ocr_min_confidence,
        structure_validator=structure_validator,
        preview_validator=preview_validator,
    )


def _ocr_provider_for_config(config, *, timeout_seconds: int = 180):
    if config.provider.startswith("fake_"):
        return FakeOCRProvider(config)
    if config.provider == "local_tesseract":
        return LocalTesseractOCRProvider(config, timeout_seconds=timeout_seconds)
    return OpenAIChatOCRProvider(config, timeout_seconds=timeout_seconds)


def _image_edit_provider_for_config(config):
    if config.provider.startswith("fake_"):
        return FakeImageEditProvider(config)
    return OpenAIChatImageEditProvider(config)


def _image_generation_provider_for_config(config):
    if config.provider.startswith("fake_"):
        return FakeImageGenerationProvider(config)
    return OpenAIChatImageGenerationProvider(config)


def _generative_editable_slide_inputs(
    *,
    image_paths: list,
    slides: list | None,
    slide_order: list[str] | None,
) -> list[GenerativeEditableSlideInput]:
    slide_records = []
    for index, image_path in enumerate(image_paths):
        slide = slides[index] if slides is not None and index < len(slides) else None
        slide_id = getattr(slide, "slide_id", None) or f"slide-{index + 1}"
        text_metadata = [
            item.model_dump() if hasattr(item, "model_dump") else dict(item)
            for item in (getattr(slide, "text_metadata", []) if slide is not None else [])
        ]
        slide_records.append(
            GenerativeEditableSlideInput(
                slide_id=slide_id,
                image_path=str(image_path),
                text_metadata=text_metadata,
            )
        )
    duplicate_slide_ids = _duplicates([slide.slide_id for slide in slide_records])
    if duplicate_slide_ids:
        raise ValueError(f"slide_id values must be unique: {', '.join(duplicate_slide_ids)}")
    if not slide_order:
        return slide_records
    duplicate_order_ids = _duplicates(slide_order)
    if duplicate_order_ids:
        raise ValueError(f"slide_order values must be unique: {', '.join(duplicate_order_ids)}")
    by_id = {slide.slide_id: slide for slide in slide_records}
    unknown_order_ids = [slide_id for slide_id in slide_order if slide_id not in by_id]
    if unknown_order_ids:
        raise ValueError(f"slide_order contains unknown slide_id: {', '.join(unknown_order_ids)}")
    ordered = [by_id[slide_id] for slide_id in slide_order]
    ordered.extend(slide for slide in slide_records if slide.slide_id not in set(slide_order))
    return ordered


def _duplicates(values: list[str]) -> list[str]:
    seen = set()
    duplicates = []
    for value in values:
        if value in seen and value not in duplicates:
            duplicates.append(value)
        seen.add(value)
    return duplicates


def _generative_editable_response_headers(export_result: object | None) -> dict[str, str]:
    headers = {}
    status = getattr(export_result, "status", "")
    fallback_used = getattr(export_result, "fallback_used", "")
    fallback_policy = getattr(export_result, "fallback_policy", "")
    if status:
        headers["X-Generative-Editable-Status"] = str(status)
    if fallback_policy:
        headers["X-Generative-Editable-Fallback-Policy"] = str(fallback_policy)
    if fallback_used:
        headers["X-Generative-Editable-Fallback-Used"] = str(fallback_used)
    return headers
