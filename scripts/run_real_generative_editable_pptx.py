#!/usr/bin/env python3
"""Run real-image generative editable PPTX verification on arbitrary slide images."""

from __future__ import annotations

import argparse
from collections import Counter
from dataclasses import replace
from datetime import datetime
from difflib import SequenceMatcher
import glob
import json
import os
import re
import signal
import subprocess
import sys
import time
from pathlib import Path
from typing import Any
from xml.etree import ElementTree as ET
import zipfile

from PIL import Image, ImageChops, ImageDraw, ImageStat
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

import api.routes.export as export_route  # noqa: E402
from src.generative_editable_config import load_generative_editable_config  # noqa: E402
from src.generative_editable_manifest import read_deck_manifest, read_page_manifest  # noqa: E402
from src.model_profiles import load_default_profiles  # noqa: E402
from src.generative_editable_pipeline import (  # noqa: E402
    GenerativeEditableFallbackError,
    GenerativeEditableSlideInput,
    GenerativeEditableValidationError,
    run_generative_editable_pipeline,
    with_provider_retries,
)
from src.generative_editable_vlm_reconstruction import (  # noqa: E402
    FakeVLMPageAnalysisProvider,
    OpenAIChatVLMPageAnalysisProvider,
    VLMEditablePipelineDependencies,
    run_vlm_editable_pptx_pipeline,
    with_vlm_provider_retries,
)
from src.generative_editable_preview_validator import (  # noqa: E402
    quality_threshold_to_preview_gates,
    render_powerpoint_preview_with_metadata,
    validate_preview_similarity,
)
from src.generative_editable_providers import (  # noqa: E402
    FakeImageEditProvider,
    ImageEditRequest,
    ImageGenerationRequest,
    ProviderConfig,
    ProviderError,
    ProviderTimeoutError,
    safe_provider_error_message,
)


DEFAULT_REPLAY_GLOB = "output/replay-assets/slide_[0-9]*.png"
GATE_CHOICES = {
    "config",
    "vlm_analysis",
    "ocr",
    "image_generation",
    "image_edit_clean",
    "image_edit_asset",
    "image_edit_repair",
}
SLIDE_POSITION_TOLERANCE_EMU = 2
MAX_TIMEOUT_AUGMENT_PPTX_BYTES = 64 * 1024 * 1024


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    subparsers = parser.add_subparsers(dest="command", required=True)

    gates_parser = subparsers.add_parser("gates", help="run provider readiness gates")
    _add_common_args(gates_parser)
    gates_parser.add_argument(
        "--provider-gate",
        action="append",
        choices=sorted(GATE_CHOICES | {"all"}),
        default=None,
        help="provider gate to run; repeatable; defaults to all gates",
    )
    gates_parser.add_argument(
        "--provider-max-attempts",
        type=int,
        default=0,
        help="override provider retry attempts for gates; 0 keeps single-attempt gate behavior",
    )
    gates_parser.add_argument(
        "--provider-retry-backoff",
        type=float,
        default=-1.0,
        help="override provider retry backoff seconds for gates; negative keeps zero backoff",
    )
    gates_parser.add_argument(
        "--gate-wall-timeout",
        type=int,
        default=0,
        help="wall-clock timeout in seconds for each provider gate; 0 runs gates in-process",
    )

    run_parser = subparsers.add_parser("run", help="run the editable PPTX pipeline")
    _add_common_args(run_parser)
    run_parser.add_argument(
        "--mode",
        choices=["legacy", "vlm_first"],
        default="vlm_first",
        help="pipeline mode; legacy uses OCR/CV, vlm_first uses VLM structured analysis",
    )
    run_parser.add_argument("--job-id", default="")
    run_parser.add_argument(
        "--provider-max-attempts",
        type=int,
        default=0,
        help="override configured provider retry attempts for run mode; 0 keeps config",
    )
    run_parser.add_argument(
        "--provider-retry-backoff",
        type=float,
        default=-1.0,
        help="override configured provider retry backoff seconds for run mode; negative keeps config",
    )
    run_parser.add_argument(
        "--pipeline-page-timeout",
        type=float,
        default=0.0,
        help="override configured per-page pipeline timeout seconds for VLM-first run mode; 0 keeps config",
    )
    run_parser.add_argument(
        "--fallback-policy",
        choices=["fail", "text_editable_background", "raster_pptx"],
        default="fail",
    )
    run_parser.add_argument(
        "--isolate-pages",
        action="store_true",
        help="run each input slide in a child process and aggregate page reports",
    )
    run_parser.add_argument(
        "--page-wall-timeout",
        type=int,
        default=0,
        help="wall-clock timeout in seconds for each isolated page child process",
    )

    args = parser.parse_args(argv)
    output_dir = Path(args.output_dir or _default_output_dir()).resolve()
    output_dir.mkdir(parents=True, exist_ok=True)
    if args.command == "run" and not args.job_id:
        args.job_id = _generated_job_id()
    try:
        if args.command == "gates":
            result = _run_gates(args, output_dir)
        elif getattr(args, "isolate_pages", False):
            result = _run_isolated_pages(args, output_dir)
        else:
            result = _run_pipeline(args, output_dir)
    except Exception as exc:
        if args.command == "run" and not getattr(args, "isolate_pages", False):
            job_id = args.job_id or ""
            try:
                input_images = _resolve_input_images(args)
            except Exception:
                input_images = []
            result = _failed_pipeline_report_from_artifacts(
                error=exc,
                output_dir=output_dir,
                output_path=output_dir / f"{job_id}.pptx" if job_id else None,
                artifact_root=output_dir / "artifacts",
                job_id=job_id,
                input_images=input_images,
            )
        else:
            result = {
                "status": "failed",
                "error_type": type(exc).__name__,
                "error": safe_provider_error_message(str(exc)),
                "output_dir": str(output_dir),
            }
        _write_json(output_dir / "report.json", result)
        print(json.dumps(result, ensure_ascii=False))
        return 1
    print(json.dumps(result, ensure_ascii=False))
    return 0 if result.get("status") == "passed" else 1


def _add_common_args(parser: argparse.ArgumentParser) -> None:
    parser.add_argument("--input-images", nargs="*", default=[])
    parser.add_argument("--input-glob", action="append", default=[])
    parser.add_argument("--slides", type=int, default=0)
    parser.add_argument("--aspect-ratio", choices=["16:9", "4:3"], default="16:9")
    parser.add_argument("--output-dir", default="")
    parser.add_argument("--use-fake", action="store_true")
    parser.add_argument(
        "--provider-timeout",
        type=int,
        default=0,
        help="provider call timeout seconds; 0 uses generative_editable_pptx.timeouts.provider_call",
    )


def _provider_timeout_for_args(args: argparse.Namespace) -> int:
    configured = int(getattr(args, "provider_timeout", 0) or 0)
    if configured > 0:
        return configured
    config = load_generative_editable_config(use_fake=bool(getattr(args, "use_fake", False)))
    return config.timeouts.provider_call


def _run_gates(args: argparse.Namespace, output_dir: Path) -> dict[str, Any]:
    images = _resolve_input_images(args)
    requested = _requested_gates(args.provider_gate)
    gate_wall_timeout = int(getattr(args, "gate_wall_timeout", 0) or 0)
    if gate_wall_timeout > 0:
        return _run_gates_with_wall_timeout(args, output_dir, images, requested, gate_wall_timeout)
    provider_timeout = _provider_timeout_for_args(args)
    dependencies = _gate_dependencies(
        use_fake=args.use_fake,
        provider_timeout=provider_timeout,
        provider_max_attempts=getattr(args, "provider_max_attempts", 0),
        provider_retry_backoff_seconds=getattr(args, "provider_retry_backoff", -1.0),
    )
    config = load_generative_editable_config(use_fake=args.use_fake)
    gates = []
    for gate in requested:
        gate_dependencies = dependencies["vlm"] if gate == "vlm_analysis" else dependencies["legacy"]
        gates.append(_run_one_gate(gate, images, gate_dependencies, output_dir, provider_timeout))
    status = "passed" if all(item["status"] == "passed" for item in gates) else "failed"
    report = {
        "status": status,
        "input_images": [str(path) for path in images],
        "output_dir": str(output_dir),
        "config": _redacted_config_summary(config),
        "gates": gates,
    }
    report_path = output_dir / "provider-gates.json"
    _write_json(report_path, report)
    return {"status": status, "report_path": str(report_path), "gates": gates}


def _run_gates_with_wall_timeout(
    args: argparse.Namespace,
    output_dir: Path,
    images: list[Path],
    requested: list[str],
    gate_wall_timeout: int,
) -> dict[str, Any]:
    config = load_generative_editable_config(use_fake=args.use_fake)
    gates: list[dict[str, Any]] = []
    for gate in requested:
        gate_dir = output_dir / f"gate-{gate}"
        gate_dir.mkdir(parents=True, exist_ok=True)
        command = [
            sys.executable,
            str(Path(__file__).resolve()),
            "gates",
            "--input-images",
            *[str(path) for path in images],
            "--output-dir",
            str(gate_dir),
            "--provider-gate",
            gate,
            "--aspect-ratio",
            args.aspect_ratio,
            "--provider-timeout",
            str(args.provider_timeout),
            "--provider-max-attempts",
            str(getattr(args, "provider_max_attempts", 0)),
            "--provider-retry-backoff",
            str(getattr(args, "provider_retry_backoff", -1.0)),
        ]
        if args.use_fake:
            command.append("--use-fake")
        child_result = _run_subprocess_json(
            command,
            timeout_seconds=gate_wall_timeout,
            report_path=gate_dir / "provider-gates.json",
            timeout_payload={
                "gate": gate,
                "status": "failed",
                "error_type": "TimeoutExpired",
                "error": f"provider gate exceeded wall-clock timeout_seconds={gate_wall_timeout}",
                "timeout_seconds": gate_wall_timeout,
            },
        )
        gates.append(_extract_gate_result(gate, child_result))
    status = "passed" if all(item["status"] == "passed" for item in gates) else "failed"
    report = {
        "status": status,
        "input_images": [str(path) for path in images],
        "output_dir": str(output_dir),
        "config": _redacted_config_summary(config),
        "gate_wall_timeout": gate_wall_timeout,
        "gates": gates,
    }
    report_path = output_dir / "provider-gates.json"
    _write_json(report_path, report)
    return {"status": status, "report_path": str(report_path), "gates": gates}


def _extract_gate_result(gate: str, child_result: dict[str, Any]) -> dict[str, Any]:
    child_gates = child_result.get("gates")
    if isinstance(child_gates, list) and child_gates:
        first = child_gates[0]
        if isinstance(first, dict):
            return first
    payload = dict(child_result)
    payload.setdefault("gate", gate)
    payload.setdefault("status", "failed")
    return payload


def _run_one_gate(
    gate: str,
    images: list[Path],
    dependencies,
    output_dir: Path,
    provider_timeout: int,
) -> dict[str, Any]:
    attempt_provider = None
    try:
        if gate == "config":
            return {"gate": gate, "status": "passed"}
        if gate == "vlm_analysis":
            attempt_provider = dependencies.vlm_provider
            analysis = dependencies.vlm_provider.analyze_page(
                str(images[0]),
                timeout_seconds=provider_timeout,
            )
            payload = {
                "gate": gate,
                "status": "passed",
                "text_region_count": len(analysis.text_regions),
                "bitmap_region_count": len(analysis.bitmap_regions),
                "shape_region_count": len(analysis.shape_regions),
                "coordinate_space": {
                    "width": analysis.coordinate_space.width,
                    "height": analysis.coordinate_space.height,
                    "unit": analysis.coordinate_space.unit,
                },
            }
            _add_provider_attempts(payload, attempt_provider)
            return payload
        if gate == "ocr":
            attempt_provider = dependencies.ocr_provider
            result = dependencies.ocr_provider.extract_text(str(images[0]))
            payload = {
                "gate": gate,
                "status": "passed",
                "provider_role": result.provider_role,
                "provider_name": result.provider_name,
                "item_count": len(result.items),
                "image_size": list(result.image_size),
            }
            _add_provider_attempts(payload, attempt_provider)
            return payload
        if gate == "image_generation":
            attempt_provider = dependencies.image_generation_provider
            output_path = output_dir / "provider-image-generation.png"
            result = dependencies.image_generation_provider.generate(
                ImageGenerationRequest(
                    prompt_id="real_gate_image_generation",
                    prompt="Generate a transparent presentation icon with no text.",
                    output_asset_path=str(output_path),
                    asset_root=str(output_dir),
                    visual_reference={"source_image_path": str(images[0])},
                    timeout_seconds=provider_timeout,
                )
            )
            payload = _provider_image_gate_result(gate, result.output_asset_path, result)
            _add_provider_attempts(payload, attempt_provider)
            return payload
        if gate.startswith("image_edit_"):
            role, provider = _image_edit_provider_for_gate(gate, dependencies)
            attempt_provider = provider
            source_path = output_dir / f"provider-{gate}-source.png"
            _write_tiny_edit_probe_image(source_path)
            output_path = output_dir / f"provider-{gate}.png"
            result = provider.edit(
                ImageEditRequest(
                    source_image_path=str(source_path),
                    prompt_id=f"real_gate_{gate}",
                    prompt="Remove the black word TEST and keep the light background. Return image only.",
                    output_asset_path=str(output_path),
                    asset_root=str(output_dir),
                    timeout_seconds=provider_timeout,
                )
            )
            payload = _provider_image_gate_result(gate, result.output_asset_path, result)
            payload["role"] = role
            _add_provider_attempts(payload, attempt_provider)
            return payload
    except (ProviderError, ProviderTimeoutError) as exc:
        payload = {
            "gate": gate,
            "status": "failed",
            "error_type": type(exc).__name__,
            "error": str(exc),
            "retryable": bool(getattr(exc, "retryable", False)),
        }
        _add_provider_error_metadata(payload, exc)
        _add_provider_attempts(payload, attempt_provider, error=exc)
        return payload
    except Exception as exc:
        return {
            "gate": gate,
            "status": "failed",
            "error_type": type(exc).__name__,
            "error": safe_provider_error_message(str(exc)),
        }
    raise ValueError(f"unknown provider gate: {gate}")


def _run_pipeline(args: argparse.Namespace, output_dir: Path) -> dict[str, Any]:
    if getattr(args, "mode", "vlm_first") == "vlm_first":
        return _run_vlm_pipeline(args, output_dir)
    images = _resolve_input_images(args)
    job_id = args.job_id or _generated_job_id()
    artifact_root = output_dir / "artifacts"
    output_path = output_dir / f"{job_id}.pptx"
    raster_fallback_path = output_dir / f"{job_id}.raster-fallback.pptx"
    provider_timeout = _provider_timeout_for_args(args)
    dependencies = replace(
        _with_timeout(_dependencies(use_fake=args.use_fake), provider_timeout),
        allow_source_crop_asset_fallback=False,
    )
    result = run_generative_editable_pipeline(
        slides=[
            GenerativeEditableSlideInput(
                slide_id=f"slide-{index + 1}",
                image_path=str(path),
                text_metadata=[],
            )
            for index, path in enumerate(images)
        ],
        output_path=str(output_path),
        artifact_root=str(artifact_root),
        job_id=job_id,
        aspect_ratio=args.aspect_ratio,
        dependencies=dependencies,
        cleanup_artifacts=False,
        max_page_concurrency=1,
        fallback_policy=args.fallback_policy,
        fallback_output_factory=(
            lambda: _export_raster_pptx_fallback(
                image_paths=images,
                output_path=raster_fallback_path,
                aspect_ratio=args.aspect_ratio,
            )
        )
        if args.fallback_policy == "raster_pptx"
        else None,
    )
    job_dir = artifact_root / job_id
    deck_path = job_dir / "deck.json"
    result_output_path = Path(result.output_path)
    object_stats = _pptx_object_stats(result_output_path)
    reconstruction_issues = _reconstruction_object_issues(
        deck_manifest_path=deck_path,
        artifact_root=job_dir,
        object_stats=object_stats,
    )
    preview_reports = _write_preview_reports(
        source_images=images,
        pptx_path=result_output_path,
        deck_manifest_path=deck_path,
        artifact_root=job_dir,
        similarity_threshold=dependencies.preview_similarity_threshold,
        output_dir=output_dir / "previews",
    )
    status = _status_with_preview_reports(
        _status_with_reconstruction_issues(result.status, reconstruction_issues),
        preview_reports,
    )
    warning_summary = _warning_summary([{"page": 1, "reconstruction_issues": reconstruction_issues}])
    report = {
        "status": status,
        "output_path": str(result_output_path),
        "generative_output_path": str(output_path),
        "artifact_root": str(job_dir),
        "deck_manifest_path": str(deck_path),
        "input_images": [str(path) for path in images],
        "slide_count": len(images),
        "fallback_used": result.fallback_used,
        "object_stats": object_stats,
        "reconstruction_issues": reconstruction_issues,
        **warning_summary,
        "preview_reports": preview_reports,
    }
    _augment_result_with_stage_events(report, job_dir)
    report_path = output_dir / "report.json"
    _write_json(report_path, report)
    return {
        "status": status,
        "output_path": str(result_output_path),
        "generative_output_path": str(output_path),
        "report_path": str(report_path),
        "artifact_root": str(job_dir),
    }


def _run_vlm_pipeline(args: argparse.Namespace, output_dir: Path) -> dict[str, Any]:
    images = _resolve_input_images(args)
    job_id = args.job_id or _generated_job_id()
    artifact_root = output_dir / "artifacts"
    output_path = output_dir / f"{job_id}.pptx"
    provider_timeout = _provider_timeout_for_args(args)
    dependencies = _vlm_dependencies(
        use_fake=args.use_fake,
        provider_timeout=provider_timeout,
        provider_max_attempts=getattr(args, "provider_max_attempts", 0),
        provider_retry_backoff_seconds=getattr(args, "provider_retry_backoff", -1.0),
        page_timeout_seconds=getattr(args, "pipeline_page_timeout", 0.0),
    )
    result = run_vlm_editable_pptx_pipeline(
        slides=[
            {"slide_id": f"slide-{index + 1}", "image_path": str(path)}
            for index, path in enumerate(images)
        ],
        output_path=str(output_path),
        artifact_root=str(artifact_root),
        job_id=job_id,
        aspect_ratio=args.aspect_ratio,
        dependencies=dependencies,
        cleanup_artifacts=False,
    )
    job_dir = artifact_root / job_id
    deck_path = job_dir / "deck.json"
    result_output_path = Path(result.output_path)
    object_stats = _pptx_object_stats(result_output_path)
    reconstruction_issues = _reconstruction_object_issues(
        deck_manifest_path=deck_path,
        artifact_root=job_dir,
        object_stats=object_stats,
    )
    preview_reports = _write_preview_reports(
        source_images=images,
        pptx_path=result_output_path,
        deck_manifest_path=deck_path,
        artifact_root=job_dir,
        similarity_threshold=0.92,
        output_dir=output_dir / "previews",
    )
    status = _status_with_preview_reports(
        _status_with_reconstruction_issues(result.status, reconstruction_issues),
        preview_reports,
    )
    warning_summary = _warning_summary([{"page": 1, "reconstruction_issues": reconstruction_issues}])
    report = {
        "status": status,
        "mode": "vlm_first",
        "output_path": str(result_output_path),
        "generative_output_path": str(output_path),
        "artifact_root": str(job_dir),
        "deck_manifest_path": str(deck_path),
        "input_images": [str(path) for path in images],
        "slide_count": len(images),
        "fallback_used": result.fallback_used,
        "object_stats": object_stats,
        "reconstruction_issues": reconstruction_issues,
        **warning_summary,
        "preview_reports": preview_reports,
    }
    _augment_result_with_stage_events(report, job_dir)
    report_path = output_dir / "report.json"
    _write_json(report_path, report)
    return {
        "status": status,
        "output_path": str(result_output_path),
        "generative_output_path": str(output_path),
        "report_path": str(report_path),
        "artifact_root": str(job_dir),
    }


def _export_raster_pptx_fallback(
    *,
    image_paths: list[Path],
    output_path: Path,
    aspect_ratio: str,
) -> str:
    export_route._export_pptx([str(path) for path in image_paths], str(output_path), aspect_ratio=aspect_ratio)
    return str(output_path.resolve())


def _generated_job_id() -> str:
    return f"real-editable-{datetime.now().strftime('%Y%m%d%H%M%S')}"


def _failed_pipeline_report_from_artifacts(
    *,
    error: Exception,
    output_dir: Path,
    output_path: Path | None,
    artifact_root: Path,
    job_id: str,
    input_images: list[Path],
) -> dict[str, Any]:
    job_dir = artifact_root / job_id if job_id else artifact_root
    report: dict[str, Any] = {
        "status": "failed",
        "error_type": type(error).__name__,
        "error": safe_provider_error_message(str(error)),
        "output_dir": str(output_dir),
        "artifact_root": str(job_dir),
        "input_images": [str(path) for path in input_images],
        "slide_count": len(input_images),
    }
    _add_provider_error_metadata(report, error)
    if isinstance(error, (GenerativeEditableValidationError, GenerativeEditableFallbackError)):
        report["validation_status"] = error.validation_report.status
        report["validation_issues"] = [issue.code for issue in error.validation_report.issues]
    deck_path = job_dir / "deck.json"
    if deck_path.is_file():
        report["deck_manifest_path"] = str(deck_path)
    if output_path is not None and output_path.is_file():
        report["output_path"] = str(output_path)
        try:
            object_stats = _pptx_object_stats(output_path)
        except Exception as exc:
            report["object_stats_error"] = safe_provider_error_message(str(exc))
        else:
            report["object_stats"] = object_stats
            if deck_path.is_file():
                try:
                    report["reconstruction_issues"] = _reconstruction_object_issues(
                        deck_manifest_path=deck_path,
                        artifact_root=job_dir,
                        object_stats=object_stats,
                    )
                except Exception as exc:
                    report["reconstruction_issues_error"] = safe_provider_error_message(str(exc))
                if input_images:
                    try:
                        config = load_generative_editable_config()
                        threshold = config.preview_similarity_threshold
                    except Exception:
                        threshold = 0.92
                    try:
                        report["preview_reports"] = _write_preview_reports(
                            source_images=input_images,
                            pptx_path=output_path,
                            deck_manifest_path=deck_path,
                            artifact_root=job_dir,
                            similarity_threshold=threshold,
                            output_dir=output_dir / "previews",
                        )
                    except Exception as exc:
                        report["preview_reports_error"] = safe_provider_error_message(str(exc))
    _augment_result_with_stage_events(report, job_dir)
    return report


def _run_isolated_pages(args: argparse.Namespace, output_dir: Path) -> dict[str, Any]:
    images = _resolve_input_images(args)
    job_id = args.job_id or f"real-editable-{datetime.now().strftime('%Y%m%d%H%M%S')}"
    page_reports = []
    provider_timeout = _provider_timeout_for_args(args)
    for index, image in enumerate(images, start=1):
        page_dir = output_dir / f"page-{index:02d}"
        page_dir.mkdir(parents=True, exist_ok=True)
        page_job_id = f"{job_id}-page-{index:02d}"
        command = [
            sys.executable,
            str(Path(__file__).resolve()),
            "run",
            "--mode",
            getattr(args, "mode", "vlm_first"),
            "--input-images",
            str(image),
            "--output-dir",
            str(page_dir),
            "--job-id",
            page_job_id,
            "--aspect-ratio",
            args.aspect_ratio,
            "--provider-timeout",
            str(provider_timeout),
            "--provider-max-attempts",
            str(getattr(args, "provider_max_attempts", 0)),
            "--provider-retry-backoff",
            str(getattr(args, "provider_retry_backoff", -1.0)),
            "--pipeline-page-timeout",
            str(getattr(args, "pipeline_page_timeout", 0.0)),
            "--fallback-policy",
            args.fallback_policy,
        ]
        if args.use_fake:
            command.append("--use-fake")
        page_report_path = page_dir / "report.json"
        page_output_path = page_dir / f"{page_job_id}.pptx"
        page_artifact_root = page_dir / "artifacts" / page_job_id
        page_result = _run_subprocess_json(
            command,
            timeout_seconds=max(0, int(args.page_wall_timeout or 0)),
            report_path=page_report_path,
            timeout_payload={
                "status": "failed",
                "page": index,
                "input_image": str(image),
                "job_id": page_job_id,
                "output_path": str(page_output_path),
                "artifact_root": str(page_artifact_root),
                "deck_manifest_path": str(page_artifact_root / "deck.json"),
            },
        )
        page_reports.append(_isolated_page_report_summary(index, image, page_report_path, page_result))
    status = _aggregate_page_status(page_reports)
    warning_summary = _warning_summary(page_reports)
    report = {
        "status": status,
        "mode": "isolated_pages",
        "output_dir": str(output_dir),
        "input_images": [str(path) for path in images],
        "slide_count": len(images),
        "page_wall_timeout": int(args.page_wall_timeout or 0),
        **warning_summary,
        "page_reports": page_reports,
    }
    report_path = output_dir / "report.json"
    _write_json(report_path, report)
    return {
        "status": status,
        "report_path": str(report_path),
        "output_dir": str(output_dir),
    }


def _run_subprocess_json(
    command: list[str],
    *,
    timeout_seconds: int,
    report_path: Path,
    timeout_payload: dict[str, Any],
) -> dict[str, Any]:
    process = subprocess.Popen(
        command,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        text=True,
        start_new_session=True,
    )
    try:
        stdout, stderr = process.communicate(
            timeout=timeout_seconds if timeout_seconds > 0 else None,
        )
        result = _last_json_line(stdout)
        if not result:
            result = {
                "status": "failed",
                "error_type": "SubprocessOutputError",
                "error": "child process did not print a JSON result",
                "returncode": process.returncode,
            }
            _add_subprocess_output_tails(result, stdout=stdout, stderr=stderr)
        if process.returncode != 0 and result.get("status") == "passed":
            result = {
                **result,
                "status": "failed",
                "error_type": "SubprocessFailed",
                "returncode": process.returncode,
            }
            _add_subprocess_output_tails(result, stdout=stdout, stderr=stderr)
        if not report_path.exists():
            _write_json(report_path, result)
        return result
    except subprocess.TimeoutExpired as exc:
        _kill_process_tree(process.pid)
        _finish_timed_out_process(process)
        existing_result = _read_existing_report(report_path)
        if existing_result is not None:
            result = {
                **existing_result,
                "subprocess_cleanup": {
                    "error_type": "TimeoutExpired",
                    "error": f"child process stdout/stderr did not close before timeout_seconds={timeout_seconds}",
                    "timeout_seconds": timeout_seconds,
                },
            }
            _write_json(report_path, result)
            return result
        result = {
            **timeout_payload,
            "status": "failed",
            "error_type": "TimeoutExpired",
            "error": timeout_payload.get(
                "error",
                f"child process exceeded wall-clock timeout_seconds={timeout_seconds}",
            ),
            "timeout_seconds": timeout_seconds,
        }
        result = _augment_timeout_result_from_artifacts(result)
        _write_json(report_path, result)
        return result


def _add_subprocess_output_tails(
    result: dict[str, Any],
    *,
    stdout: str,
    stderr: str,
    max_chars: int = 2000,
) -> None:
    if stdout:
        result["stdout_tail"] = safe_provider_error_message(stdout[-max_chars:])
    if stderr:
        result["stderr_tail"] = safe_provider_error_message(stderr[-max_chars:])


def _read_existing_report(report_path: Path) -> dict[str, Any] | None:
    if not report_path.exists():
        return None
    try:
        value = json.loads(report_path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError):
        return None
    return value if isinstance(value, dict) else None


def _finish_timed_out_process(process: subprocess.Popen[str]) -> None:
    _close_process_pipe(process.stdout)
    _close_process_pipe(process.stderr)
    try:
        process.wait(timeout=2)
        return
    except subprocess.TimeoutExpired:
        _kill_process_tree(process.pid)
        try:
            process.kill()
        except ProcessLookupError:
            pass
    try:
        process.wait(timeout=1)
    except subprocess.TimeoutExpired:
        pass


def _close_process_pipe(pipe: Any) -> None:
    if pipe is None:
        return
    try:
        pipe.close()
    except OSError:
        pass


def _isolated_page_report_summary(
    index: int,
    image: Path,
    page_report_path: Path,
    page_result: dict[str, Any],
) -> dict[str, Any]:
    summary = {
        "page": index,
        "input_image": str(image),
        "status": page_result.get("status", "failed"),
        "report_path": str(page_report_path),
        "output_path": page_result.get("output_path", ""),
        "artifact_root": page_result.get("artifact_root", ""),
        "error_type": page_result.get("error_type", ""),
        "error": page_result.get("error", ""),
    }
    for key in (
        "object_stats",
        "object_stats_error",
        "reconstruction_issues",
        "reconstruction_issues_error",
        "validation_status",
        "validation_issues",
        "last_stage_event",
        "active_stage",
        "stage_events_error",
    ):
        if key in page_result:
            summary[key] = page_result[key]
    return summary


def _kill_process_tree(pid: int | None) -> None:
    if not pid:
        return
    try:
        os.killpg(pid, signal.SIGKILL)
    except ProcessLookupError:
        return
    except OSError:
        try:
            os.kill(pid, signal.SIGKILL)
        except OSError:
            return


def _augment_timeout_result_from_artifacts(result: dict[str, Any]) -> dict[str, Any]:
    output_path_value = str(result.get("output_path") or "")
    output_path = Path(output_path_value) if output_path_value else None
    if output_path is not None and not output_path.is_file():
        result.pop("output_path", None)
        output_path = None
    if output_path is not None:
        try:
            if output_path.stat().st_size > MAX_TIMEOUT_AUGMENT_PPTX_BYTES:
                raise ValueError("pptx exceeds timeout augmentation size cap")
            if not _file_size_is_stable(output_path):
                raise ValueError("pptx is still being written")
            object_stats = _pptx_object_stats(output_path)
        except Exception as exc:
            result["object_stats_error"] = safe_provider_error_message(str(exc))
        else:
            result["object_stats"] = object_stats
    artifact_root = Path(str(result.get("artifact_root") or ""))
    _augment_result_with_stage_events(result, artifact_root)
    deck_path = artifact_root / "deck.json" if artifact_root else None
    if output_path is not None and deck_path is not None and deck_path.is_file() and "object_stats" in result:
        try:
            result["reconstruction_issues"] = _reconstruction_object_issues(
                deck_manifest_path=deck_path,
                artifact_root=artifact_root,
                object_stats=result["object_stats"],
            )
        except Exception as exc:
            result["reconstruction_issues_error"] = safe_provider_error_message(str(exc))
    return result


def _augment_result_with_stage_events(result: dict[str, Any], artifact_root: Path) -> None:
    if not artifact_root:
        return
    stage_events_path = artifact_root / "stage-events.jsonl"
    if not stage_events_path.is_file():
        return
    events, skipped_bad_line = _read_stage_events(stage_events_path)
    if skipped_bad_line:
        result["stage_events_error"] = "malformed stage-events line skipped"
    if not events:
        return
    last_event = events[-1]
    result["last_stage_event"] = last_event
    if last_event.get("status") in {"started", "failed"} and last_event.get("stage"):
        result["active_stage"] = last_event["stage"]


def _read_stage_events(path: Path) -> tuple[list[dict[str, Any]], bool]:
    events: list[dict[str, Any]] = []
    skipped_bad_line = False
    for line in path.read_text(encoding="utf-8").splitlines():
        if not line.strip():
            continue
        try:
            parsed = json.loads(line)
        except json.JSONDecodeError:
            skipped_bad_line = True
            continue
        if isinstance(parsed, dict):
            events.append(parsed)
    return events, skipped_bad_line


def _add_provider_error_metadata(payload: dict[str, Any], error: Exception) -> None:
    if hasattr(error, "retryable"):
        payload["retryable"] = bool(getattr(error, "retryable"))
    status_code = getattr(error, "status_code", None)
    provider_error_code = getattr(error, "provider_error_code", "")
    if status_code is not None:
        payload["status_code"] = status_code
    if provider_error_code:
        payload["provider_error_code"] = provider_error_code


def _add_provider_attempts(
    payload: dict[str, Any],
    provider: Any | None,
    *,
    error: Exception | None = None,
) -> None:
    attempts = getattr(provider, "last_attempts", None)
    if not attempts and error is not None:
        attempts = getattr(error, "attempts", None)
    if attempts:
        payload["provider_attempts"] = [dict(item) for item in attempts if isinstance(item, dict)]


def _file_size_is_stable(path: Path, *, interval_seconds: float = 0.2) -> bool:
    first = path.stat().st_size
    time.sleep(interval_seconds)
    try:
        second = path.stat().st_size
    except FileNotFoundError:
        return False
    return first == second


def _last_json_line(stdout: str) -> dict[str, Any] | None:
    for line in reversed(stdout.splitlines()):
        line = line.strip()
        if not line:
            continue
        try:
            parsed = json.loads(line)
        except json.JSONDecodeError:
            continue
        return parsed if isinstance(parsed, dict) else None
    return None


def _dependencies(*, use_fake: bool):
    if use_fake:
        return export_route._build_fake_generative_editable_pipeline_dependencies()
    return export_route._build_generative_editable_pipeline_dependencies()


def _gate_dependencies(
    *,
    use_fake: bool,
    provider_timeout: int,
    provider_max_attempts: int = 0,
    provider_retry_backoff_seconds: float = -1.0,
) -> dict[str, Any]:
    retry_attempts = provider_max_attempts if provider_max_attempts and provider_max_attempts > 0 else 1
    retry_backoff = provider_retry_backoff_seconds if provider_retry_backoff_seconds >= 0 else 0
    legacy = _with_timeout(_dependencies(use_fake=use_fake), provider_timeout)
    legacy = replace(
        legacy,
        provider_max_attempts=retry_attempts,
        provider_retry_backoff_seconds=retry_backoff,
    )
    vlm = _vlm_dependencies(
        use_fake=use_fake,
        provider_timeout=provider_timeout,
        provider_max_attempts=retry_attempts,
        provider_retry_backoff_seconds=retry_backoff,
    )
    vlm = replace(
        vlm,
        provider_max_attempts=retry_attempts,
        provider_retry_backoff_seconds=retry_backoff,
    )
    return {
        "legacy": with_provider_retries(legacy),
        "vlm": with_vlm_provider_retries(vlm),
    }


def _vlm_dependencies(
    *,
    use_fake: bool,
    provider_timeout: int,
    provider_max_attempts: int = 0,
    provider_retry_backoff_seconds: float = -1.0,
    page_timeout_seconds: float = 0.0,
) -> VLMEditablePipelineDependencies:
    config = load_generative_editable_config(use_fake=use_fake)
    retry_attempts = (
        provider_max_attempts
        if provider_max_attempts and provider_max_attempts > 0
        else config.retries.provider_max_attempts
    )
    retry_backoff = (
        provider_retry_backoff_seconds
        if provider_retry_backoff_seconds >= 0
        else config.retries.backoff_seconds
    )
    page_timeout = page_timeout_seconds if page_timeout_seconds > 0 else config.timeouts.page
    if use_fake:
        return VLMEditablePipelineDependencies(
            vlm_provider=FakeVLMPageAnalysisProvider(),
            image_edit_provider=FakeImageEditProvider(config.clean_base_model),
            asset_sheet_image_edit_provider=FakeImageEditProvider(config.asset_sheet_model),
            ocr_provider=export_route._ocr_provider_for_config(
                config.ocr,
                timeout_seconds=provider_timeout,
            ),
            provider_timeout_seconds=provider_timeout,
            provider_max_attempts=retry_attempts,
            provider_retry_backoff_seconds=retry_backoff,
            page_timeout_seconds=page_timeout,
            ocr_min_confidence=config.ocr_min_confidence,
        )
    profiles = load_default_profiles()
    if profiles is None or profiles.vlm is None:
        raise RuntimeError("VLM model profile is required for --mode vlm_first")
    vlm_config = ProviderConfig(
        role="VLM",
        model=profiles.vlm.model,
        base_url=profiles.vlm.base_url,
        api_key=profiles.vlm.api_key,
        adapter=profiles.vlm.adapter,
    )
    return VLMEditablePipelineDependencies(
        vlm_provider=OpenAIChatVLMPageAnalysisProvider(vlm_config),
        image_edit_provider=export_route._image_edit_provider_for_config(config.clean_base_model),
        asset_sheet_image_edit_provider=export_route._image_edit_provider_for_config(config.asset_sheet_model),
        ocr_provider=export_route._ocr_provider_for_config(
            config.ocr,
            timeout_seconds=provider_timeout,
        ),
        provider_timeout_seconds=provider_timeout,
        provider_max_attempts=retry_attempts,
        provider_retry_backoff_seconds=retry_backoff,
        page_timeout_seconds=page_timeout,
        ocr_min_confidence=config.ocr_min_confidence,
    )


def _with_timeout(dependencies, provider_timeout: int):
    if provider_timeout <= 0:
        return dependencies
    if hasattr(dependencies.ocr_provider, "timeout_seconds"):
        dependencies.ocr_provider.timeout_seconds = provider_timeout
    return replace(dependencies, provider_timeout_seconds=provider_timeout)


def _resolve_input_images(args: argparse.Namespace) -> list[Path]:
    paths: list[Path] = [Path(value).resolve() for value in args.input_images]
    for pattern in args.input_glob:
        paths.extend(Path(value).resolve() for value in sorted(glob.glob(pattern), key=_natural_key))
    paths = _filter_replay_slide_images(paths)
    if not paths:
        paths = _default_replay_input_images()
    if args.slides and args.slides > 0:
        paths = paths[: args.slides]
    missing = [str(path) for path in paths if not path.is_file()]
    if missing:
        raise FileNotFoundError("input slide images are missing: " + ", ".join(missing))
    if not paths:
        raise FileNotFoundError(
            f"No input images found. Expected replay assets matching {DEFAULT_REPLAY_GLOB}."
        )
    return paths


def _default_replay_input_images() -> list[Path]:
    matches = [Path(value).resolve() for value in sorted(glob.glob(DEFAULT_REPLAY_GLOB), key=_natural_key)]
    numeric_slide_images = _filter_replay_slide_images(matches)
    return numeric_slide_images or matches


def _filter_replay_slide_images(paths: list[Path]) -> list[Path]:
    replay_like = [path for path in paths if path.name.startswith("slide_")]
    if not replay_like:
        return paths
    numeric_slide_images = [
        path for path in paths if re.fullmatch(r"slide_\d+\.png", path.name)
    ]
    return numeric_slide_images or paths


def _requested_gates(values: list[str] | None) -> list[str]:
    if not values or "all" in values:
        return [
            "config",
            "vlm_analysis",
            "ocr",
            "image_edit_clean",
            "image_edit_asset",
            "image_edit_repair",
            "image_generation",
        ]
    return values


def _image_edit_provider_for_gate(gate: str, dependencies):
    if gate == "image_edit_clean":
        return "clean_base_model", dependencies.image_edit_provider
    if gate == "image_edit_asset":
        return "asset_sheet_model", dependencies.asset_sheet_image_edit_provider or dependencies.image_edit_provider
    if gate == "image_edit_repair":
        return "repair_model", dependencies.repair_image_edit_provider or dependencies.image_edit_provider
    raise ValueError(f"unknown image edit gate: {gate}")


def _provider_image_gate_result(gate: str, output_path: str, result) -> dict[str, Any]:
    return {
        "gate": gate,
        "status": "passed",
        "output_path": output_path,
        "provider_role": result.provider_role,
        "provider_name": result.provider_name,
    }


def _pptx_object_stats(pptx_path: Path) -> dict[str, Any]:
    presentation = Presentation(str(pptx_path))
    slides = []
    totals: Counter[str] = Counter()
    background_picture_counts = _pptx_slide_background_picture_counts(pptx_path)
    for index, slide in enumerate(presentation.slides, start=1):
        counts: Counter[str] = Counter()
        background_picture_count = background_picture_counts.get(index, 0)
        if background_picture_count:
            totals["BACKGROUND_PICTURE"] += background_picture_count
        full_slide_picture_count = 0
        non_full_slide_picture_count = 0
        for shape in slide.shapes:
            name = _shape_type_name(shape)
            counts[name] += 1
            totals[name] += 1
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                if _is_full_slide_shape(shape, presentation.slide_width, presentation.slide_height):
                    full_slide_picture_count += 1
                else:
                    non_full_slide_picture_count += 1
        slides.append(
            {
                "page": index,
                "shape_counts": dict(sorted(counts.items())),
                "background_picture_count": background_picture_count,
                "full_slide_picture_count": full_slide_picture_count,
                "non_full_slide_picture_count": non_full_slide_picture_count,
                "text_box_count": counts.get("TEXT_BOX", 0),
            }
        )
    return {
        "slide_count": len(presentation.slides),
        "slides": slides,
        "totals": dict(sorted(totals.items())),
    }


def _pptx_slide_background_picture_counts(pptx_path: Path) -> dict[int, int]:
    counts: dict[int, int] = {}
    namespaces = {
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    }
    try:
        with zipfile.ZipFile(pptx_path) as archive:
            slide_parts = sorted(
                (
                    name
                    for name in archive.namelist()
                    if re.fullmatch(r"ppt/slides/slide\d+\.xml", name)
                ),
                key=lambda value: int(re.search(r"slide(\d+)\.xml$", value).group(1)),
            )
            for slide_index, slide_part in enumerate(slide_parts, start=1):
                root = ET.fromstring(archive.read(slide_part))
                background = root.find("p:cSld/p:bg", namespaces)
                if background is not None and background.findall(".//a:blip", namespaces):
                    counts[slide_index] = 1
    except Exception:
        return counts
    return counts


def _is_full_slide_shape(shape, slide_width: int, slide_height: int) -> bool:
    return (
        abs(shape.left) <= SLIDE_POSITION_TOLERANCE_EMU
        and abs(shape.top) <= SLIDE_POSITION_TOLERANCE_EMU
        and abs(shape.width - slide_width) <= SLIDE_POSITION_TOLERANCE_EMU
        and abs(shape.height - slide_height) <= SLIDE_POSITION_TOLERANCE_EMU
    )


def _reconstruction_object_issues(
    *,
    deck_manifest_path: Path,
    artifact_root: Path,
    object_stats: dict[str, Any],
) -> list[dict[str, Any]]:
    issues: list[dict[str, Any]] = []
    deck = read_deck_manifest(deck_manifest_path)
    slide_stats = object_stats.get("slides", [])
    for index, page_ref in enumerate(deck.page_manifest_paths):
        page = read_page_manifest(artifact_root / page_ref)
        stats = slide_stats[index] if index < len(slide_stats) else {}
        shape_counts = stats.get("shape_counts", {})
        decomposed_visual_count = len(page.bitmap_assets) + len(page.native_shapes)
        picture_count = int(shape_counts.get("PICTURE", 0))
        background_picture_count = int(stats.get("background_picture_count", 0))
        full_slide_picture_count = int(stats.get("full_slide_picture_count", 0))
        non_full_slide_picture_count = int(stats.get("non_full_slide_picture_count", 0))
        text_box_count = int(stats.get("text_box_count", 0))
        if page.provenance.get("chosen_background_kind") == "source_raster_guardrail":
            issues.append(
                {
                    "code": "source_raster_guardrail_degraded",
                    "severity": "error",
                    "page": index + 1,
                    "slide_id": page.slide_id,
                    "message": "Slide used source raster guardrail and is not an editable reconstruction",
                    "details": {
                        "bitmap_asset_count": len(page.bitmap_assets),
                        "native_shape_count": len(page.native_shapes),
                        "text_box_count": int(stats.get("text_box_count", 0)),
                        "full_slide_picture_count": full_slide_picture_count,
                        "source_raster_guardrail": page.provenance.get("source_raster_guardrail", {}),
                    },
                }
            )
        issues.extend(
            _bitmap_asset_coverage_issues(
                page=page,
                page_number=index + 1,
                stats=stats,
            )
        )
        issues.extend(
            _overlapping_duplicate_text_issues(
                page=page,
                page_number=index + 1,
            )
        )
        filtered_ocr_issues = [
            issue
            for issue in page.provenance.get("text_issues", [])
            if isinstance(issue, dict) and issue.get("code") == "ignored_spurious_ocr"
        ]
        if int(stats.get("text_box_count", 0)) == 0 and filtered_ocr_issues:
            filtered_reason_counts = Counter(
                str(issue.get("code", "unknown")) for issue in filtered_ocr_issues
            )
            ocr_returned_count = _ocr_returned_item_count(page, artifact_root)
            issues.append(
                {
                    "code": "no_editable_text_after_ocr_filtering",
                    "severity": "warning",
                    "page": index + 1,
                    "slide_id": page.slide_id,
                    "message": "OCR returned text-like content but all items were filtered as spurious",
                    "details": {
                        "filtered_ocr_issue_count": len(filtered_ocr_issues),
                        "filtered_reason_counts": dict(sorted(filtered_reason_counts.items())),
                        "ocr_returned_count": ocr_returned_count,
                        "sample_filtered_text": [
                            str(issue.get("ocr_text", ""))[:120]
                            for issue in filtered_ocr_issues[:3]
                        ],
                    },
                }
            )
        if (
            page.provenance.get("chosen_background_kind")
            == "source_preserving_low_opacity_text_overlay"
            and full_slide_picture_count >= 1
            and decomposed_visual_count == 0
        ):
            issues.append(
                {
                    "code": "source_preserving_low_opacity_overlay_degraded",
                    "severity": "error",
                    "page": index + 1,
                    "slide_id": page.slide_id,
                    "message": "Slide preserved source visual fidelity with low-opacity editable OCR text overlay",
                    "details": {
                        "bitmap_asset_count": len(page.bitmap_assets),
                        "native_shape_count": len(page.native_shapes),
                        "full_slide_picture_count": full_slide_picture_count,
                        "text_box_count": int(stats.get("text_box_count", 0)),
                        "text_overlay_opacity": page.provenance.get("text_overlay_opacity"),
                    },
                }
            )
            continue
        if (
            decomposed_visual_count == 0
            and text_box_count == 0
            and (
                (picture_count > 0 and full_slide_picture_count == picture_count)
                or background_picture_count > 0
            )
            and non_full_slide_picture_count == 0
        ):
            issues.append(
                {
                    "code": "no_decomposed_visual_elements",
                    "severity": "error",
                    "page": index + 1,
                    "slide_id": page.slide_id,
                    "message": "PPTX slide contains only full-slide picture objects and no decomposed visual elements",
                    "details": {
                        "bitmap_asset_count": len(page.bitmap_assets),
                        "native_shape_count": len(page.native_shapes),
                        "background_picture_count": background_picture_count,
                        "full_slide_picture_count": full_slide_picture_count,
                        "text_box_count": text_box_count,
                    },
                }
            )
        if (
            page.provenance.get("chosen_background_kind") == "source_preserving_text_fast_path"
            and full_slide_picture_count >= 1
            and int(shape_counts.get("AUTO_SHAPE", 0)) == 0
            and int(shape_counts.get("LINE", 0)) == 0
            and all(
                asset.provenance.get("asset_strategy") == "source_preserving_anchor"
                for asset in page.bitmap_assets
            )
        ):
            issues.append(
                {
                    "code": "source_preserving_fast_path_degraded",
                    "severity": "error",
                    "page": index + 1,
                    "slide_id": page.slide_id,
                    "message": "Slide used source-preserving fast path with only anchor bitmap assets",
                    "details": {
                        "bitmap_asset_count": len(page.bitmap_assets),
                        "native_shape_count": len(page.native_shapes),
                        "text_box_count": int(stats.get("text_box_count", 0)),
                        "full_slide_picture_count": full_slide_picture_count,
                    },
                }
            )
    return issues


def _overlapping_duplicate_text_issues(
    *,
    page,
    page_number: int,
) -> list[dict[str, Any]]:
    issues: list[dict[str, Any]] = []
    text_boxes = list(page.text_boxes)
    for left_index, left in enumerate(text_boxes):
        for right_index in range(left_index + 1, len(text_boxes)):
            right = text_boxes[right_index]
            if not _text_boxes_are_overlapping_duplicates(left, right):
                continue
            issues.append(
                {
                    "code": "overlapping_duplicate_text_boxes",
                    "severity": "warning",
                    "page": page_number,
                    "slide_id": page.slide_id,
                    "message": "Two editable text boxes contain duplicate text and overlap geometrically",
                    "details": {
                        "duplicate_text": left.text,
                        "left_text": left.text,
                        "right_text": right.text,
                        "left_bbox": list(left.source_pixel_bbox),
                        "right_bbox": list(right.source_pixel_bbox),
                    },
                }
            )
    return issues


def _text_boxes_are_overlapping_duplicates(left: Any, right: Any) -> bool:
    if not _normalized_texts_are_exact_duplicate(str(left.text), str(right.text)):
        return False
    overlap = _bbox_intersection_area(left.source_pixel_bbox, right.source_pixel_bbox)
    if overlap <= 0:
        return False
    smaller_area = min(_bbox_area(left.source_pixel_bbox), _bbox_area(right.source_pixel_bbox))
    larger_area = max(_bbox_area(left.source_pixel_bbox), _bbox_area(right.source_pixel_bbox))
    return (
        overlap / max(1, smaller_area) >= 0.70
        and overlap / max(1, larger_area) >= 0.45
    )


def _text_similarity(left: str, right: str) -> float:
    normalized_left = _normalize_text_for_match(left)
    normalized_right = _normalize_text_for_match(right)
    if not normalized_left or not normalized_right:
        return 0.0
    if normalized_left in normalized_right or normalized_right in normalized_left:
        return 1.0
    return SequenceMatcher(None, normalized_left, normalized_right).ratio()


def _normalize_text_for_match(value: str) -> str:
    return re.sub(r"\s+", "", str(value)).lower()


def _normalized_texts_are_exact_duplicate(left: str, right: str) -> bool:
    normalized_left = _normalize_text_for_match(left)
    normalized_right = _normalize_text_for_match(right)
    if normalized_left != normalized_right:
        return False
    return len(normalized_left) >= 4


def _bbox_intersection_area(
    left_bbox: tuple[int, int, int, int],
    right_bbox: tuple[int, int, int, int],
) -> int:
    left = max(int(left_bbox[0]), int(right_bbox[0]))
    top = max(int(left_bbox[1]), int(right_bbox[1]))
    right = min(int(left_bbox[2]), int(right_bbox[2]))
    bottom = min(int(left_bbox[3]), int(right_bbox[3]))
    if right <= left or bottom <= top:
        return 0
    return (right - left) * (bottom - top)


def _bitmap_asset_coverage_issues(
    *,
    page,
    page_number: int,
    stats: dict[str, Any],
) -> list[dict[str, Any]]:
    source_width, source_height = page.source_image_size
    page_area = max(1, int(source_width) * int(source_height))
    if not page.bitmap_assets or page_area <= 1:
        return []
    source_preserved_assets = [
        asset
        for asset in page.bitmap_assets
        if _is_source_preserved_bitmap_asset(asset)
    ]
    if not source_preserved_assets:
        return []
    issues: list[dict[str, Any]] = []
    largest_asset = max(
        source_preserved_assets,
        key=lambda asset: _source_preserved_asset_effective_area_ratio(asset, (int(source_width), int(source_height))),
    )
    largest_ratio = _source_preserved_asset_effective_area_ratio(
        largest_asset,
        (int(source_width), int(source_height)),
    )
    largest_bbox_ratio = _bbox_area(largest_asset.source_pixel_bbox) / float(page_area)
    combined_bbox_ratio = _combined_bbox_coverage_ratio(
        [asset.source_pixel_bbox for asset in source_preserved_assets],
        (int(source_width), int(source_height)),
    )
    combined_visible_ratio = min(
        1.0,
        sum(
            _source_preserved_asset_effective_area_ratio(asset, (int(source_width), int(source_height)))
            for asset in source_preserved_assets
        ),
    )
    combined_ratio = min(combined_bbox_ratio, combined_visible_ratio)
    text_box_count = int(stats.get("text_box_count", 0))
    native_shape_count = len(page.native_shapes)
    details = {
        "largest_asset_id": largest_asset.asset_id,
        "largest_asset_area_ratio": round(largest_ratio, 4),
        "largest_asset_bbox_area_ratio": round(largest_bbox_ratio, 4),
        "combined_bitmap_asset_coverage_ratio": round(combined_ratio, 4),
        "combined_bitmap_asset_bbox_coverage_ratio": round(combined_bbox_ratio, 4),
        "combined_bitmap_asset_visible_area_ratio": round(combined_visible_ratio, 4),
        "bitmap_asset_count": len(page.bitmap_assets),
        "source_preserved_bitmap_asset_count": len(source_preserved_assets),
        "native_shape_count": native_shape_count,
        "text_box_count": text_box_count,
    }
    if largest_ratio >= 0.80:
        issues.append(
            {
                "code": "oversized_bitmap_asset_coverage",
                "severity": "error",
                "page": page_number,
                "slide_id": page.slide_id,
                "message": "A single source-preserved bitmap asset covers most of the slide",
                "details": details,
            }
        )
        return issues
    if combined_ratio >= 0.85 and text_box_count + native_shape_count < 3:
        issues.append(
            {
                "code": "excessive_bitmap_asset_coverage",
                "severity": "error",
                "page": page_number,
                "slide_id": page.slide_id,
                "message": "Source-preserved bitmap assets cover most of the slide without enough editable structure",
                "details": details,
            }
        )
        return issues
    if combined_ratio >= 0.65 and not _has_split_row_level_bitmap_structure(
        source_preserved_assets,
        largest_ratio=largest_ratio,
        structure_count=text_box_count + native_shape_count,
    ):
        issues.append(
            {
                "code": "high_bitmap_asset_coverage",
                "severity": "error",
                "page": page_number,
                "slide_id": page.slide_id,
                "message": "Source-preserved bitmap assets cover a large part of the slide",
                "details": details,
            }
        )
    return issues


def _is_source_preserved_bitmap_asset(asset: Any) -> bool:
    return asset.provenance.get("asset_strategy") in {
        "masked_source_element",
        "source_preserved_crop",
        "source_preserving_anchor",
    } or asset.provenance.get("source_type") == "vlm_source_crop"


def _has_split_row_level_bitmap_structure(
    source_preserved_assets: list[Any],
    *,
    largest_ratio: float,
    structure_count: int,
) -> bool:
    split_row_structure = (
        len(source_preserved_assets) >= 3
        and largest_ratio <= 0.35
        and structure_count >= 3
    )
    dense_infographic_structure = (
        len(source_preserved_assets) >= 3
        and largest_ratio <= 0.60
        and structure_count >= 12
    )
    return split_row_structure or dense_infographic_structure


def _source_preserved_asset_effective_area_ratio(
    asset: Any,
    source_size: tuple[int, int],
) -> float:
    visible_ratio = asset.provenance.get("alpha_visible_area_ratio")
    if (
        asset.provenance.get("background_difference_alpha") is True
        and isinstance(visible_ratio, (int, float))
        and 0.0 <= float(visible_ratio) <= 1.0
    ):
        return float(visible_ratio)
    page_area = max(1, source_size[0] * source_size[1])
    return _bbox_area(asset.source_pixel_bbox) / float(page_area)


def _combined_bbox_coverage_ratio(
    bboxes: list[tuple[int, int, int, int]],
    source_size: tuple[int, int],
) -> float:
    width, height = source_size
    if width <= 0 or height <= 0:
        return 0.0
    mask = Image.new("L", (width, height), 0)
    draw = ImageDraw.Draw(mask)
    for bbox in bboxes:
        left, top, right, bottom = _clamp_bbox(bbox, source_size)
        if right > left and bottom > top:
            draw.rectangle((left, top, right, bottom), fill=255)
    histogram = mask.histogram()
    covered = width * height - histogram[0]
    return covered / float(width * height)


def _bbox_area(bbox: tuple[int, int, int, int]) -> int:
    return max(0, int(bbox[2]) - int(bbox[0])) * max(0, int(bbox[3]) - int(bbox[1]))


def _clamp_bbox(
    bbox: tuple[int, int, int, int],
    source_size: tuple[int, int],
) -> tuple[int, int, int, int]:
    width, height = source_size
    left, top, right, bottom = bbox
    return (
        max(0, min(width, int(left))),
        max(0, min(height, int(top))),
        max(0, min(width, int(right))),
        max(0, min(height, int(bottom))),
    )


def _ocr_returned_item_count(page, artifact_root: Path) -> int | None:
    ocr_ref = page.provider_output_paths.get("ocr")
    if not ocr_ref:
        return None
    try:
        path = (artifact_root / ocr_ref).resolve()
        if not path.is_relative_to(artifact_root.resolve()) or not path.is_file():
            return None
        payload = json.loads(path.read_text(encoding="utf-8"))
    except Exception:
        return None
    items = payload.get("items") if isinstance(payload, dict) else None
    return len(items) if isinstance(items, list) else None


def _status_with_reconstruction_issues(status: str, issues: list[dict[str, Any]]) -> str:
    if status == "failed" or _has_error_reconstruction_issue(issues):
        return "failed"
    if _has_warning_reconstruction_issue(issues):
        return "degraded"
    return status


def _status_with_preview_reports(status: str, preview_reports: list[dict[str, Any]]) -> str:
    if status == "failed":
        return "failed"
    if any(str(report.get("status")) == "failed" for report in preview_reports):
        return "failed"
    return status


def _aggregate_page_status(page_reports: list[dict[str, Any]]) -> str:
    statuses = [str(report.get("status", "failed")) for report in page_reports]
    if any(status == "failed" for status in statuses):
        return "failed"
    if any(
        _has_error_reconstruction_issue(report.get("reconstruction_issues", []) or [])
        for report in page_reports
    ):
        return "failed"
    if any(status == "degraded" for status in statuses) or any(
        _has_warning_reconstruction_issue(report.get("reconstruction_issues", []) or [])
        for report in page_reports
    ):
        return "degraded"
    return "passed"


def _has_error_reconstruction_issue(issues: list[dict[str, Any]]) -> bool:
    return any(issue.get("severity", "error") == "error" for issue in issues)


def _has_warning_reconstruction_issue(issues: list[dict[str, Any]]) -> bool:
    return any(issue.get("severity") == "warning" for issue in issues)


def _warning_summary(page_reports: list[dict[str, Any]]) -> dict[str, Any]:
    warning_codes: Counter[str] = Counter()
    warning_pages: set[int] = set()
    warning_count = 0
    for report in page_reports:
        report_page = int(report.get("page", 1))
        for issue in report.get("reconstruction_issues", []) or []:
            if not isinstance(issue, dict) or issue.get("severity") != "warning":
                continue
            page = int(issue.get("page", report_page))
            warning_count += 1
            warning_pages.add(page)
            warning_codes[str(issue.get("code", "unknown_warning"))] += 1
    return {
        "warning_count": warning_count,
        "warning_pages": sorted(warning_pages),
        "warning_codes": dict(sorted(warning_codes.items())),
    }


def _write_preview_reports(
    *,
    source_images: list[Path],
    pptx_path: Path,
    deck_manifest_path: Path,
    artifact_root: Path,
    similarity_threshold: float,
    output_dir: Path,
) -> list[dict[str, Any]]:
    output_dir.mkdir(parents=True, exist_ok=True)
    deck = read_deck_manifest(deck_manifest_path)
    gates = quality_threshold_to_preview_gates(similarity_threshold)
    reports = []
    for page_index, page_ref in enumerate(deck.page_manifest_paths):
        page = read_page_manifest(artifact_root / page_ref)
        source_path = source_images[page_index]
        with Image.open(source_path) as source:
            output_size = source.size
        preview = render_powerpoint_preview_with_metadata(
            page,
            artifact_root,
            pptx_path=pptx_path,
            output_size=output_size,
        )
        preview_path = output_dir / f"{page_index + 1:02d}-preview.png"
        preview.image.save(preview_path)
        diff_path = output_dir / f"{page_index + 1:02d}-diff.png"
        diff_metrics = _write_diff_image(source_path, preview_path, diff_path)
        validation = validate_preview_similarity(
            source_image_path=source_path,
            preview=preview,
            slide_id=page.slide_id,
            page_index=page_index,
            max_mean_abs_delta=gates.max_mean_abs_delta,
            max_changed_pixel_ratio=gates.max_changed_pixel_ratio,
            require_powerpoint_render=True,
        )
        reports.append(
            {
                "page": page_index + 1,
                "slide_id": page.slide_id,
                "status": validation.status,
                "preview_path": str(preview_path),
                "diff_path": str(diff_path),
                "renderer": preview.metadata.get("renderer"),
                "is_powerpoint_render": bool(preview.metadata.get("is_powerpoint_render")),
                "diff_metrics": diff_metrics,
                "issues": [issue.to_dict() for issue in validation.issues],
            }
        )
    return reports


def _write_diff_image(source_path: Path, preview_path: Path, diff_path: Path) -> dict[str, float]:
    with Image.open(source_path) as source_image, Image.open(preview_path) as preview_image:
        source = source_image.convert("RGB")
        preview = preview_image.convert("RGB")
        if preview.size != source.size:
            preview = preview.resize(source.size)
        diff = ImageChops.difference(source, preview)
        diff.save(diff_path)
        mean_abs_delta = sum(ImageStat.Stat(diff).mean) / 3.0
        changed = diff.convert("L").point(lambda value: 255 if value > 16 else 0)
        changed_ratio = sum(ImageStat.Stat(changed).sum) / 255.0 / (source.width * source.height)
        return {
            "mean_abs_delta": round(mean_abs_delta, 4),
            "changed_pixel_ratio": round(changed_ratio, 6),
        }


def _shape_type_name(shape) -> str:
    raw = str(shape.shape_type)
    return raw.split(" ", 1)[0]


def _write_tiny_edit_probe_image(path: Path) -> None:
    image = Image.new("RGB", (160, 90), "white")
    from PIL import ImageDraw

    draw = ImageDraw.Draw(image)
    draw.rectangle((10, 10, 150, 80), fill="#eef2ff")
    draw.text((44, 35), "TEST", fill="black")
    image.save(path)


def _redacted_config_summary(config) -> dict[str, Any]:
    return {
        "ocr": _provider_summary(config.ocr),
        "clean_base_model": _provider_summary(config.clean_base_model),
        "asset_sheet_model": _provider_summary(config.asset_sheet_model),
        "repair_model": _provider_summary(config.repair_model),
        "generation_model": _provider_summary(config.generation_model),
    }


def _provider_summary(provider) -> dict[str, str]:
    return {
        "role": provider.role,
        "provider": provider.provider or "openai_chat",
        "model": provider.model,
        "base_url": "SET" if provider.base_url else "EMPTY",
        "api_key": "SET" if provider.api_key else "EMPTY",
    }


def _write_json(path: Path, payload: dict[str, Any]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(payload, ensure_ascii=False, indent=2, sort_keys=True), encoding="utf-8")


def _default_output_dir() -> str:
    return f"/private/tmp/aippt-real-editable-run-{datetime.now().strftime('%Y%m%d-%H%M%S')}"


def _natural_key(value: str) -> list[Any]:
    return [int(part) if part.isdigit() else part.lower() for part in re.split(r"(\d+)", value)]


if __name__ == "__main__":
    raise SystemExit(main())
