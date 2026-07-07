"""Deterministic PPTX composition from generative editable manifests."""

from __future__ import annotations

from dataclasses import dataclass
import tempfile
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

from .generative_editable_manifest import (
    DeckManifest,
    NativeShapeSpec,
    PageManifest,
    TextBoxSpec,
    read_deck_manifest,
    read_page_manifest,
)


@dataclass(frozen=True)
class SlideRect:
    left: int
    top: int
    width: int
    height: int


class CompositionArtifactError(RuntimeError):
    pass


class CompositionGeometryError(RuntimeError):
    pass


def slide_rect_from_source_pixels(
    bbox: tuple[int, int, int, int],
    *,
    source_image_size: tuple[int, int],
    aspect_ratio: str,
) -> SlideRect:
    slide_width, slide_height = slide_dimensions(aspect_ratio)
    source_width, source_height = source_image_size
    left, top, right, bottom = bbox
    return SlideRect(
        left=round(left / source_width * slide_width),
        top=round(top / source_height * slide_height),
        width=round((right - left) / source_width * slide_width),
        height=round((bottom - top) / source_height * slide_height),
    )


def slide_dimensions(aspect_ratio: str) -> tuple[int, int]:
    if aspect_ratio == "4:3":
        return Inches(10), Inches(7.5)
    return Inches(10), Inches(5.625)


def compose_deck_from_manifests(
    deck_manifest_path: str | Path,
    artifact_root: str | Path,
    output_path: str | Path,
) -> None:
    root = Path(artifact_root)
    deck = read_deck_manifest(deck_manifest_path)
    prs = Presentation()
    width, height = slide_dimensions(deck.aspect_ratio)
    prs.slide_width = width
    prs.slide_height = height
    blank_layout = prs.slide_layouts[6]
    while len(prs.slides) > 0:
        r_id = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[0]

    for page_manifest_ref in deck.page_manifest_paths:
        page = read_page_manifest(root / page_manifest_ref)
        _preflight_page_artifacts(page, root)
        _validate_page_geometry(page, deck.aspect_ratio)
        slide = prs.slides.add_slide(blank_layout)
        _compose_page(slide, prs, page, deck.aspect_ratio, root)

    prs.save(output_path)
    _normalize_pptx_zip(output_path)


def _compose_page(slide, prs: Presentation, page: PageManifest, aspect_ratio: str, root: Path) -> None:
    if page.chosen_background:
        _set_slide_picture_background(slide, _artifact_path(root, page.chosen_background))
    for shape in page.native_shapes:
        _add_native_shape(slide, shape, page, aspect_ratio)
    for asset in sorted(page.bitmap_assets, key=lambda item: item.z_order):
        rect = slide_rect_from_source_pixels(
            asset.source_pixel_bbox,
            source_image_size=page.source_image_size,
            aspect_ratio=aspect_ratio,
        )
        slide.shapes.add_picture(
            str(_artifact_path(root, asset.asset_path)),
            rect.left,
            rect.top,
            width=rect.width,
            height=rect.height,
        )
    for text_box in page.text_boxes:
        _add_text_box(slide, text_box, page, aspect_ratio)


def _set_slide_picture_background(slide, image_path: Path) -> None:
    image_part = slide.part._package.get_or_add_image_part(str(image_path))
    relationship_id = slide.part.relate_to(image_part, RT.IMAGE)
    common_slide_data = slide._element.cSld
    existing_background = common_slide_data.find(qn("p:bg"))
    if existing_background is not None:
        common_slide_data.remove(existing_background)

    background = OxmlElement("p:bg")
    background_properties = OxmlElement("p:bgPr")
    blip_fill = OxmlElement("a:blipFill")
    blip = OxmlElement("a:blip")
    blip.set(qn("r:embed"), relationship_id)
    stretch = OxmlElement("a:stretch")
    fill_rect = OxmlElement("a:fillRect")
    stretch.append(fill_rect)
    blip_fill.append(blip)
    blip_fill.append(stretch)
    background_properties.append(blip_fill)
    background_properties.append(OxmlElement("a:effectLst"))
    background.append(background_properties)
    common_slide_data.insert(0, background)


def _add_native_shape(slide, shape: NativeShapeSpec, page: PageManifest, aspect_ratio: str) -> None:
    rect = slide_rect_from_source_pixels(
        shape.source_pixel_bbox,
        source_image_size=page.source_image_size,
        aspect_ratio=aspect_ratio,
    )
    if shape.shape_type == "line":
        start = shape.line_start or (shape.source_pixel_bbox[0], shape.source_pixel_bbox[1])
        end = shape.line_end or (shape.source_pixel_bbox[2], shape.source_pixel_bbox[3])
        start_rect = slide_rect_from_source_pixels(
            (start[0], start[1], start[0] + 1, start[1] + 1),
            source_image_size=page.source_image_size,
            aspect_ratio=aspect_ratio,
        )
        end_rect = slide_rect_from_source_pixels(
            (end[0], end[1], end[0] + 1, end[1] + 1),
            source_image_size=page.source_image_size,
            aspect_ratio=aspect_ratio,
        )
        line = slide.shapes.add_connector(
            1,
            start_rect.left,
            start_rect.top,
            end_rect.left,
            end_rect.top,
        )
        if shape.line_color:
            line.line.color.rgb = _rgb(shape.line_color)
            _set_line_alpha(line, shape.opacity)
        else:
            line.line.fill.background()
        if shape.stroke_width:
            line.line.width = _stroke_width_to_emu(shape.stroke_width, page, aspect_ratio)
        return

    ppt_shape = MSO_SHAPE.RECTANGLE
    if shape.shape_type == "rounded_rectangle":
        ppt_shape = MSO_SHAPE.ROUNDED_RECTANGLE
    elif shape.shape_type == "ellipse":
        ppt_shape = MSO_SHAPE.OVAL
    native = slide.shapes.add_shape(ppt_shape, rect.left, rect.top, rect.width, rect.height)
    if shape.fill_color:
        native.fill.solid()
        native.fill.fore_color.rgb = _rgb(shape.fill_color)
        _set_fill_alpha(native, shape.opacity)
    else:
        native.fill.background()
    if shape.line_color:
        native.line.color.rgb = _rgb(shape.line_color)
        _set_line_alpha(native, shape.opacity)
    else:
        native.line.fill.background()
    if shape.stroke_width:
        native.line.width = _stroke_width_to_emu(shape.stroke_width, page, aspect_ratio)


def _add_text_box(slide, text_box: TextBoxSpec, page: PageManifest, aspect_ratio: str) -> None:
    rect = slide_rect_from_source_pixels(
        text_box.source_pixel_bbox,
        source_image_size=page.source_image_size,
        aspect_ratio=aspect_ratio,
    )
    rect = _expand_approximate_text_rect(rect, text_box, aspect_ratio)
    box = slide.shapes.add_textbox(rect.left, rect.top, rect.width, rect.height)
    frame = box.text_frame
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.word_wrap = False
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.alignment = _alignment(text_box.alignment)
    run = paragraph.add_run()
    run.text = text_box.text
    if text_box.font_family:
        run.font.name = text_box.font_family
    if text_box.style_hints.get("bold") is True:
        run.font.bold = True
    if text_box.font_size:
        run.font.size = Pt(text_box.font_size)
    if text_box.color_hex:
        run.font.color.rgb = _rgb(text_box.color_hex)
        _set_text_alpha(run, text_box.opacity)


def _expand_approximate_text_rect(
    rect: SlideRect,
    text_box: TextBoxSpec,
    aspect_ratio: str,
) -> SlideRect:
    if not text_box.style_hints.get("approximate_layout"):
        return rect
    font_size = float(text_box.font_size or 18)
    estimated_width = round(_text_visual_units(text_box.text) * Pt(font_size) * 3.2)
    if estimated_width <= rect.width:
        return rect
    slide_width, _slide_height = slide_dimensions(aspect_ratio)
    width = min(estimated_width, max(rect.width, slide_width - rect.left))
    return SlideRect(left=rect.left, top=rect.top, width=width, height=rect.height)


def _text_visual_units(text: str) -> float:
    units = 0.0
    for char in text.strip():
        if char.isspace():
            continue
        if "\u4e00" <= char <= "\u9fff":
            units += 1.0
        elif char.isascii():
            units += 0.58
        else:
            units += 0.8
    return max(1.0, units)


def _alignment(value: str):
    return {
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }.get(value, PP_ALIGN.LEFT)


def _artifact_path(root: Path, ref: str) -> Path:
    path = (root / ref).resolve()
    if not path.is_relative_to(root.resolve()):
        raise ValueError("artifact path must be inside artifact_root")
    return path


def _preflight_page_artifacts(page: PageManifest, root: Path) -> None:
    if page.chosen_background:
        _require_artifact(root, page.chosen_background, "background", page.slide_id)
    for asset in page.bitmap_assets:
        _require_artifact(root, asset.asset_path, f"bitmap asset {asset.asset_id}", page.slide_id)


def _require_artifact(root: Path, ref: str, label: str, slide_id: str) -> None:
    path = _artifact_path(root, ref)
    if not path.exists():
        raise CompositionArtifactError(f"missing {label} for slide {slide_id}: {ref}")


def _validate_page_geometry(page: PageManifest, aspect_ratio: str) -> None:
    expected_width, expected_height = _slide_size_inches(aspect_ratio)
    actual_width, actual_height = page.slide_size
    expected_ratio = expected_width / expected_height
    slide_ratio = actual_width / actual_height
    source_ratio = page.source_image_size[0] / page.source_image_size[1]
    if not _ratio_close(slide_ratio, expected_ratio) or not _ratio_close(source_ratio, expected_ratio):
        raise CompositionGeometryError(
            f"page {page.slide_id} geometry does not match deck aspect ratio {aspect_ratio}"
        )


def _slide_size_inches(aspect_ratio: str) -> tuple[float, float]:
    if aspect_ratio == "4:3":
        return 10.0, 7.5
    return 10.0, 5.625


def _ratio_close(actual: float, expected: float, tolerance: float = 0.01) -> bool:
    return abs(actual - expected) <= tolerance


def _stroke_width_to_emu(stroke_width_px: float, page: PageManifest, aspect_ratio: str) -> int:
    slide_width, slide_height = slide_dimensions(aspect_ratio)
    source_width, source_height = page.source_image_size
    px_to_emu = min(slide_width / source_width, slide_height / source_height)
    return round(stroke_width_px * px_to_emu)


def _set_fill_alpha(shape, opacity: float) -> None:
    _append_alpha(shape._element.spPr.find(qn("a:solidFill")), opacity)


def _set_line_alpha(shape, opacity: float) -> None:
    line = shape._element.spPr.find(qn("a:ln"))
    if line is not None:
        _append_alpha(line.find(qn("a:solidFill")), opacity)


def _set_text_alpha(run, opacity: float) -> None:
    if opacity >= 1:
        return
    r_pr = run._r.get_or_add_rPr()
    _append_alpha(r_pr.find(qn("a:solidFill")), opacity)


def _append_alpha(solid_fill, opacity: float) -> None:
    if solid_fill is None or opacity >= 1:
        return
    color = solid_fill.find(qn("a:srgbClr"))
    if color is None:
        color = solid_fill.find(qn("a:schemeClr"))
    if color is None:
        return
    for existing in color.findall(qn("a:alpha")):
        color.remove(existing)
    alpha = OxmlElement("a:alpha")
    alpha.set("val", str(round(opacity * 100000)))
    color.append(alpha)


def _normalize_pptx_zip(path: str | Path) -> None:
    pptx_path = Path(path)
    fixed_date = (1980, 1, 1, 0, 0, 0)
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
        tmp_path = Path(tmp.name)
    try:
        with zipfile.ZipFile(pptx_path, "r") as source, zipfile.ZipFile(
            tmp_path, "w", compression=zipfile.ZIP_DEFLATED
        ) as target:
            for name in sorted(source.namelist()):
                info = zipfile.ZipInfo(name, fixed_date)
                original = source.getinfo(name)
                info.compress_type = zipfile.ZIP_DEFLATED
                info.external_attr = original.external_attr
                target.writestr(info, source.read(name))
        tmp_path.replace(pptx_path)
    finally:
        if tmp_path.exists():
            tmp_path.unlink()


def _rgb(value: str) -> RGBColor:
    cleaned = value.lstrip("#")
    return RGBColor(int(cleaned[0:2], 16), int(cleaned[2:4], 16), int(cleaned[4:6], 16))
