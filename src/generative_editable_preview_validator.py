"""Preview rendering and validation helpers for generative editable PPTX."""

from __future__ import annotations

from dataclasses import asdict, dataclass
import hashlib
import math
import posixpath
from pathlib import Path
import shutil
import subprocess
import tempfile
from typing import NamedTuple
import zipfile
from xml.etree import ElementTree as ET

from PIL import Image, ImageChops, ImageDraw, ImageStat
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .generative_editable_composer import slide_dimensions
from .generative_editable_manifest import PageManifest, read_deck_manifest, read_page_manifest


VALIDATION_PASSED = "passed"
VALIDATION_FAILED = "failed"
SLIDE_POSITION_TOLERANCE_EMU = 24_000
SOURCE_BACKGROUND_MEAN_DELTA_THRESHOLD = 3.0
SOURCE_BACKGROUND_CHANGED_RATIO_THRESHOLD = 0.02
MAX_PREVIEW_PIXELS = 8_000_000
POWERPOINT_RENDER_TIMEOUT_SECONDS = 30
FORBIDDEN_BITMAP_ASSET_PROVENANCE_TERMS = {
    "source_crop",
    "source-crop",
    "source crop",
    "source crops",
    "source_crops",
    "direct_source",
    "direct-source",
    "direct source",
    "direct crop",
    "source snippet",
    "cropped source",
    "source-cropped",
}


class PreviewGateThresholds(NamedTuple):
    max_mean_abs_delta: float
    max_changed_pixel_ratio: float


@dataclass(frozen=True)
class PreviewRenderResult:
    image: Image.Image
    metadata: dict[str, object]


@dataclass(frozen=True)
class ValidationIssue:
    code: str
    message: str
    slide_id: str = ""
    severity: str = "error"
    details: dict[str, object] | None = None

    def to_dict(self) -> dict[str, object]:
        payload = _json_safe(asdict(self))
        if self.details is None:
            payload.pop("details")
        return payload


@dataclass(frozen=True)
class ValidationReport:
    status: str
    checked_pages: int
    issues: list[ValidationIssue]

    def to_dict(self) -> dict[str, object]:
        return {
            "status": self.status,
            "checked_pages": self.checked_pages,
            "issues": [issue.to_dict() for issue in self.issues],
        }


def validate_composed_deck_structure(
    *,
    deck_manifest_path: str | Path,
    artifact_root: str | Path,
    pptx_path: str | Path,
) -> ValidationReport:
    root = Path(artifact_root).resolve()
    issues: list[ValidationIssue] = []
    try:
        deck = read_deck_manifest(deck_manifest_path)
    except Exception as exc:
        return ValidationReport(
            status=VALIDATION_FAILED,
            checked_pages=0,
            issues=[
                ValidationIssue(
                    code="invalid_deck_manifest",
                    message="deck manifest cannot be read for validation",
                    details={"error": str(exc)},
                )
            ],
        )

    try:
        presentation = Presentation(str(pptx_path))
    except Exception as exc:
        presentation = None
        issues.append(
            ValidationIssue(
                code="invalid_pptx",
                message="composed PPTX cannot be opened for structural validation",
                details={"pptx_path": str(pptx_path), "error": str(exc)},
            )
        )

    if presentation is not None and len(presentation.slides) != len(deck.page_manifest_paths):
        issues.append(
            ValidationIssue(
                code="page_count_mismatch",
                message="composed PPTX page count does not match deck manifest",
                details={
                    "expected": len(deck.page_manifest_paths),
                    "actual": len(presentation.slides),
                },
            )
        )

    expected_width, expected_height = slide_dimensions(deck.aspect_ratio)
    if (
        presentation is not None
        and (presentation.slide_width != expected_width or presentation.slide_height != expected_height)
    ):
        issues.append(
            ValidationIssue(
                code="slide_dimensions_mismatch",
                message="composed PPTX slide dimensions do not match requested aspect ratio",
                details={
                    "expected": [expected_width, expected_height],
                    "actual": [presentation.slide_width, presentation.slide_height],
                },
            )
        )

    for page_index, page_ref in enumerate(deck.page_manifest_paths):
        try:
            page_path = _artifact_path(root, page_ref)
        except ValueError as exc:
            issues.append(
                ValidationIssue(
                    code="invalid_page_manifest_path",
                    message="page manifest path is outside artifact root",
                    details={"page_index": page_index, "page_ref": page_ref, "error": str(exc)},
                )
            )
            continue
        if not page_path.exists():
            issues.append(
                ValidationIssue(
                    code="missing_page_manifest",
                    message="page manifest referenced by deck is missing",
                    details={"page_index": page_index, "page_ref": page_ref},
                )
            )
            continue
        try:
            page = read_page_manifest(page_path)
        except Exception as exc:
            issues.append(
                ValidationIssue(
                    code="invalid_page_manifest",
                    message="page manifest cannot be read for validation",
                    details={"page_index": page_index, "page_ref": page_ref, "error": str(exc)},
                )
            )
            continue
        _validate_page_assets(page, root, issues)
        _validate_slide_background_identity(page, root, pptx_path, page_index, issues)
        source_background_used = _uses_full_slide_source_background(page, root)
        editable_object_count = len(page.text_boxes) + len(page.native_shapes) + len(page.bitmap_assets)
        if source_background_used and editable_object_count == 0:
            issues.append(
                ValidationIssue(
                    code="full_slide_source_background_only",
                    message="slide background is the full source image and no editable or decomposed objects were produced",
                    slide_id=page.slide_id,
                    details={
                        "page_index": page.page_index,
                        "target_kind": "background",
                        "background_ref": page.chosen_background,
                        "editable_object_count": editable_object_count,
                    },
                )
            )
        if (
            _has_non_empty_text(page)
            and (
                source_background_used
                or (
                    presentation is not None
                    and page_index < len(presentation.slides)
                    and _slide_contains_full_slide_source_picture(
                        page,
                        presentation.slides[page_index],
                        root,
                        presentation.slide_width,
                        presentation.slide_height,
                    )
                )
            )
        ):
            issues.append(
                ValidationIssue(
                    code="unsafe_full_slide_source_with_text",
                    message="full-slide source image would leave baked text underneath editable text",
                    slide_id=page.slide_id,
                    details={
                        "page_index": page.page_index,
                        "target_kind": "background",
                        "background_ref": page.chosen_background,
                    },
                )
            )
        if presentation is None or page_index >= len(presentation.slides):
            continue
        slide = presentation.slides[page_index]
        _validate_required_text(page, slide, issues, deck.aspect_ratio)
        _validate_object_order(
            page,
            slide,
            root,
            issues,
            deck.aspect_ratio,
            presentation.slide_width,
            presentation.slide_height,
        )

    return ValidationReport(
        status=VALIDATION_FAILED if issues else VALIDATION_PASSED,
        checked_pages=len(deck.page_manifest_paths),
        issues=issues,
    )


def validate_preview_similarity(
    *,
    source_image_path: str | Path,
    preview: PreviewRenderResult,
    slide_id: str,
    page_index: int,
    max_mean_abs_delta: float,
    max_changed_pixel_ratio: float,
    require_powerpoint_render: bool = True,
    changed_pixel_delta_threshold: int = 16,
) -> ValidationReport:
    renderer = str(preview.metadata.get("renderer", ""))
    if require_powerpoint_render and (
        preview.metadata.get("is_powerpoint_render") is not True or renderer == "manifest_stub"
    ):
        return ValidationReport(
            status=VALIDATION_FAILED,
            checked_pages=1,
            issues=[
                ValidationIssue(
                    code="preview_renderer_not_powerpoint",
                    message="preview similarity gate requires a real PowerPoint render, not a manifest stub",
                    slide_id=slide_id,
                    details={
                        "page_index": page_index,
                        "target_kind": "preview",
                        "renderer": renderer,
                        "is_powerpoint_render": bool(preview.metadata.get("is_powerpoint_render")),
                        "repairable": False,
                    },
                )
            ],
        )
    try:
        with Image.open(source_image_path) as source:
            source_rgb = source.convert("RGB")
    except Exception as exc:
        return ValidationReport(
            status=VALIDATION_FAILED,
            checked_pages=1,
            issues=[
                ValidationIssue(
                    code="missing_source_preview_reference",
                    message="source slide image cannot be opened for preview comparison",
                    slide_id=slide_id,
                    details={
                        "page_index": page_index,
                        "target_kind": "source",
                        "source_image_path": str(source_image_path),
                        "error": str(exc),
                        "repairable": False,
                    },
                )
            ],
        )
    if source_rgb.width * source_rgb.height > MAX_PREVIEW_PIXELS:
        return _preview_too_large_report(slide_id, page_index)
    try:
        preview_rgb = preview.image.convert("RGB")
    except Exception as exc:
        return ValidationReport(
            status=VALIDATION_FAILED,
            checked_pages=1,
            issues=[
                ValidationIssue(
                    code="invalid_preview_image",
                    message="preview image cannot be opened for similarity comparison",
                    slide_id=slide_id,
                    details={
                        "page_index": page_index,
                        "target_kind": "preview",
                        "renderer": renderer,
                        "error": str(exc),
                        "repairable": True,
                    },
                )
            ],
        )
    if preview_rgb.width * preview_rgb.height > MAX_PREVIEW_PIXELS:
        return _preview_too_large_report(slide_id, page_index)
    if preview_rgb.size != source_rgb.size:
        if not _same_aspect_ratio(preview_rgb.size, source_rgb.size):
            return ValidationReport(
                status=VALIDATION_FAILED,
                checked_pages=1,
                issues=[
                    ValidationIssue(
                        code="preview_dimensions_mismatch",
                        message="preview aspect ratio does not match source slide",
                        slide_id=slide_id,
                        details={
                            "page_index": page_index,
                            "target_kind": "preview",
                            "source_size": list(source_rgb.size),
                            "preview_size": list(preview_rgb.size),
                            "repairable": True,
                        },
                    )
                ],
            )
        preview_rgb = preview_rgb.resize(source_rgb.size)
    mean_abs_delta, changed_pixel_ratio = _image_delta_metrics(
        source_rgb,
        preview_rgb,
        changed_pixel_delta_threshold=changed_pixel_delta_threshold,
    )
    details = {
        "page_index": page_index,
        "target_kind": "preview",
        "mean_abs_delta": mean_abs_delta,
        "max_mean_abs_delta": max_mean_abs_delta,
        "changed_pixel_ratio": changed_pixel_ratio,
        "max_changed_pixel_ratio": max_changed_pixel_ratio,
        "changed_pixel_delta_threshold": changed_pixel_delta_threshold,
        "renderer": renderer,
        "repairable": True,
        "fallback_candidate": "text_editable_background",
    }
    if mean_abs_delta > max_mean_abs_delta or changed_pixel_ratio > max_changed_pixel_ratio:
        return ValidationReport(
            status=VALIDATION_FAILED,
            checked_pages=1,
            issues=[
                ValidationIssue(
                    code="preview_similarity_failed",
                    message="rendered preview differs from source slide beyond configured thresholds",
                    slide_id=slide_id,
                    details=details,
                )
            ],
        )
    return ValidationReport(status=VALIDATION_PASSED, checked_pages=1, issues=[])


def quality_threshold_to_preview_gates(similarity_threshold: float) -> PreviewGateThresholds:
    similarity = min(1.0, max(0.0, float(similarity_threshold)))
    allowed_difference = 1.0 - similarity
    return PreviewGateThresholds(
        max_mean_abs_delta=max(1.0, round(255.0 * allowed_difference, 3)),
        max_changed_pixel_ratio=max(0.005, round(allowed_difference * 1.65, 4)),
    )


def render_manifest_preview_with_metadata(
    page_manifest: PageManifest,
    artifact_root: str | Path,
    *,
    pptx_path: str | Path | None = None,
    output_size: tuple[int, int] | None = None,
    **_: object,
) -> PreviewRenderResult:
    return render_powerpoint_preview_with_metadata(
        page_manifest,
        artifact_root,
        pptx_path=pptx_path,
        output_size=output_size,
    )


def render_powerpoint_preview_with_metadata(
    page_manifest: PageManifest,
    artifact_root: str | Path,
    *,
    pptx_path: str | Path | None = None,
    output_size: tuple[int, int] | None = None,
) -> PreviewRenderResult:
    if pptx_path is None:
        return _manifest_stub_preview(page_manifest, artifact_root, output_size=output_size)
    soffice = shutil.which("soffice")
    pdftoppm = shutil.which("pdftoppm")
    if not soffice or not pdftoppm:
        return _manifest_stub_preview(page_manifest, artifact_root, output_size=output_size)
    with tempfile.TemporaryDirectory() as tmp:
        tmp_path = Path(tmp)
        user_profile = tmp_path / "lo-profile"
        pptx = Path(pptx_path).resolve()
        try:
            subprocess.run(
                [
                    soffice,
                    "--headless",
                    f"-env:UserInstallation=file://{user_profile}",
                    "--convert-to",
                    "pdf",
                    "--outdir",
                    str(tmp_path),
                    str(pptx),
                ],
                check=True,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                timeout=POWERPOINT_RENDER_TIMEOUT_SECONDS,
            )
        except (OSError, subprocess.CalledProcessError, subprocess.TimeoutExpired):
            return _manifest_stub_preview(page_manifest, artifact_root, output_size=output_size)
        pdf_path = tmp_path / f"{pptx.stem}.pdf"
        if not pdf_path.exists():
            return _manifest_stub_preview(page_manifest, artifact_root, output_size=output_size)
        output_prefix = tmp_path / "preview"
        try:
            subprocess.run(
                [
                    pdftoppm,
                    "-f",
                    str(page_manifest.page_index + 1),
                    "-l",
                    str(page_manifest.page_index + 1),
                    "-singlefile",
                    "-png",
                    str(pdf_path),
                    str(output_prefix),
                ],
                check=True,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                timeout=POWERPOINT_RENDER_TIMEOUT_SECONDS,
            )
        except (OSError, subprocess.CalledProcessError, subprocess.TimeoutExpired):
            return _manifest_stub_preview(page_manifest, artifact_root, output_size=output_size)
        preview_path = output_prefix.with_suffix(".png")
        if not preview_path.exists():
            return _manifest_stub_preview(page_manifest, artifact_root, output_size=output_size)
        with Image.open(preview_path) as image:
            preview = image.convert("RGB")
            if output_size is not None:
                preview = preview.resize(output_size)
            else:
                preview = preview.copy()
    return PreviewRenderResult(
        image=preview,
        metadata={
            "renderer": "soffice_pdf_pdftoppm",
            "is_powerpoint_render": True,
        },
    )


def _manifest_stub_preview(
    page_manifest: PageManifest,
    artifact_root: str | Path,
    *,
    output_size: tuple[int, int] | None = None,
) -> PreviewRenderResult:
    return PreviewRenderResult(
        image=render_manifest_preview(page_manifest, artifact_root, output_size=output_size),
        metadata={
            "renderer": "manifest_stub",
            "is_powerpoint_render": False,
        },
    )


def render_manifest_preview(
    page_manifest: PageManifest,
    artifact_root: str | Path,
    *,
    output_size: tuple[int, int] | None = None,
) -> Image.Image:
    root = Path(artifact_root)
    source_size = page_manifest.source_image_size
    size = output_size or source_size
    if page_manifest.chosen_background:
        with Image.open(_artifact_path(root, page_manifest.chosen_background)) as background:
            preview = background.convert("RGB").resize(size)
    else:
        preview = Image.new("RGB", size, "white")
    draw = ImageDraw.Draw(preview)
    for shape in page_manifest.native_shapes:
        bbox = _scale_bbox(shape.source_pixel_bbox, source_size, size)
        fill = shape.fill_color or shape.line_color or "#000000"
        if shape.shape_type == "ellipse":
            draw.ellipse(bbox, fill=fill)
        elif shape.shape_type == "line":
            start = shape.line_start or (shape.source_pixel_bbox[0], shape.source_pixel_bbox[1])
            end = shape.line_end or (shape.source_pixel_bbox[2], shape.source_pixel_bbox[3])
            draw.line(
                [_scale_point(start, source_size, size), _scale_point(end, source_size, size)],
                fill=shape.line_color or fill,
                width=max(1, round(shape.stroke_width or 1)),
            )
        else:
            draw.rectangle(bbox, fill=fill)
    for asset in sorted(page_manifest.bitmap_assets, key=lambda item: item.z_order):
        with Image.open(_artifact_path(root, asset.asset_path)) as image:
            rect = _scale_bbox(asset.source_pixel_bbox, source_size, size)
            resized = image.convert("RGBA").resize((rect[2] - rect[0], rect[3] - rect[1]))
            preview.paste(resized, rect[:2], resized)
    for text_box in page_manifest.text_boxes:
        bbox = _scale_bbox(text_box.source_pixel_bbox, source_size, size)
        draw.text((bbox[0], bbox[1]), text_box.text, fill=text_box.color_hex or "#000000")
    return preview


def _artifact_path(root: Path, ref: str) -> Path:
    if not ref:
        raise ValueError("artifact path must be non-empty")
    path = (root / ref).resolve()
    if not path.is_relative_to(root.resolve()):
        raise ValueError("artifact path must be inside artifact_root")
    return path


def _preview_too_large_report(slide_id: str, page_index: int) -> ValidationReport:
    return ValidationReport(
        status=VALIDATION_FAILED,
        checked_pages=1,
        issues=[
            ValidationIssue(
                code="preview_too_large",
                message="preview comparison image exceeds validation pixel limit",
                slide_id=slide_id,
                details={"page_index": page_index, "target_kind": "preview", "repairable": False},
            )
        ],
    )


def _validate_page_assets(page: PageManifest, root: Path, issues: list[ValidationIssue]) -> None:
    if page.chosen_background:
        _validate_image_asset(
            root=root,
            ref=page.chosen_background,
            slide_id=page.slide_id,
            issues=issues,
            details={
                "page_index": page.page_index,
                "target_kind": "background",
                "asset_ref": page.chosen_background,
                "repairable": True,
                "fallback_candidate": "text_editable_background",
            },
        )
    _validate_image_asset(
        root=root,
        ref=page.source_image_path,
        slide_id=page.slide_id,
        issues=issues,
        details={
            "page_index": page.page_index,
            "target_kind": "source",
            "asset_ref": page.source_image_path,
            "repairable": False,
        },
    )
    for asset in page.bitmap_assets:
        _validate_image_asset(
            root=root,
            ref=asset.asset_path,
            slide_id=page.slide_id,
            issues=issues,
            details={
                "page_index": page.page_index,
                "target_kind": "bitmap_asset",
                "target_id": asset.asset_id,
                "asset_ref": asset.asset_path,
                "repairable": True,
            },
        )
        _validate_bitmap_asset_provenance(page, asset, issues)


def _validate_bitmap_asset_provenance(
    page: PageManifest,
    asset,
    issues: list[ValidationIssue],
) -> None:
    provenance_text = _compact_text(asset.provenance)
    asset_path_text = _compact_text(asset.asset_path)
    if not (
        any(term in provenance_text for term in FORBIDDEN_BITMAP_ASSET_PROVENANCE_TERMS)
        or ".source-crop." in asset_path_text
        or asset_path_text.endswith(".source-crop.png")
    ):
        return
    issues.append(
        ValidationIssue(
            code="forbidden_source_crop_bitmap_asset",
            message="foreground bitmap asset uses source-crop or direct-source fallback instead of asset-sheet separation",
            slide_id=page.slide_id,
            details={
                "page_index": page.page_index,
                "target_kind": "bitmap_asset",
                "target_id": asset.asset_id,
                "asset_ref": asset.asset_path,
                "provenance": asset.provenance,
                "repairable": True,
            },
        )
    )


def _compact_text(value: object) -> str:
    if value is None:
        return ""
    if isinstance(value, str):
        return value.lower()
    if isinstance(value, (int, float, bool)):
        return str(value).lower()
    if isinstance(value, dict):
        return " ".join(_compact_text(item) for item in value.values())
    if isinstance(value, (list, tuple, set)):
        return " ".join(_compact_text(item) for item in value)
    return str(value).lower()


def _validate_image_asset(
    *,
    root: Path,
    ref: str,
    slide_id: str,
    issues: list[ValidationIssue],
    details: dict[str, object],
    required: bool = True,
) -> None:
    try:
        path = _artifact_path(root, ref)
    except ValueError as exc:
        issues.append(
            ValidationIssue(
                code="invalid_asset_path",
                message="asset path is outside artifact root",
                slide_id=slide_id,
                details={**details, "error": str(exc)},
            )
        )
        return
    if not path.is_file():
        if required:
            issues.append(
                ValidationIssue(
                    code="missing_asset",
                    message="required image asset is missing",
                    slide_id=slide_id,
                    details=details,
                )
            )
        return
    try:
        with Image.open(path) as image:
            image.verify()
    except Exception as exc:
        issues.append(
            ValidationIssue(
                code="invalid_asset",
                message="image asset cannot be opened",
                slide_id=slide_id,
                details={**details, "error": str(exc)},
            )
        )


def _uses_full_slide_source_background(page: PageManifest, root: Path) -> bool:
    value = str(page.provenance.get("chosen_background_kind", "")).lower()
    if value in {"source_full_slide", "full_slide_source", "uncleaned_source"}:
        return True
    if not page.chosen_background or not page.source_image_path:
        return False
    try:
        background_path = _artifact_path(root, page.chosen_background)
        source_path = _artifact_path(root, page.source_image_path)
    except ValueError:
        return False
    if not background_path.is_file() or not source_path.is_file():
        return False
    try:
        with Image.open(background_path) as background, Image.open(source_path) as source:
            bg = background.convert("RGB")
            src = source.convert("RGB")
            if src.size != page.source_image_size:
                return False
            if bg.size != src.size:
                bg = bg.resize(src.size)
            if _chosen_background_is_text_clean(page) and _text_mask_region_changed(page, root, src, bg):
                return False
            mean_delta, changed_ratio = _image_delta_metrics(
                src,
                bg,
                changed_pixel_delta_threshold=16,
            )
            return (
                mean_delta <= SOURCE_BACKGROUND_MEAN_DELTA_THRESHOLD
                and changed_ratio <= SOURCE_BACKGROUND_CHANGED_RATIO_THRESHOLD
            )
    except Exception:
        return False


def _chosen_background_is_text_clean(page: PageManifest) -> bool:
    if str(page.provenance.get("chosen_background_kind", "")).lower() == "source_preserving_text_clean":
        return True
    if page.chosen_background and page.chosen_background == page.text_clean_background:
        return True
    backgrounds = page.provenance.get("backgrounds")
    if not isinstance(backgrounds, dict):
        return False
    for key in ("base_clean", "text_clean", "source_preserving"):
        record = backgrounds.get(key)
        if not isinstance(record, dict):
            continue
        if record.get("output_asset_ref") and record.get("output_asset_ref") != page.chosen_background:
            continue
        if record.get("provider_role") == "local" and record.get("prompt_id") in {
            "local_text_cleanup",
            "source_preserving_text_background",
        }:
            return True
    return False


def _text_mask_region_changed(
    page: PageManifest,
    root: Path,
    source: Image.Image,
    background: Image.Image,
    *,
    mean_delta_threshold: float = 3.0,
    changed_ratio_threshold: float = 0.08,
    min_changed_pixels: int = 20,
) -> bool:
    text_mask_ref = page.provenance.get("text_mask_path")
    if not isinstance(text_mask_ref, str) or not text_mask_ref:
        return False
    try:
        mask_path = _artifact_path(root, text_mask_ref)
    except ValueError:
        return False
    if not mask_path.is_file():
        return False
    with Image.open(mask_path) as mask_image:
        mask = mask_image.convert("L").resize(source.size)
    diff = ImageChops.difference(source, background).convert("L")
    mask_pixels = changed_pixels = total_delta = 0
    for y in range(diff.height):
        for x in range(diff.width):
            if mask.getpixel((x, y)) <= 0:
                continue
            mask_pixels += 1
            delta = diff.getpixel((x, y))
            total_delta += delta
            if delta > 16:
                changed_pixels += 1
    if mask_pixels <= 0:
        return False
    return (
        changed_pixels >= min_changed_pixels
        or total_delta / float(mask_pixels) > mean_delta_threshold
        or changed_pixels / float(mask_pixels) > changed_ratio_threshold
    )


def _has_non_empty_text(page: PageManifest) -> bool:
    return any(text_box.text.strip() for text_box in page.text_boxes)


def _allows_low_opacity_source_text_overlay(page: PageManifest) -> bool:
    # Kept for reading legacy manifests only. A full-slide source image plus
    # editable OCR text is not an editable reconstruction, even at low opacity.
    return False


def _slide_contains_full_slide_source_picture(
    page: PageManifest,
    slide,
    root: Path,
    slide_width: int,
    slide_height: int,
) -> bool:
    try:
        source_path = _artifact_path(root, page.source_image_path)
    except ValueError:
        return False
    source_sha = _image_sha1(source_path)
    if not source_sha:
        return False
    for shape in slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue
        if getattr(shape.image, "sha1", "") != source_sha:
            continue
        if (
            abs(shape.left) <= SLIDE_POSITION_TOLERANCE_EMU
            and abs(shape.top) <= SLIDE_POSITION_TOLERANCE_EMU
            and abs(shape.width - slide_width) <= SLIDE_POSITION_TOLERANCE_EMU
            and abs(shape.height - slide_height) <= SLIDE_POSITION_TOLERANCE_EMU
        ):
            return True
    return False


def _validate_required_text(page: PageManifest, slide, issues: list[ValidationIssue], aspect_ratio: str) -> None:
    actual_counts: dict[str, int] = {}
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            normalized = _normalize_text(shape.text)
            if normalized:
                actual_counts[normalized] = actual_counts.get(normalized, 0) + 1
    seen_expected: dict[str, int] = {}
    for index, text_box in enumerate(page.text_boxes):
        normalized = _normalize_text(text_box.text)
        if not normalized:
            continue
        seen_expected[normalized] = seen_expected.get(normalized, 0) + 1
        matching_shapes = [
            shape
            for shape in slide.shapes
            if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX and _normalize_text(shape.text) == normalized
        ]
        if len(matching_shapes) < seen_expected[normalized]:
            issues.append(
                ValidationIssue(
                    code="missing_required_text",
                    message="required text from manifest is absent from composed PPTX",
                    slide_id=page.slide_id,
                    details={
                        "page_index": page.page_index,
                        "target_kind": "text",
                        "target_index": index,
                        "text": text_box.text,
                        "source_pixel_bbox": list(text_box.source_pixel_bbox),
                    },
                )
            )
            continue
        matched_shape = matching_shapes[seen_expected[normalized] - 1]
        expected_rect = _slide_rect_for_bbox(page, text_box.source_pixel_bbox, aspect_ratio)
        actual_rect = (matched_shape.left, matched_shape.top, matched_shape.width, matched_shape.height)
        if not _text_rect_close(actual_rect, expected_rect, text_box):
            issues.append(
                ValidationIssue(
                    code="text_position_mismatch",
                    message="required text is present but not positioned at the manifest location",
                    slide_id=page.slide_id,
                    details={
                        "page_index": page.page_index,
                        "target_kind": "text",
                        "target_index": index,
                        "text": text_box.text,
                        "expected_rect": list(expected_rect),
                        "actual_rect": list(actual_rect),
                    },
                )
            )


def _validate_object_order(
    page: PageManifest,
    slide,
    root: Path,
    issues: list[ValidationIssue],
    aspect_ratio: str,
    slide_width: int,
    slide_height: int,
) -> None:
    expected: list[str] = []
    expected.extend(["native_shape"] * len(page.native_shapes))
    expected.extend(["picture"] * len(page.bitmap_assets))
    expected.extend(["text"] * len(page.text_boxes))
    actual = [_shape_role(shape) for shape in slide.shapes]
    actual_relevant = [role for role in actual if role in {"picture", "native_shape", "text"}]
    if actual_relevant != expected:
        issues.append(
            ValidationIssue(
                code="object_order_mismatch",
                message="composed PPTX objects are not ordered as background, native shapes, bitmap assets, text",
                slide_id=page.slide_id,
                details={
                    "page_index": page.page_index,
                    "target_kind": "slide_object_order",
                    "expected": expected,
                    "actual": actual_relevant,
                    "repairable": False,
                },
            )
        )
        return
    _validate_object_identity(page, slide, root, issues, aspect_ratio, slide_width, slide_height)


def _validate_object_identity(
    page: PageManifest,
    slide,
    root: Path,
    issues: list[ValidationIssue],
    aspect_ratio: str,
    slide_width: int,
    slide_height: int,
) -> None:
    relevant_shapes = [shape for shape in slide.shapes if _shape_role(shape) in {"picture", "native_shape", "text"}]
    expected_index = 0
    for shape_index, shape_spec in enumerate(page.native_shapes):
        actual_shape = relevant_shapes[expected_index]
        if shape_spec.shape_type == "line":
            if actual_shape.shape_type != MSO_SHAPE_TYPE.LINE or not _line_close(page, shape_spec, actual_shape, aspect_ratio):
                _append_identity_issue(page, issues, "native_shape", str(shape_index), expected_index)
        elif actual_shape.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
            _append_identity_issue(page, issues, "native_shape", str(shape_index), expected_index)
        else:
            expected_rect = _slide_rect_for_bbox(page, shape_spec.source_pixel_bbox, aspect_ratio)
            actual_rect = (actual_shape.left, actual_shape.top, actual_shape.width, actual_shape.height)
            if not _native_auto_shape_type_matches(shape_spec, actual_shape) or not _rect_close(actual_rect, expected_rect):
                _append_identity_issue(page, issues, "native_shape", str(shape_index), expected_index)
        expected_index += 1
    for asset in sorted(page.bitmap_assets, key=lambda item: item.z_order):
        _validate_picture_identity(
            page=page,
            root=root,
            root_shape=relevant_shapes[expected_index],
            expected_ref=asset.asset_path,
            target_kind="bitmap_asset",
            target_id=asset.asset_id,
            issues=issues,
            index=expected_index,
            expected_bbox=asset.source_pixel_bbox,
            aspect_ratio=aspect_ratio,
        )
        expected_index += 1


def _validate_slide_background_identity(
    page: PageManifest,
    root: Path,
    pptx_path: str | Path,
    page_index: int,
    issues: list[ValidationIssue],
) -> None:
    if not page.chosen_background:
        return
    try:
        expected_path = _artifact_path(root, page.chosen_background)
    except Exception:
        expected_path = None
    expected_sha = _image_sha1(expected_path)
    if not expected_sha:
        return

    slide_part = f"ppt/slides/slide{page_index + 1}.xml"
    rels_part = f"ppt/slides/_rels/slide{page_index + 1}.xml.rels"
    try:
        with zipfile.ZipFile(pptx_path) as pptx:
            slide_xml = pptx.read(slide_part)
            rels_xml = pptx.read(rels_part)
            embed_id = _slide_background_embed_id(slide_xml)
            if not embed_id:
                _append_identity_issue(
                    page,
                    issues,
                    "background",
                    "background",
                    -1,
                    expected_sha=expected_sha,
                    actual_sha="",
                )
                return
            media_part = _relationship_target_part(rels_xml, embed_id, slide_part)
            actual_sha = hashlib.sha1(pptx.read(media_part)).hexdigest() if media_part else ""
    except Exception:
        actual_sha = ""
    if actual_sha != expected_sha:
        _append_identity_issue(
            page,
            issues,
            "background",
            "background",
            -1,
            expected_sha=expected_sha,
            actual_sha=actual_sha,
        )


def _slide_background_embed_id(slide_xml: bytes) -> str:
    namespaces = {
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    }
    root = ET.fromstring(slide_xml)
    blip = root.find(".//p:bg//a:blip", namespaces)
    if blip is None:
        return ""
    return blip.attrib.get(f"{{{namespaces['r']}}}embed", "")


def _relationship_target_part(rels_xml: bytes, relationship_id: str, source_part: str) -> str:
    namespace = {"rel": "http://schemas.openxmlformats.org/package/2006/relationships"}
    root = ET.fromstring(rels_xml)
    for relationship in root.findall("rel:Relationship", namespace):
        if relationship.attrib.get("Id") != relationship_id:
            continue
        target = relationship.attrib.get("Target", "")
        if not target:
            return ""
        if target.startswith("/"):
            return target.lstrip("/")
        return posixpath.normpath(posixpath.join(posixpath.dirname(source_part), target))
    return ""


def _validate_picture_identity(
    *,
    page: PageManifest,
    root: Path,
    root_shape,
    expected_ref: str,
    target_kind: str,
    target_id: str,
    issues: list[ValidationIssue],
    index: int,
    expected_bbox: tuple[int, int, int, int] | None = None,
    expected_rect: tuple[int, int, int, int] | None = None,
    aspect_ratio: str | None = None,
) -> None:
    if root_shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
        _append_identity_issue(page, issues, target_kind, target_id, index)
        return
    try:
        expected_path = _artifact_path(root, expected_ref)
    except Exception:
        expected_path = None
    expected_sha = _image_sha1(expected_path) if expected_path is not None else ""
    actual_sha = getattr(root_shape.image, "sha1", "")
    if expected_sha and actual_sha != expected_sha:
        _append_identity_issue(page, issues, target_kind, target_id, index, expected_sha=expected_sha, actual_sha=actual_sha)
        return
    if expected_rect is None and expected_bbox is not None:
        expected_rect = _slide_rect_for_bbox(page, expected_bbox, aspect_ratio)
    if expected_rect is not None:
        actual_rect = (root_shape.left, root_shape.top, root_shape.width, root_shape.height)
        if not _rect_close(actual_rect, expected_rect):
            _append_identity_issue(page, issues, target_kind, target_id, index)


def _append_identity_issue(
    page: PageManifest,
    issues: list[ValidationIssue],
    target_kind: str,
    target_id: str,
    index: int,
    *,
    expected_sha: str = "",
    actual_sha: str = "",
) -> None:
    issues.append(
        ValidationIssue(
            code="object_identity_mismatch",
            message="composed PPTX object identity does not match the manifest object at this position",
            slide_id=page.slide_id,
            details={
                "page_index": page.page_index,
                "target_kind": target_kind,
                "target_id": target_id,
                "object_index": index,
                "expected_sha1": expected_sha,
                "actual_sha1": actual_sha,
                "repairable": False,
            },
        )
    )


def _shape_role(shape) -> str:
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return "picture"
    if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
        return "text"
    if shape.shape_type in {MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.LINE}:
        return "native_shape"
    return "other"


def _normalize_text(value: str) -> str:
    return " ".join(value.split()).strip()


def _slide_rect_for_bbox(
    page: PageManifest,
    bbox: tuple[int, int, int, int],
    aspect_ratio: str | None = None,
) -> tuple[int, int, int, int]:
    from .generative_editable_composer import slide_rect_from_source_pixels

    slide_ratio = aspect_ratio or _aspect_ratio_from_page(page)
    rect = slide_rect_from_source_pixels(
        bbox,
        source_image_size=page.source_image_size,
        aspect_ratio=slide_ratio,
    )
    return rect.left, rect.top, rect.width, rect.height


def _aspect_ratio_from_page(page: PageManifest) -> str:
    width, height = page.slide_size
    return "4:3" if abs((width / height) - (4 / 3)) < abs((width / height) - (16 / 9)) else "16:9"


def _rect_close(actual: tuple[int, int, int, int], expected: tuple[int, int, int, int]) -> bool:
    return all(abs(actual[index] - expected[index]) <= SLIDE_POSITION_TOLERANCE_EMU for index in range(4))


def _text_rect_close(actual: tuple[int, int, int, int], expected: tuple[int, int, int, int], text_box) -> bool:
    if not getattr(text_box, "style_hints", {}).get("approximate_layout"):
        return _rect_close(actual, expected)
    left_close = abs(actual[0] - expected[0]) <= SLIDE_POSITION_TOLERANCE_EMU
    top_close = abs(actual[1] - expected[1]) <= SLIDE_POSITION_TOLERANCE_EMU
    height_close = abs(actual[3] - expected[3]) <= SLIDE_POSITION_TOLERANCE_EMU
    width_not_smaller = actual[2] + SLIDE_POSITION_TOLERANCE_EMU >= expected[2]
    return left_close and top_close and height_close and width_not_smaller


def _native_auto_shape_type_matches(shape_spec, actual_shape) -> bool:
    expected = {
        "rectangle": MSO_SHAPE.RECTANGLE,
        "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
        "ellipse": MSO_SHAPE.OVAL,
    }.get(shape_spec.shape_type)
    return expected is not None and getattr(actual_shape, "auto_shape_type", None) == expected


def _line_close(page: PageManifest, shape_spec, actual_shape, aspect_ratio: str | None = None) -> bool:
    start = shape_spec.line_start or (shape_spec.source_pixel_bbox[0], shape_spec.source_pixel_bbox[1])
    end = shape_spec.line_end or (shape_spec.source_pixel_bbox[2], shape_spec.source_pixel_bbox[3])
    expected_start = _slide_point_for_source_pixel(page, start, aspect_ratio)
    expected_end = _slide_point_for_source_pixel(page, end, aspect_ratio)
    actual_start = (actual_shape.left, actual_shape.top)
    actual_end = (actual_shape.left + actual_shape.width, actual_shape.top + actual_shape.height)
    return (
        _point_close(actual_start, expected_start)
        and _point_close(actual_end, expected_end)
    ) or (
        _point_close(actual_start, expected_end)
        and _point_close(actual_end, expected_start)
    ) or _rect_close(
        (actual_shape.left, actual_shape.top, actual_shape.width, actual_shape.height),
        _slide_rect_for_bbox(
            page,
            (
                min(start[0], end[0]),
                min(start[1], end[1]),
                max(start[0], end[0]),
                max(start[1], end[1]),
            ),
            aspect_ratio,
        ),
    )


def _slide_point_for_source_pixel(
    page: PageManifest,
    point: tuple[int, int],
    aspect_ratio: str | None = None,
) -> tuple[int, int]:
    left, top, _, _ = _slide_rect_for_bbox(page, (point[0], point[1], point[0] + 1, point[1] + 1), aspect_ratio)
    return left, top


def _point_close(actual: tuple[int, int], expected: tuple[int, int]) -> bool:
    return all(abs(actual[index] - expected[index]) <= SLIDE_POSITION_TOLERANCE_EMU for index in range(2))


def _same_aspect_ratio(first: tuple[int, int], second: tuple[int, int], tolerance: float = 0.01) -> bool:
    return abs((first[0] / first[1]) - (second[0] / second[1])) <= tolerance


def _image_sha1(path: Path | None) -> str:
    if path is None or not path.exists():
        return ""
    try:
        return hashlib.sha1(path.read_bytes()).hexdigest()
    except Exception:
        return ""


def _image_delta_metrics(
    source: Image.Image,
    preview: Image.Image,
    *,
    changed_pixel_delta_threshold: int,
) -> tuple[float, float]:
    diff = ImageChops.difference(source, preview)
    mean_abs_delta = sum(ImageStat.Stat(diff).mean) / 3.0
    grayscale = diff.convert("L").point(
        lambda value: 255 if value > changed_pixel_delta_threshold else 0,
        mode="L",
    )
    changed_pixel_ratio = sum(ImageStat.Stat(grayscale).sum) / (255.0 * source.width * source.height)
    return mean_abs_delta, changed_pixel_ratio


def _json_safe(value):
    if isinstance(value, dict):
        return {str(key): _json_safe(nested) for key, nested in value.items()}
    if isinstance(value, list):
        return [_json_safe(item) for item in value]
    if isinstance(value, tuple):
        return [_json_safe(item) for item in value]
    if isinstance(value, Path):
        return str(value)
    if isinstance(value, float) and not math.isfinite(value):
        return None
    if isinstance(value, str | int | float | bool) or value is None:
        return value
    return str(value)


def _scale_bbox(
    bbox: tuple[int, int, int, int],
    source_size: tuple[int, int],
    output_size: tuple[int, int],
) -> tuple[int, int, int, int]:
    left, top, right, bottom = bbox
    source_width, source_height = source_size
    output_width, output_height = output_size
    return (
        round(left / source_width * output_width),
        round(top / source_height * output_height),
        round(right / source_width * output_width),
        round(bottom / source_height * output_height),
    )


def _scale_point(
    point: tuple[int, int],
    source_size: tuple[int, int],
    output_size: tuple[int, int],
) -> tuple[int, int]:
    source_width, source_height = source_size
    output_width, output_height = output_size
    return (
        round(point[0] / source_width * output_width),
        round(point[1] / source_height * output_height),
    )
