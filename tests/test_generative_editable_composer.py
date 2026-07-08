import tempfile
import time
import unittest
import zipfile
from dataclasses import replace
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

from src.generative_editable_manifest import (
    BitmapAssetSpec,
    DeckManifest,
    NativeShapeSpec,
    PageManifest,
    TextBoxSpec,
    write_manifest,
)


class GenerativeEditableComposerTest(unittest.TestCase):
    def test_converts_source_pixels_to_slide_coordinates_for_16_9_and_4_3(self):
        from src.generative_editable_composer import slide_rect_from_source_pixels

        rect_16_9 = slide_rect_from_source_pixels(
            (80, 45, 400, 225),
            source_image_size=(800, 450),
            aspect_ratio="16:9",
        )
        rect_4_3 = slide_rect_from_source_pixels(
            (80, 60, 400, 300),
            source_image_size=(800, 600),
            aspect_ratio="4:3",
        )

        self.assertEqual(rect_16_9.left, Inches(1))
        self.assertEqual(rect_16_9.top, Inches(0.5625))
        self.assertEqual(rect_16_9.width, Inches(4))
        self.assertEqual(rect_16_9.height, Inches(2.25))
        self.assertEqual(rect_4_3.left, Inches(1))
        self.assertEqual(rect_4_3.top, Inches(0.75))
        self.assertEqual(rect_4_3.width, Inches(4))
        self.assertEqual(rect_4_3.height, Inches(3))

    def test_composes_page_in_background_shape_asset_text_z_order(self):
        from src.generative_editable_composer import compose_deck_from_manifests

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            self._write_fixture_assets(root)
            page = self._page_manifest(root)
            deck = DeckManifest(
                job_id="job-001",
                slide_order=["slide-a"],
                aspect_ratio="16:9",
                provider_roles={},
                quality_settings={},
                fallback_policy="fail",
                page_manifest_paths=["pages/0000-slide-a.json"],
            )
            page_path = root / "pages" / "0000-slide-a.json"
            write_manifest(page_path, page)
            write_manifest(root / "deck.json", deck)
            output = root / "out.pptx"

            compose_deck_from_manifests(root / "deck.json", root, output)
            prs = Presentation(str(output))
            slide = prs.slides[0]
            shape_types = [shape.shape_type for shape in slide.shapes]

        self.assertEqual(
            shape_types,
            [
                MSO_SHAPE_TYPE.AUTO_SHAPE,
                MSO_SHAPE_TYPE.PICTURE,
                MSO_SHAPE_TYPE.TEXT_BOX,
            ],
        )

    def test_composes_chosen_background_as_slide_background_not_selectable_picture(self):
        from src.generative_editable_composer import compose_deck_from_manifests

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            self._write_fixture_assets(root)
            page = self._page_manifest(root)
            deck = DeckManifest(
                job_id="job-001",
                slide_order=["slide-a"],
                aspect_ratio="16:9",
                provider_roles={},
                quality_settings={},
                fallback_policy="fail",
                page_manifest_paths=["pages/0000-slide-a.json"],
            )
            page_path = root / "pages" / "0000-slide-a.json"
            write_manifest(page_path, page)
            write_manifest(root / "deck.json", deck)
            output = root / "out.pptx"

            compose_deck_from_manifests(root / "deck.json", root, output)
            prs = Presentation(str(output))
            slide = prs.slides[0]
            shape_types = [shape.shape_type for shape in slide.shapes]
            slide_xml = self._slide_xml(output, 1)

        self.assertEqual(
            shape_types,
            [
                MSO_SHAPE_TYPE.AUTO_SHAPE,
                MSO_SHAPE_TYPE.PICTURE,
                MSO_SHAPE_TYPE.TEXT_BOX,
            ],
        )
        self.assertIn("<p:bg>", slide_xml)
        self.assertIn("<a:blipFill>", slide_xml)

    def test_composes_approximate_layout_text_as_visible_editable_text(self):
        from src.generative_editable_composer import compose_deck_from_manifests

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            self._write_fixture_assets(root)
            page = self._page_manifest(
                root,
                title="Structured title",
            )
            page = replace(
                page,
                text_boxes=[
                    TextBoxSpec(
                        text="Approximate OCR text",
                        source_pixel_bbox=(70, 60, 730, 120),
                        source_pixel_polygon=((70, 60), (730, 60), (730, 120), (70, 120)),
                        font_family="Arial",
                        font_size=18,
                        color_hex="#FFFFFF",
                        alignment="left",
                        style_hints={"approximate_layout": True},
                    ),
                    *page.text_boxes,
                ],
            )
            deck = DeckManifest(
                job_id="job-001",
                slide_order=["slide-a"],
                aspect_ratio="16:9",
                provider_roles={},
                quality_settings={},
                fallback_policy="fail",
                page_manifest_paths=["pages/0000-slide-a.json"],
            )
            write_manifest(root / "pages" / "0000-slide-a.json", page)
            write_manifest(root / "deck.json", deck)
            output = root / "out.pptx"

            compose_deck_from_manifests(root / "deck.json", root, output)
            prs = Presentation(str(output))
            slide = prs.slides[0]
            shape_types = [shape.shape_type for shape in slide.shapes]
            texts = [
                shape.text for shape in slide.shapes if shape.has_text_frame and shape.text.strip()
            ]
            slide_xml = self._slide_xml(output, 1)

        self.assertEqual(
            shape_types,
            [
                MSO_SHAPE_TYPE.AUTO_SHAPE,
                MSO_SHAPE_TYPE.PICTURE,
                MSO_SHAPE_TYPE.TEXT_BOX,
                MSO_SHAPE_TYPE.TEXT_BOX,
            ],
        )
        self.assertEqual(texts, ["Approximate OCR text", "Structured title"])
        self.assertNotIn('<a:alpha val="0"/>', slide_xml)
        self.assertIn('wrap="none"', slide_xml)

    def test_composes_low_opacity_text_as_native_text_box_with_alpha(self):
        from src.generative_editable_composer import compose_deck_from_manifests

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            self._write_fixture_assets(root)
            page = replace(
                self._page_manifest(root),
                text_boxes=[
                    TextBoxSpec(
                        text="Editable OCR overlay",
                        source_pixel_bbox=(70, 60, 730, 120),
                        source_pixel_polygon=((70, 60), (730, 60), (730, 120), (70, 120)),
                        font_family="Arial",
                        font_size=18,
                        color_hex="#FFFFFF",
                        opacity=0.1,
                    )
                ],
                native_shapes=[],
                bitmap_assets=[],
            )
            deck = DeckManifest(
                job_id="job-001",
                slide_order=["slide-a"],
                aspect_ratio="16:9",
                provider_roles={},
                quality_settings={},
                fallback_policy="fail",
                page_manifest_paths=["pages/0000-slide-a.json"],
            )
            write_manifest(root / "pages" / "0000-slide-a.json", page)
            write_manifest(root / "deck.json", deck)
            output = root / "out.pptx"

            compose_deck_from_manifests(root / "deck.json", root, output)
            prs = Presentation(str(output))
            slide = prs.slides[0]
            text_shapes = [shape for shape in slide.shapes if shape.has_text_frame]
            slide_xml = self._slide_xml(output, 1)

        self.assertEqual([shape.text for shape in text_shapes], ["Editable OCR overlay"])
        self.assertEqual(text_shapes[0].shape_type, MSO_SHAPE_TYPE.TEXT_BOX)
        self.assertIn('<a:alpha val="10000"/>', slide_xml)

    def test_applies_bold_text_style_hint(self):
        from src.generative_editable_composer import compose_deck_from_manifests

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            self._write_fixture_assets(root)
            page = replace(
                self._page_manifest(root),
                text_boxes=[
                    TextBoxSpec(
                        text="Bold heading",
                        source_pixel_bbox=(70, 60, 730, 120),
                        source_pixel_polygon=((70, 60), (730, 60), (730, 120), (70, 120)),
                        font_family="Arial",
                        font_size=24,
                        color_hex="#FFFFFF",
                        style_hints={"bold": True},
                    )
                ],
                native_shapes=[],
                bitmap_assets=[],
            )
            deck = DeckManifest(
                job_id="job-001",
                slide_order=["slide-a"],
                aspect_ratio="16:9",
                provider_roles={},
                quality_settings={},
                fallback_policy="fail",
                page_manifest_paths=["pages/0000-slide-a.json"],
            )
            write_manifest(root / "pages" / "0000-slide-a.json", page)
            write_manifest(root / "deck.json", deck)
            output = root / "out.pptx"

            compose_deck_from_manifests(root / "deck.json", root, output)
            prs = Presentation(str(output))
            text_shape = next(shape for shape in prs.slides[0].shapes if shape.has_text_frame)

        self.assertTrue(text_shape.text_frame.paragraphs[0].runs[0].font.bold)

    def test_expands_narrow_approximate_text_boxes_to_prevent_wrapping(self):
        from src.generative_editable_composer import compose_deck_from_manifests

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            self._write_fixture_assets(root)
            page = replace(
                self._page_manifest(root),
                text_boxes=[
                    TextBoxSpec(
                        text="EMB",
                        source_pixel_bbox=(100, 80, 112, 96),
                        source_pixel_polygon=((100, 80), (112, 80), (112, 96), (100, 96)),
                        font_size=8,
                        color_hex="#FFFFFF",
                        style_hints={"approximate_layout": True},
                    )
                ],
                native_shapes=[],
                bitmap_assets=[],
            )
            deck = DeckManifest(
                job_id="job-001",
                slide_order=["slide-a"],
                aspect_ratio="16:9",
                provider_roles={},
                quality_settings={},
                fallback_policy="fail",
                page_manifest_paths=["pages/0000-slide-a.json"],
            )
            write_manifest(root / "pages" / "0000-slide-a.json", page)
            write_manifest(root / "deck.json", deck)
            output = root / "out.pptx"

            compose_deck_from_manifests(root / "deck.json", root, output)
            prs = Presentation(str(output))
            shape = next(shape for shape in prs.slides[0].shapes if shape.has_text_frame)
            raw_width = Inches(10) * 12 // 800

        self.assertGreater(shape.width, raw_width * 4)

    def test_text_boxes_and_simple_shapes_are_native_powerpoint_objects(self):
        from src.generative_editable_composer import compose_deck_from_manifests

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            self._write_fixture_assets(root)
            page = self._page_manifest(root)
            deck = DeckManifest(
                job_id="job-001",
                slide_order=["slide-a"],
                aspect_ratio="16:9",
                provider_roles={},
                quality_settings={},
                fallback_policy="fail",
                page_manifest_paths=["pages/0000-slide-a.json"],
            )
            write_manifest(root / "pages" / "0000-slide-a.json", page)
            write_manifest(root / "deck.json", deck)
            output = root / "out.pptx"

            compose_deck_from_manifests(root / "deck.json", root, output)
            prs = Presentation(str(output))
            slide = prs.slides[0]
            texts = [shape.text for shape in slide.shapes if shape.has_text_frame]
            auto_shapes = [
                shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
            ]

        self.assertIn("Quarterly Plan", texts)
        self.assertEqual(len(auto_shapes), 1)
        self.assertEqual(auto_shapes[0].fill.fore_color.rgb.__str__(), "2563EB")

    def test_composes_multi_page_deck_order_dimensions_media_and_rebuild_deterministically(self):
        from src.generative_editable_composer import compose_deck_from_manifests

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            self._write_fixture_assets(
                root, slide_ids=("slide-b", "slide-a"), image_size=(800, 600)
            )
            page_b = self._page_manifest(
                root,
                slide_id="slide-b",
                page_index=0,
                title="Slide B",
                source_image_size=(800, 600),
                slide_size=(10.0, 7.5),
            )
            page_a = self._page_manifest(
                root,
                slide_id="slide-a",
                page_index=1,
                title="Slide A",
                source_image_size=(800, 600),
                slide_size=(10.0, 7.5),
            )
            deck = DeckManifest(
                job_id="job-001",
                slide_order=["slide-b", "slide-a"],
                aspect_ratio="4:3",
                provider_roles={},
                quality_settings={},
                fallback_policy="fail",
                page_manifest_paths=["pages/0000-slide-b.json", "pages/0001-slide-a.json"],
            )
            write_manifest(root / "pages" / "0000-slide-b.json", page_b)
            write_manifest(root / "pages" / "0001-slide-a.json", page_a)
            write_manifest(root / "deck.json", deck)
            output_one = root / "one.pptx"
            output_two = root / "two.pptx"

            compose_deck_from_manifests(root / "deck.json", root, output_one)
            time.sleep(2)
            compose_deck_from_manifests(root / "deck.json", root, output_two)
            prs = Presentation(str(output_one))
            media_names = self._media_names(output_one)
            output_one_bytes = output_one.read_bytes()
            output_two_bytes = output_two.read_bytes()

        self.assertEqual(len(prs.slides), 2)
        self.assertEqual(prs.slide_width, Inches(10))
        self.assertEqual(prs.slide_height, Inches(7.5))
        self.assertEqual(prs.slides[0].shapes[-1].text, "Slide B")
        self.assertEqual(prs.slides[1].shapes[-1].text, "Slide A")
        self.assertTrue(any(name.endswith(".png") for name in media_names))
        self.assertEqual(output_one_bytes, output_two_bytes)

    def test_rejects_page_geometry_that_disagrees_with_deck_aspect_ratio(self):
        from src.generative_editable_composer import (
            CompositionGeometryError,
            compose_deck_from_manifests,
        )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            self._write_fixture_assets(root)
            page = self._page_manifest(root)
            deck = DeckManifest(
                job_id="job-001",
                slide_order=["slide-a"],
                aspect_ratio="4:3",
                provider_roles={},
                quality_settings={},
                fallback_policy="fail",
                page_manifest_paths=["pages/0000-slide-a.json"],
            )
            write_manifest(root / "pages" / "0000-slide-a.json", page)
            write_manifest(root / "deck.json", deck)

            with self.assertRaisesRegex(CompositionGeometryError, "aspect ratio"):
                compose_deck_from_manifests(root / "deck.json", root, root / "out.pptx")

    def test_handles_empty_text_boxes_and_scales_line_stroke_width(self):
        from src.generative_editable_composer import compose_deck_from_manifests

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            self._write_fixture_assets(root)
            page = self._page_manifest(root, title="")
            page = replace(
                page,
                native_shapes=[
                    NativeShapeSpec(
                        shape_type="line",
                        source_pixel_bbox=(100, 100, 180, 104),
                        line_start=(100, 102),
                        line_end=(180, 102),
                        line_color="#111827",
                        stroke_width=4,
                        confidence=0.95,
                    )
                ],
            )
            deck = DeckManifest(
                job_id="job-001",
                slide_order=["slide-a"],
                aspect_ratio="16:9",
                provider_roles={},
                quality_settings={},
                fallback_policy="fail",
                page_manifest_paths=["pages/0000-slide-a.json"],
            )
            write_manifest(root / "pages" / "0000-slide-a.json", page)
            write_manifest(root / "deck.json", deck)
            output = root / "out.pptx"

            compose_deck_from_manifests(root / "deck.json", root, output)
            xml = self._slide_xml(output, 1)
            prs = Presentation(str(output))

        self.assertEqual(prs.slides[0].shapes[-1].text, "")
        self.assertIn('w="45720"', xml)
        self.assertIn("<p:cxnSp>", xml)

    def test_applies_explicit_native_shape_fill_line_and_width_defaults(self):
        from src.generative_editable_composer import compose_deck_from_manifests

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            self._write_fixture_assets(root)
            page = self._page_manifest(root)
            page = replace(
                page,
                native_shapes=[
                    NativeShapeSpec(
                        shape_type="rectangle",
                        source_pixel_bbox=(60, 140, 300, 220),
                        fill_color="",
                        line_color="",
                        confidence=0.95,
                    ),
                    NativeShapeSpec(
                        shape_type="ellipse",
                        source_pixel_bbox=(340, 140, 460, 220),
                        fill_color="#22C55E",
                        line_color="#14532D",
                        stroke_width=4,
                        opacity=0.5,
                        confidence=0.96,
                    ),
                ],
            )
            deck = DeckManifest(
                job_id="job-001",
                slide_order=["slide-a"],
                aspect_ratio="16:9",
                provider_roles={},
                quality_settings={},
                fallback_policy="fail",
                page_manifest_paths=["pages/0000-slide-a.json"],
            )
            write_manifest(root / "pages" / "0000-slide-a.json", page)
            write_manifest(root / "deck.json", deck)
            output = root / "out.pptx"

            compose_deck_from_manifests(root / "deck.json", root, output)
            xml = self._slide_xml(output, 1)

        self.assertIn("<a:noFill/>", xml)
        self.assertIn('w="45720"', xml)
        self.assertIn('val="50000"', xml)

    def test_missing_artifacts_raise_structured_composition_error(self):
        from src.generative_editable_composer import (
            CompositionArtifactError,
            compose_deck_from_manifests,
        )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            page = self._page_manifest(root)
            deck = DeckManifest(
                job_id="job-001",
                slide_order=["slide-a"],
                aspect_ratio="16:9",
                provider_roles={},
                quality_settings={},
                fallback_policy="fail",
                page_manifest_paths=["pages/0000-slide-a.json"],
            )
            write_manifest(root / "pages" / "0000-slide-a.json", page)
            write_manifest(root / "deck.json", deck)

            with self.assertRaisesRegex(CompositionArtifactError, "background"):
                compose_deck_from_manifests(root / "deck.json", root, root / "out.pptx")

    def _write_fixture_assets(self, root: Path, slide_ids=("slide-a",), image_size=(800, 450)):
        for index, slide_id in enumerate(slide_ids):
            suffix = f"{index:04d}-{slide_id}"
            (root / "backgrounds" / suffix).mkdir(parents=True, exist_ok=True)
            (root / "assets" / suffix).mkdir(parents=True, exist_ok=True)
            Image.new("RGB", image_size, "#F8FAFC").save(root / "backgrounds" / suffix / "base.png")
            asset = Image.new("RGBA", (120, 80), (0, 0, 0, 0))
            for x in range(20, 100):
                for y in range(18, 62):
                    asset.putpixel((x, y), (16, 185, 129, 255))
            asset.save(root / "assets" / suffix / "asset.png")

    def _page_manifest(
        self,
        root: Path,
        *,
        slide_id="slide-a",
        page_index=0,
        title="Quarterly Plan",
        source_image_size=(800, 450),
        slide_size=(10.0, 5.625),
    ):
        suffix = f"{page_index:04d}-{slide_id}"
        return PageManifest(
            slide_id=slide_id,
            page_index=page_index,
            source_image_path=f"sources/{suffix}/source.png",
            source_image_size=source_image_size,
            slide_size=slide_size,
            chosen_background=f"backgrounds/{suffix}/base.png",
            text_boxes=[
                TextBoxSpec(
                    text=title,
                    source_pixel_bbox=(80, 54, 420, 102),
                    source_pixel_polygon=((80, 54), (420, 54), (420, 102), (80, 102)),
                    font_family="Arial",
                    font_size=28,
                    color_hex="#111827",
                    alignment="center",
                )
            ],
            native_shapes=[
                NativeShapeSpec(
                    shape_type="rectangle",
                    source_pixel_bbox=(60, 140, 300, 220),
                    fill_color="#2563EB",
                    line_color="#1D4ED8",
                    confidence=0.95,
                )
            ],
            bitmap_assets=[
                BitmapAssetSpec(
                    asset_id="asset-1",
                    source_pixel_bbox=(340, 120, 460, 200),
                    asset_path=f"assets/{suffix}/asset.png",
                    z_order=1,
                )
            ],
        )

    def _media_names(self, pptx_path: Path):
        with zipfile.ZipFile(pptx_path) as package:
            return sorted(name for name in package.namelist() if name.startswith("ppt/media/"))

    def _slide_xml(self, pptx_path: Path, slide_number: int):
        with zipfile.ZipFile(pptx_path) as package:
            return package.read(f"ppt/slides/slide{slide_number}.xml").decode("utf-8")


if __name__ == "__main__":
    unittest.main()
