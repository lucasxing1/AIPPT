import asyncio
import base64
import tempfile
import unittest
from pathlib import Path
from types import SimpleNamespace
from unittest.mock import patch

from PIL import Image, ImageDraw
from fastapi import HTTPException

from api.models import ExportRequest, ExportSlide
import api.routes.export as export_route


def _slide_base64() -> str:
    with tempfile.TemporaryDirectory() as tmp:
        image_path = Path(tmp) / "slide.png"
        _write_fake_vlm_source(image_path)
        return base64.b64encode(image_path.read_bytes()).decode()


def _write_fake_vlm_source(path: Path, *, size: tuple[int, int] = (800, 450)) -> None:
    image = Image.new("RGB", size, "white")
    width, height = size
    ImageDraw.Draw(image).rectangle(
        (
            round(width * 0.08),
            round(height * 0.08),
            round(width * 0.48),
            round(height * 0.18),
        ),
        fill="#2563EB",
    )
    image.save(path)


def _fake_route_dependencies(
    *,
    validation_status="passed",
    validation_code="preview_similarity_failed",
    provider_exception=None,
    repair_limit_failure=False,
):
    dependencies = export_route._build_fake_generative_editable_pipeline_dependencies(
        validation_status=validation_status,
        validation_code=validation_code,
    )
    if provider_exception is not None:
        dependencies = _with_raising_ocr_provider(dependencies, provider_exception)
    if repair_limit_failure:
        dependencies = _with_repair_limit_asset_builder(dependencies)
    return dependencies


def _fake_vlm_route_dependencies(
    *,
    validation_status="passed",
    validation_code="preview_similarity_failed",
    provider_exception=None,
):
    from src.generative_editable_config import ProviderConfig
    from src.generative_editable_preview_validator import ValidationIssue, ValidationReport
    from src.generative_editable_providers import FakeImageEditProvider, FakeOCRProvider
    from src.generative_editable_vlm_reconstruction import (
        FakeVLMPageAnalysisProvider,
        VLMEditablePipelineDependencies,
    )

    def preview_validator(**kwargs):
        if validation_status == "failed":
            return ValidationReport(
                status="failed",
                checked_pages=1,
                issues=[
                    ValidationIssue(
                        code=validation_code,
                        message="forced preview failure",
                        slide_id=kwargs["slide_id"],
                    )
                ],
            )
        return ValidationReport(status="passed", checked_pages=1, issues=[])

    provider = ProviderConfig(
        role="edit_model",
        provider="fake_image_edit",
        model="fake-image-edit",
        base_url="fake",
        api_key="fake",
    )
    ocr_provider = FakeOCRProvider(
        ProviderConfig(
            role="ocr_model",
            provider="fake_ocr",
            model="fake-ocr",
            base_url="fake",
            api_key="fake",
        )
    )
    dependencies = VLMEditablePipelineDependencies(
        vlm_provider=FakeVLMPageAnalysisProvider(),
        image_edit_provider=FakeImageEditProvider(provider),
        asset_sheet_image_edit_provider=FakeImageEditProvider(provider),
        ocr_provider=ocr_provider,
        preview_validator=preview_validator,
    )
    if provider_exception is not None:
        dependencies = _with_raising_ocr_provider(dependencies, provider_exception)
    return dependencies


def _legacy_generative_config():
    from dataclasses import replace

    from src.generative_editable_config import load_generative_editable_config

    return replace(load_generative_editable_config(use_fake=True), reconstruction_mode="generative")


def _vlm_first_config():
    from dataclasses import replace

    from src.generative_editable_config import load_generative_editable_config

    return replace(load_generative_editable_config(use_fake=True), reconstruction_mode="vlm_first")


def _with_raising_ocr_provider(dependencies, exception):
    class RaisingOCRProvider:
        config = dependencies.ocr_provider.config

        def extract_text(self, image_path):
            raise exception

    from dataclasses import replace

    return replace(dependencies, ocr_provider=RaisingOCRProvider())


def _with_repair_limit_asset_builder(dependencies):
    from dataclasses import replace

    from src.generative_editable_manifest import BitmapAssetSpec, RepairAttempt
    from src.generative_editable_pipeline import AssetBuildResult

    def asset_builder(**kwargs):
        asset = BitmapAssetSpec(
            asset_id="fg-001",
            source_pixel_bbox=(500, 120, 620, 220),
            asset_path=f"assets/{kwargs['page_index']:04d}-{kwargs['slide_id']}/fg-001.png",
            z_order=1,
        )
        Path(kwargs["asset_root"], asset.asset_path).parent.mkdir(parents=True, exist_ok=True)
        Image.new("RGBA", (120, 100), (0, 0, 0, 0)).save(
            Path(kwargs["asset_root"], asset.asset_path)
        )
        return AssetBuildResult(
            bitmap_assets=[asset],
            repair_attempts=[
                RepairAttempt(
                    target_id="fg-001",
                    attempt_index=1,
                    reason="repair_limit_exceeded:edge_touch",
                    provider_role="edit_model",
                    status="failed",
                )
            ],
        )

    return replace(dependencies, asset_builder=asset_builder)


class GenerativeEditableExportRouteTest(unittest.TestCase):
    def setUp(self):
        self._config_patcher = patch(
            "api.routes.export.load_generative_editable_config",
            return_value=_vlm_first_config(),
        )
        self._config_patcher.start()

    def tearDown(self):
        self._config_patcher.stop()

    def test_export_uses_vlm_first_pipeline_when_configured(self):
        from src.generative_editable_config import (
            GenerativeEditableConfig,
            ProviderConfig,
            QualityConfig,
            RetryConfig,
            TimeoutConfig,
        )

        config = GenerativeEditableConfig(
            reconstruction_mode="vlm_first",
            ocr=ProviderConfig(role="ocr_model", model="ocr", base_url="fake", api_key="fake"),
            clean_base_model=ProviderConfig(
                role="edit_model",
                provider="fake_image_edit",
                model="edit",
                base_url="fake",
                api_key="fake",
            ),
            asset_sheet_model=ProviderConfig(
                role="asset_sheet_edit_model",
                provider="fake_image_edit",
                model="asset-sheet",
                base_url="fake",
                api_key="fake",
            ),
            repair_model=ProviderConfig(
                role="edit_model",
                provider="fake_image_edit",
                model="edit",
                base_url="fake",
                api_key="fake",
            ),
            generation_model=ProviderConfig(
                role="image_model", model="image", base_url="fake", api_key="fake"
            ),
            use_aippt_metadata_first=False,
            ocr_min_confidence=0.75,
            quality=QualityConfig(),
            retries=RetryConfig(),
            timeouts=TimeoutConfig(),
        )
        captured = {}
        dependency_sentinel = object()

        def fake_vlm_pipeline(**kwargs):
            captured.update(kwargs)
            Path(kwargs["output_path"]).write_bytes(b"pptx")
            return type(
                "Result",
                (),
                {
                    "status": "passed",
                    "output_path": kwargs["output_path"],
                    "fallback_policy": "fail",
                    "fallback_used": "",
                    "validation_report": None,
                },
            )()

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "slide.png"
            Image.new("RGB", (800, 450), "#FFFFFF").save(source)
            output = root / "out.pptx"
            with (
                patch("api.routes.export.load_generative_editable_config", return_value=config),
                patch(
                    "api.routes.export.run_vlm_editable_pptx_pipeline",
                    side_effect=fake_vlm_pipeline,
                ),
                patch(
                    "api.routes.export._build_vlm_editable_pipeline_dependencies",
                    return_value=dependency_sentinel,
                ),
            ):
                result = export_route._export_generative_editable_pptx(
                    [str(source)],
                    str(output),
                    aspect_ratio="16:9",
                )

        self.assertEqual(result.status, "passed")
        self.assertIs(captured["dependencies"], dependency_sentinel)
        self.assertEqual(captured["slides"][0].slide_id, "slide-1")

    def test_public_route_maps_vlm_page_timeout_from_worker_thread_to_504(self):
        import time as clock
        from dataclasses import replace

        from src.generative_editable_config import TimeoutConfig, load_generative_editable_config
        from src.generative_editable_providers import FakeImageEditProvider
        from src.generative_editable_vlm_reconstruction import (
            VLMEditablePipelineDependencies,
            VLMPageAnalysisProvider,
            coerce_vlm_analysis,
        )

        class SlowVLMProvider(VLMPageAnalysisProvider):
            def analyze_page(self, image_path: str, *, timeout_seconds: int = 180):
                clock.sleep(0.2)
                return coerce_vlm_analysis(
                    {
                        "coordinate_space": {"width": 160, "height": 90, "unit": "px"},
                        "text_regions": [],
                        "bitmap_regions": [],
                        "shape_regions": [
                            {"id": "line-a", "type": "line", "bbox": [20, 60, 140, 61]}
                        ],
                    }
                )

        config = replace(
            load_generative_editable_config(use_fake=True),
            reconstruction_mode="vlm_first",
            timeouts=TimeoutConfig(provider_call=5, page=0.05),
        )
        dependencies = VLMEditablePipelineDependencies(
            vlm_provider=SlowVLMProvider(),
            image_edit_provider=FakeImageEditProvider(config.clean_base_model),
            page_timeout_seconds=0.05,
        )
        request = ExportRequest(
            slides=[ExportSlide(image_base64=_slide_base64())],
            format="generative_editable_pptx",
        )

        with (
            patch("api.routes.export.load_generative_editable_config", return_value=config),
            patch(
                "api.routes.export._build_vlm_editable_pipeline_dependencies",
                return_value=dependencies,
            ),
        ):
            with self.assertRaises(HTTPException) as ctx:
                asyncio.run(export_route.export_presentation(request))

        self.assertEqual(ctx.exception.status_code, 504)
        self.assertIn("vlm_first.build_page_manifest", ctx.exception.detail)

    def test_live_dependency_builder_does_not_require_provider_adapter_fields(self):
        from src.generative_editable_config import (
            GenerativeEditableConfig,
            ProviderConfig,
            QualityConfig,
            RetryConfig,
            TimeoutConfig,
        )

        config = GenerativeEditableConfig(
            ocr=ProviderConfig(
                role="ocr_model",
                model="ocr",
                base_url="https://provider.example/v1",
                api_key="secret",
            ),
            clean_base_model=ProviderConfig(
                role="edit_model",
                model="edit",
                base_url="https://provider.example/v1",
                api_key="secret",
                adapter="multipart_images",
            ),
            asset_sheet_model=ProviderConfig(
                role="edit_model",
                model="edit",
                base_url="https://provider.example/v1",
                api_key="secret",
                adapter="multipart_images",
            ),
            repair_model=ProviderConfig(
                role="edit_model",
                model="edit",
                base_url="https://provider.example/v1",
                api_key="secret",
                adapter="multipart_images",
            ),
            generation_model=ProviderConfig(
                role="image_model",
                model="image",
                base_url="https://provider.example/v1",
                api_key="secret",
                adapter="unexpected",
            ),
            use_aippt_metadata_first=False,
            ocr_min_confidence=0.75,
            quality=QualityConfig(),
            retries=RetryConfig(),
            timeouts=TimeoutConfig(),
        )

        with patch("api.routes.export.load_generative_editable_config", return_value=config):
            dependencies = export_route._build_generative_editable_pipeline_dependencies()

        self.assertEqual(dependencies.ocr_provider.config.role, "ocr_model")
        self.assertEqual(dependencies.image_edit_provider.config.role, "edit_model")
        self.assertEqual(dependencies.image_generation_provider.config.role, "image_model")

    def test_vlm_dependency_builder_wires_vlm_edit_and_ocr_providers(self):
        from src.generative_editable_config import (
            GenerativeEditableConfig,
            ProviderConfig,
            QualityConfig,
            RetryConfig,
            TimeoutConfig,
        )

        config = GenerativeEditableConfig(
            reconstruction_mode="vlm_first",
            ocr=ProviderConfig(
                role="ocr_model",
                provider="fake_ocr",
                model="ocr",
                base_url="fake",
                api_key="fake",
            ),
            clean_base_model=ProviderConfig(
                role="edit_model",
                provider="fake_image_edit",
                model="edit",
                base_url="fake",
                api_key="fake",
            ),
            asset_sheet_model=ProviderConfig(
                role="asset_sheet_edit_model",
                provider="fake_image_edit",
                model="asset-sheet",
                base_url="fake",
                api_key="fake",
            ),
            repair_model=ProviderConfig(
                role="edit_model",
                provider="fake_image_edit",
                model="edit",
                base_url="fake",
                api_key="fake",
            ),
            generation_model=ProviderConfig(
                role="image_model",
                provider="fake_image_generation",
                model="image",
                base_url="fake",
                api_key="fake",
            ),
            use_aippt_metadata_first=False,
            ocr_min_confidence=0.83,
            quality=QualityConfig(
                preview_similarity_threshold=0.86, require_preview_validation=True
            ),
            retries=RetryConfig(provider_max_attempts=4, backoff_seconds=0.5),
            timeouts=TimeoutConfig(provider_call=77, page=500),
        )
        vlm_profile = SimpleNamespace(
            model="vlm",
            base_url="https://vlm.example/v1",
            api_key="vlm-secret",
            adapter="openai_chat",
        )

        with patch(
            "api.routes.export.load_default_profiles", return_value=SimpleNamespace(vlm=vlm_profile)
        ):
            dependencies = export_route._build_vlm_editable_pipeline_dependencies(config)

        self.assertEqual(dependencies.vlm_provider.config.model, "vlm")
        self.assertEqual(dependencies.image_edit_provider.config.role, "edit_model")
        self.assertEqual(
            dependencies.asset_sheet_image_edit_provider.config.role, "asset_sheet_edit_model"
        )
        self.assertEqual(dependencies.ocr_provider.config.role, "ocr_model")
        self.assertEqual(dependencies.ocr_min_confidence, 0.83)
        self.assertEqual(dependencies.provider_timeout_seconds, 77)
        self.assertEqual(dependencies.page_timeout_seconds, 500)
        self.assertEqual(dependencies.provider_max_attempts, 4)
        self.assertEqual(dependencies.provider_retry_backoff_seconds, 0.5)
        self.assertEqual(dependencies.preview_similarity_threshold, 0.86)
        self.assertTrue(dependencies.require_preview_validation)

    def test_live_dependency_builder_wires_distinct_configured_providers(self):
        from src.generative_editable_config import (
            GenerativeEditableConfig,
            ProviderConfig,
            QualityConfig,
            RetryConfig,
            TimeoutConfig,
        )

        config = GenerativeEditableConfig(
            ocr=ProviderConfig(
                role="ocr_model",
                provider="fake_ocr",
                model="ocr",
                base_url="fake",
                api_key="fake",
            ),
            clean_base_model=ProviderConfig(
                role="clean_edit_model",
                provider="fake_image_edit",
                model="clean",
                base_url="fake",
                api_key="fake",
            ),
            asset_sheet_model=ProviderConfig(
                role="asset_sheet_edit_model",
                provider="fake_image_edit",
                model="asset-sheet",
                base_url="fake",
                api_key="fake",
            ),
            repair_model=ProviderConfig(
                role="repair_edit_model",
                provider="fake_image_edit",
                model="repair",
                base_url="fake",
                api_key="fake",
            ),
            generation_model=ProviderConfig(
                role="image_model",
                provider="fake_image_generation",
                model="image",
                base_url="fake",
                api_key="fake",
            ),
            use_aippt_metadata_first=False,
            ocr_min_confidence=0.88,
            quality=QualityConfig(max_repair_attempts=3, preview_similarity_threshold=0.87),
            retries=RetryConfig(provider_max_attempts=2, repair_max_attempts=4, backoff_seconds=0),
            timeouts=TimeoutConfig(provider_call=77, page=500),
        )

        with patch("api.routes.export.load_generative_editable_config", return_value=config):
            dependencies = export_route._build_generative_editable_pipeline_dependencies()

        self.assertEqual(dependencies.image_edit_provider.config.role, "clean_edit_model")
        self.assertEqual(
            dependencies.asset_sheet_image_edit_provider.config.role, "asset_sheet_edit_model"
        )
        self.assertEqual(dependencies.repair_image_edit_provider.config.role, "repair_edit_model")
        self.assertEqual(dependencies.preview_similarity_threshold, 0.87)
        self.assertEqual(dependencies.provider_timeout_seconds, 77)
        self.assertEqual(dependencies.provider_max_attempts, 2)
        self.assertEqual(dependencies.provider_retry_backoff_seconds, 0)
        self.assertEqual(dependencies.max_repair_attempts, 3)
        self.assertFalse(dependencies.use_aippt_metadata_first)
        self.assertEqual(dependencies.ocr_min_confidence, 0.88)

    def test_public_route_maps_missing_provider_config_to_400(self):
        from fastapi import HTTPException

        from api.routes.export import export_presentation
        from src.generative_editable_config import GenerativeEditableConfigError

        async def fake_to_thread(func, *args, **kwargs):
            return func(*args, **kwargs)

        original_to_thread = asyncio.to_thread
        asyncio.to_thread = fake_to_thread
        try:
            with patch(
                "api.routes.export._build_vlm_editable_pipeline_dependencies",
                side_effect=GenerativeEditableConfigError("Missing provider configuration"),
            ):
                with self.assertRaises(HTTPException) as ctx:
                    asyncio.run(
                        export_presentation(
                            ExportRequest(
                                slides=[ExportSlide(image_base64=_slide_base64())],
                                format="generative_editable_pptx",
                            )
                        )
                    )

            self.assertEqual(ctx.exception.status_code, 400)
            self.assertIn("provider configuration", str(ctx.exception.detail))
        finally:
            asyncio.to_thread = original_to_thread

    def test_public_route_maps_provider_timeout_and_failure(self):
        from fastapi import HTTPException

        from api.routes.export import export_presentation
        from src.generative_editable_providers import ProviderError, ProviderTimeoutError

        async def fake_to_thread(func, *args, **kwargs):
            return func(*args, **kwargs)

        timeout = ProviderTimeoutError(
            provider_role="ocr_model",
            operation="extract_text",
            message="timeout calling provider",
            retryable=True,
            timeout_seconds=30,
        )
        provider_error = ProviderError(
            provider_role="edit_model",
            operation="edit",
            message="upstream failed",
            retryable=False,
        )

        original_to_thread = asyncio.to_thread
        asyncio.to_thread = fake_to_thread
        try:
            for raised, expected_status in ((timeout, 504), (provider_error, 502)):
                with patch(
                    "api.routes.export._build_vlm_editable_pipeline_dependencies",
                    lambda config, raised=raised: _fake_vlm_route_dependencies(
                        provider_exception=raised
                    ),
                ):
                    with self.assertRaises(HTTPException) as ctx:
                        asyncio.run(
                            export_presentation(
                                ExportRequest(
                                    slides=[ExportSlide(image_base64=_slide_base64())],
                                    format="generative_editable_pptx",
                                )
                            )
                        )

                self.assertEqual(ctx.exception.status_code, expected_status)
                self.assertIn("导出失败", str(ctx.exception.detail))
        finally:
            asyncio.to_thread = original_to_thread

    def test_public_route_maps_validation_failure_to_422(self):
        from fastapi import HTTPException

        from api.routes.export import export_presentation

        async def fake_to_thread(func, *args, **kwargs):
            return func(*args, **kwargs)

        original_to_thread = asyncio.to_thread
        asyncio.to_thread = fake_to_thread
        try:
            with patch(
                "api.routes.export._build_vlm_editable_pipeline_dependencies",
                lambda config: _fake_vlm_route_dependencies(validation_status="failed"),
            ):
                with self.assertRaises(HTTPException) as ctx:
                    asyncio.run(
                        export_presentation(
                            ExportRequest(
                                slides=[ExportSlide(image_base64=_slide_base64())],
                                format="generative_editable_pptx",
                            )
                        )
                    )

            self.assertEqual(ctx.exception.status_code, 422)
            self.assertIn("preview_similarity_failed", str(ctx.exception.detail))
        finally:
            asyncio.to_thread = original_to_thread

    def test_public_route_maps_real_repair_limit_failure_to_422(self):
        from fastapi import HTTPException

        from api.routes.export import export_presentation

        async def fake_to_thread(func, *args, **kwargs):
            return func(*args, **kwargs)

        original_to_thread = asyncio.to_thread
        asyncio.to_thread = fake_to_thread
        try:
            with patch(
                "api.routes.export._build_vlm_editable_pipeline_dependencies",
                lambda config: _fake_vlm_route_dependencies(
                    validation_status="failed",
                    validation_code="repair_limit_exceeded",
                ),
            ):
                with self.assertRaises(HTTPException) as ctx:
                    asyncio.run(
                        export_presentation(
                            ExportRequest(
                                slides=[ExportSlide(image_base64=_slide_base64())],
                                format="generative_editable_pptx",
                            )
                        )
                    )

            self.assertEqual(ctx.exception.status_code, 422)
            self.assertIn("repair_limit_exceeded", str(ctx.exception.detail))
        finally:
            asyncio.to_thread = original_to_thread

    def test_public_route_returns_explicit_raster_fallback_response_on_validation_failure(self):
        from api.routes.export import export_presentation

        async def fake_to_thread(func, *args, **kwargs):
            return func(*args, **kwargs)

        original_to_thread = asyncio.to_thread
        asyncio.to_thread = fake_to_thread
        try:
            with patch(
                "api.routes.export._build_vlm_editable_pipeline_dependencies",
                lambda config: _fake_vlm_route_dependencies(validation_status="failed"),
            ):
                response = asyncio.run(
                    export_presentation(
                        ExportRequest(
                            slides=[ExportSlide(image_base64=_slide_base64())],
                            format="generative_editable_pptx",
                            editable_options={"fallback_policy": "raster_pptx"},
                        )
                    )
                )

            self.assertEqual(response.headers["X-Generative-Editable-Status"], "fallback_used")
            self.assertEqual(response.headers["X-Generative-Editable-Fallback-Used"], "raster_pptx")
            asyncio.run(response.background())
        finally:
            asyncio.to_thread = original_to_thread

    def test_public_export_route_returns_valid_generative_editable_pptx_with_fake_providers(self):
        from pptx import Presentation

        from api.routes.export import export_presentation

        async def fake_to_thread(func, *args, **kwargs):
            return func(*args, **kwargs)

        original_to_thread = asyncio.to_thread
        asyncio.to_thread = fake_to_thread
        try:
            with patch(
                "api.routes.export._build_vlm_editable_pipeline_dependencies",
                lambda config: _fake_vlm_route_dependencies(),
            ):
                response = asyncio.run(
                    export_presentation(
                        ExportRequest(
                            slides=[
                                ExportSlide(
                                    image_base64=_slide_base64(),
                                    slide_id="slide-a",
                                    text_metadata=[
                                        {
                                            "text": "Quarterly Plan",
                                            "role": "title",
                                            "order": 1,
                                        }
                                    ],
                                )
                            ],
                            format="generative_editable_pptx",
                        )
                    )
                )

            presentation = Presentation(str(response.path))
            self.assertEqual(len(presentation.slides), 1)
            self.assertEqual(response.headers["X-Generative-Editable-Status"], "passed")
            self.assertEqual(response.headers["X-Generative-Editable-Fallback-Policy"], "fail")
            asyncio.run(response.background())
            self.assertFalse(Path(response.path).exists())
        finally:
            asyncio.to_thread = original_to_thread

    def test_generative_editable_route_returns_valid_pptx_with_fake_providers(self):
        from pptx import Presentation

        with tempfile.TemporaryDirectory() as tmp:
            image_path = Path(tmp) / "slide.png"
            output_path = Path(tmp) / "out.pptx"
            _write_fake_vlm_source(image_path)

            with patch(
                "api.routes.export._build_vlm_editable_pipeline_dependencies",
                lambda config: _fake_vlm_route_dependencies(),
            ):
                result = export_route._export_generative_editable_pptx(
                    [str(image_path)],
                    str(output_path),
                    slides=[
                        ExportSlide(
                            image_base64=_slide_base64(),
                            slide_id="slide-a",
                            text_metadata=[
                                {
                                    "text": "Quarterly Plan",
                                    "role": "title",
                                    "order": 1,
                                }
                            ],
                        )
                    ],
                )

            presentation = Presentation(str(output_path))
            self.assertEqual(len(presentation.slides), 1)
            self.assertEqual(result.status, "passed")
            self.assertEqual(result.output_path, str(output_path))

    def test_slide_order_rejects_duplicate_and_unknown_ids(self):
        with tempfile.TemporaryDirectory() as tmp:
            first = Path(tmp) / "first.png"
            second = Path(tmp) / "second.png"
            Image.new("RGB", (800, 450), "white").save(first)
            Image.new("RGB", (800, 450), "white").save(second)

            with self.assertRaisesRegex(ValueError, "slide_id values must be unique"):
                export_route._generative_editable_slide_inputs(
                    image_paths=[str(first), str(second)],
                    slides=[
                        ExportSlide(image_base64=_slide_base64(), slide_id="slide-a"),
                        ExportSlide(image_base64=_slide_base64(), slide_id="slide-a"),
                    ],
                    slide_order=None,
                )

            with self.assertRaisesRegex(ValueError, "slide_order values must be unique"):
                export_route._generative_editable_slide_inputs(
                    image_paths=[str(first), str(second)],
                    slides=[
                        ExportSlide(image_base64=_slide_base64(), slide_id="slide-a"),
                        ExportSlide(image_base64=_slide_base64(), slide_id="slide-b"),
                    ],
                    slide_order=["slide-a", "slide-a"],
                )

            with self.assertRaisesRegex(ValueError, "unknown slide_id"):
                export_route._generative_editable_slide_inputs(
                    image_paths=[str(first), str(second)],
                    slides=[
                        ExportSlide(image_base64=_slide_base64(), slide_id="slide-a"),
                        ExportSlide(image_base64=_slide_base64(), slide_id="slide-b"),
                    ],
                    slide_order=["slide-c"],
                )

    def test_explicit_raster_fallback_uses_existing_raster_pptx_exporter(self):
        calls = []

        def fake_raster_export(image_paths, output_path, aspect_ratio="16:9"):
            calls.append((list(image_paths), output_path, aspect_ratio))
            Path(output_path).write_bytes(b"raster-fallback")

        with tempfile.TemporaryDirectory() as tmp:
            image_path = Path(tmp) / "slide.png"
            output_path = Path(tmp) / "out.pptx"
            _write_fake_vlm_source(image_path, size=(800, 600))

            with (
                patch("api.routes.export._export_pptx", fake_raster_export),
                patch(
                    "api.routes.export._build_vlm_editable_pipeline_dependencies",
                    lambda config: _fake_vlm_route_dependencies(validation_status="failed"),
                ),
            ):
                result = export_route._export_generative_editable_pptx(
                    [str(image_path)],
                    str(output_path),
                    aspect_ratio="4:3",
                    editable_options=ExportRequest(
                        slides=[ExportSlide(image_base64=_slide_base64())],
                        format="generative_editable_pptx",
                        editable_options={"fallback_policy": "raster_pptx"},
                    ).editable_options,
                )

            self.assertEqual(output_path.read_bytes(), b"raster-fallback")

        self.assertEqual(calls[0][2], "4:3")
        self.assertEqual(result.status, "fallback_used")
        self.assertEqual(result.fallback_used, "raster_pptx")

    def test_explicit_raster_fallback_handles_vlm_provider_failure(self):
        from src.generative_editable_providers import ProviderError

        calls = []
        provider_error = ProviderError(
            provider_role="ocr_model",
            operation="extract_text",
            message="upstream failed",
            retryable=False,
        )

        def fake_raster_export(image_paths, output_path, aspect_ratio="16:9"):
            calls.append((list(image_paths), output_path, aspect_ratio))
            Path(output_path).write_bytes(b"raster-fallback")

        with tempfile.TemporaryDirectory() as tmp:
            image_path = Path(tmp) / "slide.png"
            output_path = Path(tmp) / "out.pptx"
            _write_fake_vlm_source(image_path)

            with (
                patch("api.routes.export._export_pptx", fake_raster_export),
                patch(
                    "api.routes.export._build_vlm_editable_pipeline_dependencies",
                    lambda config: _fake_vlm_route_dependencies(provider_exception=provider_error),
                ),
            ):
                result = export_route._export_generative_editable_pptx(
                    [str(image_path)],
                    str(output_path),
                    aspect_ratio="16:9",
                    editable_options=ExportRequest(
                        slides=[ExportSlide(image_base64=_slide_base64())],
                        format="generative_editable_pptx",
                        editable_options={"fallback_policy": "raster_pptx"},
                    ).editable_options,
                )

            self.assertEqual(output_path.read_bytes(), b"raster-fallback")

        self.assertEqual(calls[0][2], "16:9")
        self.assertEqual(result.status, "fallback_used")
        self.assertEqual(result.fallback_used, "raster_pptx")

    def test_text_editable_background_fallback_returns_text_editable_background_deck(self):
        from pptx import Presentation

        with tempfile.TemporaryDirectory() as tmp:
            image_path = Path(tmp) / "slide.png"
            output_path = Path(tmp) / "out.pptx"
            _write_fake_vlm_source(image_path)

            with patch(
                "api.routes.export._build_vlm_editable_pipeline_dependencies",
                lambda config: _fake_vlm_route_dependencies(validation_status="failed"),
            ):
                result = export_route._export_generative_editable_pptx(
                    [str(image_path)],
                    str(output_path),
                    editable_options={
                        "fallback_policy": "text_editable_background",
                    },
                )

            presentation = Presentation(str(output_path))
            self.assertEqual(len(presentation.slides), 1)
            self.assertEqual(result.status, "fallback_used")
            self.assertEqual(result.fallback_used, "text_editable_background")
            self.assertTrue(any(shape.has_text_frame for shape in presentation.slides[0].shapes))

    def test_text_editable_background_fallback_requires_deck_manifest(self):
        from src.generative_editable_pipeline import GenerativeEditableFallbackError

        with tempfile.TemporaryDirectory() as tmp:
            with self.assertRaises(GenerativeEditableFallbackError):
                export_route._export_text_editable_background_fallback(
                    artifact_root=Path(tmp),
                    job_id="export",
                    output_path=str(Path(tmp) / "out.pptx"),
                )

    def test_text_editable_background_fallback_reports_missing_background_artifact(self):
        from src.generative_editable_manifest import DeckManifest, PageManifest, write_manifest
        from src.generative_editable_pipeline import GenerativeEditableFallbackError

        with tempfile.TemporaryDirectory() as tmp:
            artifact_root = Path(tmp)
            job_dir = artifact_root / "export"
            page_path = job_dir / "pages" / "0000-slide-1.json"
            write_manifest(
                page_path,
                PageManifest(
                    slide_id="slide-1",
                    page_index=0,
                    source_image_path="sources/0000-slide-1.png",
                    source_image_size=(800, 450),
                    slide_size=(13.333, 7.5),
                ),
            )
            write_manifest(
                job_dir / "deck.json",
                DeckManifest(
                    job_id="export",
                    slide_order=["slide-1"],
                    aspect_ratio="16:9",
                    provider_roles={},
                    quality_settings={},
                    fallback_policy="fail",
                    page_manifest_paths=["pages/0000-slide-1.json"],
                ),
            )

            with self.assertRaises(GenerativeEditableFallbackError):
                export_route._export_text_editable_background_fallback(
                    artifact_root=artifact_root,
                    job_id="export",
                    output_path=str(artifact_root / "out.pptx"),
                )


if __name__ == "__main__":
    unittest.main()
