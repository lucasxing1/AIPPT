import tempfile
import unittest
import shutil
from dataclasses import replace
from pathlib import Path
from unittest.mock import patch

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Inches

from src.generative_editable_manifest import (
    BitmapAssetSpec,
    DeckManifest,
    NativeShapeSpec,
    PageManifest,
    TextBoxSpec,
    write_manifest,
)


class GenerativeEditablePreviewValidatorTest(unittest.TestCase):
    def test_renders_deterministic_manifest_preview_stub_from_fake_assets(self):
        from src.generative_editable_preview_validator import (
            render_manifest_preview,
            render_manifest_preview_with_metadata,
        )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            (root / "backgrounds" / "0000-slide-a").mkdir(parents=True)
            (root / "assets" / "0000-slide-a").mkdir(parents=True)
            Image.new("RGB", (800, 450), "#FFFFFF").save(
                root / "backgrounds" / "0000-slide-a" / "base.png"
            )
            Image.new("RGBA", (80, 60), (16, 185, 129, 255)).save(
                root / "assets" / "0000-slide-a" / "asset.png"
            )
            page = PageManifest(
                slide_id="slide-a",
                page_index=0,
                source_image_path="sources/0000-slide-a/source.png",
                source_image_size=(800, 450),
                slide_size=(10.0, 5.625),
                chosen_background="backgrounds/0000-slide-a/base.png",
                native_shapes=[
                    NativeShapeSpec(
                        shape_type="rectangle",
                        source_pixel_bbox=(80, 80, 240, 140),
                        fill_color="#2563EB",
                    ),
                    NativeShapeSpec(
                        shape_type="line",
                        source_pixel_bbox=(280, 80, 380, 84),
                        line_color="#111827",
                        line_start=(280, 82),
                        line_end=(380, 82),
                        stroke_width=4,
                    )
                ],
                bitmap_assets=[
                    BitmapAssetSpec(
                        asset_id="asset",
                        source_pixel_bbox=(300, 90, 380, 150),
                        asset_path="assets/0000-slide-a/asset.png",
                        z_order=1,
                    )
                ],
                text_boxes=[
                    TextBoxSpec(
                        text="Preview",
                        source_pixel_bbox=(100, 210, 260, 250),
                        source_pixel_polygon=((100, 210), (260, 210), (260, 250), (100, 250)),
                    )
                ],
            )

            preview_one = render_manifest_preview(page, root, output_size=(800, 450))
            preview_two = render_manifest_preview(page, root, output_size=(800, 450))
            preview_result = render_manifest_preview_with_metadata(page, root, output_size=(800, 450))

        self.assertEqual(preview_one.size, (800, 450))
        self.assertEqual(preview_one.tobytes(), preview_two.tobytes())
        self.assertEqual(preview_one.getpixel((100, 100)), (37, 99, 235))
        self.assertEqual(preview_one.getpixel((320, 110)), (16, 185, 129))
        text_region = preview_one.crop((100, 210, 260, 250))
        histogram = text_region.convert("L").histogram()
        self.assertGreater(sum(histogram[:255]), 0)
        self.assertEqual(preview_result.metadata["renderer"], "manifest_stub")
        self.assertFalse(preview_result.metadata["is_powerpoint_render"])

    def test_structural_validation_passes_for_composed_manifest_deck(self):
        from src.generative_editable_composer import compose_deck_from_manifests
        from src.generative_editable_preview_validator import validate_composed_deck_structure

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            deck, _ = self._write_validation_fixture(root)
            output = root / "out.pptx"
            compose_deck_from_manifests(root / "deck.json", root, output)

            report = validate_composed_deck_structure(
                deck_manifest_path=root / "deck.json",
                artifact_root=root,
                pptx_path=output,
            )

        self.assertEqual(report.status, "passed")
        self.assertEqual(report.issues, [])
        self.assertEqual(report.checked_pages, len(deck.slide_order))
        self.assertEqual(
            report.to_dict(),
            {
                "status": "passed",
                "checked_pages": 1,
                "issues": [],
            },
        )

    def test_structural_validation_reports_wrong_slide_background_identity(self):
        from src.generative_editable_composer import compose_deck_from_manifests
        from src.generative_editable_preview_validator import validate_composed_deck_structure

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            self._write_validation_fixture(root)
            output = root / "out.pptx"
            compose_deck_from_manifests(root / "deck.json", root, output)
            Image.new("RGB", (800, 450), "#111827").save(
                root / "backgrounds" / "0000-slide-a" / "base.png"
            )

            report = validate_composed_deck_structure(
                deck_manifest_path=root / "deck.json",
                artifact_root=root,
                pptx_path=output,
            )

        self.assertEqual(report.status, "failed")
        background_issues = [
            issue for issue in report.issues
            if issue.code == "object_identity_mismatch"
            and issue.details["target_kind"] == "background"
        ]
        self.assertEqual(len(background_issues), 1)
        self.assertNotEqual(
            background_issues[0].details["expected_sha1"],
            background_issues[0].details["actual_sha1"],
        )

    def test_structural_validation_accepts_negative_slope_native_line_bbox(self):
        from src.generative_editable_composer import compose_deck_from_manifests
        from src.generative_editable_preview_validator import validate_composed_deck_structure

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            _, page = self._write_validation_fixture(root)
            page = replace(
                page,
                native_shapes=[
                    NativeShapeSpec(
                        shape_type="line",
                        source_pixel_bbox=(572, 186, 720, 295),
                        line_color="#7C3AED",
                        line_start=(572, 290),
                        line_end=(719, 190),
                        stroke_width=13,
                    )
                ],
                bitmap_assets=[],
                text_boxes=[],
            )
            write_manifest(root / "pages" / "0000-slide-a.json", page)
            output = root / "negative-line.pptx"
            compose_deck_from_manifests(root / "deck.json", root, output)

            report = validate_composed_deck_structure(
                deck_manifest_path=root / "deck.json",
                artifact_root=root,
                pptx_path=output,
            )

        self.assertEqual(report.status, "passed")
        self.assertEqual(report.issues, [])

    def test_structural_validation_reports_page_count_and_dimension_failures(self):
        from src.generative_editable_preview_validator import validate_composed_deck_structure

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            self._write_validation_fixture(root)
            wrong = Presentation()
            wrong.slide_width = Inches(10)
            wrong.slide_height = Inches(7.5)
            wrong.slides.add_slide(wrong.slide_layouts[6])
            wrong.slides.add_slide(wrong.slide_layouts[6])
            output = root / "wrong.pptx"
            wrong.save(output)

            report = validate_composed_deck_structure(
                deck_manifest_path=root / "deck.json",
                artifact_root=root,
                pptx_path=output,
            )

        self.assertEqual(report.status, "failed")
        self.assertIn("page_count_mismatch", [issue.code for issue in report.issues])
        self.assertIn("slide_dimensions_mismatch", [issue.code for issue in report.issues])

    def test_structural_validation_reports_missing_assets_and_required_text(self):
        from src.generative_editable_composer import compose_deck_from_manifests
        from src.generative_editable_preview_validator import validate_composed_deck_structure

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            _, page = self._write_validation_fixture(root)
            output = root / "out.pptx"
            compose_deck_from_manifests(root / "deck.json", root, output)
            missing_text_page = PageManifest(
                **{
                    **page.__dict__,
                    "bitmap_assets": [
                        BitmapAssetSpec(
                            asset_id="missing",
                            source_pixel_bbox=(300, 90, 380, 150),
                            asset_path="assets/0000-slide-a/missing.png",
                            z_order=1,
                        )
                    ],
                    "text_boxes": [
                        TextBoxSpec(
                            text="Missing Required Text",
                            source_pixel_bbox=(100, 210, 330, 250),
                            source_pixel_polygon=((100, 210), (330, 210), (330, 250), (100, 250)),
                        )
                    ],
                }
            )
            write_manifest(root / "pages" / "0000-slide-a.json", missing_text_page)

            report = validate_composed_deck_structure(
                deck_manifest_path=root / "deck.json",
                artifact_root=root,
                pptx_path=output,
            )

        codes = [issue.code for issue in report.issues]
        self.assertIn("missing_asset", codes)
        self.assertIn("missing_required_text", codes)
        failed_payload = report.to_dict()
        self.assertEqual(failed_payload["status"], "failed")
        self.assertEqual(failed_payload["checked_pages"], 1)
        text_issue = next(issue for issue in failed_payload["issues"] if issue["code"] == "missing_required_text")
        self.assertEqual(text_issue["severity"], "error")
        self.assertEqual(text_issue["slide_id"], "slide-a")
        self.assertEqual(text_issue["details"]["target_kind"], "text")
        self.assertEqual(text_issue["details"]["page_index"], 0)

    def test_structural_validation_rejects_source_crop_bitmap_asset_provenance(self):
        from src.generative_editable_composer import compose_deck_from_manifests
        from src.generative_editable_preview_validator import validate_composed_deck_structure

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            _, page = self._write_validation_fixture(root)
            unsafe_page = replace(
                page,
                bitmap_assets=[
                    replace(
                        page.bitmap_assets[0],
                        provenance={
                            "fallback": "source_crop_after_asset_sheet_failure",
                            "failure": "asset sheet cannot be sliced",
                        },
                    )
                ],
            )
            write_manifest(root / "pages" / "0000-slide-a.json", unsafe_page)
            output = root / "source-crop-asset.pptx"
            compose_deck_from_manifests(root / "deck.json", root, output)

            report = validate_composed_deck_structure(
                deck_manifest_path=root / "deck.json",
                artifact_root=root,
                pptx_path=output,
            )

        self.assertEqual(report.status, "failed")
        issue = next(
            issue
            for issue in report.issues
            if issue.code == "forbidden_source_crop_bitmap_asset"
        )
        self.assertEqual(issue.details["target_kind"], "bitmap_asset")
        self.assertEqual(issue.details["target_id"], "asset")
        self.assertEqual(issue.details["asset_ref"], "assets/0000-slide-a/asset.png")

    def test_structural_validation_allows_vlm_bbox_source_preserved_crop_assets(self):
        from src.generative_editable_composer import compose_deck_from_manifests
        from src.generative_editable_preview_validator import validate_composed_deck_structure

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            _, page = self._write_validation_fixture(root)
            source_preserved_path = root / "assets" / "0000-slide-a" / "asset.source-preserved.png"
            source_preserved_path.write_bytes(
                (root / "assets" / "0000-slide-a" / "asset.png").read_bytes()
            )
            unsafe_page = replace(
                page,
                bitmap_assets=[
                    replace(
                        page.bitmap_assets[0],
                        asset_path="assets/0000-slide-a/asset.source-preserved.png",
                        provenance={
                            "asset_strategy": "source_preserved_crop",
                            "alpha_strategy": "opaque_source_crop",
                            "candidate_classification": "complex_whole_visual",
                            "asset_sheet_skipped_reason": "complex_bitmap_region",
                            "candidate_provenance": {"source": "vlm_bitmap_region"},
                            "original_source_pixel_bbox": [10, 10, 80, 80],
                        },
                    )
                ],
            )
            write_manifest(root / "pages" / "0000-slide-a.json", unsafe_page)
            output = root / "source-preserved-provider-failed.pptx"
            compose_deck_from_manifests(root / "deck.json", root, output)

            report = validate_composed_deck_structure(
                deck_manifest_path=root / "deck.json",
                artifact_root=root,
                pptx_path=output,
            )

        self.assertEqual(report.status, "passed")
        self.assertNotIn(
            "forbidden_source_crop_bitmap_asset",
            [issue.code for issue in report.issues],
        )

    def test_structural_validation_allows_clean_fallback_icon_source_preserved_crop_assets(self):
        from src.generative_editable_composer import compose_deck_from_manifests
        from src.generative_editable_preview_validator import validate_composed_deck_structure

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            _, page = self._write_validation_fixture(root)
            source_preserved_path = root / "assets" / "0000-slide-a" / "icon.source-preserved.png"
            source_preserved_path.write_bytes(
                (root / "assets" / "0000-slide-a" / "asset.png").read_bytes()
            )
            clean_fallback_page = replace(
                page,
                provenance={**page.provenance, "clean_background_provider_failed": True},
                bitmap_assets=[
                    replace(
                        page.bitmap_assets[0],
                        asset_path="assets/0000-slide-a/icon.source-preserved.png",
                        provenance={
                            "asset_strategy": "source_preserved_crop",
                            "alpha_strategy": "opaque_source_crop",
                            "candidate_classification": "bitmap_asset_candidate",
                            "asset_sheet_skipped_reason": "icon_source_preserved",
                            "candidate_provenance": {
                                "source": "vlm_bitmap_region",
                                "vlm_type": "icon",
                            },
                            "original_source_pixel_bbox": [10, 10, 80, 80],
                        },
                    )
                ],
            )
            write_manifest(root / "pages" / "0000-slide-a.json", clean_fallback_page)
            output = root / "source-preserved-clean-fallback-icon.pptx"
            compose_deck_from_manifests(root / "deck.json", root, output)

            report = validate_composed_deck_structure(
                deck_manifest_path=root / "deck.json",
                artifact_root=root,
                pptx_path=output,
            )

        self.assertEqual(report.status, "passed")
        self.assertNotIn(
            "forbidden_source_crop_bitmap_asset",
            [issue.code for issue in report.issues],
        )

    def test_structural_validation_rejects_oversized_clean_fallback_icon_source_preserved_crop_assets(self):
        from src.generative_editable_composer import compose_deck_from_manifests
        from src.generative_editable_preview_validator import validate_composed_deck_structure

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            _, page = self._write_validation_fixture(root)
            source_preserved_path = root / "assets" / "0000-slide-a" / "icon.source-preserved.png"
            source_preserved_path.write_bytes(
                (root / "assets" / "0000-slide-a" / "asset.png").read_bytes()
            )
            oversized_icon_page = replace(
                page,
                provenance={**page.provenance, "clean_background_provider_failed": True},
                bitmap_assets=[
                    replace(
                        page.bitmap_assets[0],
                        asset_path="assets/0000-slide-a/icon.source-preserved.png",
                        source_pixel_bbox=(20, 20, 720, 380),
                        provenance={
                            "asset_strategy": "source_preserved_crop",
                            "alpha_strategy": "opaque_source_crop",
                            "candidate_classification": "bitmap_asset_candidate",
                            "asset_sheet_skipped_reason": "icon_source_preserved",
                            "candidate_provenance": {
                                "source": "vlm_bitmap_region",
                                "vlm_type": "icon",
                            },
                            "original_source_pixel_bbox": [20, 20, 720, 380],
                        },
                    )
                ],
            )
            write_manifest(root / "pages" / "0000-slide-a.json", oversized_icon_page)
            output = root / "source-preserved-oversized-clean-fallback-icon.pptx"
            compose_deck_from_manifests(root / "deck.json", root, output)

            report = validate_composed_deck_structure(
                deck_manifest_path=root / "deck.json",
                artifact_root=root,
                pptx_path=output,
            )

        self.assertEqual(report.status, "failed")
        self.assertIn(
            "forbidden_source_crop_bitmap_asset",
            [issue.code for issue in report.issues],
        )

    def test_structural_validation_rejects_incomplete_opaque_source_crop_provenance(self):
        from src.generative_editable_composer import compose_deck_from_manifests
        from src.generative_editable_preview_validator import validate_composed_deck_structure

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            _, page = self._write_validation_fixture(root)
            source_preserved_path = root / "assets" / "0000-slide-a" / "asset.source-preserved.png"
            source_preserved_path.write_bytes(
                (root / "assets" / "0000-slide-a" / "asset.png").read_bytes()
            )
            unsafe_page = replace(
                page,
                bitmap_assets=[
                    replace(
                        page.bitmap_assets[0],
                        asset_path="assets/0000-slide-a/asset.source-preserved.png",
                        provenance={
                            "asset_strategy": "source_preserved_crop",
                            "alpha_strategy": "opaque_source_crop",
                            "candidate_classification": "complex_whole_visual",
                        },
                    )
                ],
            )
            write_manifest(root / "pages" / "0000-slide-a.json", unsafe_page)
            output = root / "source-preserved-incomplete.pptx"
            compose_deck_from_manifests(root / "deck.json", root, output)

            report = validate_composed_deck_structure(
                deck_manifest_path=root / "deck.json",
                artifact_root=root,
                pptx_path=output,
            )

        self.assertEqual(report.status, "failed")
        self.assertIn(
            "forbidden_source_crop_bitmap_asset",
            [issue.code for issue in report.issues],
        )

    def test_structural_validation_reports_unsafe_full_slide_source_with_text(self):
        from src.generative_editable_composer import compose_deck_from_manifests
        from src.generative_editable_preview_validator import validate_composed_deck_structure

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            _, page = self._write_validation_fixture(root)
            unsafe_page = PageManifest(
                **{
                    **page.__dict__,
                    "provenance": {"chosen_background_kind": "source_full_slide"},
                }
            )
            write_manifest(root / "pages" / "0000-slide-a.json", unsafe_page)
            output = root / "out.pptx"
            compose_deck_from_manifests(root / "deck.json", root, output)

            report = validate_composed_deck_structure(
                deck_manifest_path=root / "deck.json",
                artifact_root=root,
                pptx_path=output,
            )

        self.assertIn("unsafe_full_slide_source_with_text", [issue.code for issue in report.issues])
        self.assertIn("baked text", report.issues[0].message)

    def test_structural_validation_rejects_source_background_without_editable_structure(self):
        from src.generative_editable_composer import compose_deck_from_manifests
        from src.generative_editable_preview_validator import validate_composed_deck_structure

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            _, page = self._write_validation_fixture(root)
            source_dir = root / "sources" / "0000-slide-a"
            source_dir.mkdir(parents=True, exist_ok=True)
            source = source_dir / "source.png"
            Image.new("RGB", (800, 450), "#FFFFFF").save(source)
            source_background = root / "backgrounds" / "0000-slide-a" / "source-copy.png"
            source_background.parent.mkdir(parents=True, exist_ok=True)
            shutil.copyfile(source, source_background)
            source_page = replace(
                page,
                source_image_path="sources/0000-slide-a/source.png",
                chosen_background="backgrounds/0000-slide-a/source-copy.png",
                text_clean_background="backgrounds/0000-slide-a/source-copy.png",
                base_clean_background="backgrounds/0000-slide-a/source-copy.png",
                native_shapes=[],
                bitmap_assets=[],
                text_boxes=[],
            )
            write_manifest(root / "pages" / "0000-slide-a.json", source_page)
            output = root / "source-background-only.pptx"
            compose_deck_from_manifests(root / "deck.json", root, output)

            report = validate_composed_deck_structure(
                deck_manifest_path=root / "deck.json",
                artifact_root=root,
                pptx_path=output,
            )

        self.assertEqual(report.status, "failed")
        issue = next(issue for issue in report.issues if issue.code == "full_slide_source_background_only")
        self.assertEqual(issue.details["target_kind"], "background")
        self.assertEqual(issue.details["editable_object_count"], 0)

    def test_structural_validation_reports_full_slide_source_with_approximate_text(self):
        from src.generative_editable_composer import compose_deck_from_manifests
        from src.generative_editable_preview_validator import validate_composed_deck_structure

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            _, page = self._write_validation_fixture(root)
            approximate_page = replace(
                page,
                provenance={
                    **page.provenance,
                    "chosen_background_kind": "source_full_slide",
                    "reconstruction_mode": "source_visual_background",
                },
                text_boxes=[
                    replace(
                        text_box,
                        style_hints={**text_box.style_hints, "approximate_layout": True},
                    )
                    for text_box in page.text_boxes
                ],
            )
            write_manifest(root / "pages" / "0000-slide-a.json", approximate_page)
            output = root / "out.pptx"
            compose_deck_from_manifests(root / "deck.json", root, output)

            report = validate_composed_deck_structure(
                deck_manifest_path=root / "deck.json",
                artifact_root=root,
                pptx_path=output,
            )

        self.assertIn("unsafe_full_slide_source_with_text", [issue.code for issue in report.issues])

    def test_structural_validation_rejects_declared_low_opacity_source_overlay(self):
        from src.generative_editable_composer import compose_deck_from_manifests
        from src.generative_editable_preview_validator import validate_composed_deck_structure

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            _, page = self._write_validation_fixture(root)
            source_dir = root / "sources" / "0000-slide-a"
            source_dir.mkdir(parents=True, exist_ok=True)
            source = source_dir / "source.png"
            Image.new("RGB", (800, 450), "#FFFFFF").save(source)
            background = root / "backgrounds" / "0000-slide-a" / "source-preserving.png"
            background.parent.mkdir(parents=True, exist_ok=True)
            Image.new("RGB", (800, 450), "#FFFFFF").save(background)
            overlay_page = replace(
                page,
                source_image_path="sources/0000-slide-a/source.png",
                chosen_background="backgrounds/0000-slide-a/source-preserving.png",
                text_clean_background="backgrounds/0000-slide-a/source-preserving.png",
                base_clean_background="backgrounds/0000-slide-a/source-preserving.png",
                text_boxes=[replace(text_box, opacity=0.1) for text_box in page.text_boxes],
                provenance={
                    **page.provenance,
                    "chosen_background_kind": "source_preserving_low_opacity_text_overlay",
                    "text_overlay_opacity": 0.1,
                },
            )
            write_manifest(root / "pages" / "0000-slide-a.json", overlay_page)
            output = root / "out.pptx"
            compose_deck_from_manifests(root / "deck.json", root, output)

            report = validate_composed_deck_structure(
                deck_manifest_path=root / "deck.json",
                artifact_root=root,
                pptx_path=output,
            )

        self.assertIn("unsafe_full_slide_source_with_text", [issue.code for issue in report.issues])

    def test_structural_validation_detects_actual_full_slide_source_picture_with_text(self):
        from src.generative_editable_preview_validator import validate_composed_deck_structure

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            _, page = self._write_validation_fixture(root)
            source_dir = root / "sources" / "0000-slide-a"
            source_dir.mkdir(parents=True, exist_ok=True)
            source = source_dir / "source.png"
            Image.new("RGB", (800, 450), "#FFFFFF").save(source)
            source_page = replace(page, source_image_path="sources/0000-slide-a/source.png")
            write_manifest(root / "pages" / "0000-slide-a.json", source_page)
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(5.625)
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(str(source), 0, 0, width=prs.slide_width, height=prs.slide_height)
            text = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(0.5))
            text.text_frame.text = "Preview"
            output = root / "unsafe.pptx"
            prs.save(output)

            report = validate_composed_deck_structure(
                deck_manifest_path=root / "deck.json",
                artifact_root=root,
                pptx_path=output,
            )

        self.assertIn("unsafe_full_slide_source_with_text", [issue.code for issue in report.issues])

    def test_structural_validation_detects_near_source_full_slide_background(self):
        from src.generative_editable_preview_validator import validate_composed_deck_structure

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            _, page = self._write_validation_fixture(root)
            source_dir = root / "sources" / "0000-slide-a"
            source_dir.mkdir(parents=True, exist_ok=True)
            source = source_dir / "source.png"
            Image.new("RGB", (800, 450), "#FFFFFF").save(source)
            near_source = root / "backgrounds" / "0000-slide-a" / "near-source.png"
            Image.new("RGB", (800, 450), "#FEFEFE").save(near_source)
            source_page = replace(
                page,
                source_image_path="sources/0000-slide-a/source.png",
                chosen_background="backgrounds/0000-slide-a/near-source.png",
            )
            write_manifest(root / "pages" / "0000-slide-a.json", source_page)
            output = root / "out.pptx"
            from src.generative_editable_composer import compose_deck_from_manifests

            compose_deck_from_manifests(root / "deck.json", root, output)

            report = validate_composed_deck_structure(
                deck_manifest_path=root / "deck.json",
                artifact_root=root,
                pptx_path=output,
            )

        self.assertIn("unsafe_full_slide_source_with_text", [issue.code for issue in report.issues])

    def test_structural_validation_detects_resized_source_full_slide_background(self):
        from src.generative_editable_composer import compose_deck_from_manifests
        from src.generative_editable_preview_validator import validate_composed_deck_structure

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            _, page = self._write_validation_fixture(root)
            source_dir = root / "sources" / "0000-slide-a"
            source_dir.mkdir(parents=True, exist_ok=True)
            source = source_dir / "source.png"
            source_image = Image.new("RGB", (800, 450), "#FFFFFF")
            ImageDraw.Draw(source_image).text((100, 120), "Preview", fill="#111827")
            source_image.save(source)
            resized = root / "backgrounds" / "0000-slide-a" / "resized-source.png"
            source_image.resize((1600, 900)).save(resized)
            source_page = replace(
                page,
                source_image_path="sources/0000-slide-a/source.png",
                chosen_background="backgrounds/0000-slide-a/resized-source.png",
            )
            write_manifest(root / "pages" / "0000-slide-a.json", source_page)
            output = root / "out.pptx"
            compose_deck_from_manifests(root / "deck.json", root, output)

            report = validate_composed_deck_structure(
                deck_manifest_path=root / "deck.json",
                artifact_root=root,
                pptx_path=output,
            )

        self.assertIn("unsafe_full_slide_source_with_text", [issue.code for issue in report.issues])

    def test_structural_validation_allows_text_clean_local_background_with_editable_text(self):
        from src.generative_editable_composer import compose_deck_from_manifests
        from src.generative_editable_preview_validator import validate_composed_deck_structure

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            _, page = self._write_validation_fixture(root)
            source_dir = root / "sources" / "0000-slide-a"
            source_dir.mkdir(parents=True, exist_ok=True)
            source = source_dir / "source.png"
            Image.new("RGB", (800, 450), "#FFFFFF").save(source)
            cleaned = root / "backgrounds" / "0000-slide-a" / "text-clean.png"
            Image.new("RGB", (800, 450), "#F0F0F0").save(cleaned)
            source_page = replace(
                page,
                source_image_path="sources/0000-slide-a/source.png",
                text_clean_background="backgrounds/0000-slide-a/text-clean.png",
                base_clean_background="backgrounds/0000-slide-a/text-clean.png",
                chosen_background="backgrounds/0000-slide-a/text-clean.png",
                provenance={
                    "backgrounds": {
                        "base_clean": {
                            "prompt_id": "local_text_cleanup",
                            "provider_role": "local",
                        }
                    }
                },
            )
            write_manifest(root / "pages" / "0000-slide-a.json", source_page)
            output = root / "out.pptx"
            compose_deck_from_manifests(root / "deck.json", root, output)

            report = validate_composed_deck_structure(
                deck_manifest_path=root / "deck.json",
                artifact_root=root,
                pptx_path=output,
            )

        self.assertNotIn("unsafe_full_slide_source_with_text", [issue.code for issue in report.issues])

    def test_structural_validation_reports_object_order_failures(self):
        from src.generative_editable_preview_validator import validate_composed_deck_structure

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            self._write_validation_fixture(root)
            wrong = Presentation()
            wrong.slide_width = Inches(10)
            wrong.slide_height = Inches(5.625)
            slide = wrong.slides.add_slide(wrong.slide_layouts[6])
            text = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(0.5))
            text.text_frame.text = "Preview"
            slide.shapes.add_picture(str(root / "backgrounds" / "0000-slide-a" / "base.png"), 0, 0)
            output = root / "wrong-order.pptx"
            wrong.save(output)

            report = validate_composed_deck_structure(
                deck_manifest_path=root / "deck.json",
                artifact_root=root,
                pptx_path=output,
            )

        self.assertIn("object_order_mismatch", [issue.code for issue in report.issues])

    def test_structural_validation_reports_extra_unmanifested_relevant_objects(self):
        from src.generative_editable_composer import compose_deck_from_manifests
        from src.generative_editable_preview_validator import validate_composed_deck_structure

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            self._write_validation_fixture(root)
            (root / "sources" / "0000-slide-a" / "source.png").unlink()
            output = root / "out.pptx"
            compose_deck_from_manifests(root / "deck.json", root, output)
            prs = Presentation(str(output))
            slide = prs.slides[0]
            slide.shapes.add_picture(
                str(root / "assets" / "0000-slide-a" / "asset.png"),
                Inches(8),
                Inches(4),
                width=Inches(1),
                height=Inches(1),
            )
            prs.save(output)

            report = validate_composed_deck_structure(
                deck_manifest_path=root / "deck.json",
                artifact_root=root,
                pptx_path=output,
            )

        self.assertIn("object_order_mismatch", [issue.code for issue in report.issues])

    def test_structural_validation_reports_wrong_bitmap_asset_identity_at_expected_position(self):
        from src.generative_editable_composer import _set_slide_picture_background
        from src.generative_editable_preview_validator import validate_composed_deck_structure

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            self._write_validation_fixture(root)
            wrong_asset = root / "assets" / "0000-slide-a" / "wrong.png"
            Image.new("RGBA", (80, 60), (249, 115, 22, 255)).save(wrong_asset)
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(5.625)
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            _set_slide_picture_background(slide, root / "backgrounds" / "0000-slide-a" / "base.png")
            slide.shapes.add_shape(1, Inches(1), Inches(1), Inches(2), Inches(1))
            slide.shapes.add_connector(1, Inches(3.5), Inches(1.0125), Inches(4.75), Inches(1.0125))
            slide.shapes.add_picture(str(wrong_asset), Inches(3.75), Inches(1.125), width=Inches(1), height=Inches(0.75))
            text = slide.shapes.add_textbox(Inches(1.25), Inches(2.625), Inches(2), Inches(0.5))
            text.text_frame.text = "Preview"
            output = root / "wrong-identity.pptx"
            prs.save(output)

            report = validate_composed_deck_structure(
                deck_manifest_path=root / "deck.json",
                artifact_root=root,
                pptx_path=output,
            )

        self.assertIn("object_identity_mismatch", [issue.code for issue in report.issues])

    def test_structural_validation_rejects_text_clean_background_that_matches_source(self):
        from src.generative_editable_composer import compose_deck_from_manifests
        from src.generative_editable_preview_validator import validate_composed_deck_structure

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            _, page = self._write_validation_fixture(root)
            text_clean = root / "backgrounds" / "0000-slide-a" / "text-clean.png"
            with Image.open(root / page.source_image_path) as source:
                source.save(text_clean)
            page = replace(
                page,
                chosen_background="backgrounds/0000-slide-a/text-clean.png",
                text_clean_background="backgrounds/0000-slide-a/text-clean.png",
            )
            write_manifest(root / "pages" / "0000-slide-a.json", page)
            output = root / "source-as-text-clean.pptx"
            compose_deck_from_manifests(root / "deck.json", root, output)

            report = validate_composed_deck_structure(
                deck_manifest_path=root / "deck.json",
                artifact_root=root,
                pptx_path=output,
            )

        self.assertIn("unsafe_full_slide_source_with_text", [issue.code for issue in report.issues])

    def test_structural_validation_allows_sparse_text_clean_region(self):
        from src.generative_editable_composer import compose_deck_from_manifests
        from src.generative_editable_preview_validator import validate_composed_deck_structure

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            _, page = self._write_validation_fixture(root)
            source = root / page.source_image_path
            cleaned = root / "backgrounds" / "0000-slide-a" / "text-clean.png"
            mask = root / "assets" / "0000-slide-a" / "text-mask.png"
            mask.parent.mkdir(parents=True, exist_ok=True)
            with Image.open(source) as source_image:
                source_with_text = source_image.convert("RGB")
            for x in range(100, 140):
                for y in range(210, 216):
                    source_with_text.putpixel((x, y), (17, 24, 39))
            source_with_text.save(source)
            cleaned_image = source_with_text.copy()
            for x in range(100, 140):
                for y in range(210, 216):
                    cleaned_image.putpixel((x, y), (248, 250, 252))
            cleaned_image.save(cleaned)
            mask_image = Image.new("L", (800, 450), 0)
            for x in range(96, 265):
                for y in range(206, 254):
                    mask_image.putpixel((x, y), 255)
            mask_image.save(mask)
            page = replace(
                page,
                chosen_background="backgrounds/0000-slide-a/text-clean.png",
                text_clean_background="backgrounds/0000-slide-a/text-clean.png",
                provenance={
                    **page.provenance,
                    "text_mask_path": "assets/0000-slide-a/text-mask.png",
                    "backgrounds": {
                        "base_clean": {
                            "output_asset_ref": "backgrounds/0000-slide-a/text-clean.png",
                            "prompt_id": "local_text_cleanup",
                            "provider_role": "local",
                        }
                    },
                },
            )
            write_manifest(root / "pages" / "0000-slide-a.json", page)
            output = root / "sparse-text-clean.pptx"
            compose_deck_from_manifests(root / "deck.json", root, output)

            report = validate_composed_deck_structure(
                deck_manifest_path=root / "deck.json",
                artifact_root=root,
                pptx_path=output,
            )

        self.assertNotIn("unsafe_full_slide_source_with_text", [issue.code for issue in report.issues])

    def test_structural_validation_allows_source_preserving_background_when_text_region_changed(self):
        from src.generative_editable_composer import compose_deck_from_manifests
        from src.generative_editable_preview_validator import validate_composed_deck_structure

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            _, page = self._write_validation_fixture(root)
            source = root / page.source_image_path
            cleaned = root / "backgrounds" / "0000-slide-a" / "source-preserving.png"
            mask = root / "assets" / "0000-slide-a" / "text-mask.png"
            mask.parent.mkdir(parents=True, exist_ok=True)
            with Image.open(source) as source_image:
                source_with_text = source_image.convert("RGB")
            for x in range(100, 140):
                for y in range(210, 216):
                    source_with_text.putpixel((x, y), (17, 24, 39))
            source_with_text.save(source)
            cleaned_image = source_with_text.copy()
            for x in range(100, 140):
                for y in range(210, 216):
                    cleaned_image.putpixel((x, y), (248, 250, 252))
            cleaned_image.save(cleaned)
            mask_image = Image.new("L", (800, 450), 0)
            for x in range(96, 265):
                for y in range(206, 254):
                    mask_image.putpixel((x, y), 255)
            mask_image.save(mask)
            page = replace(
                page,
                chosen_background="backgrounds/0000-slide-a/source-preserving.png",
                provenance={
                    **page.provenance,
                    "chosen_background_kind": "source_preserving_text_clean",
                    "text_mask_path": "assets/0000-slide-a/text-mask.png",
                    "backgrounds": {
                        "source_preserving": {
                            "output_asset_ref": "backgrounds/0000-slide-a/source-preserving.png",
                            "prompt_id": "source_preserving_text_background",
                            "provider_role": "local",
                        }
                    },
                },
            )
            write_manifest(root / "pages" / "0000-slide-a.json", page)
            output = root / "source-preserving.pptx"
            compose_deck_from_manifests(root / "deck.json", root, output)

            report = validate_composed_deck_structure(
                deck_manifest_path=root / "deck.json",
                artifact_root=root,
                pptx_path=output,
            )

        self.assertNotIn("unsafe_full_slide_source_with_text", [issue.code for issue in report.issues])

    def test_structural_validation_reports_wrong_native_shape_type_at_expected_position(self):
        from src.generative_editable_composer import _set_slide_picture_background
        from src.generative_editable_preview_validator import validate_composed_deck_structure

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            _, page = self._write_validation_fixture(root)
            page = replace(
                page,
                native_shapes=[
                    NativeShapeSpec(
                        shape_type="ellipse",
                        source_pixel_bbox=(80, 80, 240, 140),
                        fill_color="#2563EB",
                    )
                ],
                bitmap_assets=[],
            )
            write_manifest(root / "pages" / "0000-slide-a.json", page)
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(5.625)
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            _set_slide_picture_background(slide, root / "backgrounds" / "0000-slide-a" / "base.png")
            slide.shapes.add_shape(1, Inches(1), Inches(1), Inches(2), Inches(0.75))
            text = slide.shapes.add_textbox(Inches(1.25), Inches(2.625), Inches(2), Inches(0.5))
            text.text_frame.text = "Preview"
            output = root / "wrong-native-shape.pptx"
            prs.save(output)

            report = validate_composed_deck_structure(
                deck_manifest_path=root / "deck.json",
                artifact_root=root,
                pptx_path=output,
            )

        self.assertIn("object_identity_mismatch", [issue.code for issue in report.issues])

    def test_structural_validation_rejects_selectable_background_picture(self):
        from src.generative_editable_preview_validator import validate_composed_deck_structure

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            self._write_validation_fixture(root)
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(5.625)
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(
                str(root / "backgrounds" / "0000-slide-a" / "base.png"),
                Inches(0.5),
                Inches(0),
                width=Inches(9.5),
                height=prs.slide_height,
            )
            slide.shapes.add_shape(1, Inches(1), Inches(1), Inches(2), Inches(0.75))
            slide.shapes.add_connector(1, Inches(3.5), Inches(1.0125), Inches(4.75), Inches(1.0125))
            slide.shapes.add_picture(
                str(root / "assets" / "0000-slide-a" / "asset.png"),
                Inches(3.75),
                Inches(1.125),
                width=Inches(1),
                height=Inches(0.75),
            )
            text = slide.shapes.add_textbox(Inches(1.25), Inches(2.625), Inches(2), Inches(0.5))
            text.text_frame.text = "Preview"
            output = root / "wrong-background-geometry.pptx"
            prs.save(output)

            report = validate_composed_deck_structure(
                deck_manifest_path=root / "deck.json",
                artifact_root=root,
                pptx_path=output,
            )

        self.assertIn("object_order_mismatch", [issue.code for issue in report.issues])
        issue = next(issue for issue in report.issues if issue.code == "object_order_mismatch")
        self.assertEqual(issue.details["target_kind"], "slide_object_order")

    def test_structural_validation_uses_deck_aspect_ratio_for_equivalent_wide_slide_sizes(self):
        from src.generative_editable_composer import compose_deck_from_manifests
        from src.generative_editable_preview_validator import validate_composed_deck_structure

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            _, page = self._write_validation_fixture(root)
            write_manifest(
                root / "pages" / "0000-slide-a.json",
                replace(page, slide_size=(13.333, 7.5)),
            )
            output = root / "wide-16-9.pptx"
            compose_deck_from_manifests(root / "deck.json", root, output)

            report = validate_composed_deck_structure(
                deck_manifest_path=root / "deck.json",
                artifact_root=root,
                pptx_path=output,
            )

        self.assertEqual(report.status, "passed")

    def test_structural_validation_counts_required_text_occurrences_exactly(self):
        from src.generative_editable_preview_validator import validate_composed_deck_structure

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            _, page = self._write_validation_fixture(root)
            page = replace(
                page,
                text_boxes=[
                    TextBoxSpec(
                        text="Plan",
                        source_pixel_bbox=(100, 210, 170, 250),
                        source_pixel_polygon=((100, 210), (170, 210), (170, 250), (100, 250)),
                    ),
                    TextBoxSpec(
                        text="Plan",
                        source_pixel_bbox=(180, 210, 250, 250),
                        source_pixel_polygon=((180, 210), (250, 210), (250, 250), (180, 250)),
                    ),
                ],
                native_shapes=[],
                bitmap_assets=[],
            )
            write_manifest(root / "pages" / "0000-slide-a.json", page)
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(5.625)
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(
                str(root / "backgrounds" / "0000-slide-a" / "base.png"),
                0,
                0,
                width=prs.slide_width,
                height=prs.slide_height,
            )
            text = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(0.5))
            text.text_frame.text = "Planning"
            output = root / "text-mismatch.pptx"
            prs.save(output)

            report = validate_composed_deck_structure(
                deck_manifest_path=root / "deck.json",
                artifact_root=root,
                pptx_path=output,
            )

        issues = [issue for issue in report.issues if issue.code == "missing_required_text"]
        self.assertEqual(len(issues), 2)
        self.assertEqual(issues[0].details["target_index"], 0)
        self.assertEqual(issues[1].details["target_index"], 1)

    def test_structural_validation_reports_required_text_position_mismatch(self):
        from src.generative_editable_preview_validator import validate_composed_deck_structure

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            _, page = self._write_validation_fixture(root)
            page = replace(page, native_shapes=[], bitmap_assets=[])
            write_manifest(root / "pages" / "0000-slide-a.json", page)
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(5.625)
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(
                str(root / "backgrounds" / "0000-slide-a" / "base.png"),
                0,
                0,
                width=prs.slide_width,
                height=prs.slide_height,
            )
            text = slide.shapes.add_textbox(Inches(7), Inches(4), Inches(2), Inches(0.5))
            text.text_frame.text = "Preview"
            output = root / "text-wrong-position.pptx"
            prs.save(output)

            report = validate_composed_deck_structure(
                deck_manifest_path=root / "deck.json",
                artifact_root=root,
                pptx_path=output,
            )

        self.assertIn("text_position_mismatch", [issue.code for issue in report.issues])

    def test_structural_validation_allows_expanded_approximate_text_width(self):
        from src.generative_editable_composer import compose_deck_from_manifests
        from src.generative_editable_preview_validator import validate_composed_deck_structure

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            _, page = self._write_validation_fixture(root)
            page = replace(
                page,
                native_shapes=[],
                bitmap_assets=[],
                text_boxes=[
                    replace(
                        page.text_boxes[0],
                        text="EMB",
                        source_pixel_bbox=(100, 210, 112, 226),
                        source_pixel_polygon=((100, 210), (112, 210), (112, 226), (100, 226)),
                        font_size=8,
                        style_hints={"approximate_layout": True},
                    )
                ],
            )
            write_manifest(root / "pages" / "0000-slide-a.json", page)
            output = root / "expanded-approx-text.pptx"
            compose_deck_from_manifests(root / "deck.json", root, output)

            report = validate_composed_deck_structure(
                deck_manifest_path=root / "deck.json",
                artifact_root=root,
                pptx_path=output,
            )

        self.assertNotIn("text_position_mismatch", [issue.code for issue in report.issues])

    def test_structural_validation_reports_invalid_pptx_and_missing_page_manifest_as_issues(self):
        from src.generative_editable_preview_validator import validate_composed_deck_structure

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            self._write_validation_fixture(root)
            (root / "pages" / "0000-slide-a.json").unlink()
            invalid = root / "invalid.pptx"
            invalid.write_text("not a pptx", encoding="utf-8")

            report = validate_composed_deck_structure(
                deck_manifest_path=root / "deck.json",
                artifact_root=root,
                pptx_path=invalid,
            )

        codes = [issue.code for issue in report.issues]
        self.assertIn("invalid_pptx", codes)
        self.assertIn("missing_page_manifest", codes)

    def test_structural_validation_reports_missing_source_asset(self):
        from src.generative_editable_composer import compose_deck_from_manifests
        from src.generative_editable_preview_validator import validate_composed_deck_structure

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            self._write_validation_fixture(root)
            output = root / "out.pptx"
            compose_deck_from_manifests(root / "deck.json", root, output)
            (root / "sources" / "0000-slide-a" / "source.png").unlink()

            report = validate_composed_deck_structure(
                deck_manifest_path=root / "deck.json",
                artifact_root=root,
                pptx_path=output,
            )

        self.assertIn("missing_asset", [issue.code for issue in report.issues])
        source_issue = next(issue for issue in report.issues if issue.details.get("target_kind") == "source")
        self.assertFalse(source_issue.details["repairable"])

    def test_structural_validation_reports_empty_source_asset_path(self):
        from src.generative_editable_composer import compose_deck_from_manifests
        from src.generative_editable_preview_validator import validate_composed_deck_structure

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            _, page = self._write_validation_fixture(root)
            write_manifest(root / "pages" / "0000-slide-a.json", replace(page, source_image_path=""))
            output = root / "out.pptx"
            compose_deck_from_manifests(root / "deck.json", root, output)

            report = validate_composed_deck_structure(
                deck_manifest_path=root / "deck.json",
                artifact_root=root,
                pptx_path=output,
            )

        self.assertIn("invalid_asset_path", [issue.code for issue in report.issues])
        source_issue = next(issue for issue in report.issues if issue.details.get("target_kind") == "source")
        self.assertEqual(source_issue.details["asset_ref"], "")

    def test_preview_similarity_passes_for_deterministic_matching_fixture_images(self):
        from src.generative_editable_preview_validator import PreviewRenderResult, validate_preview_similarity

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "source.png"
            Image.new("RGB", (120, 80), "#F8FAFC").save(source)
            preview = Image.new("RGB", (120, 80), "#F8FAFC")

            report = validate_preview_similarity(
                source_image_path=source,
                preview=PreviewRenderResult(
                    image=preview,
                    metadata={"renderer": "fake_powerpoint", "is_powerpoint_render": True},
                ),
                slide_id="slide-a",
                page_index=0,
                max_mean_abs_delta=1.0,
                max_changed_pixel_ratio=0.01,
            )

        self.assertEqual(report.status, "passed")
        self.assertEqual(report.issues, [])

    def test_preview_similarity_fails_when_fixture_images_exceed_threshold(self):
        from src.generative_editable_preview_validator import PreviewRenderResult, validate_preview_similarity

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "source.png"
            Image.new("RGB", (120, 80), "#F8FAFC").save(source)
            preview = Image.new("RGB", (120, 80), "#0F172A")

            report = validate_preview_similarity(
                source_image_path=source,
                preview=PreviewRenderResult(
                    image=preview,
                    metadata={"renderer": "fake_powerpoint", "is_powerpoint_render": True},
                ),
                slide_id="slide-a",
                page_index=0,
                max_mean_abs_delta=1.0,
                max_changed_pixel_ratio=0.01,
            )

        self.assertEqual(report.status, "failed")
        self.assertEqual(report.issues[0].code, "preview_similarity_failed")
        self.assertGreater(report.issues[0].details["mean_abs_delta"], 1.0)
        self.assertGreater(report.issues[0].details["changed_pixel_ratio"], 0.01)

    def test_preview_similarity_allows_text_dense_editable_render_drift(self):
        from src.generative_editable_preview_validator import PreviewRenderResult, validate_preview_similarity

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            _, page = self._write_validation_fixture(root)
            source = root / "source.png"
            source_image = Image.new("RGB", (100, 100), (128, 128, 128))
            source_image.save(source)
            preview = source_image.copy()
            draw = ImageDraw.Draw(preview)
            draw.rectangle((0, 0, 99, 13), fill=(160, 160, 160))
            dense_page = replace(
                page,
                source_image_size=(100, 100),
                native_shapes=[],
                text_boxes=[
                    replace(
                        page.text_boxes[0],
                        text=f"Line {index}",
                        source_pixel_bbox=(0, index, 80, index + 1),
                        source_pixel_polygon=((0, index), (80, index), (80, index + 1), (0, index + 1)),
                    )
                    for index in range(24)
                ],
                bitmap_assets=[
                    replace(page.bitmap_assets[0], asset_id=f"asset-{index}", source_pixel_bbox=(0, 0, 10, 10))
                    for index in range(6)
                ],
            )

            report = validate_preview_similarity(
                source_image_path=source,
                preview=PreviewRenderResult(
                    image=preview,
                    metadata={"renderer": "fake_powerpoint", "is_powerpoint_render": True},
                ),
                slide_id="slide-a",
                page_index=0,
                max_mean_abs_delta=20.4,
                max_changed_pixel_ratio=0.132,
                page_manifest=dense_page,
            )

        self.assertEqual(report.status, "passed")
        self.assertEqual(report.issues, [])

    def test_preview_similarity_rejects_text_dense_render_drift_outside_text_regions(self):
        from src.generative_editable_preview_validator import PreviewRenderResult, validate_preview_similarity

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            _, page = self._write_validation_fixture(root)
            source = root / "source.png"
            source_image = Image.new("RGB", (100, 100), (128, 128, 128))
            source_image.save(source)
            preview = source_image.copy()
            draw = ImageDraw.Draw(preview)
            draw.rectangle((0, 80, 99, 93), fill=(160, 160, 160))
            dense_page = replace(
                page,
                source_image_size=(100, 100),
                native_shapes=[],
                text_boxes=[
                    replace(
                        page.text_boxes[0],
                        text=f"Line {index}",
                        source_pixel_bbox=(0, index, 80, index + 1),
                        source_pixel_polygon=((0, index), (80, index), (80, index + 1), (0, index + 1)),
                    )
                    for index in range(24)
                ],
                bitmap_assets=[
                    replace(page.bitmap_assets[0], asset_id=f"asset-{index}", source_pixel_bbox=(0, 0, 10, 10))
                    for index in range(6)
                ],
            )

            report = validate_preview_similarity(
                source_image_path=source,
                preview=PreviewRenderResult(
                    image=preview,
                    metadata={"renderer": "fake_powerpoint", "is_powerpoint_render": True},
                ),
                slide_id="slide-a",
                page_index=0,
                max_mean_abs_delta=20.4,
                max_changed_pixel_ratio=0.132,
                page_manifest=dense_page,
            )

        self.assertEqual(report.status, "failed")
        self.assertEqual(report.issues[0].code, "preview_similarity_failed")

    def test_preview_similarity_rejects_text_dense_mixed_drift_with_too_much_non_text_change(self):
        from src.generative_editable_preview_validator import PreviewRenderResult, validate_preview_similarity

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            _, page = self._write_validation_fixture(root)
            source = root / "source.png"
            source_image = Image.new("RGB", (100, 100), (128, 128, 128))
            source_image.save(source)
            preview = source_image.copy()
            draw = ImageDraw.Draw(preview)
            for y in range(0, 11):
                for x in range(0, 100):
                    draw.point((x, y), fill=(160, 160, 160))
            for y in range(80, 85):
                for x in range(0, 100):
                    draw.point((x, y), fill=(160, 160, 160))
            for x in range(0, 10):
                draw.point((x, 85), fill=(160, 160, 160))
            dense_page = replace(
                page,
                source_image_size=(100, 100),
                native_shapes=[],
                text_boxes=[
                    replace(
                        page.text_boxes[0],
                        text=f"Line {index}",
                        source_pixel_bbox=(0, index, 80, index + 1),
                        source_pixel_polygon=((0, index), (80, index), (80, index + 1), (0, index + 1)),
                    )
                    for index in range(24)
                ],
                bitmap_assets=[
                    replace(page.bitmap_assets[0], asset_id=f"asset-{index}", source_pixel_bbox=(0, 0, 10, 10))
                    for index in range(6)
                ],
            )

            report = validate_preview_similarity(
                source_image_path=source,
                preview=PreviewRenderResult(
                    image=preview,
                    metadata={"renderer": "fake_powerpoint", "is_powerpoint_render": True},
                ),
                slide_id="slide-a",
                page_index=0,
                max_mean_abs_delta=20.4,
                max_changed_pixel_ratio=0.132,
                page_manifest=dense_page,
            )

        self.assertEqual(report.status, "failed")
        self.assertEqual(report.issues[0].code, "preview_similarity_failed")
        self.assertGreater(report.issues[0].details["outside_text_changed_pixel_ratio"], 0.05)

    def test_preview_similarity_still_rejects_sparse_render_drift_over_changed_area_threshold(self):
        from src.generative_editable_preview_validator import PreviewRenderResult, validate_preview_similarity

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            _, page = self._write_validation_fixture(root)
            source = root / "source.png"
            source_image = Image.new("RGB", (100, 100), (128, 128, 128))
            source_image.save(source)
            preview = source_image.copy()
            draw = ImageDraw.Draw(preview)
            draw.rectangle((0, 0, 99, 13), fill=(160, 160, 160))

            report = validate_preview_similarity(
                source_image_path=source,
                preview=PreviewRenderResult(
                    image=preview,
                    metadata={"renderer": "fake_powerpoint", "is_powerpoint_render": True},
                ),
                slide_id="slide-a",
                page_index=0,
                max_mean_abs_delta=20.4,
                max_changed_pixel_ratio=0.132,
                page_manifest=page,
            )

        self.assertEqual(report.status, "failed")
        self.assertEqual(report.issues[0].code, "preview_similarity_failed")

    def test_preview_similarity_rejects_manifest_stub_as_real_powerpoint_quality_gate(self):
        from src.generative_editable_preview_validator import (
            render_manifest_preview_with_metadata,
            validate_preview_similarity,
        )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            _, page = self._write_validation_fixture(root)
            source = root / "sources" / "0000-slide-a" / "source.png"
            source.parent.mkdir(parents=True, exist_ok=True)
            Image.new("RGB", (800, 450), "#FFFFFF").save(source)
            preview = render_manifest_preview_with_metadata(page, root, output_size=(800, 450))

            report = validate_preview_similarity(
                source_image_path=source,
                preview=preview,
                slide_id="slide-a",
                page_index=0,
                max_mean_abs_delta=12.0,
                max_changed_pixel_ratio=0.10,
                require_powerpoint_render=True,
            )

        self.assertEqual(report.status, "failed")
        self.assertEqual(report.issues[0].code, "preview_renderer_not_powerpoint")

    @unittest.skipIf(
        not shutil.which("soffice") or not shutil.which("pdftoppm"),
        "PowerPoint preview render tools are unavailable",
    )
    def test_powerpoint_preview_renderer_returns_real_render_metadata(self):
        from src.generative_editable_preview_validator import (
            render_powerpoint_preview_with_metadata,
        )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            _, page = self._write_validation_fixture(root)
            pptx_path = root / "deck.pptx"
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(5.625)
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(
                str(root / page.chosen_background),
                0,
                0,
                width=prs.slide_width,
                height=prs.slide_height,
            )
            prs.save(pptx_path)

            preview = render_powerpoint_preview_with_metadata(
                page,
                root,
                pptx_path=pptx_path,
                output_size=(800, 450),
            )

            self.assertEqual(preview.image.size, (800, 450))
            self.assertTrue(preview.metadata["is_powerpoint_render"])
            self.assertEqual(preview.metadata["renderer"], "soffice_pdf_pdftoppm")

    def test_powerpoint_preview_renderer_falls_back_when_renderer_times_out(self):
        from subprocess import TimeoutExpired

        from src.generative_editable_preview_validator import (
            render_powerpoint_preview_with_metadata,
        )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            _, page = self._write_validation_fixture(root)
            pptx_path = root / "deck.pptx"
            pptx_path.write_bytes(b"not a real pptx")
            with patch("src.generative_editable_preview_validator.shutil.which") as which:
                which.side_effect = lambda name: f"/usr/bin/{name}"
                with patch("src.generative_editable_preview_validator.subprocess.run") as run:
                    run.side_effect = TimeoutExpired(cmd="soffice", timeout=30)
                    preview = render_powerpoint_preview_with_metadata(
                        page,
                        root,
                        pptx_path=pptx_path,
                        output_size=(800, 450),
                    )

        self.assertFalse(preview.metadata["is_powerpoint_render"])
        self.assertEqual(preview.metadata["renderer"], "manifest_stub")
        self.assertIn("timeout", run.call_args.kwargs)

    def test_preview_similarity_rejects_manifest_stub_even_if_metadata_claims_powerpoint(self):
        from src.generative_editable_preview_validator import PreviewRenderResult, validate_preview_similarity

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "source.png"
            Image.new("RGB", (120, 80), "#F8FAFC").save(source)

            report = validate_preview_similarity(
                source_image_path=source,
                preview=PreviewRenderResult(
                    image=Image.new("RGB", (120, 80), "#F8FAFC"),
                    metadata={"renderer": "manifest_stub", "is_powerpoint_render": True},
                ),
                slide_id="slide-a",
                page_index=0,
                max_mean_abs_delta=1.0,
                max_changed_pixel_ratio=0.01,
            )

        self.assertEqual(report.status, "failed")
        self.assertEqual(report.issues[0].code, "preview_renderer_not_powerpoint")

    def test_preview_similarity_reports_wrong_aspect_preview_without_resizing(self):
        from src.generative_editable_preview_validator import PreviewRenderResult, validate_preview_similarity

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "source.png"
            Image.new("RGB", (120, 80), "#F8FAFC").save(source)

            report = validate_preview_similarity(
                source_image_path=source,
                preview=PreviewRenderResult(
                    image=Image.new("RGB", (120, 120), "#F8FAFC"),
                    metadata={"renderer": "fake_powerpoint", "is_powerpoint_render": True},
                ),
                slide_id="slide-a",
                page_index=0,
                max_mean_abs_delta=1.0,
                max_changed_pixel_ratio=0.01,
            )

        self.assertEqual(report.status, "failed")
        self.assertEqual(report.issues[0].code, "preview_dimensions_mismatch")

    def test_preview_similarity_returns_report_for_invalid_preview_image(self):
        from src.generative_editable_preview_validator import PreviewRenderResult, validate_preview_similarity

        class BrokenPreview:
            size = (120, 80)

            def convert(self, mode):
                raise OSError("broken preview")

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "source.png"
            Image.new("RGB", (120, 80), "#F8FAFC").save(source)

            report = validate_preview_similarity(
                source_image_path=source,
                preview=PreviewRenderResult(
                    image=BrokenPreview(),
                    metadata={"renderer": "fake_powerpoint", "is_powerpoint_render": True},
                ),
                slide_id="slide-a",
                page_index=0,
                max_mean_abs_delta=1.0,
                max_changed_pixel_ratio=0.01,
            )

        self.assertEqual(report.status, "failed")
        self.assertEqual(report.issues[0].code, "invalid_preview_image")

    def test_preview_similarity_caps_source_and_preview_pixel_counts_before_comparison(self):
        from src.generative_editable_preview_validator import PreviewRenderResult, validate_preview_similarity

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "source.png"
            Image.new("RGB", (11, 10), "#F8FAFC").save(source)

            with patch("src.generative_editable_preview_validator.MAX_PREVIEW_PIXELS", 100):
                source_report = validate_preview_similarity(
                    source_image_path=source,
                    preview=PreviewRenderResult(
                        image=Image.new("RGB", (10, 10), "#F8FAFC"),
                        metadata={"renderer": "fake_powerpoint", "is_powerpoint_render": True},
                    ),
                    slide_id="slide-a",
                    page_index=0,
                    max_mean_abs_delta=1.0,
                    max_changed_pixel_ratio=0.01,
                )
                Image.new("RGB", (10, 10), "#F8FAFC").save(source)
                preview_report = validate_preview_similarity(
                    source_image_path=source,
                    preview=PreviewRenderResult(
                        image=Image.new("RGB", (11, 10), "#F8FAFC"),
                        metadata={"renderer": "fake_powerpoint", "is_powerpoint_render": True},
                    ),
                    slide_id="slide-a",
                    page_index=0,
                    max_mean_abs_delta=1.0,
                    max_changed_pixel_ratio=0.01,
                )

        self.assertEqual(source_report.issues[0].code, "preview_too_large")
        self.assertEqual(preview_report.issues[0].code, "preview_too_large")

    def test_quality_similarity_threshold_maps_to_delta_gate_thresholds(self):
        from src.generative_editable_preview_validator import quality_threshold_to_preview_gates

        strict = quality_threshold_to_preview_gates(0.95)
        loose = quality_threshold_to_preview_gates(0.80)

        self.assertLess(strict.max_mean_abs_delta, loose.max_mean_abs_delta)
        self.assertLess(strict.max_changed_pixel_ratio, loose.max_changed_pixel_ratio)

    def test_quality_similarity_threshold_allows_render_edge_headroom(self):
        from src.generative_editable_preview_validator import quality_threshold_to_preview_gates

        gates = quality_threshold_to_preview_gates(0.92)

        self.assertEqual(gates.max_mean_abs_delta, 20.4)
        self.assertEqual(gates.max_changed_pixel_ratio, 0.132)

    def _write_validation_fixture(self, root: Path):
        (root / "backgrounds" / "0000-slide-a").mkdir(parents=True)
        (root / "assets" / "0000-slide-a").mkdir(parents=True)
        (root / "sources" / "0000-slide-a").mkdir(parents=True)
        Image.new("RGB", (800, 450), "#F8FAFC").save(
            root / "sources" / "0000-slide-a" / "source.png"
        )
        Image.new("RGB", (800, 450), "#FFFFFF").save(
            root / "backgrounds" / "0000-slide-a" / "base.png"
        )
        Image.new("RGBA", (80, 60), (16, 185, 129, 255)).save(
            root / "assets" / "0000-slide-a" / "asset.png"
        )
        page = PageManifest(
            slide_id="slide-a",
            page_index=0,
            source_image_path="sources/0000-slide-a/source.png",
            source_image_size=(800, 450),
            slide_size=(10.0, 5.625),
            chosen_background="backgrounds/0000-slide-a/base.png",
            native_shapes=[
                NativeShapeSpec(
                    shape_type="rectangle",
                    source_pixel_bbox=(80, 80, 240, 140),
                    fill_color="#2563EB",
                ),
                NativeShapeSpec(
                    shape_type="line",
                    source_pixel_bbox=(280, 80, 380, 83),
                    line_color="#111827",
                    line_start=(280, 81),
                    line_end=(380, 81),
                    stroke_width=3,
                )
            ],
            bitmap_assets=[
                BitmapAssetSpec(
                    asset_id="asset",
                    source_pixel_bbox=(300, 90, 380, 150),
                    asset_path="assets/0000-slide-a/asset.png",
                    z_order=1,
                )
            ],
            text_boxes=[
                TextBoxSpec(
                    text="Preview",
                    source_pixel_bbox=(100, 210, 260, 250),
                    source_pixel_polygon=((100, 210), (260, 210), (260, 250), (100, 250)),
                )
            ],
        )
        deck = DeckManifest(
            job_id="job-001",
            slide_order=["slide-a"],
            aspect_ratio="16:9",
            provider_roles={},
            quality_settings={},
            fallback_policy="fail",
            page_manifest_paths=["pages/0000-slide-a.json"],
        )
        write_manifest(root / "pages" / "0000-slide-a.json", page)
        write_manifest(root / "deck.json", deck)
        return deck, page


if __name__ == "__main__":
    unittest.main()
