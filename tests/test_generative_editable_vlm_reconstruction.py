import tempfile
import time
import unittest
from concurrent.futures import ThreadPoolExecutor
from pathlib import Path
from unittest.mock import patch

from PIL import Image, ImageDraw

from src.generative_editable_config import ProviderConfig
from src.generative_editable_providers import FakeImageEditProvider
from src.generative_editable_providers import ImageEditProvider
from src.generative_editable_providers import OCRResult, OCRTextItem
from src.generative_editable_providers import ProviderError


class GenerativeEditableVLMReconstructionTest(unittest.TestCase):
    def test_vlm_validation_fails_large_source_preserved_bitmap_coverage_without_runner(self):
        from src.generative_editable_manifest import BitmapAssetSpec, PageManifest
        from src.generative_editable_vlm_reconstruction import _validate_vlm_source_preserved_bitmap_coverage

        page = PageManifest(
            slide_id="slide-a",
            page_index=0,
            source_image_path="sources/slide-a.png",
            source_image_size=(1000, 500),
            slide_size=(10.0, 5.0),
            bitmap_assets=[
                BitmapAssetSpec(
                    asset_id="huge-source-preserved-crop",
                    source_pixel_bbox=(0, 0, 850, 500),
                    asset_path="assets/huge-source-preserved-crop.png",
                    z_order=1,
                    provenance={"asset_strategy": "source_preserved_crop"},
                )
            ],
        )

        report = _validate_vlm_source_preserved_bitmap_coverage([page])

        self.assertEqual(report.status, "failed")
        self.assertEqual(report.issues[0].code, "oversized_bitmap_asset_coverage")

    def test_vlm_validation_allows_split_row_source_preserved_bitmap_structure(self):
        from src.generative_editable_manifest import BitmapAssetSpec, NativeShapeSpec, PageManifest, TextBoxSpec
        from src.generative_editable_vlm_reconstruction import _validate_vlm_source_preserved_bitmap_coverage

        page = PageManifest(
            slide_id="slide-a",
            page_index=0,
            source_image_path="sources/slide-a.png",
            source_image_size=(1000, 500),
            slide_size=(10.0, 5.0),
            text_boxes=[
                TextBoxSpec(
                    text="标题",
                    source_pixel_bbox=(40, 20, 160, 60),
                    source_pixel_polygon=((40, 20), (160, 20), (160, 60), (40, 60)),
                )
            ],
            native_shapes=[
                NativeShapeSpec(
                    shape_type="line",
                    source_pixel_bbox=(20, 80, 980, 82),
                    line_start=(20, 80),
                    line_end=(980, 80),
                ),
                NativeShapeSpec(
                    shape_type="line",
                    source_pixel_bbox=(20, 240, 980, 242),
                    line_start=(20, 240),
                    line_end=(980, 240),
                ),
            ],
            bitmap_assets=[
                BitmapAssetSpec(
                    asset_id="row-a",
                    source_pixel_bbox=(0, 0, 1000, 150),
                    asset_path="assets/row-a.png",
                    z_order=1,
                    provenance={"asset_strategy": "source_preserved_crop", "alpha_visible_area_ratio": 0.12},
                ),
                BitmapAssetSpec(
                    asset_id="row-b",
                    source_pixel_bbox=(0, 170, 1000, 320),
                    asset_path="assets/row-b.png",
                    z_order=2,
                    provenance={"asset_strategy": "source_preserved_crop", "alpha_visible_area_ratio": 0.12},
                ),
                BitmapAssetSpec(
                    asset_id="row-c",
                    source_pixel_bbox=(0, 340, 1000, 490),
                    asset_path="assets/row-c.png",
                    z_order=3,
                    provenance={"asset_strategy": "source_preserved_crop", "alpha_visible_area_ratio": 0.12},
                ),
            ],
        )

        report = _validate_vlm_source_preserved_bitmap_coverage([page])

        self.assertEqual(report.status, "passed")
        self.assertEqual(report.issues, [])

    def test_vlm_pipeline_runs_structure_and_preview_validation_before_returning(self):
        from src.generative_editable_preview_validator import ValidationReport
        from src.generative_editable_vlm_reconstruction import (
            VLMEditablePipelineDependencies,
            VLMPageAnalysisProvider,
            coerce_vlm_analysis,
            run_vlm_editable_pptx_pipeline,
        )

        class EmptyVLMProvider(VLMPageAnalysisProvider):
            def analyze_page(self, image_path: str, *, timeout_seconds: int = 180):
                return coerce_vlm_analysis(
                    {
                        "coordinate_space": {"width": 160, "height": 90, "unit": "px"},
                        "text_regions": [],
                        "bitmap_regions": [],
                        "shape_regions": [
                            {"id": "line-a", "type": "line", "bbox": [20, 60, 140, 61]}
                        ],
                    }
                )

        calls = []

        def structure_validator(**kwargs):
            calls.append(("structure", kwargs["pptx_path"], kwargs["deck_manifest_path"]))
            return ValidationReport(status="passed", checked_pages=1, issues=[])

        def preview_renderer(page, artifact_root, *, pptx_path):
            calls.append(("preview_renderer", page.slide_id, pptx_path))
            return object()

        def preview_validator(**kwargs):
            calls.append(
                (
                    "preview_validator",
                    kwargs["slide_id"],
                    kwargs["pptx_path"],
                    kwargs["preview_similarity_threshold"],
                    kwargs["require_preview_validation"],
                )
            )
            return ValidationReport(status="passed", checked_pages=1, issues=[])

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "input.png"
            output = root / "out.pptx"
            Image.new("RGB", (160, 90), "#001122").save(source)

            result = run_vlm_editable_pptx_pipeline(
                slides=[{"slide_id": "slide-a", "image_path": str(source)}],
                output_path=str(output),
                artifact_root=str(root / "artifacts"),
                job_id="job-validate",
                dependencies=VLMEditablePipelineDependencies(
                    vlm_provider=EmptyVLMProvider(),
                    image_edit_provider=FakeImageEditProvider(
                        ProviderConfig(
                            role="edit_model",
                            provider="fake_image_edit",
                            model="fake-image-edit",
                            base_url="fake",
                            api_key="fake",
                        )
                    ),
                    preview_similarity_threshold=0.87,
                    require_preview_validation=True,
                    structure_validator=structure_validator,
                    preview_renderer=preview_renderer,
                    preview_validator=preview_validator,
                ),
            )

        self.assertEqual(result.validation_report.status, "passed")
        self.assertEqual([call[0] for call in calls], ["structure", "preview_renderer", "preview_validator"])
        self.assertEqual(calls[2][2], str(output))
        self.assertEqual(calls[2][3], 0.87)
        self.assertIs(calls[2][4], True)

    def test_vlm_pipeline_requires_preview_validation_by_default(self):
        from src.generative_editable_preview_validator import ValidationReport
        from src.generative_editable_vlm_reconstruction import (
            VLMEditablePipelineDependencies,
            VLMPageAnalysisProvider,
            coerce_vlm_analysis,
            run_vlm_editable_pptx_pipeline,
        )

        class EmptyVLMProvider(VLMPageAnalysisProvider):
            def analyze_page(self, image_path: str, *, timeout_seconds: int = 180):
                return coerce_vlm_analysis(
                    {
                        "coordinate_space": {"width": 160, "height": 90, "unit": "px"},
                        "text_regions": [],
                        "bitmap_regions": [],
                        "shape_regions": [
                            {"id": "line-a", "type": "line", "bbox": [20, 60, 140, 61]}
                        ],
                    }
                )

        observed = {}

        def preview_validator(**kwargs):
            observed["require_preview_validation"] = kwargs["require_preview_validation"]
            return ValidationReport(status="passed", checked_pages=1, issues=[])

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "input.png"
            output = root / "out.pptx"
            Image.new("RGB", (160, 90), "#001122").save(source)

            run_vlm_editable_pptx_pipeline(
                slides=[{"slide_id": "slide-a", "image_path": str(source)}],
                output_path=str(output),
                artifact_root=str(root / "artifacts"),
                job_id="job-default-preview",
                dependencies=VLMEditablePipelineDependencies(
                    vlm_provider=EmptyVLMProvider(),
                    image_edit_provider=FakeImageEditProvider(
                        ProviderConfig(
                            role="edit_model",
                            provider="fake_image_edit",
                            model="fake-image-edit",
                            base_url="fake",
                            api_key="fake",
                        )
                    ),
                    structure_validator=lambda **kwargs: ValidationReport(
                        status="passed", checked_pages=1, issues=[]
                    ),
                    preview_renderer=lambda page, artifact_root, *, pptx_path: object(),
                    preview_validator=preview_validator,
                ),
            )

        self.assertIs(observed["require_preview_validation"], True)

    def test_vlm_pipeline_raises_validation_error_when_validation_fails(self):
        from src.generative_editable_pipeline import GenerativeEditableValidationError
        from src.generative_editable_preview_validator import ValidationIssue, ValidationReport
        from src.generative_editable_vlm_reconstruction import (
            VLMEditablePipelineDependencies,
            VLMPageAnalysisProvider,
            coerce_vlm_analysis,
            run_vlm_editable_pptx_pipeline,
        )

        class EmptyVLMProvider(VLMPageAnalysisProvider):
            def analyze_page(self, image_path: str, *, timeout_seconds: int = 180):
                return coerce_vlm_analysis(
                    {
                        "coordinate_space": {"width": 160, "height": 90, "unit": "px"},
                        "text_regions": [],
                        "bitmap_regions": [],
                        "shape_regions": [
                            {"id": "line-a", "type": "line", "bbox": [20, 60, 140, 61]}
                        ],
                    }
                )

        def failing_preview_validator(**kwargs):
            return ValidationReport(
                status="failed",
                checked_pages=1,
                issues=[
                    ValidationIssue(
                        code="preview_similarity_failed",
                        message="forced preview mismatch",
                        slide_id=kwargs["slide_id"],
                    )
                ],
            )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "input.png"
            output = root / "out.pptx"
            Image.new("RGB", (160, 90), "#001122").save(source)

            with self.assertRaises(GenerativeEditableValidationError) as ctx:
                run_vlm_editable_pptx_pipeline(
                    slides=[{"slide_id": "slide-a", "image_path": str(source)}],
                    output_path=str(output),
                    artifact_root=str(root / "artifacts"),
                    job_id="job-validate-fail",
                    dependencies=VLMEditablePipelineDependencies(
                        vlm_provider=EmptyVLMProvider(),
                        image_edit_provider=FakeImageEditProvider(
                            ProviderConfig(
                                role="edit_model",
                                provider="fake_image_edit",
                                model="fake-image-edit",
                                base_url="fake",
                                api_key="fake",
                            )
                        ),
                        structure_validator=lambda **kwargs: ValidationReport(
                            status="passed", checked_pages=1, issues=[]
                        ),
                        preview_renderer=lambda page, artifact_root, *, pptx_path: object(),
                        preview_validator=failing_preview_validator,
                    ),
                )

        self.assertEqual(ctx.exception.validation_report.status, "failed")
        self.assertIn("preview_similarity_failed", str(ctx.exception))

    def test_vlm_pipeline_fails_when_page_has_no_editable_or_decomposed_objects(self):
        from src.generative_editable_pipeline import GenerativeEditableValidationError
        from src.generative_editable_preview_validator import ValidationReport
        from src.generative_editable_vlm_reconstruction import (
            VLMEditablePipelineDependencies,
            VLMPageAnalysisProvider,
            coerce_vlm_analysis,
            run_vlm_editable_pptx_pipeline,
        )

        class EmptyVLMProvider(VLMPageAnalysisProvider):
            def analyze_page(self, image_path: str, *, timeout_seconds: int = 180):
                return coerce_vlm_analysis(
                    {
                        "coordinate_space": {"width": 160, "height": 90, "unit": "px"},
                        "text_regions": [],
                        "bitmap_regions": [],
                        "shape_regions": [],
                    }
                )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "input.png"
            output = root / "out.pptx"
            Image.new("RGB", (160, 90), "#001122").save(source)

            with self.assertRaises(GenerativeEditableValidationError) as ctx:
                run_vlm_editable_pptx_pipeline(
                    slides=[{"slide_id": "slide-a", "image_path": str(source)}],
                    output_path=str(output),
                    artifact_root=str(root / "artifacts"),
                    job_id="job-empty-structure",
                    dependencies=VLMEditablePipelineDependencies(
                        vlm_provider=EmptyVLMProvider(),
                        image_edit_provider=FakeImageEditProvider(
                            ProviderConfig(
                                role="edit_model",
                                provider="fake_image_edit",
                                model="fake-image-edit",
                                base_url="fake",
                                api_key="fake",
                            )
                        ),
                        structure_validator=lambda **kwargs: ValidationReport(
                            status="passed", checked_pages=1, issues=[]
                        ),
                        preview_renderer=lambda page, artifact_root, *, pptx_path: object(),
                        preview_validator=lambda **kwargs: ValidationReport(
                            status="passed", checked_pages=1, issues=[]
                        ),
                    ),
                )

        self.assertEqual(ctx.exception.validation_report.status, "failed")
        self.assertIn("vlm_no_editable_or_decomposed_objects", str(ctx.exception))

    def test_vlm_provider_retries_retryable_errors(self):
        from src.generative_editable_vlm_reconstruction import (
            RetryingVLMPageAnalysisProvider,
            VLMPageAnalysisProvider,
            coerce_vlm_analysis,
        )

        class FlakyVLMProvider(VLMPageAnalysisProvider):
            def __init__(self):
                self.calls = 0

            def analyze_page(self, image_path: str, *, timeout_seconds: int = 180):
                self.calls += 1
                if self.calls == 1:
                    raise ProviderError(
                        provider_role="VLM",
                        operation="vlm_page_analysis",
                        message="temporary EOF",
                        retryable=True,
                    )
                return coerce_vlm_analysis(
                    {
                        "coordinate_space": {"width": 10, "height": 10, "unit": "px"},
                        "text_regions": [],
                        "bitmap_regions": [],
                        "shape_regions": [
                            {"id": "line-a", "type": "line", "bbox": [20, 60, 140, 61]}
                        ],
                    }
                )

        provider = FlakyVLMProvider()

        result = RetryingVLMPageAnalysisProvider(
            provider,
            max_attempts=2,
            backoff_seconds=0,
        ).analyze_page("source.png", timeout_seconds=5)

        self.assertEqual(provider.calls, 2)
        self.assertEqual(result.coordinate_space.width, 10)

    def test_vlm_pipeline_applies_configured_provider_retries(self):
        from src.generative_editable_vlm_reconstruction import (
            VLMEditablePipelineDependencies,
            VLMPageAnalysisProvider,
            coerce_vlm_analysis,
            run_vlm_editable_pptx_pipeline,
        )

        class FlakyPipelineVLMProvider(VLMPageAnalysisProvider):
            def __init__(self):
                self.calls = 0

            def analyze_page(self, image_path: str, *, timeout_seconds: int = 180):
                self.calls += 1
                if self.calls == 1:
                    raise ProviderError(
                        provider_role="VLM",
                        operation="vlm_page_analysis",
                        message="temporary EOF",
                        retryable=True,
                    )
                return coerce_vlm_analysis(
                    {
                        "coordinate_space": {"width": 160, "height": 90, "unit": "px"},
                        "text_regions": [],
                        "bitmap_regions": [],
                        "shape_regions": [
                            {"id": "line-a", "type": "line", "bbox": [20, 60, 140, 61]}
                        ],
                    }
                )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "input.png"
            Image.new("RGB", (160, 90), "#001122").save(source)
            output = root / "out.pptx"
            vlm_provider = FlakyPipelineVLMProvider()

            run_vlm_editable_pptx_pipeline(
                slides=[{"slide_id": "slide-a", "image_path": str(source)}],
                output_path=str(output),
                artifact_root=str(root / "artifacts"),
                job_id="job-retry",
                dependencies=VLMEditablePipelineDependencies(
                    vlm_provider=vlm_provider,
                    image_edit_provider=FakeImageEditProvider(
                        ProviderConfig(
                            role="edit_model",
                            provider="fake_image_edit",
                            model="fake-image-edit",
                            base_url="fake",
                            api_key="fake",
                        )
                    ),
                    provider_max_attempts=2,
                    provider_retry_backoff_seconds=0,
                ),
            )

        self.assertEqual(vlm_provider.calls, 2)

    def test_vlm_pipeline_writes_stage_events_for_timeout_diagnostics(self):
        from src.generative_editable_job_artifacts import GenerativeEditableJobArtifacts
        from src.generative_editable_vlm_reconstruction import (
            VLMEditablePipelineDependencies,
            VLMPageAnalysisProvider,
            coerce_vlm_analysis,
            run_vlm_editable_pptx_pipeline,
        )

        class EmptyVLMProvider(VLMPageAnalysisProvider):
            def analyze_page(self, image_path: str, *, timeout_seconds: int = 180):
                return coerce_vlm_analysis(
                    {
                        "coordinate_space": {"width": 160, "height": 90, "unit": "px"},
                        "text_regions": [],
                        "bitmap_regions": [],
                        "shape_regions": [
                            {"id": "line-a", "type": "line", "bbox": [20, 60, 140, 61]}
                        ],
                    }
                )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "input.png"
            output = root / "out.pptx"
            artifact_root = root / "artifacts"
            job_id = "job-stage-events"
            Image.new("RGB", (160, 90), "#001122").save(source)

            run_vlm_editable_pptx_pipeline(
                slides=[{"slide_id": "slide-a", "image_path": str(source)}],
                output_path=str(output),
                artifact_root=str(artifact_root),
                job_id=job_id,
                dependencies=VLMEditablePipelineDependencies(
                    vlm_provider=EmptyVLMProvider(),
                    image_edit_provider=FakeImageEditProvider(
                        ProviderConfig(
                            role="edit_model",
                            provider="fake_image_edit",
                            model="fake-image-edit",
                            base_url="fake",
                            api_key="fake",
                        )
                    ),
                ),
            )
            events = GenerativeEditableJobArtifacts(
                root_dir=artifact_root,
                job_id=job_id,
            ).read_stage_events()

        stages = [(event["stage"], event["status"]) for event in events]
        self.assertIn(("vlm_page_analysis", "started"), stages)
        self.assertIn(("vlm_page_analysis", "finished"), stages)
        self.assertIn(("vlm_clean_background", "started"), stages)
        self.assertIn(("vlm_asset_sheet", "finished"), stages)
        self.assertIn(("compose_deck", "finished"), stages)

    def test_vlm_pipeline_applies_configured_retries_to_image_edit(self):
        from src.generative_editable_providers import ImageEditProvider
        from src.generative_editable_vlm_reconstruction import (
            VLMEditablePipelineDependencies,
            VLMPageAnalysisProvider,
            coerce_vlm_analysis,
            run_vlm_editable_pptx_pipeline,
        )

        class EmptyVLMProvider(VLMPageAnalysisProvider):
            def analyze_page(self, image_path: str, *, timeout_seconds: int = 180):
                return coerce_vlm_analysis(
                    {
                        "coordinate_space": {"width": 160, "height": 90, "unit": "px"},
                        "text_regions": [],
                        "bitmap_regions": [],
                        "shape_regions": [
                            {"id": "line-a", "type": "line", "bbox": [20, 60, 140, 61]}
                        ],
                    }
                )

        class FlakyImageEditProvider(ImageEditProvider):
            def __init__(self):
                super().__init__(
                    ProviderConfig(
                        role="edit_model",
                        provider="fake_image_edit",
                        model="fake-image-edit",
                        base_url="fake",
                        api_key="fake",
                    )
                )
                self.calls = 0
                self.fake = FakeImageEditProvider(self.config)

            def edit(self, request):
                self.calls += 1
                if self.calls == 1:
                    raise ProviderError(
                        provider_role="edit_model",
                        operation=request.prompt_id,
                        message="temporary EOF",
                        retryable=True,
                    )
                return self.fake.edit(request)

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "input.png"
            Image.new("RGB", (160, 90), "#001122").save(source)
            output = root / "out.pptx"
            image_edit_provider = FlakyImageEditProvider()

            run_vlm_editable_pptx_pipeline(
                slides=[{"slide_id": "slide-a", "image_path": str(source)}],
                output_path=str(output),
                artifact_root=str(root / "artifacts"),
                job_id="job-edit-retry",
                dependencies=VLMEditablePipelineDependencies(
                    vlm_provider=EmptyVLMProvider(),
                    image_edit_provider=image_edit_provider,
                    provider_max_attempts=2,
                    provider_retry_backoff_seconds=0,
                ),
            )

        self.assertEqual(image_edit_provider.calls, 2)

    def test_vlm_pipeline_enforces_configured_page_timeout(self):
        from src.generative_editable_providers import ProviderTimeoutError
        from src.generative_editable_vlm_reconstruction import (
            VLMEditablePipelineDependencies,
            VLMPageAnalysisProvider,
            coerce_vlm_analysis,
            run_vlm_editable_pptx_pipeline,
        )

        class SlowVLMProvider(VLMPageAnalysisProvider):
            def analyze_page(self, image_path: str, *, timeout_seconds: int = 180):
                time.sleep(0.2)
                return coerce_vlm_analysis(
                    {
                        "coordinate_space": {"width": 160, "height": 90, "unit": "px"},
                        "text_regions": [],
                        "bitmap_regions": [],
                        "shape_regions": [
                            {"id": "line-a", "type": "line", "bbox": [20, 60, 140, 61]}
                        ],
                    }
                )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "input.png"
            Image.new("RGB", (160, 90), "#001122").save(source)

            with self.assertRaises(ProviderTimeoutError) as ctx:
                run_vlm_editable_pptx_pipeline(
                    slides=[{"slide_id": "slide-a", "image_path": str(source)}],
                    output_path=str(root / "out.pptx"),
                    artifact_root=str(root / "artifacts"),
                    job_id="job-page-timeout",
                    dependencies=VLMEditablePipelineDependencies(
                        vlm_provider=SlowVLMProvider(),
                        image_edit_provider=FakeImageEditProvider(
                            ProviderConfig(
                                role="edit_model",
                                provider="fake_image_edit",
                                model="fake-image-edit",
                                base_url="fake",
                                api_key="fake",
                            )
                        ),
                        page_timeout_seconds=0.05,
                    ),
                )

        self.assertEqual(ctx.exception.provider_role, "vlm_first")
        self.assertEqual(ctx.exception.operation, "build_page_manifest")
        self.assertEqual(ctx.exception.timeout_seconds, 0.05)

    def test_vlm_pipeline_enforces_page_timeout_in_worker_thread(self):
        from src.generative_editable_providers import ProviderTimeoutError
        from src.generative_editable_preview_validator import ValidationReport
        from src.generative_editable_vlm_reconstruction import (
            VLMEditablePipelineDependencies,
            VLMPageAnalysisProvider,
            coerce_vlm_analysis,
            run_vlm_editable_pptx_pipeline,
        )

        class SlowVLMProvider(VLMPageAnalysisProvider):
            def analyze_page(self, image_path: str, *, timeout_seconds: int = 180):
                time.sleep(0.2)
                return coerce_vlm_analysis(
                    {
                        "coordinate_space": {"width": 160, "height": 90, "unit": "px"},
                        "text_regions": [],
                        "bitmap_regions": [],
                        "shape_regions": [
                            {"id": "line-a", "type": "line", "bbox": [20, 60, 140, 61]}
                        ],
                    }
                )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "input.png"
            Image.new("RGB", (160, 90), "#001122").save(source)

            def run_pipeline():
                return run_vlm_editable_pptx_pipeline(
                    slides=[{"slide_id": "slide-a", "image_path": str(source)}],
                    output_path=str(root / "out.pptx"),
                    artifact_root=str(root / "artifacts"),
                    job_id="job-thread-page-timeout",
                    dependencies=VLMEditablePipelineDependencies(
                        vlm_provider=SlowVLMProvider(),
                        image_edit_provider=FakeImageEditProvider(
                            ProviderConfig(
                                role="edit_model",
                                provider="fake_image_edit",
                                model="fake-image-edit",
                                base_url="fake",
                                api_key="fake",
                            )
                        ),
                        page_timeout_seconds=0.05,
                        structure_validator=lambda **kwargs: ValidationReport(
                            status="passed", checked_pages=1, issues=[]
                        ),
                        preview_renderer=lambda page, artifact_root, *, pptx_path: object(),
                        preview_validator=lambda **kwargs: ValidationReport(
                            status="passed", checked_pages=1, issues=[]
                        ),
                    ),
                )

            with ThreadPoolExecutor(max_workers=1) as executor:
                future = executor.submit(run_pipeline)
                with self.assertRaises(ProviderTimeoutError) as ctx:
                    future.result(timeout=1)

        self.assertEqual(ctx.exception.provider_role, "vlm_first")
        self.assertEqual(ctx.exception.operation, "build_page_manifest")

    def test_vlm_pipeline_attributes_provider_timeout_to_page_when_page_budget_is_smaller(self):
        from src.generative_editable_providers import ProviderTimeoutError
        from src.generative_editable_vlm_reconstruction import (
            VLMEditablePipelineDependencies,
            VLMPageAnalysisProvider,
            run_vlm_editable_pptx_pipeline,
        )

        class TimingOutVLMProvider(VLMPageAnalysisProvider):
            def analyze_page(self, image_path: str, *, timeout_seconds: int = 180):
                raise ProviderTimeoutError(
                    provider_role="VLM",
                    operation="vlm_page_analysis",
                    message="request timed out",
                    timeout_seconds=timeout_seconds,
                )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "input.png"
            Image.new("RGB", (160, 90), "#001122").save(source)

            with self.assertRaises(ProviderTimeoutError) as ctx:
                run_vlm_editable_pptx_pipeline(
                    slides=[{"slide_id": "slide-a", "image_path": str(source)}],
                    output_path=str(root / "out.pptx"),
                    artifact_root=str(root / "artifacts"),
                    job_id="job-page-timeout-attribution",
                    dependencies=VLMEditablePipelineDependencies(
                        vlm_provider=TimingOutVLMProvider(),
                        image_edit_provider=FakeImageEditProvider(
                            ProviderConfig(
                                role="edit_model",
                                provider="fake_image_edit",
                                model="fake-image-edit",
                                base_url="fake",
                                api_key="fake",
                            )
                        ),
                        provider_timeout_seconds=30,
                        page_timeout_seconds=0.05,
                    ),
                )

        self.assertEqual(ctx.exception.provider_role, "vlm_first")
        self.assertEqual(ctx.exception.operation, "build_page_manifest")

    def test_maps_vlm_pixel_coordinates_back_to_source_pixels(self):
        from src.generative_editable_vlm_reconstruction import (
            VLMCoordinateMapper,
            coerce_vlm_analysis,
        )

        analysis = coerce_vlm_analysis(
            {
                "coordinate_space": {"width": 960, "height": 540, "unit": "px"},
                "text_regions": [
                    {
                        "id": "t1",
                        "text": "标题",
                        "bbox": [96, 54, 192, 108],
                        "role": "title",
                    }
                ],
                "bitmap_regions": [],
                "shape_regions": [],
            }
        )

        mapper = VLMCoordinateMapper(
            analysis_space=analysis.coordinate_space,
            source_image_size=(1920, 1080),
        )

        self.assertEqual(mapper.to_source_bbox(analysis.text_regions[0].bbox), (192, 108, 384, 216))

    def test_mask_uses_text_and_bitmap_regions_but_excludes_shape_regions(self):
        from src.generative_editable_vlm_reconstruction import (
            VLMCoordinateMapper,
            build_text_bitmap_mask,
            coerce_vlm_analysis,
        )

        analysis = coerce_vlm_analysis(
            {
                "coordinate_space": {"width": 100, "height": 100, "unit": "px"},
                "text_regions": [{"id": "t1", "text": "A", "bbox": [10, 10, 20, 20]}],
                "bitmap_regions": [{"id": "b1", "type": "icon", "bbox": [40, 40, 50, 50]}],
                "shape_regions": [
                    {"id": "s1", "type": "rounded_rect", "bbox": [0, 0, 100, 100]},
                ],
            }
        )
        mapper = VLMCoordinateMapper(
            analysis_space=analysis.coordinate_space,
            source_image_size=(100, 100),
        )

        mask, mask_boxes = build_text_bitmap_mask(analysis, mapper, padding=0)

        self.assertEqual(mask.getpixel((15, 15)), 255)
        self.assertEqual(mask.getpixel((45, 45)), 255)
        self.assertEqual(mask.getpixel((80, 80)), 0)
        self.assertEqual([box.kind for box in mask_boxes], ["text", "bitmap"])

    def test_mask_expands_complex_bitmap_regions_more_than_icon_regions(self):
        from src.generative_editable_vlm_reconstruction import (
            VLMCoordinateMapper,
            build_text_bitmap_mask,
            coerce_vlm_analysis,
        )

        analysis = coerce_vlm_analysis(
            {
                "coordinate_space": {"width": 160, "height": 90, "unit": "px"},
                "text_regions": [],
                "bitmap_regions": [
                    {"id": "product", "type": "product", "bbox": [80, 20, 120, 45]},
                    {"id": "icon", "type": "icon", "bbox": [10, 20, 20, 30]},
                ],
                "shape_regions": [],
            }
        )
        mapper = VLMCoordinateMapper(
            analysis_space=analysis.coordinate_space,
            source_image_size=(160, 90),
        )

        _mask, mask_boxes = build_text_bitmap_mask(analysis, mapper, padding=4)

        by_id = {box.region_id: box.bbox for box in mask_boxes}
        self.assertEqual(by_id["icon"], (6, 16, 24, 34))
        self.assertLessEqual(by_id["product"][0], 64)
        self.assertLessEqual(by_id["product"][1], 4)
        self.assertGreaterEqual(by_id["product"][2], 136)
        self.assertGreaterEqual(by_id["product"][3], 61)

    def test_mask_expansion_does_not_cover_adjacent_text_region(self):
        from src.generative_editable_vlm_reconstruction import (
            VLMCoordinateMapper,
            build_text_bitmap_mask,
            coerce_vlm_analysis,
        )

        analysis = coerce_vlm_analysis(
            {
                "coordinate_space": {"width": 160, "height": 90, "unit": "px"},
                "text_regions": [{"id": "label", "text": "文字", "bbox": [25, 25, 37, 45]}],
                "bitmap_regions": [
                    {"id": "product", "type": "product", "bbox": [40, 20, 80, 50]},
                ],
                "shape_regions": [],
            }
        )
        mapper = VLMCoordinateMapper(
            analysis_space=analysis.coordinate_space,
            source_image_size=(160, 90),
        )

        _mask, mask_boxes = build_text_bitmap_mask(analysis, mapper, padding=4)

        by_id = {box.region_id: box.bbox for box in mask_boxes}
        self.assertGreaterEqual(by_id["product"][0], 38)

    def test_builds_manifest_with_editable_text_split_assets_and_conservative_shapes(self):
        from src.generative_editable_manifest import PageManifest
        from src.generative_editable_vlm_reconstruction import (
            build_page_manifest_from_vlm_analysis,
            coerce_vlm_analysis,
        )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "sources" / "0000-slide-a" / "source.png"
            background = root / "backgrounds" / "0000-slide-a" / "clean.png"
            source.parent.mkdir(parents=True)
            background.parent.mkdir(parents=True)
            image = Image.new("RGB", (200, 100), "#001122")
            draw = ImageDraw.Draw(image)
            draw.rectangle((20, 20, 50, 40), fill="#FFFFFF")
            draw.rectangle((130, 30, 140, 40), fill="#22AAFF")
            draw.point((135, 35), fill="#001122")
            image.save(source)
            Image.new("RGB", (200, 100), "#001122").save(background)

            analysis = coerce_vlm_analysis(
                {
                    "coordinate_space": {"width": 100, "height": 50, "unit": "px"},
                    "text_regions": [
                        {
                            "id": "t1",
                            "text": "标题",
                            "bbox": [10, 10, 25, 20],
                            "role": "title",
                            "color": "#FFFFFF",
                        }
                    ],
                    "bitmap_regions": [
                        {"id": "b1", "type": "icon", "bbox": [60, 10, 75, 25]},
                    ],
                    "shape_regions": [
                        {"id": "panel", "type": "rounded_rect", "bbox": [0, 0, 100, 50]},
                        {"id": "divider", "type": "divider", "bbox": [10, 30, 90, 31]},
                    ],
                }
            )

            page = build_page_manifest_from_vlm_analysis(
                analysis=analysis,
                slide_id="slide-a",
                page_index=0,
                source_image_path=source,
                clean_background_path=background,
                artifact_root=root,
                aspect_ratio="16:9",
            )
            asset_path = root / page.bitmap_assets[0].asset_path
            asset = Image.open(asset_path).convert("RGBA")
            transparent_alpha = asset.getpixel((1, 1))[3]
            foreground_alpha = asset.getpixel((14, 15))[3]
            enclosed_detail_alpha = asset.getpixel((15, 15))[3]

        self.assertIsInstance(page, PageManifest)
        self.assertEqual(page.chosen_background, "backgrounds/0000-slide-a/clean.png")
        self.assertEqual([text.text for text in page.text_boxes], ["标题"])
        self.assertEqual(page.text_boxes[0].source_pixel_bbox, (20, 20, 50, 40))
        self.assertTrue(page.text_boxes[0].style_hints["bold"])
        self.assertEqual(len(page.bitmap_assets), 1)
        self.assertEqual(page.bitmap_assets[0].source_pixel_bbox, (120, 20, 150, 50))
        self.assertEqual(transparent_alpha, 0)
        self.assertEqual(foreground_alpha, 255)
        self.assertEqual(enclosed_detail_alpha, 255)
        self.assertEqual(len(page.native_shapes), 1)
        self.assertEqual(page.native_shapes[0].shape_type, "line")
        self.assertEqual(page.provenance["reconstruction_strategy"], "vlm_first")

    def test_selects_source_preserved_crop_when_generated_asset_drifts_from_source(self):
        from src.generative_editable_manifest import BitmapAssetSpec
        from src.generative_editable_vlm_reconstruction import _select_source_faithful_bitmap_assets

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "sources" / "0000-slide-a" / "source.png"
            background = root / "backgrounds" / "0000-slide-a" / "clean.png"
            generated_path = root / "assets" / "0000-slide-a" / "asset.png"
            source.parent.mkdir(parents=True)
            background.parent.mkdir(parents=True)
            generated_path.parent.mkdir(parents=True)
            source_image = Image.new("RGB", (120, 80), "#001122")
            ImageDraw.Draw(source_image).rectangle((30, 20, 70, 50), fill="#22AAFF")
            source_image.save(source)
            Image.new("RGB", (120, 80), "#001122").save(background)
            generated = Image.new("RGBA", (40, 30), (0, 0, 0, 0))
            ImageDraw.Draw(generated).rectangle((0, 0, 39, 29), fill="#FF0000")
            generated.save(generated_path)
            asset = BitmapAssetSpec(
                asset_id="b1",
                source_pixel_bbox=(30, 20, 70, 50),
                asset_path="assets/0000-slide-a/asset.png",
                z_order=1,
                provenance={"split_method": "connected_components"},
            )

            selected = _select_source_faithful_bitmap_assets(
                generated_assets=[asset],
                source_image_path=source,
                clean_background_path=background,
                artifact_root=root,
                slide_id="slide-a",
                page_index=0,
            )
            selected_path = root / selected[0].asset_path
            selected_pixel = Image.open(selected_path).convert("RGBA").getpixel((20, 15))

        self.assertEqual(selected[0].provenance["asset_strategy"], "source_preserved_crop")
        self.assertEqual(selected[0].provenance["original_source_pixel_bbox"], [30, 20, 70, 50])
        self.assertEqual(selected[0].source_pixel_bbox, (22, 12, 78, 58))
        self.assertEqual(selected_pixel[:3], (34, 170, 255))

    def test_complex_source_preserved_assets_use_opaque_crop_not_background_diff_alpha(self):
        from src.generative_editable_foreground_planner import ForegroundCandidate
        from src.generative_editable_vlm_reconstruction import (
            _source_preserved_bitmap_assets_from_candidates,
        )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "sources" / "0000-slide-a" / "source.png"
            background = root / "backgrounds" / "0000-slide-a" / "clean.png"
            source.parent.mkdir(parents=True)
            background.parent.mkdir(parents=True)
            source_image = Image.new("RGB", (120, 80), "#001122")
            draw = ImageDraw.Draw(source_image)
            draw.rectangle((30, 18, 105, 55), fill="#0B2A66")
            draw.ellipse((36, 46, 52, 62), fill="#000000")
            source_image.save(source)
            Image.new("RGB", (120, 80), "#001122").save(background)
            candidate = ForegroundCandidate(
                candidate_id="vehicle",
                source_pixel_bbox=(30, 18, 105, 62),
                area=3300,
                classification="complex_whole_visual",
                provenance={"vlm_type": "photo"},
            )

            assets = _source_preserved_bitmap_assets_from_candidates(
                candidates=[candidate],
                source_image_path=source,
                clean_background_path=background,
                artifact_root=root,
                slide_id="slide-a",
                page_index=0,
                reason="complex_bitmap_region",
                asset_sheet_skipped_reason="complex_bitmap_region",
            )
            crop = Image.open(root / assets[0].asset_path).convert("RGBA")

        self.assertEqual(crop.getpixel((0, 0))[3], 255)
        self.assertEqual(crop.getpixel((crop.width - 1, crop.height - 1))[3], 255)
        self.assertEqual(assets[0].provenance["alpha_strategy"], "opaque_source_crop")
        self.assertNotIn("background_difference_alpha", assets[0].provenance)

    def test_foreground_crop_preserves_edge_foreground_close_to_background(self):
        from src.generative_editable_vlm_reconstruction import _foreground_rgba_crop

        source = Image.new("RGB", (20, 12), "#001122")
        background = Image.new("RGB", (20, 12), "#001122")
        ImageDraw.Draw(source).rectangle((0, 3, 8, 9), fill="#102132")

        crop = _foreground_rgba_crop(source, background, (0, 0, 12, 12))

        self.assertEqual(crop.getpixel((0, 6))[3], 255)

    def test_foreground_crop_can_mask_overlapping_text_regions(self):
        from src.generative_editable_vlm_reconstruction import _foreground_rgba_crop

        source = Image.new("RGB", (40, 24), "#001122")
        background = Image.new("RGB", (40, 24), "#001122")
        draw = ImageDraw.Draw(source)
        draw.rectangle((4, 6, 15, 18), fill="#FFFFFF")
        draw.rectangle((26, 6, 34, 18), fill="#22AAFF")

        crop = _foreground_rgba_crop(
            source,
            background,
            (0, 0, 40, 24),
            transparent_bboxes=[(0, 0, 20, 24)],
        )

        self.assertEqual(crop.getpixel((8, 12))[3], 0)
        self.assertEqual(crop.getpixel((30, 12))[3], 255)

    def test_refines_vlm_bitmap_candidate_bbox_with_clean_background_difference(self):
        from src.generative_editable_foreground_planner import ForegroundCandidate
        from src.generative_editable_vlm_reconstruction import _refine_candidates_with_clean_background_difference

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "source.png"
            background = root / "background.png"
            source_image = Image.new("RGB", (160, 100), "#001122")
            ImageDraw.Draw(source_image).rectangle((44, 28, 104, 62), fill="#22AAFF")
            source_image.save(source)
            Image.new("RGB", (160, 100), "#001122").save(background)
            candidate = ForegroundCandidate(
                candidate_id="b1",
                source_pixel_bbox=(55, 34, 92, 55),
                area=777,
                classification="bitmap_asset_candidate",
            )

            refined = _refine_candidates_with_clean_background_difference(
                candidates=[candidate],
                source_image_path=source,
                clean_background_path=background,
            )

        self.assertEqual(refined[0].source_pixel_bbox, (44, 28, 105, 63))
        self.assertEqual(refined[0].provenance["original_vlm_bbox"], [55, 34, 92, 55])

    def test_bridges_narrow_complex_candidate_with_same_group_row_neighbor(self):
        from src.generative_editable_foreground_planner import ForegroundCandidate
        from src.generative_editable_vlm_reconstruction import (
            _bridge_narrow_same_group_complex_candidates,
        )

        candidates = [
            ForegroundCandidate(
                candidate_id="left",
                source_pixel_bbox=(457, 393, 882, 621),
                area=96900,
                classification="complex_whole_visual",
                provenance={"group_id": "g_middle"},
            ),
            ForegroundCandidate(
                candidate_id="right",
                source_pixel_bbox=(976, 418, 1180, 605),
                area=38148,
                classification="complex_whole_visual",
                provenance={"group_id": "g_middle"},
            ),
            ForegroundCandidate(
                candidate_id="bottom",
                source_pixel_bbox=(848, 649, 1173, 841),
                area=62400,
                classification="complex_whole_visual",
                provenance={"group_id": "g_bottom"},
            ),
        ]

        bridged = _bridge_narrow_same_group_complex_candidates(candidates)

        self.assertLessEqual(bridged[1].source_pixel_bbox[0], 856)
        self.assertEqual(bridged[1].source_pixel_bbox[1:], (418, 1180, 605))
        self.assertEqual(bridged[2].source_pixel_bbox, (848, 649, 1173, 841))
        self.assertEqual(bridged[1].provenance["bbox_refinement"], "same_group_row_bridge")

    def test_vlm_text_regions_gate_ocr_boxes_to_prevent_stray_text(self):
        from src.generative_editable_manifest import TextBoxSpec
        from src.generative_editable_vlm_reconstruction import (
            build_page_manifest_from_vlm_analysis,
            coerce_vlm_analysis,
        )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "sources" / "0000-slide-a" / "source.png"
            background = root / "backgrounds" / "0000-slide-a" / "clean.png"
            source.parent.mkdir(parents=True)
            background.parent.mkdir(parents=True)
            Image.new("RGB", (200, 100), "#001122").save(source)
            Image.new("RGB", (200, 100), "#001122").save(background)
            analysis = coerce_vlm_analysis(
                {
                    "coordinate_space": {"width": 200, "height": 100, "unit": "px"},
                    "text_regions": [
                        {"id": "t1", "text": "OCR 粗文", "bbox": [20, 20, 90, 38]},
                    ],
                    "bitmap_regions": [
                        {"id": "b1", "type": "icon", "bbox": [8, 18, 45, 40]},
                    ],
                    "shape_regions": [],
                }
            )
            matching_ocr = TextBoxSpec(
                text="OCR 正文",
                source_pixel_bbox=(8, 18, 92, 40),
                source_pixel_polygon=((8, 18), (92, 18), (92, 40), (8, 40)),
                provenance={"content_source": "ocr", "ocr_confidence": 0.99},
            )
            stray_ocr = TextBoxSpec(
                text="串位文本",
                source_pixel_bbox=(120, 70, 180, 88),
                source_pixel_polygon=((120, 70), (180, 70), (180, 88), (120, 88)),
                provenance={"content_source": "ocr", "ocr_confidence": 0.99},
            )

            page = build_page_manifest_from_vlm_analysis(
                analysis=analysis,
                slide_id="slide-a",
                page_index=0,
                source_image_path=source,
                clean_background_path=background,
                artifact_root=root,
                aspect_ratio="16:9",
                text_boxes=[matching_ocr, stray_ocr],
            )

        self.assertEqual([box.text for box in page.text_boxes], ["OCR 正文"])
        self.assertEqual(page.text_boxes[0].provenance["content_source"], "ocr")
        self.assertEqual(page.text_boxes[0].source_pixel_bbox, (53, 20, 123, 38))
        self.assertTrue(page.text_boxes[0].style_hints["approximate_layout"])

    def test_bitmap_overlap_avoidance_does_not_shift_minor_edge_overlap(self):
        from src.generative_editable_manifest import BitmapAssetSpec, TextBoxSpec
        from src.generative_editable_vlm_reconstruction import _avoid_bitmap_text_overlap

        text_box = TextBoxSpec(
            text="Livis版800V主动悬架",
            source_pixel_bbox=(1292, 540, 1502, 563),
            source_pixel_polygon=((1292, 540), (1502, 540), (1502, 563), (1292, 563)),
            provenance={"content_source": "ocr", "layout_source": "vlm"},
        )
        asset = BitmapAssetSpec(
            asset_id="component",
            source_pixel_bbox=(1004, 418, 1312, 605),
            asset_path="assets/component.png",
            z_order=1,
            provenance={"asset_strategy": "source_preserved_crop"},
        )

        adjusted = _avoid_bitmap_text_overlap([text_box], [asset], (1600, 900))

        self.assertEqual(adjusted[0].source_pixel_bbox, (1292, 540, 1502, 563))

    def test_vlm_text_region_can_keep_multiple_ocr_lines_without_stray_text(self):
        from src.generative_editable_manifest import TextBoxSpec
        from src.generative_editable_vlm_reconstruction import (
            build_page_manifest_from_vlm_analysis,
            coerce_vlm_analysis,
        )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "sources" / "0000-slide-a" / "source.png"
            background = root / "backgrounds" / "0000-slide-a" / "clean.png"
            source.parent.mkdir(parents=True)
            background.parent.mkdir(parents=True)
            Image.new("RGB", (240, 120), "#001122").save(source)
            Image.new("RGB", (240, 120), "#001122").save(background)
            analysis = coerce_vlm_analysis(
                {
                    "coordinate_space": {"width": 240, "height": 120, "unit": "px"},
                    "text_regions": [
                        {"id": "t1", "text": "第一行 第二行", "bbox": [20, 20, 150, 70]},
                    ],
                    "bitmap_regions": [],
                    "shape_regions": [],
                }
            )
            first_line = TextBoxSpec(
                text="第一行",
                source_pixel_bbox=(22, 22, 90, 40),
                source_pixel_polygon=((22, 22), (90, 22), (90, 40), (22, 40)),
                provenance={"content_source": "ocr", "ocr_confidence": 0.99},
            )
            second_line = TextBoxSpec(
                text="第二行",
                source_pixel_bbox=(22, 48, 90, 66),
                source_pixel_polygon=((22, 48), (90, 48), (90, 66), (22, 66)),
                provenance={"content_source": "ocr", "ocr_confidence": 0.99},
            )
            stray_line = TextBoxSpec(
                text="串位",
                source_pixel_bbox=(180, 90, 220, 108),
                source_pixel_polygon=((180, 90), (220, 90), (220, 108), (180, 108)),
                provenance={"content_source": "ocr", "ocr_confidence": 0.99},
            )

            page = build_page_manifest_from_vlm_analysis(
                analysis=analysis,
                slide_id="slide-a",
                page_index=0,
                source_image_path=source,
                clean_background_path=background,
                artifact_root=root,
                aspect_ratio="16:9",
                text_boxes=[first_line, second_line, stray_line],
            )

        self.assertEqual([box.text for box in page.text_boxes], ["第一行", "第二行"])
        self.assertEqual([box.provenance["layout_source"] for box in page.text_boxes], ["ocr", "ocr"])

    def test_vlm_text_resolution_drops_overlapping_duplicate_text(self):
        from src.generative_editable_manifest import TextBoxSpec
        from src.generative_editable_vlm_reconstruction import (
            build_page_manifest_from_vlm_analysis,
            coerce_vlm_analysis,
        )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "sources" / "0000-slide-a" / "source.png"
            background = root / "backgrounds" / "0000-slide-a" / "clean.png"
            source.parent.mkdir(parents=True)
            background.parent.mkdir(parents=True)
            Image.new("RGB", (320, 180), "#001122").save(source)
            Image.new("RGB", (320, 180), "#001122").save(background)
            analysis = coerce_vlm_analysis(
                {
                    "coordinate_space": {"width": 320, "height": 180, "unit": "px"},
                    "text_regions": [
                        {"id": "t1", "text": "21英寸4K后舱屏", "bbox": [220, 46, 298, 66]},
                        {"id": "t2", "text": "21 英寸 4K 后舱屏", "bbox": [218, 50, 302, 72]},
                    ],
                    "bitmap_regions": [],
                    "shape_regions": [],
                }
            )
            ocr_line = TextBoxSpec(
                text="21英寸4K后舱屏",
                source_pixel_bbox=(220, 46, 298, 66),
                source_pixel_polygon=((220, 46), (298, 46), (298, 66), (220, 66)),
                font_size=10,
                provenance={"content_source": "ocr", "ocr_confidence": 0.99},
            )

            page = build_page_manifest_from_vlm_analysis(
                analysis=analysis,
                slide_id="slide-a",
                page_index=0,
                source_image_path=source,
                clean_background_path=background,
                artifact_root=root,
                aspect_ratio="16:9",
                text_boxes=[ocr_line],
            )

        self.assertEqual([box.text for box in page.text_boxes], ["21英寸4K后舱屏"])
        self.assertEqual(page.text_boxes[0].provenance["content_source"], "ocr")

    def test_vlm_text_resolution_drops_nearby_ocr_vlm_duplicate_text(self):
        from src.generative_editable_manifest import TextBoxSpec
        from src.generative_editable_vlm_reconstruction import (
            build_page_manifest_from_vlm_analysis,
            coerce_vlm_analysis,
        )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "sources" / "0000-slide-a" / "source.png"
            background = root / "backgrounds" / "0000-slide-a" / "clean.png"
            source.parent.mkdir(parents=True)
            background.parent.mkdir(parents=True)
            Image.new("RGB", (1280, 720), "#001122").save(source)
            Image.new("RGB", (1280, 720), "#001122").save(background)
            analysis = coerce_vlm_analysis(
                {
                    "coordinate_space": {"width": 1280, "height": 720, "unit": "px"},
                    "text_regions": [
                        {"id": "title-vlm", "text": "理想L9:", "bbox": [67, 168, 433, 253]},
                        {
                            "id": "subtitle",
                            "text": "旗舰增程SUV的技术实验",
                            "bbox": [67, 279, 934, 362],
                        },
                    ],
                    "bitmap_regions": [],
                    "shape_regions": [],
                }
            )
            ocr_title = TextBoxSpec(
                text="理想L9:",
                source_pixel_bbox=(85, 70, 298, 94),
                source_pixel_polygon=((85, 70), (298, 70), (298, 94), (85, 94)),
                font_size=18,
                provenance={"content_source": "ocr", "layout_source": "ocr", "ocr_confidence": 0.78},
            )

            page = build_page_manifest_from_vlm_analysis(
                analysis=analysis,
                slide_id="slide-a",
                page_index=0,
                source_image_path=source,
                clean_background_path=background,
                artifact_root=root,
                aspect_ratio="16:9",
                text_boxes=[ocr_title],
            )

        self.assertEqual([box.text for box in page.text_boxes], ["理想L9:", "旗舰增程SUV的技术实验"])
        self.assertEqual(page.text_boxes[0].provenance["content_source"], "ocr")

    def test_vlm_text_resolution_drops_nearby_duplicate_after_ocr_match(self):
        from src.generative_editable_manifest import TextBoxSpec
        from src.generative_editable_vlm_reconstruction import (
            build_page_manifest_from_vlm_analysis,
            coerce_vlm_analysis,
        )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "sources" / "0000-slide-a" / "source.png"
            background = root / "backgrounds" / "0000-slide-a" / "clean.png"
            source.parent.mkdir(parents=True)
            background.parent.mkdir(parents=True)
            Image.new("RGB", (1280, 720), "#001122").save(source)
            Image.new("RGB", (1280, 720), "#001122").save(background)
            analysis = coerce_vlm_analysis(
                {
                    "coordinate_space": {"width": 1280, "height": 720, "unit": "px"},
                    "text_regions": [
                        {"id": "title-ocr-layout", "text": "理想L9:", "bbox": [85, 70, 296, 94]},
                        {"id": "title-vlm-duplicate", "text": "理想L9:", "bbox": [67, 168, 433, 250]},
                        {
                            "id": "subtitle",
                            "text": "旗舰增程SUV的技术实验",
                            "bbox": [67, 279, 934, 362],
                        },
                    ],
                    "bitmap_regions": [],
                    "shape_regions": [],
                }
            )
            ocr_title = TextBoxSpec(
                text="理想L9:",
                source_pixel_bbox=(85, 70, 296, 94),
                source_pixel_polygon=((85, 70), (296, 70), (296, 94), (85, 94)),
                font_size=18,
                provenance={"content_source": "ocr", "layout_source": "ocr", "ocr_confidence": 0.78},
            )

            page = build_page_manifest_from_vlm_analysis(
                analysis=analysis,
                slide_id="slide-a",
                page_index=0,
                source_image_path=source,
                clean_background_path=background,
                artifact_root=root,
                aspect_ratio="16:9",
                text_boxes=[ocr_title],
            )

        self.assertEqual([box.text for box in page.text_boxes], ["理想L9:", "旗舰增程SUV的技术实验"])
        self.assertEqual(page.text_boxes[0].provenance["content_source"], "ocr")

    def test_vlm_text_resolution_keeps_real_adjacent_duplicate_text(self):
        from src.generative_editable_manifest import TextBoxSpec
        from src.generative_editable_vlm_reconstruction import (
            build_page_manifest_from_vlm_analysis,
            coerce_vlm_analysis,
        )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "sources" / "0000-slide-a" / "source.png"
            background = root / "backgrounds" / "0000-slide-a" / "clean.png"
            source.parent.mkdir(parents=True)
            background.parent.mkdir(parents=True)
            Image.new("RGB", (400, 240), "#001122").save(source)
            Image.new("RGB", (400, 240), "#001122").save(background)
            analysis = coerce_vlm_analysis(
                {
                    "coordinate_space": {"width": 400, "height": 240, "unit": "px"},
                    "text_regions": [
                        {"id": "status-row-1", "text": "Pending", "bbox": [80, 80, 220, 106]},
                        {"id": "status-row-2", "text": "Pending", "bbox": [80, 125, 220, 151]},
                    ],
                    "bitmap_regions": [],
                    "shape_regions": [],
                }
            )
            ocr_row = TextBoxSpec(
                text="Pending",
                source_pixel_bbox=(80, 80, 220, 106),
                source_pixel_polygon=((80, 80), (220, 80), (220, 106), (80, 106)),
                font_size=12,
                provenance={"content_source": "ocr", "layout_source": "ocr", "ocr_confidence": 0.9},
            )

            page = build_page_manifest_from_vlm_analysis(
                analysis=analysis,
                slide_id="slide-a",
                page_index=0,
                source_image_path=source,
                clean_background_path=background,
                artifact_root=root,
                aspect_ratio="16:9",
                text_boxes=[ocr_row],
            )

        self.assertEqual([box.text for box in page.text_boxes], ["Pending", "Pending"])

    def test_vlm_text_resolution_drops_low_confidence_label_near_product_bitmap(self):
        from src.generative_editable_vlm_reconstruction import (
            build_page_manifest_from_vlm_analysis,
            coerce_vlm_analysis,
        )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "sources" / "0000-slide-a" / "source.png"
            background = root / "backgrounds" / "0000-slide-a" / "clean.png"
            source.parent.mkdir(parents=True)
            background.parent.mkdir(parents=True)
            Image.new("RGB", (1600, 900), "#001122").save(source)
            Image.new("RGB", (1600, 900), "#001122").save(background)
            analysis = coerce_vlm_analysis(
                {
                    "coordinate_space": {"width": 1600, "height": 900, "unit": "px"},
                    "text_regions": [
                        {
                            "id": "license",
                            "text": "理想L9",
                            "bbox": [702, 708, 783, 730],
                            "role": "label",
                            "confidence": 0.72,
                        },
                        {
                            "id": "caption",
                            "text": "官方帮助中心 OTA更新说明",
                            "bbox": [1424, 882, 1672, 908],
                            "role": "label",
                            "confidence": 0.89,
                        },
                    ],
                    "bitmap_regions": [
                        {
                            "id": "car",
                            "type": "product",
                            "bbox": [0, 523, 685, 927],
                            "importance": "major",
                        },
                        {
                            "id": "qr",
                            "type": "qr",
                            "bbox": [1369, 643, 1634, 905],
                            "importance": "major",
                        },
                    ],
                    "shape_regions": [],
                }
            )

            page = build_page_manifest_from_vlm_analysis(
                analysis=analysis,
                slide_id="slide-a",
                page_index=0,
                source_image_path=source,
                clean_background_path=background,
                artifact_root=root,
                aspect_ratio="16:9",
                text_boxes=[],
            )

        self.assertEqual([box.text for box in page.text_boxes], ["官方帮助中心 OTA更新说明"])

    def test_vlm_text_font_size_leaves_renderer_headroom_for_approximate_bboxes(self):
        from src.generative_editable_vlm_reconstruction import (
            build_page_manifest_from_vlm_analysis,
            coerce_vlm_analysis,
        )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "sources" / "0000-slide-a" / "source.png"
            background = root / "backgrounds" / "0000-slide-a" / "clean.png"
            source.parent.mkdir(parents=True)
            background.parent.mkdir(parents=True)
            Image.new("RGB", (1672, 941), "#001122").save(source)
            Image.new("RGB", (1672, 941), "#001122").save(background)
            analysis = coerce_vlm_analysis(
                {
                    "coordinate_space": {"width": 1672, "height": 941, "unit": "px"},
                    "text_regions": [
                        {
                            "id": "dense-body",
                            "text": "连接充电枪不可启动车辆",
                            "bbox": [160, 820, 504, 857],
                            "role": "body",
                            "confidence": 0.94,
                        }
                    ],
                    "bitmap_regions": [],
                    "shape_regions": [],
                }
            )

            page = build_page_manifest_from_vlm_analysis(
                analysis=analysis,
                slide_id="slide-a",
                page_index=0,
                source_image_path=source,
                clean_background_path=background,
                artifact_root=root,
                aspect_ratio="16:9",
            )

        self.assertLessEqual(page.text_boxes[0].font_size or 0, 10.6)
        self.assertTrue(page.text_boxes[0].style_hints["approximate_layout"])

    def test_vlm_title_font_size_preserves_large_visual_hierarchy(self):
        from src.generative_editable_vlm_reconstruction import (
            build_page_manifest_from_vlm_analysis,
            coerce_vlm_analysis,
        )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "sources" / "0000-slide-a" / "source.png"
            background = root / "backgrounds" / "0000-slide-a" / "clean.png"
            source.parent.mkdir(parents=True)
            background.parent.mkdir(parents=True)
            Image.new("RGB", (1672, 941), "#001122").save(source)
            Image.new("RGB", (1672, 941), "#001122").save(background)
            analysis = coerce_vlm_analysis(
                {
                    "coordinate_space": {"width": 1672, "height": 941, "unit": "px"},
                    "text_regions": [
                        {
                            "id": "hero-title",
                            "text": "旗舰增程SUV的技术实验",
                            "bbox": [67, 283, 923, 368],
                            "role": "title",
                            "confidence": 0.95,
                        }
                    ],
                    "bitmap_regions": [],
                    "shape_regions": [],
                }
            )

            page = build_page_manifest_from_vlm_analysis(
                analysis=analysis,
                slide_id="slide-a",
                page_index=0,
                source_image_path=source,
                clean_background_path=background,
                artifact_root=root,
                aspect_ratio="16:9",
            )

        self.assertGreaterEqual(page.text_boxes[0].font_size or 0, 32.0)
        self.assertTrue(page.text_boxes[0].style_hints["bold"])

    def test_vlm_heading_font_size_is_capped_by_text_width_density(self):
        from src.generative_editable_vlm_reconstruction import (
            build_page_manifest_from_vlm_analysis,
            coerce_vlm_analysis,
        )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "sources" / "0000-slide-a" / "source.png"
            background = root / "backgrounds" / "0000-slide-a" / "clean.png"
            source.parent.mkdir(parents=True)
            background.parent.mkdir(parents=True)
            Image.new("RGB", (1672, 941), "#001122").save(source)
            Image.new("RGB", (1672, 941), "#001122").save(background)
            analysis = coerce_vlm_analysis(
                {
                    "coordinate_space": {"width": 1672, "height": 941, "unit": "px"},
                    "text_regions": [
                        {
                            "id": "narrow-heading",
                            "text": "传统电动车\n长途焦虑",
                            "bbox": [104, 196, 224, 266],
                            "role": "heading",
                            "confidence": 0.88,
                        }
                    ],
                    "bitmap_regions": [],
                    "shape_regions": [],
                }
            )

            page = build_page_manifest_from_vlm_analysis(
                analysis=analysis,
                slide_id="slide-a",
                page_index=0,
                source_image_path=source,
                clean_background_path=background,
                artifact_root=root,
                aspect_ratio="16:9",
            )

        self.assertLessEqual(page.text_boxes[0].font_size or 0, 16.0)
        self.assertTrue(page.text_boxes[0].style_hints["bold"])

    def test_vlm_body_font_size_is_capped_by_text_width_density(self):
        from src.generative_editable_vlm_reconstruction import (
            build_page_manifest_from_vlm_analysis,
            coerce_vlm_analysis,
        )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "sources" / "0000-slide-a" / "source.png"
            background = root / "backgrounds" / "0000-slide-a" / "clean.png"
            source.parent.mkdir(parents=True)
            background.parent.mkdir(parents=True)
            Image.new("RGB", (1672, 941), "#001122").save(source)
            Image.new("RGB", (1672, 941), "#001122").save(background)
            analysis = coerce_vlm_analysis(
                {
                    "coordinate_space": {"width": 1672, "height": 941, "unit": "px"},
                    "text_regions": [
                        {
                            "id": "solution-body",
                            "text": "四零重力座椅，\n全家舒适出行",
                            "bbox": [409, 579, 594, 640],
                            "role": "body",
                            "confidence": 0.9,
                        }
                    ],
                    "bitmap_regions": [],
                    "shape_regions": [],
                }
            )

            page = build_page_manifest_from_vlm_analysis(
                analysis=analysis,
                slide_id="slide-a",
                page_index=0,
                source_image_path=source,
                clean_background_path=background,
                artifact_root=root,
                aspect_ratio="16:9",
            )

        self.assertLessEqual(page.text_boxes[0].font_size or 0, 13.0)
        self.assertFalse(page.text_boxes[0].style_hints["bold"])

    def test_vlm_page_manifest_marks_clean_background_as_text_clean(self):
        from src.generative_editable_vlm_reconstruction import (
            build_page_manifest_from_vlm_analysis,
            coerce_vlm_analysis,
        )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "sources" / "0000-slide-a" / "source.png"
            background = root / "backgrounds" / "0000-slide-a" / "clean.png"
            source.parent.mkdir(parents=True)
            background.parent.mkdir(parents=True)
            Image.new("RGB", (800, 450), "#001122").save(source)
            Image.new("RGB", (800, 450), "#001122").save(background)
            analysis = coerce_vlm_analysis(
                {
                    "coordinate_space": {"width": 800, "height": 450, "unit": "px"},
                    "text_regions": [
                        {
                            "id": "title",
                            "text": "Editable Title",
                            "bbox": [64, 36, 384, 81],
                            "role": "title",
                            "confidence": 1.0,
                        }
                    ],
                    "bitmap_regions": [],
                    "shape_regions": [],
                }
            )

            page = build_page_manifest_from_vlm_analysis(
                analysis=analysis,
                slide_id="slide-a",
                page_index=0,
                source_image_path=source,
                clean_background_path=background,
                artifact_root=root,
                aspect_ratio="16:9",
            )

        self.assertEqual(page.chosen_background, "backgrounds/0000-slide-a/clean.png")
        self.assertEqual(page.text_clean_background, page.chosen_background)
        self.assertEqual(page.base_clean_background, page.chosen_background)
        self.assertEqual(page.provenance["chosen_background_kind"], "source_preserving_text_clean")

    def test_vlm_text_resolution_keeps_high_confidence_vlm_text_over_mismatched_ocr(self):
        from src.generative_editable_manifest import TextBoxSpec
        from src.generative_editable_vlm_reconstruction import (
            build_page_manifest_from_vlm_analysis,
            coerce_vlm_analysis,
        )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "sources" / "0000-slide-a" / "source.png"
            background = root / "backgrounds" / "0000-slide-a" / "clean.png"
            source.parent.mkdir(parents=True)
            background.parent.mkdir(parents=True)
            Image.new("RGB", (1672, 941), "#001122").save(source)
            Image.new("RGB", (1672, 941), "#001122").save(background)
            analysis = coerce_vlm_analysis(
                {
                    "coordinate_space": {"width": 1672, "height": 941, "unit": "px"},
                    "text_regions": [
                        {
                            "id": "bottom-right",
                            "text": "保持视野清晰",
                            "bbox": [1419, 808, 1587, 843],
                            "role": "body",
                            "confidence": 0.95,
                        }
                    ],
                    "bitmap_regions": [],
                    "shape_regions": [],
                }
            )
            overlapping_wrong_ocr = TextBoxSpec(
                text="保持50-80%电量",
                source_pixel_bbox=(1280, 797, 1593, 867),
                source_pixel_polygon=((1280, 797), (1593, 797), (1593, 867), (1280, 867)),
                font_size=18,
                provenance={"content_source": "ocr", "layout_source": "ocr", "ocr_confidence": 0.78},
            )

            page = build_page_manifest_from_vlm_analysis(
                analysis=analysis,
                slide_id="slide-a",
                page_index=0,
                source_image_path=source,
                clean_background_path=background,
                artifact_root=root,
                aspect_ratio="16:9",
                text_boxes=[overlapping_wrong_ocr],
            )

        self.assertEqual([box.text for box in page.text_boxes], ["保持视野清晰"])
        self.assertEqual(page.text_boxes[0].provenance["provider_role"], "VLM")

    def test_vlm_text_resolution_keeps_medium_confidence_vlm_text_over_short_ocr_fragment(self):
        from src.generative_editable_manifest import TextBoxSpec
        from src.generative_editable_vlm_reconstruction import (
            build_page_manifest_from_vlm_analysis,
            coerce_vlm_analysis,
        )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "sources" / "0000-slide-a" / "source.png"
            background = root / "backgrounds" / "0000-slide-a" / "clean.png"
            source.parent.mkdir(parents=True)
            background.parent.mkdir(parents=True)
            Image.new("RGB", (160, 90), "#001122").save(source)
            Image.new("RGB", (160, 90), "#001122").save(background)
            analysis = coerce_vlm_analysis(
                {
                    "coordinate_space": {"width": 160, "height": 90, "unit": "px"},
                    "text_regions": [
                        {
                            "id": "range-anxiety",
                            "text": "续航短,\n频繁充电",
                            "bbox": [40, 24, 82, 48],
                            "role": "body",
                            "confidence": 0.86,
                        }
                    ],
                    "bitmap_regions": [],
                    "shape_regions": [],
                }
            )
            fragment_ocr = TextBoxSpec(
                text="充电慢，",
                source_pixel_bbox=(50, 30, 70, 38),
                source_pixel_polygon=((50, 30), (70, 30), (70, 38), (50, 38)),
                font_size=8,
                provenance={"content_source": "ocr", "layout_source": "ocr", "ocr_confidence": 0.78},
            )

            page = build_page_manifest_from_vlm_analysis(
                analysis=analysis,
                slide_id="slide-a",
                page_index=0,
                source_image_path=source,
                clean_background_path=background,
                artifact_root=root,
                aspect_ratio="16:9",
                text_boxes=[fragment_ocr],
            )

        self.assertEqual([box.text for box in page.text_boxes], ["续航短,\n频繁充电"])
        self.assertEqual(page.text_boxes[0].provenance["provider_role"], "VLM")

    def test_vlm_text_resolution_keeps_medium_confidence_vlm_text_over_substring_ocr_fragment(self):
        from src.generative_editable_manifest import TextBoxSpec
        from src.generative_editable_vlm_reconstruction import (
            build_page_manifest_from_vlm_analysis,
            coerce_vlm_analysis,
        )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "sources" / "0000-slide-a" / "source.png"
            background = root / "backgrounds" / "0000-slide-a" / "clean.png"
            source.parent.mkdir(parents=True)
            background.parent.mkdir(parents=True)
            Image.new("RGB", (320, 180), "#001122").save(source)
            Image.new("RGB", (320, 180), "#001122").save(background)
            analysis = coerce_vlm_analysis(
                {
                    "coordinate_space": {"width": 320, "height": 180, "unit": "px"},
                    "text_regions": [
                        {
                            "id": "range-solution",
                            "text": "纯电续航里程长",
                            "bbox": [40, 40, 170, 70],
                            "role": "body",
                            "confidence": 0.86,
                        }
                    ],
                    "bitmap_regions": [],
                    "shape_regions": [],
                }
            )
            substring_ocr = TextBoxSpec(
                text="续航里程长",
                source_pixel_bbox=(58, 42, 170, 68),
                source_pixel_polygon=((58, 42), (170, 42), (170, 68), (58, 68)),
                font_size=12,
                provenance={"content_source": "ocr", "layout_source": "ocr", "ocr_confidence": 0.78},
            )

            page = build_page_manifest_from_vlm_analysis(
                analysis=analysis,
                slide_id="slide-a",
                page_index=0,
                source_image_path=source,
                clean_background_path=background,
                artifact_root=root,
                aspect_ratio="16:9",
                text_boxes=[substring_ocr],
            )

        self.assertEqual([box.text for box in page.text_boxes], ["纯电续航里程长"])
        self.assertEqual(page.text_boxes[0].provenance["provider_role"], "VLM")

    def test_vlm_text_resolution_keeps_high_confidence_vlm_text_over_partial_ocr_substring(self):
        from src.generative_editable_manifest import TextBoxSpec
        from src.generative_editable_vlm_reconstruction import (
            build_page_manifest_from_vlm_analysis,
            coerce_vlm_analysis,
        )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "sources" / "0000-slide-a" / "source.png"
            background = root / "backgrounds" / "0000-slide-a" / "clean.png"
            source.parent.mkdir(parents=True)
            background.parent.mkdir(parents=True)
            Image.new("RGB", (800, 450), "#001122").save(source)
            Image.new("RGB", (800, 450), "#001122").save(background)
            analysis = coerce_vlm_analysis(
                {
                    "coordinate_space": {"width": 800, "height": 450, "unit": "px"},
                    "text_regions": [
                        {
                            "id": "instruction",
                            "text": "连接充电枪不可启动车辆",
                            "bbox": [100, 380, 380, 410],
                            "role": "body",
                            "confidence": 0.96,
                        }
                    ],
                    "bitmap_regions": [],
                    "shape_regions": [],
                }
            )
            partial_ocr = TextBoxSpec(
                text="连接充电枪不可启动",
                source_pixel_bbox=(100, 380, 380, 410),
                source_pixel_polygon=((100, 380), (380, 380), (380, 410), (100, 410)),
                font_size=18,
                provenance={"content_source": "ocr", "layout_source": "ocr", "ocr_confidence": 0.78},
            )

            page = build_page_manifest_from_vlm_analysis(
                analysis=analysis,
                slide_id="slide-a",
                page_index=0,
                source_image_path=source,
                clean_background_path=background,
                artifact_root=root,
                aspect_ratio="16:9",
                text_boxes=[partial_ocr],
            )

        self.assertEqual([box.text for box in page.text_boxes], ["连接充电枪不可启动车辆"])
        self.assertEqual(page.text_boxes[0].provenance["provider_role"], "VLM")

    def test_vlm_text_resolution_keeps_high_confidence_vlm_text_over_multiple_bad_ocr_fragments(self):
        from src.generative_editable_manifest import TextBoxSpec
        from src.generative_editable_vlm_reconstruction import (
            build_page_manifest_from_vlm_analysis,
            coerce_vlm_analysis,
        )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "sources" / "0000-slide-a" / "source.png"
            background = root / "backgrounds" / "0000-slide-a" / "clean.png"
            source.parent.mkdir(parents=True)
            background.parent.mkdir(parents=True)
            Image.new("RGB", (800, 450), "#001122").save(source)
            Image.new("RGB", (800, 450), "#001122").save(background)
            analysis = coerce_vlm_analysis(
                {
                    "coordinate_space": {"width": 800, "height": 450, "unit": "px"},
                    "text_regions": [
                        {
                            "id": "instruction",
                            "text": "连接充电枪不可启动车辆",
                            "bbox": [100, 380, 420, 412],
                            "role": "body",
                            "confidence": 0.96,
                        }
                    ],
                    "bitmap_regions": [],
                    "shape_regions": [],
                }
            )
            ocr_fragments = [
                TextBoxSpec(
                    text="连接充电枪不可启动",
                    source_pixel_bbox=(100, 380, 300, 397),
                    source_pixel_polygon=((100, 380), (300, 380), (300, 397), (100, 397)),
                    font_size=12,
                    provenance={"content_source": "ocr", "layout_source": "ocr", "ocr_confidence": 0.78},
                ),
                TextBoxSpec(
                    text="车",
                    source_pixel_bbox=(306, 380, 420, 397),
                    source_pixel_polygon=((306, 380), (420, 380), (420, 397), (306, 397)),
                    font_size=12,
                    provenance={"content_source": "ocr", "layout_source": "ocr", "ocr_confidence": 0.76},
                ),
            ]

            page = build_page_manifest_from_vlm_analysis(
                analysis=analysis,
                slide_id="slide-a",
                page_index=0,
                source_image_path=source,
                clean_background_path=background,
                artifact_root=root,
                aspect_ratio="16:9",
                text_boxes=ocr_fragments,
            )

        self.assertEqual([box.text for box in page.text_boxes], ["连接充电枪不可启动车辆"])
        self.assertEqual(page.text_boxes[0].provenance["provider_role"], "VLM")

    def test_vlm_text_resolution_keeps_overlapping_substring_labels(self):
        from src.generative_editable_manifest import TextBoxSpec
        from src.generative_editable_vlm_reconstruction import (
            build_page_manifest_from_vlm_analysis,
            coerce_vlm_analysis,
        )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "sources" / "0000-slide-a" / "source.png"
            background = root / "backgrounds" / "0000-slide-a" / "clean.png"
            source.parent.mkdir(parents=True)
            background.parent.mkdir(parents=True)
            Image.new("RGB", (320, 180), "#001122").save(source)
            Image.new("RGB", (320, 180), "#001122").save(background)
            analysis = coerce_vlm_analysis(
                {
                    "coordinate_space": {"width": 320, "height": 180, "unit": "px"},
                    "text_regions": [
                        {"id": "t1", "text": "Revenue", "bbox": [40, 40, 110, 62]},
                        {"id": "t2", "text": "Revenue Q1", "bbox": [42, 44, 145, 68]},
                    ],
                    "bitmap_regions": [],
                    "shape_regions": [],
                }
            )
            text_boxes = [
                TextBoxSpec(
                    text="Revenue",
                    source_pixel_bbox=(40, 40, 110, 62),
                    source_pixel_polygon=((40, 40), (110, 40), (110, 62), (40, 62)),
                    font_size=10,
                    provenance={"content_source": "ocr", "ocr_confidence": 0.99},
                ),
                TextBoxSpec(
                    text="Revenue Q1",
                    source_pixel_bbox=(42, 44, 145, 68),
                    source_pixel_polygon=((42, 44), (145, 44), (145, 68), (42, 68)),
                    font_size=10,
                    provenance={"content_source": "ocr", "ocr_confidence": 0.99},
                ),
            ]

            page = build_page_manifest_from_vlm_analysis(
                analysis=analysis,
                slide_id="slide-a",
                page_index=0,
                source_image_path=source,
                clean_background_path=background,
                artifact_root=root,
                aspect_ratio="16:9",
                text_boxes=text_boxes,
            )

        self.assertEqual([box.text for box in page.text_boxes], ["Revenue", "Revenue Q1"])

    def test_vlm_text_resolution_keeps_short_repeated_labels(self):
        from src.generative_editable_manifest import TextBoxSpec
        from src.generative_editable_vlm_reconstruction import (
            build_page_manifest_from_vlm_analysis,
            coerce_vlm_analysis,
        )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "sources" / "0000-slide-a" / "source.png"
            background = root / "backgrounds" / "0000-slide-a" / "clean.png"
            source.parent.mkdir(parents=True)
            background.parent.mkdir(parents=True)
            Image.new("RGB", (320, 180), "#001122").save(source)
            Image.new("RGB", (320, 180), "#001122").save(background)
            analysis = coerce_vlm_analysis(
                {
                    "coordinate_space": {"width": 320, "height": 180, "unit": "px"},
                    "text_regions": [
                        {"id": "t1", "text": "OK", "bbox": [40, 40, 80, 62]},
                        {"id": "t2", "text": "OK", "bbox": [44, 44, 84, 66]},
                    ],
                    "bitmap_regions": [],
                    "shape_regions": [],
                }
            )
            text_boxes = [
                TextBoxSpec(
                    text="OK",
                    source_pixel_bbox=(40, 40, 80, 62),
                    source_pixel_polygon=((40, 40), (80, 40), (80, 62), (40, 62)),
                    font_size=10,
                    provenance={"content_source": "ocr", "ocr_confidence": 0.99},
                ),
                TextBoxSpec(
                    text="OK",
                    source_pixel_bbox=(44, 44, 84, 66),
                    source_pixel_polygon=((44, 44), (84, 44), (84, 66), (44, 66)),
                    font_size=10,
                    provenance={"content_source": "ocr", "ocr_confidence": 0.99},
                ),
            ]

            page = build_page_manifest_from_vlm_analysis(
                analysis=analysis,
                slide_id="slide-a",
                page_index=0,
                source_image_path=source,
                clean_background_path=background,
                artifact_root=root,
                aspect_ratio="16:9",
                text_boxes=text_boxes,
            )

        self.assertEqual([box.text for box in page.text_boxes], ["OK", "OK"])

    def test_vlm_text_resolution_uses_ocr_layout_for_short_text(self):
        from src.generative_editable_manifest import BitmapAssetSpec, TextBoxSpec
        from src.generative_editable_vlm_reconstruction import (
            build_page_manifest_from_vlm_analysis,
            coerce_vlm_analysis,
        )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "sources" / "0000-slide-a" / "source.png"
            background = root / "backgrounds" / "0000-slide-a" / "clean.png"
            source.parent.mkdir(parents=True)
            background.parent.mkdir(parents=True)
            Image.new("RGB", (1600, 900), "#001122").save(source)
            Image.new("RGB", (1600, 900), "#001122").save(background)
            analysis = coerce_vlm_analysis(
                {
                    "coordinate_space": {"width": 1600, "height": 900, "unit": "px"},
                    "text_regions": [
                        {"id": "t1", "text": "EMB", "bbox": [1235, 587, 1287, 613]},
                    ],
                    "bitmap_regions": [],
                    "shape_regions": [],
                }
            )
            ocr_box = TextBoxSpec(
                text="EMB",
                source_pixel_bbox=(1292, 575, 1340, 595),
                source_pixel_polygon=((1292, 575), (1340, 575), (1340, 595), (1292, 595)),
                provenance={"content_source": "ocr", "ocr_confidence": 0.99},
            )
            asset = BitmapAssetSpec(
                asset_id="component",
                source_pixel_bbox=(1004, 418, 1312, 605),
                asset_path="assets/component.png",
                z_order=1,
                provenance={"asset_strategy": "source_preserved_crop"},
            )

            page = build_page_manifest_from_vlm_analysis(
                analysis=analysis,
                slide_id="slide-a",
                page_index=0,
                source_image_path=source,
                clean_background_path=background,
                artifact_root=root,
                aspect_ratio="16:9",
                text_boxes=[ocr_box],
                bitmap_assets=[asset],
            )

        self.assertEqual(page.text_boxes[0].source_pixel_bbox, (1292, 575, 1340, 595))
        self.assertEqual(page.text_boxes[0].provenance["layout_source"], "ocr")

    def test_vlm_fallback_does_not_duplicate_nearby_ocr_text_line(self):
        from src.generative_editable_manifest import TextBoxSpec
        from src.generative_editable_vlm_reconstruction import (
            build_page_manifest_from_vlm_analysis,
            coerce_vlm_analysis,
        )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "sources" / "0000-slide-a" / "source.png"
            background = root / "backgrounds" / "0000-slide-a" / "clean.png"
            source.parent.mkdir(parents=True)
            background.parent.mkdir(parents=True)
            Image.new("RGB", (1600, 900), "#001122").save(source)
            Image.new("RGB", (1600, 900), "#001122").save(background)
            analysis = coerce_vlm_analysis(
                {
                    "coordinate_space": {"width": 1600, "height": 900, "unit": "px"},
                    "text_regions": [
                        {
                            "id": "t1",
                            "text": "Livis版800V主动悬架",
                            "bbox": [1291, 558, 1501, 589],
                        },
                    ],
                    "bitmap_regions": [],
                    "shape_regions": [],
                }
            )
            ocr_box = TextBoxSpec(
                text="Livis版800V主动悬架",
                source_pixel_bbox=(1292, 540, 1502, 563),
                source_pixel_polygon=((1292, 540), (1502, 540), (1502, 563), (1292, 563)),
                provenance={"content_source": "ocr", "ocr_confidence": 0.99},
            )

            page = build_page_manifest_from_vlm_analysis(
                analysis=analysis,
                slide_id="slide-a",
                page_index=0,
                source_image_path=source,
                clean_background_path=background,
                artifact_root=root,
                aspect_ratio="16:9",
                text_boxes=[ocr_box],
            )

        self.assertEqual([box.text for box in page.text_boxes], ["Livis版800V主动悬架"])
        self.assertEqual(page.text_boxes[0].source_pixel_bbox, (1292, 540, 1502, 563))

    def test_without_vlm_text_regions_uses_high_confidence_ocr_boxes(self):
        from src.generative_editable_manifest import TextBoxSpec
        from src.generative_editable_vlm_reconstruction import (
            build_page_manifest_from_vlm_analysis,
            coerce_vlm_analysis,
        )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "sources" / "0000-slide-a" / "source.png"
            background = root / "backgrounds" / "0000-slide-a" / "clean.png"
            source.parent.mkdir(parents=True)
            background.parent.mkdir(parents=True)
            Image.new("RGB", (240, 120), "#001122").save(source)
            Image.new("RGB", (240, 120), "#001122").save(background)
            analysis = coerce_vlm_analysis(
                {
                    "coordinate_space": {"width": 240, "height": 120, "unit": "px"},
                    "text_regions": [],
                    "bitmap_regions": [],
                    "shape_regions": [],
                }
            )
            ocr_boxes = [
                TextBoxSpec(
                    text=f"OCR {index}",
                    source_pixel_bbox=(10, 10 + index * 20, 80, 25 + index * 20),
                    source_pixel_polygon=(
                        (10, 10 + index * 20),
                        (80, 10 + index * 20),
                        (80, 25 + index * 20),
                        (10, 25 + index * 20),
                    ),
                    provenance={"content_source": "ocr", "ocr_confidence": 0.99},
                )
                for index in range(5)
            ]

            page = build_page_manifest_from_vlm_analysis(
                analysis=analysis,
                slide_id="slide-a",
                page_index=0,
                source_image_path=source,
                clean_background_path=background,
                artifact_root=root,
                aspect_ratio="16:9",
                text_boxes=ocr_boxes,
            )

        self.assertEqual([box.text for box in page.text_boxes], [f"OCR {index}" for index in range(5)])
        self.assertTrue(
            all(box.provenance.get("layout_source") == "ocr" for box in page.text_boxes)
        )
        self.assertEqual(page.provenance["vlm_counts"]["text_regions"], 0)

    def test_without_vlm_text_regions_rejects_dense_short_ocr_noise(self):
        from src.generative_editable_manifest import TextBoxSpec
        from src.generative_editable_vlm_reconstruction import (
            build_page_manifest_from_vlm_analysis,
            coerce_vlm_analysis,
        )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "sources" / "0000-slide-a" / "source.png"
            background = root / "backgrounds" / "0000-slide-a" / "clean.png"
            source.parent.mkdir(parents=True)
            background.parent.mkdir(parents=True)
            Image.new("RGB", (240, 120), "#001122").save(source)
            Image.new("RGB", (240, 120), "#001122").save(background)
            analysis = coerce_vlm_analysis(
                {
                    "coordinate_space": {"width": 240, "height": 120, "unit": "px"},
                    "text_regions": [],
                    "bitmap_regions": [],
                    "shape_regions": [],
                }
            )
            ocr_boxes = [
                TextBoxSpec(
                    text="x",
                    source_pixel_bbox=(5 + (index % 20) * 11, 5 + (index // 20) * 10, 12 + (index % 20) * 11, 12 + (index // 20) * 10),
                    source_pixel_polygon=(
                        (5 + (index % 20) * 11, 5 + (index // 20) * 10),
                        (12 + (index % 20) * 11, 5 + (index // 20) * 10),
                        (12 + (index % 20) * 11, 12 + (index // 20) * 10),
                        (5 + (index % 20) * 11, 12 + (index // 20) * 10),
                    ),
                    provenance={"content_source": "ocr", "ocr_confidence": 0.99},
                )
                for index in range(60)
            ]

            page = build_page_manifest_from_vlm_analysis(
                analysis=analysis,
                slide_id="slide-a",
                page_index=0,
                source_image_path=source,
                clean_background_path=background,
                artifact_root=root,
                aspect_ratio="16:9",
                text_boxes=ocr_boxes,
            )

        self.assertEqual(page.text_boxes, [])

    def test_openai_chat_vlm_provider_requests_compact_json_from_gateway_friendly_image(self):
        from src.generative_editable_vlm_reconstruction import OpenAIChatVLMPageAnalysisProvider

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "source.png"
            Image.new("RGB", (1600, 900), "#001122").save(source)
            payloads = []
            provider = OpenAIChatVLMPageAnalysisProvider(
                ProviderConfig(
                    role="VLM",
                    provider="openai_chat",
                    model="gpt-5.5",
                    base_url="https://example.invalid/v1",
                    api_key="secret",
                )
            )

            def fake_post(config, payload, operation, timeout_seconds):
                payloads.append(payload)
                return {
                    "choices": [
                        {
                            "message": {
                                "content": (
                                    '{"coordinate_space":{"width":768,"height":432,"unit":"px"},'
                                    '"text_regions":[{"id":"t1","text":"标题","bbox":[10,10,100,40]}],'
                                    '"bitmap_regions":[],"shape_regions":[]}'
                                )
                            }
                        }
                    ]
                }

            with patch(
                "src.generative_editable_vlm_reconstruction._post_openai_chat",
                side_effect=fake_post,
            ):
                analysis = provider.analyze_page(str(source), timeout_seconds=99)

        self.assertEqual(analysis.coordinate_space.width, 768)
        self.assertEqual(analysis.text_regions[0].text, "标题")
        self.assertEqual(payloads[0]["model"], "gpt-5.5")
        self.assertEqual(payloads[0]["response_format"], {"type": "json_object"})
        self.assertLessEqual(payloads[0]["max_tokens"], 3000)
        content = payloads[0]["messages"][0]["content"]
        self.assertEqual(content[0]["type"], "text")
        self.assertIn("只输出 JSON", content[0]["text"])
        self.assertIn("768x432", content[0]["text"])
        self.assertEqual(content[1]["image_url"]["detail"], "low")
        self.assertTrue(content[1]["image_url"]["url"].startswith("data:image/jpeg;base64,"))

    def test_vlm_provider_falls_back_to_alternate_payload_shape_after_retryable_error(self):
        from src.generative_editable_vlm_reconstruction import OpenAIChatVLMPageAnalysisProvider

        with tempfile.TemporaryDirectory() as tmp:
            source = Path(tmp) / "source.png"
            Image.new("RGB", (1600, 900), "#001122").save(source)
            payloads = []
            provider = OpenAIChatVLMPageAnalysisProvider(
                ProviderConfig(
                    role="VLM",
                    provider="openai_chat",
                    model="gpt-5.5",
                    base_url="https://example.invalid/v1",
                    api_key="secret",
                )
            )

            def fake_post(config, payload, operation, timeout_seconds):
                payloads.append(payload)
                if len(payloads) == 1:
                    raise ProviderError(
                        provider_role="VLM",
                        operation="vlm_page_analysis",
                        message="503 no healthy upstream",
                        retryable=True,
                    )
                return {
                    "choices": [
                        {
                            "message": {
                                "content": (
                                    '{"coordinate_space":{"width":768,"height":432,"unit":"px"},'
                                    '"text_regions":[{"id":"t1","text":"标题","bbox":[10,10,100,40]}],'
                                    '"bitmap_regions":[],"shape_regions":[]}'
                                )
                            }
                        }
                    ]
                }

            with patch(
                "src.generative_editable_vlm_reconstruction._post_openai_chat",
                side_effect=fake_post,
            ):
                analysis = provider.analyze_page(str(source), timeout_seconds=99)

        self.assertEqual(analysis.text_regions[0].text, "标题")
        self.assertEqual(len(payloads), 2)
        first_image = payloads[0]["messages"][0]["content"][1]["image_url"]
        second_image = payloads[1]["messages"][0]["content"][1]["image_url"]
        self.assertEqual(first_image["detail"], "low")
        self.assertNotIn("detail", second_image)

    def test_vlm_provider_records_payload_variant_attempts_on_success_after_error(self):
        from src.generative_editable_vlm_reconstruction import OpenAIChatVLMPageAnalysisProvider

        with tempfile.TemporaryDirectory() as tmp:
            source = Path(tmp) / "source.png"
            Image.new("RGB", (1600, 900), "#001122").save(source)
            calls = 0
            provider = OpenAIChatVLMPageAnalysisProvider(
                ProviderConfig(
                    role="VLM",
                    provider="openai_chat",
                    model="gpt-5.5",
                    base_url="https://example.invalid/v1",
                    api_key="secret",
                )
            )

            def fake_post(config, payload, operation, timeout_seconds):
                nonlocal calls
                calls += 1
                if calls == 1:
                    raise ProviderError(
                        provider_role="VLM",
                        operation="vlm_page_analysis",
                        message="503 no healthy upstream api_key=secret",
                        retryable=True,
                        status_code=503,
                    )
                return {
                    "choices": [
                        {
                            "message": {
                                "content": (
                                    '{"coordinate_space":{"width":768,"height":432,"unit":"px"},'
                                    '"text_regions":[{"id":"t1","text":"标题","bbox":[10,10,100,40]}],'
                                    '"bitmap_regions":[],"shape_regions":[]}'
                                )
                            }
                        }
                    ]
                }

            with patch(
                "src.generative_editable_vlm_reconstruction._post_openai_chat",
                side_effect=fake_post,
            ):
                provider.analyze_page(str(source), timeout_seconds=99)

        self.assertEqual(
            [(item["attempt"], item["status"]) for item in provider.last_payload_attempts],
            [(1, "failed"), (2, "passed")],
        )
        self.assertEqual(provider.last_payload_attempts[0]["status_code"], 503)
        self.assertNotIn("secret", provider.last_payload_attempts[0]["error"])

    def test_vlm_provider_falls_back_after_payload_shape_compatibility_error(self):
        from src.generative_editable_vlm_reconstruction import OpenAIChatVLMPageAnalysisProvider

        with tempfile.TemporaryDirectory() as tmp:
            source = Path(tmp) / "source.png"
            Image.new("RGB", (1600, 900), "#001122").save(source)
            payloads = []
            provider = OpenAIChatVLMPageAnalysisProvider(
                ProviderConfig(
                    role="VLM",
                    provider="openai_chat",
                    model="gpt-5.5",
                    base_url="https://example.invalid/v1",
                    api_key="secret",
                )
            )

            def fake_post(config, payload, operation, timeout_seconds):
                payloads.append(payload)
                if len(payloads) == 1:
                    raise ProviderError(
                        provider_role="VLM",
                        operation="vlm_page_analysis",
                        message="400 Bad Request: unsupported field image_url.detail",
                        retryable=False,
                    )
                return {
                    "choices": [
                        {
                            "message": {
                                "content": (
                                    '{"coordinate_space":{"width":768,"height":432,"unit":"px"},'
                                    '"text_regions":[{"id":"t1","text":"标题","bbox":[10,10,100,40]}],'
                                    '"bitmap_regions":[],"shape_regions":[]}'
                                )
                            }
                        }
                    ]
                }

            with patch(
                "src.generative_editable_vlm_reconstruction._post_openai_chat",
                side_effect=fake_post,
            ):
                analysis = provider.analyze_page(str(source), timeout_seconds=99)

        self.assertEqual(analysis.text_regions[0].text, "标题")
        self.assertEqual(len(payloads), 2)
        self.assertNotIn("detail", payloads[1]["messages"][0]["content"][1]["image_url"])

    def test_vlm_provider_rejects_non_json_response(self):
        from src.generative_editable_providers import ProviderError
        from src.generative_editable_vlm_reconstruction import OpenAIChatVLMPageAnalysisProvider

        with tempfile.TemporaryDirectory() as tmp:
            source = Path(tmp) / "source.png"
            Image.new("RGB", (200, 100), "#001122").save(source)
            provider = OpenAIChatVLMPageAnalysisProvider(
                ProviderConfig(
                    role="VLM",
                    provider="openai_chat",
                    model="gpt-5.5",
                    base_url="https://example.invalid/v1",
                    api_key="secret",
                )
            )

            def fake_post(config, payload, operation, timeout_seconds):
                return {"choices": [{"message": {"content": "I cannot analyze this image."}}]}

            with patch(
                "src.generative_editable_vlm_reconstruction._post_openai_chat",
                side_effect=fake_post,
            ):
                with self.assertRaisesRegex(ProviderError, "valid JSON"):
                    provider.analyze_page(str(source))

    def test_vlm_pipeline_writes_manifests_and_composes_editable_pptx(self):
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        from src.generative_editable_manifest import read_deck_manifest, read_page_manifest
        from src.generative_editable_preview_validator import ValidationReport
        from src.generative_editable_vlm_reconstruction import (
            VLMEditablePipelineDependencies,
            VLMPageAnalysisProvider,
            coerce_vlm_analysis,
            run_vlm_editable_pptx_pipeline,
        )

        class FakeVLMProvider(VLMPageAnalysisProvider):
            def __init__(self):
                self.calls = []

            def analyze_page(self, image_path: str, *, timeout_seconds: int = 180):
                self.calls.append((image_path, timeout_seconds))
                return coerce_vlm_analysis(
                    {
                        "coordinate_space": {"width": 160, "height": 90, "unit": "px"},
                        "text_regions": [
                            {
                                "id": "title",
                                "text": "标题",
                                "bbox": [20, 20, 70, 36],
                                "role": "title",
                                "color": "#FFFFFF",
                            }
                        ],
                        "bitmap_regions": [
                            {"id": "logo", "type": "logo", "bbox": [120, 20, 150, 50]},
                        ],
                        "shape_regions": [
                            {"id": "divider", "type": "divider", "bbox": [20, 60, 150, 61]},
                        ],
                    }
                )

        class RecordingImageEditProvider(ImageEditProvider):
            def __init__(self):
                super().__init__(
                    ProviderConfig(
                        role="edit_model",
                        provider="fake_image_edit",
                        model="fake-image-edit",
                        base_url="",
                        api_key="",
                    )
                )
                self.calls = []
                self.fake = FakeImageEditProvider(self.config)

            def edit(self, request):
                self.calls.append(request)
                if request.prompt_id == "asset_sheet":
                    output = Path(request.output_asset_path)
                    output.parent.mkdir(parents=True, exist_ok=True)
                    sheet = Image.new("RGBA", (240, 120), (0, 0, 0, 0))
                    ImageDraw.Draw(sheet).rectangle((30, 24, 90, 84), fill="#22AAFF")
                    sheet.save(output)
                    return type(
                        "Result",
                        (),
                        {
                            "output_asset_path": str(output),
                            "source_image_path": request.source_image_path,
                            "prompt_id": request.prompt_id,
                            "provider_role": self.config.role,
                            "provider_name": self.config.provider,
                            "model": self.config.model,
                            "timeout_seconds": request.timeout_seconds,
                            "mask_path": request.mask_path,
                            "provenance": {},
                        },
                    )()
                return self.fake.edit(request)

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "input.png"
            image = Image.new("RGB", (160, 90), "#001122")
            draw = ImageDraw.Draw(image)
            draw.rectangle((20, 20, 70, 36), fill="#FFFFFF")
            draw.rectangle((120, 20, 150, 50), fill="#22AAFF")
            image.save(source)
            output = root / "out.pptx"
            vlm_provider = FakeVLMProvider()
            image_edit_provider = RecordingImageEditProvider()

            result = run_vlm_editable_pptx_pipeline(
                slides=[{"slide_id": "slide-a", "image_path": str(source)}],
                output_path=str(output),
                artifact_root=str(root / "artifacts"),
                job_id="job-001",
                dependencies=VLMEditablePipelineDependencies(
                    vlm_provider=vlm_provider,
                    image_edit_provider=image_edit_provider,
                    structure_validator=lambda **kwargs: ValidationReport(
                        status="passed", checked_pages=1, issues=[]
                    ),
                    preview_renderer=lambda page, artifact_root, *, pptx_path: object(),
                    preview_validator=lambda **kwargs: ValidationReport(
                        status="passed", checked_pages=1, issues=[]
                    ),
                ),
            )
            deck = read_deck_manifest(root / "artifacts" / "job-001" / "deck.json")
            page = read_page_manifest(root / "artifacts" / "job-001" / deck.page_manifest_paths[0])
            prs = Presentation(str(output))
            shape_types = [shape.shape_type for shape in prs.slides[0].shapes]
            mask = Image.open(root / "artifacts" / "job-001" / "assets" / "0000-slide-a" / "vlm-text-bitmap-mask.png")

        self.assertEqual(result.status, "passed")
        self.assertEqual(len(vlm_provider.calls), 1)
        self.assertIn("Remove only the masked text, icons", image_edit_provider.calls[0].prompt)
        self.assertIn("no readable residual glyphs", image_edit_provider.calls[0].prompt)
        self.assertIn("Preserve every unmasked pixel", image_edit_provider.calls[0].prompt)
        self.assertIn("asset_sheet", [request.prompt_id for request in image_edit_provider.calls])
        asset_sheet_request = next(request for request in image_edit_provider.calls if request.prompt_id == "asset_sheet")
        self.assertEqual(asset_sheet_request.timeout_seconds, 120)
        self.assertEqual(mask.getpixel((75, 40)), 255)
        self.assertIn("vlm", page.provider_output_paths)
        self.assertIn("asset_sheet", page.provider_output_paths)
        self.assertEqual([sheet.sheet_id for sheet in page.asset_sheets], ["vlm-asset-sheet-0000"])
        self.assertEqual(page.asset_sheets[0].candidate_ids, ["logo"])
        self.assertEqual(page.bitmap_assets[0].provenance["split_method"], "connected_components")
        self.assertEqual(page.provenance["reconstruction_strategy"], "vlm_first")
        self.assertEqual(page.provenance["text_mask_path"], "assets/0000-slide-a/vlm-text-bitmap-mask.png")
        self.assertEqual([box.text for box in page.text_boxes], ["标题"])
        self.assertGreaterEqual(shape_types.count(MSO_SHAPE_TYPE.TEXT_BOX), 1)
        self.assertGreaterEqual(shape_types.count(MSO_SHAPE_TYPE.PICTURE), 1)
        self.assertIn(MSO_SHAPE_TYPE.LINE, shape_types)

    def test_vlm_pipeline_rejects_asset_sheet_provider_failure_source_crops(self):
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        from src.generative_editable_manifest import read_deck_manifest, read_page_manifest
        from src.generative_editable_vlm_reconstruction import (
            VLMEditablePipelineDependencies,
            VLMPageAnalysisProvider,
            coerce_vlm_analysis,
            run_vlm_editable_pptx_pipeline,
        )

        class BitmapVLMProvider(VLMPageAnalysisProvider):
            def analyze_page(self, image_path: str, *, timeout_seconds: int = 180):
                return coerce_vlm_analysis(
                    {
                        "coordinate_space": {"width": 160, "height": 90, "unit": "px"},
                        "text_regions": [],
                        "bitmap_regions": [
                            {"id": "logo", "type": "logo", "bbox": [40, 20, 90, 60]},
                        ],
                        "shape_regions": [
                            {"id": "line-a", "type": "line", "bbox": [20, 60, 140, 61]}
                        ],
                    }
                )

        class FailingAssetSheetProvider(ImageEditProvider):
            def __init__(self):
                super().__init__(
                    ProviderConfig(
                        role="asset_sheet_model",
                        provider="fake_image_edit",
                        model="fake-image-edit",
                        base_url="",
                        api_key="",
                    )
                )

            def edit(self, request):
                raise ProviderError(
                    provider_role=self.config.role,
                    operation=request.prompt_id,
                    message="503 circuit breaker",
                    retryable=True,
                )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "input.png"
            image = Image.new("RGB", (160, 90), "#001122")
            ImageDraw.Draw(image).rectangle((40, 20, 90, 60), fill="#22AAFF")
            image.save(source)
            output = root / "out.pptx"

            run_vlm_editable_pptx_pipeline(
                slides=[{"slide_id": "slide-a", "image_path": str(source)}],
                output_path=str(output),
                artifact_root=str(root / "artifacts"),
                job_id="job-asset-fallback",
                dependencies=VLMEditablePipelineDependencies(
                    vlm_provider=BitmapVLMProvider(),
                    image_edit_provider=FakeImageEditProvider(
                        ProviderConfig(
                            role="edit_model",
                            provider="fake_image_edit",
                            model="fake-image-edit",
                            base_url="",
                            api_key="",
                        )
                    ),
                    asset_sheet_image_edit_provider=FailingAssetSheetProvider(),
                ),
            )
            job_dir = root / "artifacts" / "job-asset-fallback"
            deck = read_deck_manifest(job_dir / "deck.json")
            page = read_page_manifest(job_dir / deck.page_manifest_paths[0])
            prs = Presentation(str(output))
            shape_types = [shape.shape_type for shape in prs.slides[0].shapes]

        self.assertEqual(page.asset_sheets, [])
        self.assertEqual(len(page.bitmap_assets), 1)
        self.assertEqual(page.bitmap_assets[0].provenance["asset_strategy"], "source_preserved_crop")
        self.assertTrue(page.bitmap_assets[0].provenance["asset_sheet_provider_failed"])
        self.assertIn(MSO_SHAPE_TYPE.PICTURE, shape_types)

    def test_vlm_pipeline_uses_local_clean_background_when_clean_provider_fails(self):
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        from src.generative_editable_manifest import read_deck_manifest, read_page_manifest
        from src.generative_editable_preview_validator import ValidationReport
        from src.generative_editable_vlm_reconstruction import (
            VLMEditablePipelineDependencies,
            VLMPageAnalysisProvider,
            coerce_vlm_analysis,
            run_vlm_editable_pptx_pipeline,
        )

        class MixedVLMProvider(VLMPageAnalysisProvider):
            def analyze_page(self, image_path: str, *, timeout_seconds: int = 180):
                return coerce_vlm_analysis(
                    {
                        "coordinate_space": {"width": 160, "height": 90, "unit": "px"},
                        "text_regions": [
                            {
                                "id": "title",
                                "text": "标题",
                                "bbox": [14, 8, 58, 20],
                                "confidence": 0.96,
                            }
                        ],
                        "bitmap_regions": [
                            {"id": "product", "type": "product", "bbox": [44, 24, 102, 64]},
                        ],
                        "shape_regions": [
                            {"id": "divider", "type": "line", "bbox": [10, 72, 150, 73]},
                        ],
                    }
                )

        class FailingCleanProvider(ImageEditProvider):
            def __init__(self):
                super().__init__(
                    ProviderConfig(
                        role="edit_model",
                        provider="fake_image_edit",
                        model="fake-image-edit",
                        base_url="",
                        api_key="",
                    )
                )

            def edit(self, request):
                if request.prompt_id == "vlm_clean_background":
                    raise ProviderError(
                        provider_role=self.config.role,
                        operation=request.prompt_id,
                        message="524 timeout",
                        retryable=True,
                    )
                return FakeImageEditProvider(self.config).edit(request)

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "input.png"
            image = Image.new("RGB", (160, 90), "#001122")
            draw = ImageDraw.Draw(image)
            draw.text((14, 8), "标题", fill="#FFFFFF")
            draw.rectangle((44, 24, 102, 64), fill="#22AAFF")
            draw.line((10, 72, 150, 72), fill="#55CCFF", width=1)
            image.save(source)
            output = root / "out.pptx"

            result = run_vlm_editable_pptx_pipeline(
                slides=[{"slide_id": "slide-a", "image_path": str(source)}],
                output_path=str(output),
                artifact_root=str(root / "artifacts"),
                job_id="job-clean-fallback",
                dependencies=VLMEditablePipelineDependencies(
                    vlm_provider=MixedVLMProvider(),
                    image_edit_provider=FailingCleanProvider(),
                    structure_validator=lambda **kwargs: ValidationReport(
                        status="passed", checked_pages=1, issues=[]
                    ),
                    preview_renderer=lambda page, artifact_root, *, pptx_path: object(),
                    preview_validator=lambda **kwargs: ValidationReport(
                        status="passed", checked_pages=1, issues=[]
                    ),
                ),
            )
            job_dir = root / "artifacts" / "job-clean-fallback"
            deck = read_deck_manifest(job_dir / "deck.json")
            page = read_page_manifest(job_dir / deck.page_manifest_paths[0])
            prs = Presentation(str(output))
            shape_types = [shape.shape_type for shape in prs.slides[0].shapes]

        self.assertEqual(result.status, "passed")
        self.assertEqual(result.fallback_used, "clean_background_local")
        self.assertEqual(page.provider_output_paths["image_edit"], "provider_outputs/image_edit/0000-slide-a/backgrounds.json")
        self.assertTrue(page.provenance["clean_background_provider_failed"])
        self.assertEqual(page.provenance["clean_background_strategy"], "local_fill")
        self.assertGreaterEqual(len(page.text_boxes), 1)
        self.assertGreaterEqual(len(page.bitmap_assets), 1)
        self.assertIn(MSO_SHAPE_TYPE.TEXT_BOX, shape_types)
        self.assertIn(MSO_SHAPE_TYPE.PICTURE, shape_types)
        self.assertNotEqual(shape_types, [MSO_SHAPE_TYPE.PICTURE])

    def test_vlm_pipeline_does_not_turn_near_full_slide_bitmap_region_into_picture_asset(self):
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        from src.generative_editable_manifest import read_deck_manifest, read_page_manifest
        from src.generative_editable_preview_validator import ValidationReport
        from src.generative_editable_vlm_reconstruction import (
            VLMEditablePipelineDependencies,
            VLMPageAnalysisProvider,
            coerce_vlm_analysis,
            run_vlm_editable_pptx_pipeline,
        )

        class FullSlideBitmapVLMProvider(VLMPageAnalysisProvider):
            def analyze_page(self, image_path: str, *, timeout_seconds: int = 180):
                return coerce_vlm_analysis(
                    {
                        "coordinate_space": {"width": 160, "height": 90, "unit": "px"},
                        "text_regions": [
                            {
                                "id": "title",
                                "text": "标题",
                                "bbox": [12, 8, 52, 20],
                                "confidence": 0.96,
                            }
                        ],
                        "bitmap_regions": [
                            {"id": "whole", "type": "product", "bbox": [2, 2, 158, 88]},
                        ],
                        "shape_regions": [
                            {"id": "line-a", "type": "line", "bbox": [20, 60, 140, 61]}
                        ],
                    }
                )

        class FailingCleanProvider(ImageEditProvider):
            def __init__(self):
                super().__init__(
                    ProviderConfig(
                        role="edit_model",
                        provider="fake_image_edit",
                        model="fake-image-edit",
                        base_url="",
                        api_key="",
                    )
                )

            def edit(self, request):
                raise ProviderError(
                    provider_role=self.config.role,
                    operation=request.prompt_id,
                    message="524 timeout",
                    retryable=True,
                )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "input.png"
            image = Image.new("RGB", (160, 90), "#001122")
            draw = ImageDraw.Draw(image)
            draw.rectangle((2, 2, 158, 88), outline="#22AAFF", width=2)
            draw.text((12, 8), "标题", fill="#FFFFFF")
            image.save(source)
            output = root / "out.pptx"

            run_vlm_editable_pptx_pipeline(
                slides=[{"slide_id": "slide-a", "image_path": str(source)}],
                output_path=str(output),
                artifact_root=str(root / "artifacts"),
                job_id="job-full-slide-guard",
                dependencies=VLMEditablePipelineDependencies(
                    vlm_provider=FullSlideBitmapVLMProvider(),
                    image_edit_provider=FailingCleanProvider(),
                    structure_validator=lambda **kwargs: ValidationReport(
                        status="passed", checked_pages=1, issues=[]
                    ),
                    preview_renderer=lambda page, artifact_root, *, pptx_path: object(),
                    preview_validator=lambda **kwargs: ValidationReport(
                        status="passed", checked_pages=1, issues=[]
                    ),
                ),
            )
            job_dir = root / "artifacts" / "job-full-slide-guard"
            deck = read_deck_manifest(job_dir / "deck.json")
            page = read_page_manifest(job_dir / deck.page_manifest_paths[0])
            mask = Image.open(job_dir / "assets" / "0000-slide-a" / "vlm-text-bitmap-mask.png")
            prs = Presentation(str(output))
            shape_types = [shape.shape_type for shape in prs.slides[0].shapes]

        self.assertEqual(page.bitmap_assets, [])
        self.assertEqual(mask.getpixel((80, 45)), 0)
        self.assertEqual(mask.getpixel((18, 14)), 255)
        self.assertIn(MSO_SHAPE_TYPE.TEXT_BOX, shape_types)
        self.assertNotIn(MSO_SHAPE_TYPE.PICTURE, shape_types)

    def test_vlm_bitmap_mask_and_candidate_filter_use_same_unpadded_ignore_decision(self):
        from src.generative_editable_vlm_reconstruction import (
            VLMCoordinateMapper,
            _foreground_candidates_from_vlm_bitmap_regions,
            build_text_bitmap_mask,
            coerce_vlm_analysis,
        )

        analysis = coerce_vlm_analysis(
            {
                "coordinate_space": {"width": 160, "height": 90, "unit": "px"},
                "text_regions": [],
                "bitmap_regions": [
                    {"id": "large-panel", "type": "icon", "bbox": [8, 18, 152, 72]},
                ],
                "shape_regions": [],
            }
        )
        mapper = VLMCoordinateMapper(
            analysis_space=analysis.coordinate_space,
            source_image_size=(160, 90),
        )

        mask, mask_boxes = build_text_bitmap_mask(analysis, mapper, padding=10)
        candidates = _foreground_candidates_from_vlm_bitmap_regions(analysis, mapper)

        self.assertEqual([box.region_id for box in mask_boxes], ["large-panel"])
        self.assertEqual([candidate.candidate_id for candidate in candidates], ["large-panel"])
        self.assertEqual(mask.getpixel((1, 10)), 255)

    def test_vlm_pipeline_skips_asset_sheet_for_complex_bitmap_regions(self):
        from src.generative_editable_manifest import read_deck_manifest, read_page_manifest
        from src.generative_editable_vlm_reconstruction import (
            VLMEditablePipelineDependencies,
            VLMPageAnalysisProvider,
            coerce_vlm_analysis,
            run_vlm_editable_pptx_pipeline,
        )

        class ProductOnlyVLMProvider(VLMPageAnalysisProvider):
            def analyze_page(self, image_path: str, *, timeout_seconds: int = 180):
                return coerce_vlm_analysis(
                    {
                        "coordinate_space": {"width": 160, "height": 90, "unit": "px"},
                        "text_regions": [],
                        "bitmap_regions": [
                            {"id": "product", "type": "product", "bbox": [40, 20, 95, 62]},
                        ],
                        "shape_regions": [],
                    }
                )

        class RecordingAssetSheetProvider(ImageEditProvider):
            def __init__(self):
                super().__init__(
                    ProviderConfig(
                        role="asset_sheet_model",
                        provider="fake_image_edit",
                        model="fake-image-edit",
                        base_url="",
                        api_key="",
                    )
                )
                self.calls = []

            def edit(self, request):
                self.calls.append(request)
                raise AssertionError("complex bitmap regions should not call asset sheet provider")

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "input.png"
            image = Image.new("RGB", (160, 90), "#001122")
            ImageDraw.Draw(image).rectangle((40, 20, 95, 62), fill="#22AAFF")
            image.save(source)
            asset_sheet_provider = RecordingAssetSheetProvider()

            run_vlm_editable_pptx_pipeline(
                slides=[{"slide_id": "slide-a", "image_path": str(source)}],
                output_path=str(root / "out.pptx"),
                artifact_root=str(root / "artifacts"),
                job_id="job-complex-source",
                dependencies=VLMEditablePipelineDependencies(
                    vlm_provider=ProductOnlyVLMProvider(),
                    image_edit_provider=FakeImageEditProvider(
                        ProviderConfig(
                            role="edit_model",
                            provider="fake_image_edit",
                            model="fake-image-edit",
                            base_url="",
                            api_key="",
                        )
                    ),
                    asset_sheet_image_edit_provider=asset_sheet_provider,
                ),
            )
            job_dir = root / "artifacts" / "job-complex-source"
            deck = read_deck_manifest(job_dir / "deck.json")
            page = read_page_manifest(job_dir / deck.page_manifest_paths[0])

        self.assertEqual(asset_sheet_provider.calls, [])
        self.assertEqual(page.asset_sheets, [])
        self.assertEqual(len(page.bitmap_assets), 1)
        self.assertEqual(page.bitmap_assets[0].provenance["asset_strategy"], "source_preserved_crop")
        self.assertEqual(page.bitmap_assets[0].provenance["asset_sheet_skipped_reason"], "complex_bitmap_region")
        self.assertEqual(page.bitmap_assets[0].provenance["alpha_strategy"], "opaque_source_crop")
        self.assertNotIn("background_difference_alpha", page.bitmap_assets[0].provenance)
        self.assertIn("alpha_visible_area_ratio", page.bitmap_assets[0].provenance)

    def test_vlm_pipeline_asset_sheet_failure_does_not_duplicate_complex_source_crops(self):
        from src.generative_editable_manifest import read_deck_manifest, read_page_manifest
        from src.generative_editable_vlm_reconstruction import (
            VLMEditablePipelineDependencies,
            VLMPageAnalysisProvider,
            coerce_vlm_analysis,
            run_vlm_editable_pptx_pipeline,
        )

        class MixedVLMProvider(VLMPageAnalysisProvider):
            def analyze_page(self, image_path: str, *, timeout_seconds: int = 180):
                return coerce_vlm_analysis(
                    {
                        "coordinate_space": {"width": 160, "height": 90, "unit": "px"},
                        "text_regions": [],
                        "bitmap_regions": [
                            {"id": "product", "type": "product", "bbox": [20, 20, 70, 60]},
                            {"id": "logo", "type": "logo", "bbox": [100, 20, 130, 50]},
                        ],
                        "shape_regions": [],
                    }
                )

        class FailingAssetSheetProvider(ImageEditProvider):
            def __init__(self):
                super().__init__(
                    ProviderConfig(
                        role="asset_sheet_model",
                        provider="fake_image_edit",
                        model="fake-image-edit",
                        base_url="",
                        api_key="",
                    )
                )

            def edit(self, request):
                raise ProviderError(
                    provider_role=self.config.role,
                    operation=request.prompt_id,
                    message="503 circuit breaker",
                    retryable=True,
                )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "input.png"
            image = Image.new("RGB", (160, 90), "#001122")
            draw = ImageDraw.Draw(image)
            draw.rectangle((20, 20, 70, 60), fill="#22AAFF")
            draw.rectangle((100, 20, 130, 50), fill="#FFAA00")
            image.save(source)

            run_vlm_editable_pptx_pipeline(
                slides=[{"slide_id": "slide-a", "image_path": str(source)}],
                output_path=str(root / "out.pptx"),
                artifact_root=str(root / "artifacts"),
                job_id="job-mixed-fallback",
                dependencies=VLMEditablePipelineDependencies(
                    vlm_provider=MixedVLMProvider(),
                    image_edit_provider=FakeImageEditProvider(
                        ProviderConfig(
                            role="edit_model",
                            provider="fake_image_edit",
                            model="fake-image-edit",
                            base_url="",
                            api_key="",
                        )
                    ),
                    asset_sheet_image_edit_provider=FailingAssetSheetProvider(),
                ),
            )
            job_dir = root / "artifacts" / "job-mixed-fallback"
            deck = read_deck_manifest(job_dir / "deck.json")
            page = read_page_manifest(job_dir / deck.page_manifest_paths[0])

        self.assertEqual([asset.asset_id for asset in page.bitmap_assets], ["product", "logo"])
        self.assertEqual(page.bitmap_assets[0].provenance["asset_sheet_skipped_reason"], "complex_bitmap_region")
        self.assertTrue(page.bitmap_assets[1].provenance["asset_sheet_provider_failed"])

    def test_vlm_pipeline_falls_back_when_asset_sheet_component_count_is_short(self):
        from src.generative_editable_manifest import read_deck_manifest, read_page_manifest
        from src.generative_editable_vlm_reconstruction import (
            VLMEditablePipelineDependencies,
            VLMPageAnalysisProvider,
            coerce_vlm_analysis,
            run_vlm_editable_pptx_pipeline,
        )

        class TwoLogoVLMProvider(VLMPageAnalysisProvider):
            def analyze_page(self, image_path: str, *, timeout_seconds: int = 180):
                return coerce_vlm_analysis(
                    {
                        "coordinate_space": {"width": 160, "height": 90, "unit": "px"},
                        "text_regions": [],
                        "bitmap_regions": [
                            {"id": "logo-a", "type": "logo", "bbox": [30, 20, 55, 45]},
                            {"id": "logo-b", "type": "logo", "bbox": [95, 20, 125, 50]},
                        ],
                        "shape_regions": [],
                    }
                )

        class ShortAssetSheetProvider(ImageEditProvider):
            def __init__(self):
                super().__init__(
                    ProviderConfig(
                        role="asset_sheet_model",
                        provider="fake_image_edit",
                        model="fake-image-edit",
                        base_url="",
                        api_key="",
                    )
                )

            def edit(self, request):
                output = Path(request.output_asset_path)
                output.parent.mkdir(parents=True, exist_ok=True)
                sheet = Image.new("RGBA", (160, 90), (0, 0, 0, 0))
                ImageDraw.Draw(sheet).rectangle((30, 20, 55, 45), fill="#FFAA00")
                sheet.save(output)
                return type(
                    "Result",
                    (),
                    {
                        "output_asset_path": str(output),
                        "source_image_path": request.source_image_path,
                        "prompt_id": request.prompt_id,
                        "provider_role": self.config.role,
                        "provider_name": self.config.provider,
                        "model": self.config.model,
                        "timeout_seconds": request.timeout_seconds,
                        "mask_path": request.mask_path,
                        "provenance": {},
                    },
                )()

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "input.png"
            image = Image.new("RGB", (160, 90), "#001122")
            draw = ImageDraw.Draw(image)
            draw.rectangle((30, 20, 55, 45), fill="#FFAA00")
            draw.rectangle((95, 20, 125, 50), fill="#22AAFF")
            image.save(source)

            run_vlm_editable_pptx_pipeline(
                slides=[{"slide_id": "slide-a", "image_path": str(source)}],
                output_path=str(root / "out.pptx"),
                artifact_root=str(root / "artifacts"),
                job_id="job-short-sheet",
                dependencies=VLMEditablePipelineDependencies(
                    vlm_provider=TwoLogoVLMProvider(),
                    image_edit_provider=FakeImageEditProvider(
                        ProviderConfig(
                            role="edit_model",
                            provider="fake_image_edit",
                            model="fake-image-edit",
                            base_url="",
                            api_key="",
                        )
                    ),
                    asset_sheet_image_edit_provider=ShortAssetSheetProvider(),
                ),
            )
            job_dir = root / "artifacts" / "job-short-sheet"
            deck = read_deck_manifest(job_dir / "deck.json")
            page = read_page_manifest(job_dir / deck.page_manifest_paths[0])

        self.assertEqual(page.asset_sheets, [])
        self.assertEqual([asset.asset_id for asset in page.bitmap_assets], ["logo-a", "logo-b"])
        self.assertTrue(all(asset.provenance["asset_sheet_slicing_failed"] for asset in page.bitmap_assets))

    def test_vlm_pipeline_skips_asset_sheet_for_icon_bitmap_regions(self):
        from src.generative_editable_manifest import read_deck_manifest, read_page_manifest
        from src.generative_editable_vlm_reconstruction import (
            VLMEditablePipelineDependencies,
            VLMPageAnalysisProvider,
            coerce_vlm_analysis,
            run_vlm_editable_pptx_pipeline,
        )

        class IconOnlyVLMProvider(VLMPageAnalysisProvider):
            def analyze_page(self, image_path: str, *, timeout_seconds: int = 180):
                return coerce_vlm_analysis(
                    {
                        "coordinate_space": {"width": 160, "height": 90, "unit": "px"},
                        "text_regions": [],
                        "bitmap_regions": [
                            {"id": "icon-a", "type": "icon", "bbox": [30, 20, 55, 45]},
                            {"id": "icon-b", "type": "icon", "bbox": [95, 20, 125, 50]},
                        ],
                        "shape_regions": [],
                    }
                )

        class RecordingAssetSheetProvider(ImageEditProvider):
            def __init__(self):
                super().__init__(
                    ProviderConfig(
                        role="asset_sheet_model",
                        provider="fake_image_edit",
                        model="fake-image-edit",
                        base_url="",
                        api_key="",
                    )
                )
                self.calls = []

            def edit(self, request):
                self.calls.append(request)
                raise AssertionError("icon bitmap regions should not request asset sheet")

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "input.png"
            image = Image.new("RGB", (160, 90), "#001122")
            draw = ImageDraw.Draw(image)
            draw.rectangle((30, 20, 55, 45), fill="#FFAA00")
            draw.rectangle((95, 20, 125, 50), fill="#22AAFF")
            image.save(source)
            asset_sheet_provider = RecordingAssetSheetProvider()

            run_vlm_editable_pptx_pipeline(
                slides=[{"slide_id": "slide-a", "image_path": str(source)}],
                output_path=str(root / "out.pptx"),
                artifact_root=str(root / "artifacts"),
                job_id="job-icon-source-preserved",
                dependencies=VLMEditablePipelineDependencies(
                    vlm_provider=IconOnlyVLMProvider(),
                    image_edit_provider=FakeImageEditProvider(
                        ProviderConfig(
                            role="edit_model",
                            provider="fake_image_edit",
                            model="fake-image-edit",
                            base_url="",
                            api_key="",
                        )
                    ),
                    asset_sheet_image_edit_provider=asset_sheet_provider,
                ),
            )
            job_dir = root / "artifacts" / "job-icon-source-preserved"
            deck = read_deck_manifest(job_dir / "deck.json")
            page = read_page_manifest(job_dir / deck.page_manifest_paths[0])

        self.assertEqual(asset_sheet_provider.calls, [])
        self.assertEqual(page.asset_sheets, [])
        self.assertEqual([asset.asset_id for asset in page.bitmap_assets], ["icon-a", "icon-b"])
        self.assertTrue(
            all(
                asset.provenance.get("asset_sheet_skipped_reason") == "icon_source_preserved"
                for asset in page.bitmap_assets
            )
        )

    def test_vlm_pipeline_uses_ocr_for_mask_without_ungated_editable_text(self):
        from src.generative_editable_manifest import read_deck_manifest, read_page_manifest
        from src.generative_editable_preview_validator import ValidationReport
        from src.generative_editable_providers import OCRProvider
        from src.generative_editable_vlm_reconstruction import (
            VLMEditablePipelineDependencies,
            VLMPageAnalysisProvider,
            coerce_vlm_analysis,
            run_vlm_editable_pptx_pipeline,
        )

        class MissingTextVLMProvider(VLMPageAnalysisProvider):
            def analyze_page(self, image_path, *, timeout_seconds=180):
                return coerce_vlm_analysis(
                    {
                        "coordinate_space": {"width": 160, "height": 90, "unit": "px"},
                        "text_regions": [],
                        "bitmap_regions": [],
                        "shape_regions": [
                            {"id": "line-a", "type": "line", "bbox": [20, 60, 140, 61]}
                        ],
                    }
                )

        class ExactOCRProvider(OCRProvider):
            def __init__(self):
                super().__init__(
                    ProviderConfig(
                        role="ocr_model",
                        provider="fake_ocr",
                        model="fake-ocr",
                        base_url="fake",
                        api_key="fake",
                    )
                )
                self.calls = []

            def extract_text(self, image_path):
                self.calls.append(image_path)
                return OCRResult(
                    source_image_path=image_path,
                    image_size=(160, 90),
                    provider_role=self.config.role,
                    provider_name="fake_ocr",
                    model="fake-ocr",
                    items=[
                        OCRTextItem(
                            text="OCR 正文",
                            bbox=(20, 20, 90, 38),
                            polygon=((20, 20), (90, 20), (90, 38), (20, 38)),
                            confidence=0.99,
                            color_hex="#FFFFFF",
                            font_size_hint=18,
                            provenance={"item_id": "ocr-1"},
                        )
                    ],
                )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "input.png"
            image = Image.new("RGB", (160, 90), "#001122")
            ImageDraw.Draw(image).rectangle((20, 20, 90, 38), fill="#FFFFFF")
            image.save(source)
            output = root / "out.pptx"
            ocr_provider = ExactOCRProvider()

            run_vlm_editable_pptx_pipeline(
                slides=[{"slide_id": "slide-a", "image_path": str(source)}],
                output_path=str(output),
                artifact_root=str(root / "artifacts"),
                job_id="job-ocr",
                dependencies=VLMEditablePipelineDependencies(
                    vlm_provider=MissingTextVLMProvider(),
                    image_edit_provider=FakeImageEditProvider(
                        ProviderConfig(
                            role="edit_model",
                            provider="fake_image_edit",
                            model="fake-image-edit",
                            base_url="fake",
                            api_key="fake",
                        )
                    ),
                    ocr_provider=ocr_provider,
                    structure_validator=lambda **kwargs: ValidationReport(
                        status="passed", checked_pages=1, issues=[]
                    ),
                    preview_renderer=lambda page, artifact_root, *, pptx_path: object(),
                    preview_validator=lambda **kwargs: ValidationReport(
                        status="passed", checked_pages=1, issues=[]
                    ),
                ),
            )
            job_dir = root / "artifacts" / "job-ocr"
            deck = read_deck_manifest(job_dir / "deck.json")
            page = read_page_manifest(job_dir / deck.page_manifest_paths[0])
            mask = Image.open(job_dir / "assets" / "0000-slide-a" / "vlm-text-bitmap-mask.png")

        self.assertEqual(len(ocr_provider.calls), 1)
        self.assertEqual([box.text for box in page.text_boxes], ["OCR 正文"])
        self.assertEqual(page.text_boxes[0].provenance["layout_source"], "ocr")
        self.assertIn("ocr", page.provider_output_paths)
        self.assertEqual(mask.getpixel((25, 25)), 255)

    def test_vlm_pipeline_clears_ocr_only_text_from_complex_source_preserved_crop(self):
        from src.generative_editable_manifest import read_deck_manifest, read_page_manifest
        from src.generative_editable_preview_validator import ValidationReport
        from src.generative_editable_providers import OCRProvider
        from src.generative_editable_vlm_reconstruction import (
            VLMEditablePipelineDependencies,
            VLMPageAnalysisProvider,
            coerce_vlm_analysis,
            run_vlm_editable_pptx_pipeline,
        )

        class ProductOnlyVLMProvider(VLMPageAnalysisProvider):
            def analyze_page(self, image_path, *, timeout_seconds=180):
                return coerce_vlm_analysis(
                    {
                        "coordinate_space": {"width": 160, "height": 90, "unit": "px"},
                        "text_regions": [],
                        "bitmap_regions": [
                            {
                                "id": "product",
                                "type": "product",
                                "bbox": [20, 20, 110, 70],
                                "importance": "major",
                            }
                        ],
                        "shape_regions": [],
                    }
                )

        class OverlappingOCRProvider(OCRProvider):
            def __init__(self):
                super().__init__(
                    ProviderConfig(
                        role="ocr_model",
                        provider="fake_ocr",
                        model="fake-ocr",
                        base_url="fake",
                        api_key="fake",
                    )
                )

            def extract_text(self, image_path):
                return OCRResult(
                    source_image_path=image_path,
                    image_size=(160, 90),
                    provider_role=self.config.role,
                    provider_name="fake_ocr",
                    model="fake-ocr",
                    items=[
                        OCRTextItem(
                            text="OCR",
                            bbox=(34, 30, 58, 42),
                            polygon=((34, 30), (58, 30), (58, 42), (34, 42)),
                            confidence=0.99,
                            color_hex="#FFFFFF",
                            font_size_hint=12,
                            provenance={"item_id": "ocr-overlap"},
                        )
                    ],
                )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "input.png"
            image = Image.new("RGB", (160, 90), "#001122")
            draw = ImageDraw.Draw(image)
            draw.rectangle((20, 20, 110, 70), fill="#0B2A66")
            draw.rectangle((34, 30, 58, 42), fill="#FFFFFF")
            image.save(source)

            run_vlm_editable_pptx_pipeline(
                slides=[{"slide_id": "slide-a", "image_path": str(source)}],
                output_path=str(root / "out.pptx"),
                artifact_root=str(root / "artifacts"),
                job_id="job-ocr-crop",
                dependencies=VLMEditablePipelineDependencies(
                    vlm_provider=ProductOnlyVLMProvider(),
                    image_edit_provider=FakeImageEditProvider(
                        ProviderConfig(
                            role="edit_model",
                            provider="fake_image_edit",
                            model="fake-image-edit",
                            base_url="fake",
                            api_key="fake",
                        )
                    ),
                    ocr_provider=OverlappingOCRProvider(),
                    structure_validator=lambda **kwargs: ValidationReport(
                        status="passed", checked_pages=1, issues=[]
                    ),
                    preview_renderer=lambda page, artifact_root, *, pptx_path: object(),
                    preview_validator=lambda **kwargs: ValidationReport(
                        status="passed", checked_pages=1, issues=[]
                    ),
                ),
            )
            job_dir = root / "artifacts" / "job-ocr-crop"
            deck = read_deck_manifest(job_dir / "deck.json")
            page = read_page_manifest(job_dir / deck.page_manifest_paths[0])
            asset = page.bitmap_assets[0]
            crop = Image.open(job_dir / asset.asset_path).convert("RGBA")
            crop_x = 36 - asset.source_pixel_bbox[0]
            crop_y = 32 - asset.source_pixel_bbox[1]

        self.assertEqual(asset.provenance["alpha_strategy"], "opaque_source_crop")
        self.assertEqual(crop.getpixel((crop_x, crop_y))[3], 0)


if __name__ == "__main__":
    unittest.main()
