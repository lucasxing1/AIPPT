import io
import json
import argparse
import tempfile
import sys
import time
import unittest
from contextlib import redirect_stdout
from pathlib import Path
from unittest.mock import patch

from PIL import Image, ImageDraw

from scripts.run_real_generative_editable_pptx import main
from src.generative_editable_config import ProviderConfig, load_generative_editable_config
from src.generative_editable_pipeline import GenerativeEditablePipelineDependencies
from src.generative_editable_providers import (
    FakeImageEditProvider,
    FakeImageGenerationProvider,
    FakeOCRProvider,
    ImageGenerationProvider,
    ProviderError,
)


class RealGenerativeEditableRunnerTest(unittest.TestCase):
    def test_default_replay_inputs_exclude_edited_derivatives(self):
        import scripts.run_real_generative_editable_pptx as runner

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            Image.new("RGB", (800, 450), "#FFFFFF").save(root / "slide_1.png")
            Image.new("RGB", (800, 450), "#000000").save(root / "slide_1_edited.png")
            Image.new("RGB", (800, 450), "#F8FAFC").save(root / "slide_2.png")
            args = argparse.Namespace(
                input_images=[],
                input_glob=[],
                slides=0,
            )

            with patch.object(runner, "DEFAULT_REPLAY_GLOB", str(root / "slide_*.png")):
                paths = runner._resolve_input_images(args)

        self.assertEqual([path.name for path in paths], ["slide_1.png", "slide_2.png"])

    def test_explicit_input_glob_excludes_edited_derivatives(self):
        import scripts.run_real_generative_editable_pptx as runner

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            Image.new("RGB", (800, 450), "#FFFFFF").save(root / "slide_1.png")
            Image.new("RGB", (800, 450), "#000000").save(root / "slide_1_edited.png")
            Image.new("RGB", (800, 450), "#F8FAFC").save(root / "slide_2.png")
            args = argparse.Namespace(
                input_images=[],
                input_glob=[str(root / "slide_[0-9]*.png")],
                slides=0,
            )

            paths = runner._resolve_input_images(args)

        self.assertEqual([path.name for path in paths], ["slide_1.png", "slide_2.png"])

    def test_run_mode_accepts_sorted_input_glob_and_writes_report(self):
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            input_dir = root / "inputs"
            input_dir.mkdir()
            for name, color in (("slide_2.png", "#FFFFFF"), ("slide_1.png", "#F8FAFC")):
                image = Image.new("RGB", (800, 450), color)
                ImageDraw.Draw(image).rectangle((80, 80, 240, 180), fill="#2563EB")
                image.save(input_dir / name)

            stdout = io.StringIO()
            with redirect_stdout(stdout):
                exit_code = main(
                    [
                        "run",
                        "--use-fake",
                        "--input-glob",
                        str(input_dir / "slide_*.png"),
                        "--slides",
                        "2",
                        "--output-dir",
                        str(root / "out"),
                        "--job-id",
                        "real-test",
                    ]
                )

            self.assertEqual(exit_code, 0)
            line = json.loads(stdout.getvalue().splitlines()[-1])
            report_path = Path(line["report_path"])
            report = json.loads(report_path.read_text(encoding="utf-8"))
            self.assertEqual(report["status"], "passed")
            self.assertEqual(
                [Path(path).name for path in report["input_images"]],
                ["slide_1.png", "slide_2.png"],
            )
            self.assertTrue(Path(report["output_path"]).exists())
            self.assertEqual(len(report["object_stats"]["slides"]), 2)
            self.assertIn("TEXT_BOX", report["object_stats"]["slides"][0]["shape_counts"])

    def test_run_mode_disables_source_crop_asset_fallback_by_default(self):
        from api.routes import export as export_route

        captured = {}
        dependency_sentinel = object()

        def fake_pipeline(**kwargs):
            from pptx import Presentation

            from src.generative_editable_manifest import DeckManifest, write_manifest

            captured["dependencies"] = kwargs["dependencies"]
            output_path = Path(kwargs["output_path"])
            output_path.parent.mkdir(parents=True, exist_ok=True)
            Presentation().save(output_path)
            job_dir = Path(kwargs["artifact_root"]) / kwargs["job_id"]
            job_dir.mkdir(parents=True, exist_ok=True)
            (job_dir / "stage-events.jsonl").write_text(
                json.dumps({"stage": "compose_deck", "status": "finished"}) + "\n",
                encoding="utf-8",
            )
            write_manifest(
                job_dir / "deck.json",
                DeckManifest(
                    job_id=kwargs["job_id"],
                    slide_order=[],
                    aspect_ratio="16:9",
                    provider_roles={},
                    quality_settings={},
                    fallback_policy="fail",
                    page_manifest_paths=[],
                    validation_status="passed",
                ),
            )
            return type(
                "Result",
                (),
                {
                    "status": "passed",
                    "output_path": str(output_path),
                    "fallback_used": False,
                },
            )()

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "slide.png"
            Image.new("RGB", (800, 450), "#FFFFFF").save(source)
            stdout = io.StringIO()
            with (
                patch(
                    "scripts.run_real_generative_editable_pptx.export_route"
                    "._build_generative_editable_pipeline_dependencies",
                    return_value=export_route._build_fake_generative_editable_pipeline_dependencies(),
                ),
                patch(
                    "scripts.run_real_generative_editable_pptx.run_generative_editable_pipeline",
                    side_effect=fake_pipeline,
                ),
                patch(
                    "scripts.run_real_generative_editable_pptx._write_preview_reports",
                    return_value=[],
                ),
                redirect_stdout(stdout),
            ):
                exit_code = main(
                    [
                        "run",
                        "--input-images",
                        str(source),
                        "--output-dir",
                        str(root / "out"),
                        "--job-id",
                        "fallback-test",
                    ]
                )

        self.assertEqual(exit_code, 0)
        self.assertFalse(captured["dependencies"].allow_source_crop_asset_fallback)

    def test_run_mode_passes_runtime_fallback_policy_to_pipeline(self):
        from api.routes import export as export_route

        captured = {}

        def fake_pipeline(**kwargs):
            from pptx import Presentation

            from src.generative_editable_manifest import DeckManifest, write_manifest

            captured["fallback_policy"] = kwargs.get("fallback_policy")
            captured["fallback_output_factory"] = kwargs.get("fallback_output_factory")
            output_path = Path(kwargs["output_path"])
            output_path.parent.mkdir(parents=True, exist_ok=True)
            Presentation().save(output_path)
            job_dir = Path(kwargs["artifact_root"]) / kwargs["job_id"]
            job_dir.mkdir(parents=True, exist_ok=True)
            (job_dir / "stage-events.jsonl").write_text(
                json.dumps({"stage": "compose_deck", "status": "finished"}) + "\n",
                encoding="utf-8",
            )
            write_manifest(
                job_dir / "deck.json",
                DeckManifest(
                    job_id=kwargs["job_id"],
                    slide_order=[],
                    aspect_ratio="16:9",
                    provider_roles={},
                    quality_settings={},
                    fallback_policy="raster_pptx",
                    page_manifest_paths=[],
                    validation_status="passed",
                ),
            )
            return type(
                "Result",
                (),
                {
                    "status": "passed",
                    "output_path": str(output_path),
                    "fallback_used": "",
                },
            )()

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "slide.png"
            Image.new("RGB", (800, 450), "#FFFFFF").save(source)
            stdout = io.StringIO()
            with (
                patch(
                    "scripts.run_real_generative_editable_pptx.export_route"
                    "._build_generative_editable_pipeline_dependencies",
                    return_value=export_route._build_fake_generative_editable_pipeline_dependencies(),
                ),
                patch(
                    "scripts.run_real_generative_editable_pptx.run_generative_editable_pipeline",
                    side_effect=fake_pipeline,
                ),
                patch(
                    "scripts.run_real_generative_editable_pptx._write_preview_reports",
                    return_value=[],
                ),
                redirect_stdout(stdout),
            ):
                exit_code = main(
                    [
                        "run",
                        "--input-images",
                        str(source),
                        "--output-dir",
                        str(root / "out"),
                        "--job-id",
                        "fallback-policy",
                        "--fallback-policy",
                        "raster_pptx",
                    ]
                )
            fallback_path = captured["fallback_output_factory"]()

        self.assertEqual(exit_code, 0)
        self.assertEqual(captured["fallback_policy"], "raster_pptx")
        self.assertEqual(
            fallback_path,
            str((root / "out" / "fallback-policy.raster-fallback.pptx").resolve()),
        )

    def test_run_mode_can_route_to_vlm_first_pipeline(self):
        import scripts.run_real_generative_editable_pptx as runner

        captured = {}
        dependency_sentinel = object()

        def fake_vlm_pipeline(**kwargs):
            from pptx import Presentation

            from src.generative_editable_manifest import DeckManifest, write_manifest

            captured.update(kwargs)
            output_path = Path(kwargs["output_path"])
            output_path.parent.mkdir(parents=True, exist_ok=True)
            Presentation().save(output_path)
            job_dir = Path(kwargs["artifact_root"]) / kwargs["job_id"]
            job_dir.mkdir(parents=True, exist_ok=True)
            (job_dir / "stage-events.jsonl").write_text(
                json.dumps({"stage": "compose_deck", "status": "finished"}) + "\n",
                encoding="utf-8",
            )
            write_manifest(
                job_dir / "deck.json",
                DeckManifest(
                    job_id=kwargs["job_id"],
                    slide_order=[],
                    aspect_ratio="16:9",
                    provider_roles={"vlm": "VLM", "image_edit": "edit_model"},
                    quality_settings={},
                    fallback_policy="fail",
                    page_manifest_paths=[],
                    validation_status="passed",
                ),
            )
            return type(
                "Result",
                (),
                {
                    "status": "passed",
                    "output_path": str(output_path),
                    "fallback_used": "",
                },
            )()

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "slide.png"
            Image.new("RGB", (800, 450), "#FFFFFF").save(source)
            stdout = io.StringIO()
            with (
                patch.object(runner, "_vlm_dependencies", return_value=dependency_sentinel),
                patch.object(runner, "run_vlm_editable_pptx_pipeline", side_effect=fake_vlm_pipeline),
                patch.object(runner, "_write_preview_reports", return_value=[]),
                redirect_stdout(stdout),
            ):
                exit_code = main(
                    [
                        "run",
                        "--mode",
                        "vlm_first",
                        "--input-images",
                        str(source),
                        "--output-dir",
                        str(root / "out"),
                        "--job-id",
                        "vlm-test",
                    ]
                )
            line = json.loads(stdout.getvalue().splitlines()[-1])
            report = json.loads(Path(line["report_path"]).read_text(encoding="utf-8"))

        self.assertEqual(exit_code, 0)
        self.assertEqual(line["status"], "passed")
        self.assertIs(captured["dependencies"], dependency_sentinel)
        self.assertEqual(captured["slides"][0]["slide_id"], "slide-1")
        self.assertEqual(captured["aspect_ratio"], "16:9")
        self.assertEqual(report["last_stage_event"]["stage"], "compose_deck")
        self.assertEqual(report["last_stage_event"]["status"], "finished")

    def test_vlm_first_run_can_override_provider_retry_settings(self):
        import scripts.run_real_generative_editable_pptx as runner

        captured = {}
        dependency_sentinel = object()

        def fake_vlm_dependencies(**kwargs):
            captured.update(kwargs)
            return dependency_sentinel

        def fake_vlm_pipeline(**kwargs):
            from pptx import Presentation

            from src.generative_editable_manifest import DeckManifest, write_manifest

            output_path = Path(kwargs["output_path"])
            output_path.parent.mkdir(parents=True, exist_ok=True)
            Presentation().save(output_path)
            job_dir = Path(kwargs["artifact_root"]) / kwargs["job_id"]
            job_dir.mkdir(parents=True, exist_ok=True)
            write_manifest(
                job_dir / "deck.json",
                DeckManifest(
                    job_id=kwargs["job_id"],
                    slide_order=[],
                    aspect_ratio="16:9",
                    provider_roles={"vlm": "VLM", "image_edit": "edit_model"},
                    quality_settings={},
                    fallback_policy="fail",
                    page_manifest_paths=[],
                    validation_status="passed",
                ),
            )
            return type(
                "Result",
                (),
                {
                    "status": "passed",
                    "output_path": str(output_path),
                    "fallback_used": "",
                },
            )()

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "slide.png"
            Image.new("RGB", (800, 450), "#FFFFFF").save(source)
            stdout = io.StringIO()
            with (
                patch.object(runner, "_vlm_dependencies", side_effect=fake_vlm_dependencies),
                patch.object(runner, "run_vlm_editable_pptx_pipeline", side_effect=fake_vlm_pipeline),
                patch.object(runner, "_write_preview_reports", return_value=[]),
                redirect_stdout(stdout),
            ):
                exit_code = main(
                    [
                        "run",
                        "--mode",
                        "vlm_first",
                        "--input-images",
                        str(source),
                        "--output-dir",
                        str(root / "out"),
                        "--job-id",
                        "vlm-retry",
                        "--provider-max-attempts",
                        "6",
                        "--provider-retry-backoff",
                        "12.5",
                        "--pipeline-page-timeout",
                        "240",
                    ]
                )

        self.assertEqual(exit_code, 0)
        self.assertEqual(captured["provider_max_attempts"], 6)
        self.assertEqual(captured["provider_retry_backoff_seconds"], 12.5)
        self.assertEqual(captured["page_timeout_seconds"], 240)

    def test_vlm_first_run_marks_preview_failure_as_failed(self):
        import scripts.run_real_generative_editable_pptx as runner

        def fake_vlm_pipeline(**kwargs):
            from pptx import Presentation

            from src.generative_editable_manifest import DeckManifest, write_manifest

            output_path = Path(kwargs["output_path"])
            output_path.parent.mkdir(parents=True, exist_ok=True)
            Presentation().save(output_path)
            job_dir = Path(kwargs["artifact_root"]) / kwargs["job_id"]
            job_dir.mkdir(parents=True, exist_ok=True)
            write_manifest(
                job_dir / "deck.json",
                DeckManifest(
                    job_id=kwargs["job_id"],
                    slide_order=[],
                    aspect_ratio="16:9",
                    provider_roles={"vlm": "VLM"},
                    quality_settings={},
                    fallback_policy="fail",
                    page_manifest_paths=[],
                    validation_status="passed",
                ),
            )
            return type(
                "Result",
                (),
                {
                    "status": "passed",
                    "output_path": str(output_path),
                    "fallback_used": "",
                },
            )()

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "slide.png"
            Image.new("RGB", (800, 450), "#FFFFFF").save(source)
            stdout = io.StringIO()
            with (
                patch.object(runner, "_vlm_dependencies", return_value=object()),
                patch.object(runner, "run_vlm_editable_pptx_pipeline", side_effect=fake_vlm_pipeline),
                patch.object(
                    runner,
                    "_write_preview_reports",
                    return_value=[{"status": "failed", "issues": [{"code": "preview_similarity_failed"}]}],
                ),
                redirect_stdout(stdout),
            ):
                exit_code = main(
                    [
                        "run",
                        "--mode",
                        "vlm_first",
                        "--input-images",
                        str(source),
                        "--output-dir",
                        str(root / "out"),
                        "--job-id",
                        "vlm-preview-fail",
                    ]
                )

        self.assertEqual(exit_code, 1)
        line = json.loads(stdout.getvalue().splitlines()[-1])
        self.assertEqual(line["status"], "failed")

    def test_run_mode_fails_when_output_is_only_full_slide_picture(self):
        import shutil

        from api.routes import export as export_route
        from pptx import Presentation

        from src.generative_editable_manifest import DeckManifest, PageManifest, TextBoxSpec, write_manifest

        def fake_pipeline(**kwargs):
            source_path = Path(kwargs["slides"][0].image_path)
            output_path = Path(kwargs["output_path"])
            output_path.parent.mkdir(parents=True, exist_ok=True)
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(
                str(source_path),
                0,
                0,
                width=prs.slide_width,
                height=prs.slide_height,
            )
            prs.save(output_path)

            job_dir = Path(kwargs["artifact_root"]) / kwargs["job_id"]
            sources_dir = job_dir / "sources"
            pages_dir = job_dir / "pages"
            sources_dir.mkdir(parents=True, exist_ok=True)
            pages_dir.mkdir(parents=True, exist_ok=True)
            shutil.copyfile(source_path, sources_dir / "source.png")
            page = PageManifest(
                slide_id="slide-1",
                page_index=0,
                source_image_path="sources/source.png",
                source_image_size=(800, 450),
                slide_size=(10.0, 5.625),
                text_boxes=[
                    TextBoxSpec(
                        text="Editable text",
                        source_pixel_bbox=(80, 40, 360, 90),
                        source_pixel_polygon=((80, 40), (360, 40), (360, 90), (80, 90)),
                    )
                ],
                validation_status="passed",
            )
            write_manifest(pages_dir / "0000-slide-1.json", page)
            write_manifest(
                job_dir / "deck.json",
                DeckManifest(
                    job_id=kwargs["job_id"],
                    slide_order=["slide-1"],
                    aspect_ratio="16:9",
                    provider_roles={},
                    quality_settings={},
                    fallback_policy="fail",
                    page_manifest_paths=["pages/0000-slide-1.json"],
                    validation_status="passed",
                ),
            )
            return type(
                "Result",
                (),
                {
                    "status": "passed",
                    "output_path": str(output_path),
                    "fallback_used": False,
                },
            )()

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "slide.png"
            Image.new("RGB", (800, 450), "#FFFFFF").save(source)
            stdout = io.StringIO()
            with (
                patch(
                    "scripts.run_real_generative_editable_pptx.export_route"
                    "._build_generative_editable_pipeline_dependencies",
                    return_value=export_route._build_fake_generative_editable_pipeline_dependencies(),
                ),
                patch(
                    "scripts.run_real_generative_editable_pptx.run_generative_editable_pipeline",
                    side_effect=fake_pipeline,
                ),
                patch(
                    "scripts.run_real_generative_editable_pptx._write_preview_reports",
                    return_value=[],
                ),
                redirect_stdout(stdout),
            ):
                exit_code = main(
                    [
                        "run",
                        "--input-images",
                        str(source),
                        "--output-dir",
                        str(root / "out"),
                        "--job-id",
                        "full-slide-only",
                    ]
                )

            result = json.loads(stdout.getvalue().splitlines()[-1])
            report = json.loads(Path(result["report_path"]).read_text(encoding="utf-8"))

        self.assertEqual(exit_code, 1)
        self.assertEqual(report["status"], "failed")
        self.assertIn("no_decomposed_visual_elements", report["reconstruction_issues"][0]["code"])
        self.assertEqual(report["object_stats"]["slides"][0]["full_slide_picture_count"], 1)

    def test_single_image_edit_gate_checks_only_requested_role(self):
        config = load_generative_editable_config(use_fake=True)
        calls: list[str] = []

        class RecordingImageEditProvider(FakeImageEditProvider):
            def edit(self, request):
                calls.append(self.config.role)
                return super().edit(request)

        dependencies = GenerativeEditablePipelineDependencies(
            ocr_provider=FakeOCRProvider(config.ocr),
            image_edit_provider=RecordingImageEditProvider(
                ProviderConfig(
                    role="clean_base_model",
                    provider="fake_image_edit",
                    model="fake-edit",
                    base_url="",
                    api_key="",
                )
            ),
            asset_sheet_image_edit_provider=RecordingImageEditProvider(
                ProviderConfig(
                    role="asset_sheet_model",
                    provider="fake_image_edit",
                    model="fake-edit",
                    base_url="",
                    api_key="",
                )
            ),
            repair_image_edit_provider=RecordingImageEditProvider(
                ProviderConfig(
                    role="repair_model",
                    provider="fake_image_edit",
                    model="fake-edit",
                    base_url="",
                    api_key="",
                )
            ),
            image_generation_provider=FakeImageGenerationProvider(config.generation_model),
        )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "slide.png"
            Image.new("RGB", (800, 450), "#FFFFFF").save(source)
            stdout = io.StringIO()
            with (
                patch(
                    "scripts.run_real_generative_editable_pptx.export_route"
                    "._build_fake_generative_editable_pipeline_dependencies",
                    return_value=dependencies,
                ),
                redirect_stdout(stdout),
            ):
                exit_code = main(
                    [
                        "gates",
                        "--use-fake",
                        "--input-images",
                        str(source),
                        "--output-dir",
                        str(root / "out"),
                        "--provider-gate",
                        "image_edit_repair",
                    ]
                )

        self.assertEqual(exit_code, 0)
        self.assertEqual(calls, ["repair_model"])
        result = json.loads(stdout.getvalue().splitlines()[-1])
        self.assertEqual(result["status"], "passed")
        self.assertEqual(result["gates"][0]["gate"], "image_edit_repair")

    def test_provider_gates_do_not_retry_by_default(self):
        config = load_generative_editable_config(use_fake=True)

        class FlakyImageGenerationProvider(ImageGenerationProvider):
            def __init__(self):
                super().__init__(config.generation_model)
                self.calls = 0

            def generate(self, request):
                self.calls += 1
                raise ProviderError(
                    provider_role=self.config.role,
                    operation=request.prompt_id,
                    message="temporary provider failure",
                    retryable=True,
                )

        provider = FlakyImageGenerationProvider()
        dependencies = GenerativeEditablePipelineDependencies(
            ocr_provider=FakeOCRProvider(config.ocr),
            image_edit_provider=FakeImageEditProvider(config.clean_base_model),
            image_generation_provider=provider,
            provider_max_attempts=3,
        )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "slide.png"
            Image.new("RGB", (800, 450), "#FFFFFF").save(source)
            stdout = io.StringIO()
            with (
                patch(
                    "scripts.run_real_generative_editable_pptx.export_route"
                    "._build_fake_generative_editable_pipeline_dependencies",
                    return_value=dependencies,
                ),
                redirect_stdout(stdout),
            ):
                exit_code = main(
                    [
                        "gates",
                        "--use-fake",
                        "--input-images",
                        str(source),
                        "--output-dir",
                        str(root / "out"),
                        "--provider-gate",
                        "image_generation",
                    ]
                )

        self.assertEqual(exit_code, 1)
        self.assertEqual(provider.calls, 1)

    def test_provider_gates_include_attempt_history_without_retry_override(self):
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "slide.png"
            Image.new("RGB", (800, 450), "#FFFFFF").save(source)
            stdout = io.StringIO()
            with redirect_stdout(stdout):
                exit_code = main(
                    [
                        "gates",
                        "--use-fake",
                        "--input-images",
                        str(source),
                        "--output-dir",
                        str(root / "out"),
                        "--provider-gate",
                        "image_generation",
                    ]
                )
            report = json.loads((root / "out" / "provider-gates.json").read_text(encoding="utf-8"))

        self.assertEqual(exit_code, 0)
        self.assertEqual(report["gates"][0]["provider_attempts"][0]["attempt"], 1)
        self.assertEqual(report["gates"][0]["provider_attempts"][0]["status"], "passed")

    def test_provider_gate_report_includes_upstream_error_code(self):
        config = load_generative_editable_config(use_fake=True)

        class QuotaImageGenerationProvider(ImageGenerationProvider):
            def __init__(self):
                super().__init__(config.generation_model)

            def generate(self, request):
                raise ProviderError(
                    provider_role=self.config.role,
                    operation=request.prompt_id,
                    message="quota exhausted",
                    retryable=False,
                    status_code=403,
                    provider_error_code="insufficient_user_quota",
                )

        dependencies = GenerativeEditablePipelineDependencies(
            ocr_provider=FakeOCRProvider(config.ocr),
            image_edit_provider=FakeImageEditProvider(config.clean_base_model),
            image_generation_provider=QuotaImageGenerationProvider(),
        )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "slide.png"
            Image.new("RGB", (800, 450), "#FFFFFF").save(source)
            stdout = io.StringIO()
            with (
                patch(
                    "scripts.run_real_generative_editable_pptx.export_route"
                    "._build_fake_generative_editable_pipeline_dependencies",
                    return_value=dependencies,
                ),
                redirect_stdout(stdout),
            ):
                exit_code = main(
                    [
                        "gates",
                        "--use-fake",
                        "--input-images",
                        str(source),
                        "--output-dir",
                        str(root / "out"),
                        "--provider-gate",
                        "image_generation",
                    ]
                )
            report = json.loads((root / "out" / "provider-gates.json").read_text(encoding="utf-8"))

        self.assertEqual(exit_code, 1)
        self.assertEqual(report["gates"][0]["status_code"], 403)
        self.assertEqual(report["gates"][0]["provider_error_code"], "insufficient_user_quota")

    def test_provider_gates_use_cli_retry_override(self):
        config = load_generative_editable_config(use_fake=True)

        class FlakyImageGenerationProvider(ImageGenerationProvider):
            def __init__(self):
                super().__init__(config.generation_model)
                self.calls = 0
                self.fake = FakeImageGenerationProvider(config.generation_model)

            def generate(self, request):
                self.calls += 1
                if self.calls == 1:
                    raise ProviderError(
                        provider_role=self.config.role,
                        operation=request.prompt_id,
                        message="temporary provider failure",
                        retryable=True,
                    )
                return self.fake.generate(request)

        provider = FlakyImageGenerationProvider()
        dependencies = GenerativeEditablePipelineDependencies(
            ocr_provider=FakeOCRProvider(config.ocr),
            image_edit_provider=FakeImageEditProvider(config.clean_base_model),
            image_generation_provider=provider,
            provider_max_attempts=1,
        )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "slide.png"
            Image.new("RGB", (800, 450), "#FFFFFF").save(source)
            stdout = io.StringIO()
            with (
                patch(
                    "scripts.run_real_generative_editable_pptx.export_route"
                    "._build_fake_generative_editable_pipeline_dependencies",
                    return_value=dependencies,
                ),
                redirect_stdout(stdout),
            ):
                exit_code = main(
                    [
                        "gates",
                        "--use-fake",
                        "--input-images",
                        str(source),
                        "--output-dir",
                        str(root / "out"),
                        "--provider-gate",
                        "image_generation",
                        "--provider-max-attempts",
                        "2",
                        "--provider-retry-backoff",
                        "0",
                    ]
                )

        self.assertEqual(exit_code, 0)
        self.assertEqual(provider.calls, 2)

    def test_provider_gate_report_includes_retry_attempt_history(self):
        config = load_generative_editable_config(use_fake=True)

        class FlakyImageGenerationProvider(ImageGenerationProvider):
            def __init__(self):
                super().__init__(config.generation_model)
                self.calls = 0
                self.fake = FakeImageGenerationProvider(config.generation_model)

            def generate(self, request):
                self.calls += 1
                if self.calls == 1:
                    raise ProviderError(
                        provider_role=self.config.role,
                        operation=request.prompt_id,
                        message="temporary api_key=secret failure",
                        retryable=True,
                        status_code=500,
                        provider_error_code="upstream_timeout",
                    )
                return self.fake.generate(request)

        provider = FlakyImageGenerationProvider()
        dependencies = GenerativeEditablePipelineDependencies(
            ocr_provider=FakeOCRProvider(config.ocr),
            image_edit_provider=FakeImageEditProvider(config.clean_base_model),
            image_generation_provider=provider,
            provider_max_attempts=1,
        )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "slide.png"
            Image.new("RGB", (800, 450), "#FFFFFF").save(source)
            stdout = io.StringIO()
            with (
                patch(
                    "scripts.run_real_generative_editable_pptx.export_route"
                    "._build_fake_generative_editable_pipeline_dependencies",
                    return_value=dependencies,
                ),
                redirect_stdout(stdout),
            ):
                exit_code = main(
                    [
                        "gates",
                        "--use-fake",
                        "--input-images",
                        str(source),
                        "--output-dir",
                        str(root / "out"),
                        "--provider-gate",
                        "image_generation",
                        "--provider-max-attempts",
                        "2",
                        "--provider-retry-backoff",
                        "0",
                    ]
                )
            report = json.loads((root / "out" / "provider-gates.json").read_text(encoding="utf-8"))

        self.assertEqual(exit_code, 0)
        gate = report["gates"][0]
        self.assertEqual(
            [(item["attempt"], item["status"]) for item in gate["provider_attempts"]],
            [(1, "failed"), (2, "passed")],
        )
        self.assertTrue(gate["provider_attempts"][0]["retrying"])
        self.assertEqual(gate["provider_attempts"][0]["status_code"], 500)
        self.assertEqual(gate["provider_attempts"][0]["provider_error_code"], "upstream_timeout")
        self.assertNotIn("secret", json.dumps(gate["provider_attempts"], ensure_ascii=False))

    def test_provider_gates_can_run_vlm_analysis_gate(self):
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "slide.png"
            Image.new("RGB", (800, 450), "#FFFFFF").save(source)
            stdout = io.StringIO()
            with redirect_stdout(stdout):
                exit_code = main(
                    [
                        "gates",
                        "--use-fake",
                        "--input-images",
                        str(source),
                        "--output-dir",
                        str(root / "out"),
                        "--provider-gate",
                        "vlm_analysis",
                    ]
                )

        self.assertEqual(exit_code, 0)
        result = json.loads(stdout.getvalue().splitlines()[-1])
        self.assertEqual(result["status"], "passed")
        self.assertEqual(result["gates"][0]["gate"], "vlm_analysis")
        self.assertGreaterEqual(result["gates"][0]["text_region_count"], 1)

    def test_provider_gate_report_includes_vlm_payload_attempts(self):
        import scripts.run_real_generative_editable_pptx as runner

        from src.generative_editable_vlm_reconstruction import (
            VLMEditablePipelineDependencies,
            VLMPageAnalysisProvider,
            coerce_vlm_analysis,
            with_vlm_provider_retries,
        )

        config = load_generative_editable_config(use_fake=True)

        class PayloadFallbackVLMProvider(VLMPageAnalysisProvider):
            def __init__(self):
                self.last_payload_attempts = []

            def analyze_page(self, image_path: str, *, timeout_seconds: int = 180):
                self.last_payload_attempts = [
                    {
                        "attempt": 1,
                        "status": "failed",
                        "error": "VLM.vlm_page_analysis: 503 no healthy upstream",
                        "retryable": True,
                        "retrying": True,
                    },
                    {"attempt": 2, "status": "passed"},
                ]
                return coerce_vlm_analysis(
                    {
                        "coordinate_space": {"width": 160, "height": 90, "unit": "px"},
                        "text_regions": [{"id": "t1", "text": "标题", "bbox": [10, 10, 100, 30]}],
                        "bitmap_regions": [],
                        "shape_regions": [],
                    }
                )

        dependencies = with_vlm_provider_retries(
            VLMEditablePipelineDependencies(
                vlm_provider=PayloadFallbackVLMProvider(),
                image_edit_provider=FakeImageEditProvider(config.clean_base_model),
                provider_max_attempts=1,
            )
        )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "slide.png"
            Image.new("RGB", (800, 450), "#FFFFFF").save(source)
            stdout = io.StringIO()
            with (
                patch.object(
                    runner,
                    "_gate_dependencies",
                    return_value={
                        "legacy": object(),
                        "vlm": dependencies,
                    },
                ),
                redirect_stdout(stdout),
            ):
                exit_code = main(
                    [
                        "gates",
                        "--use-fake",
                        "--input-images",
                        str(source),
                        "--output-dir",
                        str(root / "out"),
                        "--provider-gate",
                        "vlm_analysis",
                    ]
                )
            report = json.loads((root / "out" / "provider-gates.json").read_text(encoding="utf-8"))

        self.assertEqual(exit_code, 0)
        attempts = report["gates"][0]["provider_attempts"]
        self.assertEqual(attempts[0]["status"], "passed")
        self.assertEqual(
            [(item["attempt"], item["status"]) for item in attempts[0]["payload_attempts"]],
            [(1, "failed"), (2, "passed")],
        )

    def test_subprocess_json_runner_times_out_and_writes_report(self):
        import scripts.run_real_generative_editable_pptx as runner

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            report_path = root / "report.json"

            result = runner._run_subprocess_json(
                [
                    sys.executable,
                    "-c",
                    "import time; time.sleep(2)",
                ],
                timeout_seconds=1,
                report_path=report_path,
                timeout_payload={
                    "status": "failed",
                    "stage": "test",
                },
            )

            report = json.loads(report_path.read_text(encoding="utf-8"))

        self.assertEqual(result["status"], "failed")
        self.assertEqual(result["error_type"], "TimeoutExpired")
        self.assertEqual(report["stage"], "test")
        self.assertEqual(report["error_type"], "TimeoutExpired")

    def test_subprocess_timeout_report_includes_existing_pptx_object_stats(self):
        import scripts.run_real_generative_editable_pptx as runner

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            pptx_path = root / "partial.pptx"
            report_path = root / "report.json"
            code = (
                "import sys,time;"
                "from pptx import Presentation;"
                "p=sys.argv[1];"
                "prs=Presentation();"
                "prs.slides.add_slide(prs.slide_layouts[6]);"
                "prs.save(p);"
                "time.sleep(2)"
            )

            result = runner._run_subprocess_json(
                [sys.executable, "-c", code, str(pptx_path)],
                timeout_seconds=1,
                report_path=report_path,
                timeout_payload={
                    "status": "failed",
                    "output_path": str(pptx_path),
                },
            )

            report = json.loads(report_path.read_text(encoding="utf-8"))

        self.assertEqual(result["error_type"], "TimeoutExpired")
        self.assertEqual(report["object_stats"]["slide_count"], 1)

    def test_timeout_augmentation_skips_oversized_pptx(self):
        import scripts.run_real_generative_editable_pptx as runner

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            pptx_path = root / "large.pptx"
            pptx_path.write_bytes(b"0" * (runner.MAX_TIMEOUT_AUGMENT_PPTX_BYTES + 1))

            result = runner._augment_timeout_result_from_artifacts(
                {
                    "status": "failed",
                    "output_path": str(pptx_path),
                }
            )

        self.assertNotIn("object_stats", result)
        self.assertEqual(result["object_stats_error"], "pptx exceeds timeout augmentation size cap")

    def test_timeout_augmentation_removes_missing_output_path(self):
        import scripts.run_real_generative_editable_pptx as runner

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)

            result = runner._augment_timeout_result_from_artifacts(
                {
                    "status": "failed",
                    "output_path": str(root / "missing.pptx"),
                    "artifact_root": str(root / "artifacts"),
                }
            )

        self.assertNotIn("output_path", result)

    def test_timeout_augmentation_includes_last_stage_event(self):
        import scripts.run_real_generative_editable_pptx as runner

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            artifact_root = root / "artifacts"
            artifact_root.mkdir()
            (artifact_root / "stage-events.jsonl").write_text(
                "\n".join(
                    [
                        json.dumps(
                            {
                                "stage": "ocr",
                                "status": "finished",
                                "elapsed_ms": 1000,
                            }
                        ),
                        json.dumps(
                            {
                                "stage": "text_clean_background.provider_edit",
                                "status": "started",
                                "elapsed_ms": 0,
                                "provider_role": "edit_model",
                            }
                        ),
                    ]
                )
                + "\n",
                encoding="utf-8",
            )

            result = runner._augment_timeout_result_from_artifacts(
                {
                    "status": "failed",
                    "artifact_root": str(artifact_root),
                }
            )

        self.assertEqual(
            result["last_stage_event"]["stage"],
            "text_clean_background.provider_edit",
        )
        self.assertEqual(result["active_stage"], "text_clean_background.provider_edit")

    def test_timeout_augmentation_salvages_stage_events_before_malformed_line(self):
        import scripts.run_real_generative_editable_pptx as runner

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            artifact_root = root / "artifacts"
            artifact_root.mkdir()
            (artifact_root / "stage-events.jsonl").write_text(
                json.dumps(
                    {
                        "stage": "base_clean_background",
                        "status": "started",
                    }
                )
                + "\n{bad-json",
                encoding="utf-8",
            )

            result = runner._augment_timeout_result_from_artifacts(
                {
                    "status": "failed",
                    "artifact_root": str(artifact_root),
                }
            )

        self.assertEqual(result["active_stage"], "base_clean_background")
        self.assertEqual(result["stage_events_error"], "malformed stage-events line skipped")

    def test_failed_pipeline_report_includes_existing_pptx_diagnostics(self):
        import scripts.run_real_generative_editable_pptx as runner
        from pptx import Presentation

        from src.generative_editable_manifest import (
            DeckManifest,
            PageManifest,
            TextBoxSpec,
            write_manifest,
        )
        from src.generative_editable_pipeline import GenerativeEditableValidationError
        from src.generative_editable_preview_validator import ValidationIssue, ValidationReport

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            artifact_root = root / "artifacts"
            job_dir = artifact_root / "job"
            (job_dir / "pages").mkdir(parents=True)
            ocr_dir = job_dir / "provider_outputs" / "ocr" / "0000-slide-1"
            ocr_dir.mkdir(parents=True)
            (ocr_dir / "ocr.json").write_text(
                '{"items":[{"text":"a"},{"text":"b"}]}',
                encoding="utf-8",
            )
            output_path = root / "job.pptx"
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_textbox(100, 100, 1000, 300).text = "核心架构"
            prs.save(output_path)
            page = PageManifest(
                slide_id="slide-1",
                page_index=0,
                source_image_path="sources/slide-1.png",
                source_image_size=(800, 450),
                slide_size=(10.0, 5.625),
                text_boxes=[
                    TextBoxSpec(
                        text="核心架构",
                        source_pixel_bbox=(10, 10, 120, 40),
                        source_pixel_polygon=((10, 10), (120, 10), (120, 40), (10, 40)),
                    )
                ],
            )
            write_manifest(job_dir / "pages" / "0000-slide-1.json", page)
            write_manifest(
                job_dir / "deck.json",
                DeckManifest(
                    job_id="job",
                    slide_order=["slide-1"],
                    aspect_ratio="16:9",
                    provider_roles={
                        "ocr": "ocr_model",
                        "image_edit": "edit_model",
                        "image_generation": "image_model",
                    },
                    quality_settings={},
                    fallback_policy="fail",
                    page_manifest_paths=["pages/0000-slide-1.json"],
                ),
            )
            error = GenerativeEditableValidationError(
                validation_report=ValidationReport(
                    status="failed",
                    checked_pages=1,
                    issues=[
                        ValidationIssue(
                            code="preview_similarity_failed",
                            message="preview differs",
                        )
                    ],
                ),
                fallback_policy="fail",
            )

            with patch(
                "scripts.run_real_generative_editable_pptx._write_preview_reports",
                return_value=[
                    {
                        "page": 1,
                        "status": "failed",
                        "diff_metrics": {
                            "mean_abs_delta": 23.9,
                            "changed_pixel_ratio": 0.16,
                        },
                    }
                ],
            ) as preview_reports:
                report = runner._failed_pipeline_report_from_artifacts(
                    error=error,
                    output_dir=root,
                    output_path=output_path,
                    artifact_root=artifact_root,
                    job_id="job",
                    input_images=[root / "slide_1.png"],
                )

        self.assertEqual(report["status"], "failed")
        self.assertEqual(report["object_stats"]["totals"]["TEXT_BOX"], 1)
        self.assertEqual(report["validation_issues"], ["preview_similarity_failed"])
        self.assertEqual(report["deck_manifest_path"], str(job_dir / "deck.json"))
        self.assertEqual(report["preview_reports"][0]["status"], "failed")
        preview_reports.assert_called_once()

    def test_failed_pipeline_report_includes_provider_error_metadata(self):
        import scripts.run_real_generative_editable_pptx as runner

        error = ProviderError(
            provider_role="VLM",
            operation="vlm_page_analysis",
            message="quota exhausted",
            retryable=False,
            status_code=403,
            provider_error_code="insufficient_user_quota",
        )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            report = runner._failed_pipeline_report_from_artifacts(
                error=error,
                output_dir=root,
                output_path=None,
                artifact_root=root / "artifacts",
                job_id="job",
                input_images=[root / "slide_1.png"],
            )

        self.assertEqual(report["status_code"], 403)
        self.assertEqual(report["provider_error_code"], "insufficient_user_quota")
        self.assertFalse(report["retryable"])

    def test_preview_report_rejects_manifest_stub_renderer(self):
        import scripts.run_real_generative_editable_pptx as runner

        from src.generative_editable_manifest import DeckManifest, PageManifest, write_manifest
        from src.generative_editable_preview_validator import PreviewRenderResult

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            artifact_root = root / "artifacts"
            (artifact_root / "pages").mkdir(parents=True)
            source_path = root / "slide_1.png"
            source_image = Image.new("RGB", (800, 450), "#020817")
            source_image.save(source_path)
            page = PageManifest(
                slide_id="slide-1",
                page_index=0,
                source_image_path="sources/slide-1.png",
                source_image_size=(800, 450),
                slide_size=(10.0, 5.625),
            )
            write_manifest(artifact_root / "pages" / "0000-slide-1.json", page)
            deck_path = artifact_root / "deck.json"
            write_manifest(
                deck_path,
                DeckManifest(
                    job_id="job",
                    slide_order=["slide-1"],
                    aspect_ratio="16:9",
                    provider_roles={},
                    quality_settings={},
                    fallback_policy="fail",
                    page_manifest_paths=["pages/0000-slide-1.json"],
                ),
            )

            with patch(
                "scripts.run_real_generative_editable_pptx.render_powerpoint_preview_with_metadata",
                return_value=PreviewRenderResult(
                    image=source_image.copy(),
                    metadata={"renderer": "manifest_stub", "is_powerpoint_render": False},
                ),
            ):
                reports = runner._write_preview_reports(
                    source_images=[source_path],
                    pptx_path=root / "deck.pptx",
                    deck_manifest_path=deck_path,
                    artifact_root=artifact_root,
                    similarity_threshold=0.8,
                    output_dir=root / "previews",
                )

        self.assertEqual(reports[0]["status"], "failed")
        self.assertEqual(reports[0]["renderer"], "manifest_stub")
        self.assertEqual(reports[0]["issues"][0]["code"], "preview_renderer_not_powerpoint")

    def test_reconstruction_issues_fail_for_source_preserving_fast_path_degradation(self):
        import scripts.run_real_generative_editable_pptx as runner
        from pptx import Presentation

        from src.generative_editable_manifest import (
            BitmapAssetSpec,
            DeckManifest,
            PageManifest,
            TextBoxSpec,
            write_manifest,
        )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            artifact_root = root / "artifacts"
            job_dir = artifact_root / "job"
            (job_dir / "pages").mkdir(parents=True)
            ocr_dir = job_dir / "provider_outputs" / "ocr" / "0000-slide-1"
            ocr_dir.mkdir(parents=True)
            (ocr_dir / "ocr.json").write_text(
                '{"items":[{"text":"a"},{"text":"b"}]}',
                encoding="utf-8",
            )
            output_path = root / "job.pptx"
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(str(_solid_png(root / "background.png")), 0, 0, width=prs.slide_width, height=prs.slide_height)
            slide.shapes.add_picture(str(_solid_png(root / "anchor.png")), 100, 100, width=1000, height=600)
            slide.shapes.add_textbox(100, 100, 1000, 300).text = "核心架构"
            prs.save(output_path)
            write_manifest(
                job_dir / "pages" / "0000-slide-1.json",
                PageManifest(
                    slide_id="slide-1",
                    page_index=0,
                    source_image_path="sources/slide-1.png",
                    source_image_size=(800, 450),
                    slide_size=(10.0, 5.625),
                    text_boxes=[
                        TextBoxSpec(
                            text="核心架构",
                            source_pixel_bbox=(10, 10, 120, 40),
                            source_pixel_polygon=((10, 10), (120, 10), (120, 40), (10, 40)),
                        )
                    ],
                    bitmap_assets=[
                        BitmapAssetSpec(
                            asset_id="anchor-001",
                            source_pixel_bbox=(20, 20, 220, 160),
                            asset_path="assets/anchor-001.png",
                            z_order=1,
                            provenance={"asset_strategy": "source_preserving_anchor"},
                        )
                    ],
                    provenance={"chosen_background_kind": "source_preserving_text_fast_path"},
                ),
            )
            write_manifest(
                job_dir / "deck.json",
                DeckManifest(
                    job_id="job",
                    slide_order=["slide-1"],
                    aspect_ratio="16:9",
                    provider_roles={
                        "ocr": "ocr_model",
                        "image_edit": "edit_model",
                        "image_generation": "image_model",
                    },
                    quality_settings={},
                    fallback_policy="fail",
                    page_manifest_paths=["pages/0000-slide-1.json"],
                ),
            )
            stats = runner._pptx_object_stats(output_path)

            issues = runner._reconstruction_object_issues(
                deck_manifest_path=job_dir / "deck.json",
                artifact_root=job_dir,
                object_stats=stats,
            )

        self.assertEqual(issues[0]["code"], "source_preserving_fast_path_degraded")
        self.assertEqual(issues[0]["severity"], "error")

    def test_pptx_object_stats_count_slide_background_picture(self):
        import scripts.run_real_generative_editable_pptx as runner

        from src.generative_editable_composer import compose_deck_from_manifests
        from src.generative_editable_manifest import DeckManifest, PageManifest, TextBoxSpec, write_manifest

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            job_dir = root / "artifacts" / "job"
            (job_dir / "backgrounds").mkdir(parents=True)
            (job_dir / "pages").mkdir(parents=True)
            Image.new("RGB", (800, 450), "#FFFFFF").save(job_dir / "backgrounds" / "base.png")
            write_manifest(
                job_dir / "pages" / "0000-slide-1.json",
                PageManifest(
                    slide_id="slide-1",
                    page_index=0,
                    source_image_path="sources/source.png",
                    source_image_size=(800, 450),
                    slide_size=(10.0, 5.625),
                    chosen_background="backgrounds/base.png",
                    text_boxes=[
                        TextBoxSpec(
                            text="Editable",
                            source_pixel_bbox=(100, 100, 220, 140),
                            source_pixel_polygon=((100, 100), (220, 100), (220, 140), (100, 140)),
                        )
                    ],
                ),
            )
            write_manifest(
                job_dir / "deck.json",
                DeckManifest(
                    job_id="job",
                    slide_order=["slide-1"],
                    aspect_ratio="16:9",
                    provider_roles={},
                    quality_settings={},
                    fallback_policy="fail",
                    page_manifest_paths=["pages/0000-slide-1.json"],
                ),
            )
            output_path = root / "out.pptx"
            compose_deck_from_manifests(job_dir / "deck.json", job_dir, output_path)

            stats = runner._pptx_object_stats(output_path)

        self.assertEqual(stats["slides"][0]["background_picture_count"], 1)
        self.assertEqual(stats["totals"]["BACKGROUND_PICTURE"], 1)
        self.assertEqual(stats["slides"][0]["full_slide_picture_count"], 0)

    def test_reconstruction_issues_fail_for_background_only_raster_slide(self):
        import scripts.run_real_generative_editable_pptx as runner

        from src.generative_editable_composer import compose_deck_from_manifests
        from src.generative_editable_manifest import DeckManifest, PageManifest, write_manifest

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            job_dir = root / "artifacts" / "job"
            (job_dir / "backgrounds").mkdir(parents=True)
            (job_dir / "pages").mkdir(parents=True)
            Image.new("RGB", (800, 450), "#FFFFFF").save(job_dir / "backgrounds" / "base.png")
            write_manifest(
                job_dir / "pages" / "0000-slide-1.json",
                PageManifest(
                    slide_id="slide-1",
                    page_index=0,
                    source_image_path="sources/source.png",
                    source_image_size=(800, 450),
                    slide_size=(10.0, 5.625),
                    chosen_background="backgrounds/base.png",
                ),
            )
            write_manifest(
                job_dir / "deck.json",
                DeckManifest(
                    job_id="job",
                    slide_order=["slide-1"],
                    aspect_ratio="16:9",
                    provider_roles={},
                    quality_settings={},
                    fallback_policy="fail",
                    page_manifest_paths=["pages/0000-slide-1.json"],
                ),
            )
            output_path = root / "out.pptx"
            compose_deck_from_manifests(job_dir / "deck.json", job_dir, output_path)
            stats = runner._pptx_object_stats(output_path)

            issues = runner._reconstruction_object_issues(
                deck_manifest_path=job_dir / "deck.json",
                artifact_root=job_dir,
                object_stats=stats,
            )

        self.assertIn("no_decomposed_visual_elements", [issue["code"] for issue in issues])
        issue = next(issue for issue in issues if issue["code"] == "no_decomposed_visual_elements")
        self.assertEqual(issue["details"]["background_picture_count"], 1)
        self.assertEqual(issue["details"]["full_slide_picture_count"], 0)

    def test_reconstruction_issues_fail_for_low_opacity_source_overlay(self):
        import scripts.run_real_generative_editable_pptx as runner
        from pptx import Presentation

        from src.generative_editable_manifest import (
            DeckManifest,
            PageManifest,
            TextBoxSpec,
            write_manifest,
        )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            artifact_root = root / "artifacts"
            job_dir = artifact_root / "job"
            (job_dir / "pages").mkdir(parents=True)
            output_path = root / "job.pptx"
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(str(_solid_png(root / "background.png")), 0, 0, width=prs.slide_width, height=prs.slide_height)
            slide.shapes.add_textbox(100, 100, 1000, 300).text = "核心架构"
            prs.save(output_path)
            write_manifest(
                job_dir / "pages" / "0000-slide-1.json",
                PageManifest(
                    slide_id="slide-1",
                    page_index=0,
                    source_image_path="sources/slide-1.png",
                    source_image_size=(800, 450),
                    slide_size=(10.0, 5.625),
                    text_boxes=[
                        TextBoxSpec(
                            text="核心架构",
                            source_pixel_bbox=(10, 10, 120, 40),
                            source_pixel_polygon=((10, 10), (120, 10), (120, 40), (10, 40)),
                            opacity=0.1,
                        )
                    ],
                    provenance={
                        "chosen_background_kind": "source_preserving_low_opacity_text_overlay"
                    },
                    provider_output_paths={
                        "ocr": "provider_outputs/ocr/0000-slide-1/ocr.json",
                    },
                ),
            )
            write_manifest(
                job_dir / "deck.json",
                DeckManifest(
                    job_id="job",
                    slide_order=["slide-1"],
                    aspect_ratio="16:9",
                    provider_roles={
                        "ocr": "ocr_model",
                        "image_edit": "edit_model",
                        "image_generation": "image_model",
                    },
                    quality_settings={},
                    fallback_policy="fail",
                    page_manifest_paths=["pages/0000-slide-1.json"],
                ),
            )
            stats = runner._pptx_object_stats(output_path)

            issues = runner._reconstruction_object_issues(
                deck_manifest_path=job_dir / "deck.json",
                artifact_root=job_dir,
                object_stats=stats,
            )

        self.assertEqual([issue["code"] for issue in issues], ["source_preserving_low_opacity_overlay_degraded"])
        self.assertEqual(issues[0]["severity"], "error")

    def test_reconstruction_issues_fail_for_oversized_bitmap_asset_coverage(self):
        import scripts.run_real_generative_editable_pptx as runner

        from src.generative_editable_manifest import (
            BitmapAssetSpec,
            DeckManifest,
            PageManifest,
            write_manifest,
        )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            job_dir = root / "artifacts" / "job"
            (job_dir / "pages").mkdir(parents=True)
            write_manifest(
                job_dir / "pages" / "0000-slide-1.json",
                PageManifest(
                    slide_id="slide-1",
                    page_index=0,
                    source_image_path="sources/slide-1.png",
                    source_image_size=(1000, 500),
                    slide_size=(10.0, 5.0),
                    bitmap_assets=[
                        BitmapAssetSpec(
                            asset_id="huge-source-crop",
                            source_pixel_bbox=(0, 0, 850, 500),
                            asset_path="assets/huge-source-crop.png",
                            z_order=1,
                            provenance={"asset_strategy": "masked_source_element"},
                        )
                    ],
                ),
            )
            write_manifest(
                job_dir / "deck.json",
                DeckManifest(
                    job_id="job",
                    slide_order=["slide-1"],
                    aspect_ratio="16:9",
                    provider_roles={},
                    quality_settings={},
                    fallback_policy="fail",
                    page_manifest_paths=["pages/0000-slide-1.json"],
                ),
            )

            issues = runner._reconstruction_object_issues(
                deck_manifest_path=job_dir / "deck.json",
                artifact_root=job_dir,
                object_stats={
                    "slides": [
                        {
                            "shape_counts": {"PICTURE": 2},
                            "full_slide_picture_count": 1,
                            "non_full_slide_picture_count": 1,
                            "text_box_count": 0,
                        }
                    ]
                },
            )

        self.assertEqual(issues[0]["code"], "oversized_bitmap_asset_coverage")
        self.assertEqual(issues[0]["severity"], "error")

    def test_reconstruction_issues_warn_for_overlapping_duplicate_text_boxes(self):
        import scripts.run_real_generative_editable_pptx as runner

        from src.generative_editable_manifest import (
            DeckManifest,
            PageManifest,
            TextBoxSpec,
            write_manifest,
        )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            job_dir = root / "artifacts" / "job"
            (job_dir / "pages").mkdir(parents=True)
            write_manifest(
                job_dir / "pages" / "0000-slide-1.json",
                PageManifest(
                    slide_id="slide-1",
                    page_index=0,
                    source_image_path="sources/slide-1.png",
                    source_image_size=(1000, 500),
                    slide_size=(10.0, 5.0),
                    text_boxes=[
                        TextBoxSpec(
                            text="21英寸4K后舱屏",
                            source_pixel_bbox=(700, 100, 860, 130),
                            source_pixel_polygon=((700, 100), (860, 100), (860, 130), (700, 130)),
                        ),
                        TextBoxSpec(
                            text="21 英寸 4K 后舱屏",
                            source_pixel_bbox=(705, 106, 870, 138),
                            source_pixel_polygon=((705, 106), (870, 106), (870, 138), (705, 138)),
                        ),
                    ],
                ),
            )
            write_manifest(
                job_dir / "deck.json",
                DeckManifest(
                    job_id="job",
                    slide_order=["slide-1"],
                    aspect_ratio="16:9",
                    provider_roles={},
                    quality_settings={},
                    fallback_policy="fail",
                    page_manifest_paths=["pages/0000-slide-1.json"],
                ),
            )

            issues = runner._reconstruction_object_issues(
                deck_manifest_path=job_dir / "deck.json",
                artifact_root=job_dir,
                object_stats={
                    "slides": [
                        {
                            "shape_counts": {"TEXT_BOX": 2},
                            "full_slide_picture_count": 0,
                            "non_full_slide_picture_count": 0,
                            "text_box_count": 2,
                        }
                    ]
                },
            )

        self.assertEqual(issues[0]["code"], "overlapping_duplicate_text_boxes")
        self.assertEqual(issues[0]["severity"], "warning")
        self.assertEqual(issues[0]["details"]["duplicate_text"], "21英寸4K后舱屏")

    def test_reconstruction_issues_do_not_warn_for_overlapping_substring_text_boxes(self):
        import scripts.run_real_generative_editable_pptx as runner

        from src.generative_editable_manifest import (
            DeckManifest,
            PageManifest,
            TextBoxSpec,
            write_manifest,
        )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            job_dir = root / "artifacts" / "job"
            (job_dir / "pages").mkdir(parents=True)
            write_manifest(
                job_dir / "pages" / "0000-slide-1.json",
                PageManifest(
                    slide_id="slide-1",
                    page_index=0,
                    source_image_path="sources/slide-1.png",
                    source_image_size=(1000, 500),
                    slide_size=(10.0, 5.0),
                    text_boxes=[
                        TextBoxSpec(
                            text="Revenue",
                            source_pixel_bbox=(100, 100, 170, 130),
                            source_pixel_polygon=((100, 100), (170, 100), (170, 130), (100, 130)),
                        ),
                        TextBoxSpec(
                            text="Revenue Q1",
                            source_pixel_bbox=(104, 106, 220, 138),
                            source_pixel_polygon=((104, 106), (220, 106), (220, 138), (104, 138)),
                        ),
                    ],
                ),
            )
            write_manifest(
                job_dir / "deck.json",
                DeckManifest(
                    job_id="job",
                    slide_order=["slide-1"],
                    aspect_ratio="16:9",
                    provider_roles={},
                    quality_settings={},
                    fallback_policy="fail",
                    page_manifest_paths=["pages/0000-slide-1.json"],
                ),
            )

            issues = runner._reconstruction_object_issues(
                deck_manifest_path=job_dir / "deck.json",
                artifact_root=job_dir,
                object_stats={
                    "slides": [
                        {
                            "shape_counts": {"TEXT_BOX": 2},
                            "full_slide_picture_count": 0,
                            "non_full_slide_picture_count": 0,
                            "text_box_count": 2,
                        }
                    ]
                },
            )

        self.assertEqual(issues, [])

    def test_reconstruction_issues_treat_vlm_source_crop_as_source_preserved_asset(self):
        import scripts.run_real_generative_editable_pptx as runner

        from src.generative_editable_manifest import (
            BitmapAssetSpec,
            DeckManifest,
            PageManifest,
            write_manifest,
        )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            job_dir = root / "artifacts" / "job"
            (job_dir / "pages").mkdir(parents=True)
            write_manifest(
                job_dir / "pages" / "0000-slide-1.json",
                PageManifest(
                    slide_id="slide-1",
                    page_index=0,
                    source_image_path="sources/slide-1.png",
                    source_image_size=(1000, 500),
                    slide_size=(10.0, 5.0),
                    bitmap_assets=[
                        BitmapAssetSpec(
                            asset_id="huge-vlm-crop",
                            source_pixel_bbox=(0, 0, 850, 500),
                            asset_path="assets/huge-vlm-crop.png",
                            z_order=1,
                            provenance={"source_type": "vlm_source_crop"},
                        )
                    ],
                ),
            )
            write_manifest(
                job_dir / "deck.json",
                DeckManifest(
                    job_id="job",
                    slide_order=["slide-1"],
                    aspect_ratio="16:9",
                    provider_roles={},
                    quality_settings={},
                    fallback_policy="fail",
                    page_manifest_paths=["pages/0000-slide-1.json"],
                ),
            )

            issues = runner._reconstruction_object_issues(
                deck_manifest_path=job_dir / "deck.json",
                artifact_root=job_dir,
                object_stats={
                    "slides": [
                        {
                            "shape_counts": {"PICTURE": 1},
                            "full_slide_picture_count": 0,
                            "non_full_slide_picture_count": 1,
                            "text_box_count": 0,
                        }
                    ]
                },
            )

        self.assertEqual(issues[0]["code"], "oversized_bitmap_asset_coverage")
        self.assertEqual(issues[0]["severity"], "error")

    def test_reconstruction_issues_treat_source_preserved_crop_as_source_preserved_asset(self):
        import scripts.run_real_generative_editable_pptx as runner

        from src.generative_editable_manifest import (
            BitmapAssetSpec,
            DeckManifest,
            PageManifest,
            write_manifest,
        )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            job_dir = root / "artifacts" / "job"
            (job_dir / "pages").mkdir(parents=True)
            write_manifest(
                job_dir / "pages" / "0000-slide-1.json",
                PageManifest(
                    slide_id="slide-1",
                    page_index=0,
                    source_image_path="sources/slide-1.png",
                    source_image_size=(1000, 500),
                    slide_size=(10.0, 5.0),
                    bitmap_assets=[
                        BitmapAssetSpec(
                            asset_id="huge-source-preserved-crop",
                            source_pixel_bbox=(0, 0, 850, 500),
                            asset_path="assets/huge-source-preserved-crop.png",
                            z_order=1,
                            provenance={"asset_strategy": "source_preserved_crop"},
                        )
                    ],
                ),
            )
            write_manifest(
                job_dir / "deck.json",
                DeckManifest(
                    job_id="job",
                    slide_order=["slide-1"],
                    aspect_ratio="16:9",
                    provider_roles={},
                    quality_settings={},
                    fallback_policy="fail",
                    page_manifest_paths=["pages/0000-slide-1.json"],
                ),
            )

            issues = runner._reconstruction_object_issues(
                deck_manifest_path=job_dir / "deck.json",
                artifact_root=job_dir,
                object_stats={
                    "slides": [
                        {
                            "shape_counts": {"PICTURE": 1},
                            "full_slide_picture_count": 0,
                            "non_full_slide_picture_count": 1,
                            "text_box_count": 0,
                        }
                    ]
                },
            )

        self.assertEqual(issues[0]["code"], "oversized_bitmap_asset_coverage")
        self.assertEqual(issues[0]["severity"], "error")

    def test_reconstruction_issues_degrade_for_high_combined_bitmap_coverage(self):
        import scripts.run_real_generative_editable_pptx as runner

        from src.generative_editable_manifest import (
            BitmapAssetSpec,
            DeckManifest,
            PageManifest,
            TextBoxSpec,
            write_manifest,
        )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            job_dir = root / "artifacts" / "job"
            (job_dir / "pages").mkdir(parents=True)
            write_manifest(
                job_dir / "pages" / "0000-slide-1.json",
                PageManifest(
                    slide_id="slide-1",
                    page_index=0,
                    source_image_path="sources/slide-1.png",
                    source_image_size=(1000, 500),
                    slide_size=(10.0, 5.0),
                    text_boxes=[
                        TextBoxSpec(
                            text="标题",
                            source_pixel_bbox=(40, 20, 160, 60),
                            source_pixel_polygon=((40, 20), (160, 20), (160, 60), (40, 60)),
                        )
                    ],
                    bitmap_assets=[
                        BitmapAssetSpec(
                            asset_id="top-band",
                            source_pixel_bbox=(0, 0, 1000, 180),
                            asset_path="assets/top-band.png",
                            z_order=1,
                            provenance={"asset_strategy": "masked_source_element"},
                        ),
                        BitmapAssetSpec(
                            asset_id="middle-band",
                            source_pixel_bbox=(0, 180, 1000, 180 + 170),
                            asset_path="assets/middle-band.png",
                            z_order=2,
                            provenance={"asset_strategy": "masked_source_element"},
                        ),
                    ],
                ),
            )
            write_manifest(
                job_dir / "deck.json",
                DeckManifest(
                    job_id="job",
                    slide_order=["slide-1"],
                    aspect_ratio="16:9",
                    provider_roles={},
                    quality_settings={},
                    fallback_policy="fail",
                    page_manifest_paths=["pages/0000-slide-1.json"],
                ),
            )

            issues = runner._reconstruction_object_issues(
                deck_manifest_path=job_dir / "deck.json",
                artifact_root=job_dir,
                object_stats={
                    "slides": [
                        {
                            "shape_counts": {"PICTURE": 3, "TEXT_BOX": 1},
                            "full_slide_picture_count": 1,
                            "non_full_slide_picture_count": 2,
                            "text_box_count": 1,
                        }
                    ]
                },
            )

        self.assertEqual(issues[0]["code"], "high_bitmap_asset_coverage")
        self.assertEqual(issues[0]["severity"], "warning")

    def test_reconstruction_issues_allow_split_row_level_bitmap_assets(self):
        import scripts.run_real_generative_editable_pptx as runner

        from src.generative_editable_manifest import (
            BitmapAssetSpec,
            DeckManifest,
            NativeShapeSpec,
            PageManifest,
            TextBoxSpec,
            write_manifest,
        )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            job_dir = root / "artifacts" / "job"
            (job_dir / "pages").mkdir(parents=True)
            write_manifest(
                job_dir / "pages" / "0000-slide-1.json",
                PageManifest(
                    slide_id="slide-1",
                    page_index=0,
                    source_image_path="sources/slide.png",
                    source_image_size=(1706, 960),
                    slide_size=(10.0, 5.625),
                    native_shapes=[
                        NativeShapeSpec(
                            shape_type="line",
                            source_pixel_bbox=(100, 180 + index * 30, 1500, 183 + index * 30),
                        )
                        for index in range(6)
                    ],
                    text_boxes=[
                        TextBoxSpec(
                            text="底盘域",
                            source_pixel_bbox=(160, 450, 260, 490),
                            source_pixel_polygon=((160, 450), (260, 450), (260, 490), (160, 490)),
                        )
                    ],
                    bitmap_assets=[
                        BitmapAssetSpec(
                            asset_id=f"row-{index}",
                            source_pixel_bbox=(59, top, 1672, bottom),
                            asset_path=f"assets/row-{index}.png",
                            z_order=index,
                            provenance={"asset_strategy": "masked_source_element"},
                        )
                        for index, (top, bottom) in enumerate(
                            [(168, 405), (405, 642), (642, 880)],
                            start=1,
                        )
                    ],
                ),
            )
            write_manifest(
                job_dir / "deck.json",
                DeckManifest(
                    job_id="job",
                    slide_order=["slide-1"],
                    aspect_ratio="16:9",
                    provider_roles={},
                    quality_settings={},
                    fallback_policy="fail",
                    page_manifest_paths=["pages/0000-slide-1.json"],
                ),
            )

            issues = runner._reconstruction_object_issues(
                deck_manifest_path=job_dir / "deck.json",
                artifact_root=job_dir,
                object_stats={
                    "slides": [
                        {
                            "shape_counts": {"PICTURE": 3, "TEXT_BOX": 1, "LINE": 6},
                            "full_slide_picture_count": 0,
                            "non_full_slide_picture_count": 3,
                            "text_box_count": 1,
                        }
                    ]
                },
            )

        self.assertEqual(issues, [])

    def test_reconstruction_issues_allow_dense_infographic_alpha_visible_assets(self):
        import scripts.run_real_generative_editable_pptx as runner

        from src.generative_editable_manifest import (
            BitmapAssetSpec,
            DeckManifest,
            NativeShapeSpec,
            PageManifest,
            TextBoxSpec,
            write_manifest,
        )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            job_dir = root / "artifacts" / "job"
            (job_dir / "pages").mkdir(parents=True)
            write_manifest(
                job_dir / "pages" / "0000-slide-1.json",
                PageManifest(
                    slide_id="slide-1",
                    page_index=0,
                    source_image_path="sources/slide.png",
                    source_image_size=(1672, 941),
                    slide_size=(10.0, 5.625),
                    native_shapes=[
                        NativeShapeSpec(
                            shape_type="line",
                            source_pixel_bbox=(100, 180 + index * 24, 1500, 183 + index * 24),
                        )
                        for index in range(11)
                    ],
                    text_boxes=[
                        TextBoxSpec(
                            text=f"文本{index}",
                            source_pixel_bbox=(80, 120 + index * 30, 220, 145 + index * 30),
                            source_pixel_polygon=(
                                (80, 120 + index * 30),
                                (220, 120 + index * 30),
                                (220, 145 + index * 30),
                                (80, 145 + index * 30),
                            ),
                        )
                        for index in range(15)
                    ],
                    bitmap_assets=[
                        BitmapAssetSpec(
                            asset_id="upper-band",
                            source_pixel_bbox=(95, 169, 1576, 419),
                            asset_path="assets/upper-band.png",
                            z_order=1,
                            provenance={
                                "asset_strategy": "masked_source_element",
                                "background_difference_alpha": True,
                                "alpha_visible_area_ratio": 0.097449,
                            },
                        ),
                        BitmapAssetSpec(
                            asset_id="small-detail",
                            source_pixel_bbox=(957, 375, 1002, 437),
                            asset_path="assets/small-detail.png",
                            z_order=2,
                            provenance={
                                "asset_strategy": "masked_source_element",
                                "background_difference_alpha": True,
                                "alpha_visible_area_ratio": 0.000983,
                            },
                        ),
                        BitmapAssetSpec(
                            asset_id="lower-visual",
                            source_pixel_bbox=(0, 407, 1658, 909),
                            asset_path="assets/lower-visual.png",
                            z_order=3,
                            provenance={
                                "asset_strategy": "masked_source_element",
                                "background_difference_alpha": True,
                                "alpha_visible_area_ratio": 0.21248,
                            },
                        ),
                        BitmapAssetSpec(
                            asset_id="generic",
                            source_pixel_bbox=(502, 169, 1207, 405),
                            asset_path="assets/generic.png",
                            z_order=4,
                            provenance={
                                "asset_strategy": "masked_source_element",
                                "background_difference_alpha": True,
                                "alpha_visible_area_ratio": 0.069225,
                            },
                        ),
                    ],
                ),
            )
            write_manifest(
                job_dir / "deck.json",
                DeckManifest(
                    job_id="job",
                    slide_order=["slide-1"],
                    aspect_ratio="16:9",
                    provider_roles={},
                    quality_settings={},
                    fallback_policy="fail",
                    page_manifest_paths=["pages/0000-slide-1.json"],
                ),
            )

            issues = runner._reconstruction_object_issues(
                deck_manifest_path=job_dir / "deck.json",
                artifact_root=job_dir,
                object_stats={
                    "slides": [
                        {
                            "shape_counts": {"PICTURE": 4, "TEXT_BOX": 15, "LINE": 11},
                            "full_slide_picture_count": 0,
                            "non_full_slide_picture_count": 4,
                            "text_box_count": 15,
                        }
                    ]
                },
            )

        self.assertEqual(issues, [])

    def test_reconstruction_issues_fail_for_source_raster_guardrail(self):
        import scripts.run_real_generative_editable_pptx as runner

        from src.generative_editable_manifest import (
            BitmapAssetSpec,
            DeckManifest,
            PageManifest,
            write_manifest,
        )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            job_dir = root / "artifacts" / "job"
            (job_dir / "pages").mkdir(parents=True)
            write_manifest(
                job_dir / "pages" / "0000-slide-1.json",
                PageManifest(
                    slide_id="slide-1",
                    page_index=0,
                    source_image_path="sources/slide-1.png",
                    source_image_size=(1000, 500),
                    slide_size=(10.0, 5.0),
                    bitmap_assets=[
                        BitmapAssetSpec(
                            asset_id="source-raster",
                            source_pixel_bbox=(0, 0, 1000, 500),
                            asset_path="assets/source-raster.png",
                            z_order=1,
                            provenance={"asset_strategy": "source_raster_guardrail"},
                        )
                    ],
                    provenance={
                        "chosen_background_kind": "source_raster_guardrail",
                        "source_raster_guardrail": {"reason": "unreliable_ocr_layout"},
                    },
                ),
            )
            write_manifest(
                job_dir / "deck.json",
                DeckManifest(
                    job_id="job",
                    slide_order=["slide-1"],
                    aspect_ratio="16:9",
                    provider_roles={},
                    quality_settings={},
                    fallback_policy="fail",
                    page_manifest_paths=["pages/0000-slide-1.json"],
                ),
            )

            issues = runner._reconstruction_object_issues(
                deck_manifest_path=job_dir / "deck.json",
                artifact_root=job_dir,
                object_stats={
                    "slides": [
                        {
                            "shape_counts": {"PICTURE": 1},
                            "full_slide_picture_count": 1,
                            "non_full_slide_picture_count": 0,
                            "text_box_count": 0,
                        }
                    ]
                },
            )

        self.assertEqual(issues[0]["code"], "source_raster_guardrail_degraded")
        self.assertEqual(issues[0]["severity"], "error")

    def test_error_reconstruction_issues_mark_run_failed(self):
        import scripts.run_real_generative_editable_pptx as runner

        helper = getattr(runner, "_status_with_reconstruction_issues", None)
        self.assertIsNotNone(helper)

        self.assertEqual(
            helper(
                "passed",
                [
                    {
                        "code": "source_preserving_low_opacity_overlay_degraded",
                        "severity": "error",
                    }
                ]
            ),
            "failed",
        )
        self.assertEqual(
            helper(
                "passed",
                [
                    {
                        "code": "no_decomposed_visual_elements",
                        "severity": "error",
                    }
                ]
            ),
            "failed",
        )
        self.assertEqual(helper("failed", []), "failed")

    def test_isolated_status_aggregates_failed_reconstruction_pages(self):
        import scripts.run_real_generative_editable_pptx as runner

        helper = getattr(runner, "_aggregate_page_status", None)
        self.assertIsNotNone(helper)

        self.assertEqual(
            helper(
                [
                    {"status": "passed"},
                    {
                        "status": "passed",
                        "reconstruction_issues": [
                            {
                                "code": "source_preserving_low_opacity_overlay_degraded",
                                "severity": "error",
                            }
                        ],
                    },
                    {"status": "passed"},
                ]
            ),
            "failed",
        )
        self.assertEqual(
            helper(
                [
                    {"status": "degraded"},
                    {"status": "failed"},
                ]
            ),
            "failed",
        )

    def test_reconstruction_issues_warn_when_ocr_items_are_all_filtered(self):
        import scripts.run_real_generative_editable_pptx as runner
        from pptx import Presentation

        from src.generative_editable_manifest import (
            DeckManifest,
            PageManifest,
            write_manifest,
        )

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            artifact_root = root / "artifacts"
            job_dir = artifact_root / "job"
            (job_dir / "pages").mkdir(parents=True)
            ocr_dir = job_dir / "provider_outputs" / "ocr" / "0000-slide-1"
            ocr_dir.mkdir(parents=True)
            (ocr_dir / "ocr.json").write_text(
                '{"items":[{"text":"a"},{"text":"b"}]}',
                encoding="utf-8",
            )
            output_path = root / "job.pptx"
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(str(_solid_png(root / "background.png")), 0, 0, width=prs.slide_width, height=prs.slide_height)
            prs.save(output_path)
            write_manifest(
                job_dir / "pages" / "0000-slide-1.json",
                PageManifest(
                    slide_id="slide-1",
                    page_index=0,
                    source_image_path="sources/slide-1.png",
                    source_image_size=(800, 450),
                    slide_size=(10.0, 5.625),
                    native_shapes=[],
                    provenance={
                        "text_issues": [
                            {
                                "code": "ignored_spurious_ocr",
                                "severity": "warning",
                                "ocr_text": "重复幻觉文本",
                            }
                        ],
                        "text_validation_status": "passed",
                    },
                    provider_output_paths={
                        "ocr": "provider_outputs/ocr/0000-slide-1/ocr.json",
                    },
                ),
            )
            write_manifest(
                job_dir / "deck.json",
                DeckManifest(
                    job_id="job",
                    slide_order=["slide-1"],
                    aspect_ratio="16:9",
                    provider_roles={
                        "ocr": "ocr_model",
                        "image_edit": "edit_model",
                        "image_generation": "image_model",
                    },
                    quality_settings={},
                    fallback_policy="fail",
                    page_manifest_paths=["pages/0000-slide-1.json"],
                ),
            )
            stats = runner._pptx_object_stats(output_path)

            issues = runner._reconstruction_object_issues(
                deck_manifest_path=job_dir / "deck.json",
                artifact_root=job_dir,
                object_stats=stats,
            )

        self.assertIn("no_editable_text_after_ocr_filtering", [issue["code"] for issue in issues])
        warning = next(issue for issue in issues if issue["code"] == "no_editable_text_after_ocr_filtering")
        self.assertEqual(warning["severity"], "warning")
        self.assertEqual(warning["details"]["filtered_ocr_issue_count"], 1)
        self.assertEqual(warning["details"]["ocr_returned_count"], 2)
        self.assertEqual(warning["details"]["filtered_reason_counts"], {"ignored_spurious_ocr": 1})

    def test_warning_summary_counts_pages_and_codes(self):
        import scripts.run_real_generative_editable_pptx as runner

        helper = getattr(runner, "_warning_summary", None)
        self.assertIsNotNone(helper)

        summary = helper(
            [
                {
                    "page": 1,
                    "reconstruction_issues": [
                        {"code": "source_preserving_low_opacity_overlay_degraded", "severity": "error"}
                    ],
                },
                {
                    "page": 2,
                    "reconstruction_issues": [
                        {"code": "no_editable_text_after_ocr_filtering", "severity": "warning"},
                        {"code": "no_decomposed_visual_elements", "severity": "error"},
                    ],
                },
            ]
        )

        self.assertEqual(summary["warning_count"], 1)
        self.assertEqual(summary["warning_pages"], [2])
        self.assertEqual(
            summary["warning_codes"],
            {
                "no_editable_text_after_ocr_filtering": 1,
            },
        )

    def test_warning_summary_prefers_issue_page_over_outer_report_page(self):
        import scripts.run_real_generative_editable_pptx as runner

        helper = getattr(runner, "_warning_summary", None)
        self.assertIsNotNone(helper)

        summary = helper(
            [
                {
                    "page": 1,
                    "reconstruction_issues": [
                        {
                            "page": 5,
                            "code": "source_preserving_low_opacity_overlay_degraded",
                            "severity": "error",
                        }
                    ],
                }
            ]
        )

        self.assertEqual(summary["warning_pages"], [])

    def test_main_failure_report_preserves_requested_input_images(self):
        import scripts.run_real_generative_editable_pptx as runner

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            image_path = root / "slide_1.png"
            Image.new("RGB", (800, 450), "#FFFFFF").save(image_path)
            with patch.object(runner, "_run_pipeline", side_effect=RuntimeError("boom")):
                exit_code = runner.main(
                    [
                        "run",
                        "--input-images",
                        str(image_path),
                        "--output-dir",
                        str(root / "out"),
                        "--job-id",
                        "job",
                    ]
                )
            report = json.loads((root / "out" / "report.json").read_text(encoding="utf-8"))

        self.assertEqual(exit_code, 1)
        self.assertEqual(report["input_images"], [str(image_path.resolve())])

    def test_main_failure_report_uses_generated_job_id_when_job_id_omitted(self):
        import scripts.run_real_generative_editable_pptx as runner

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            image_path = root / "slide_1.png"
            Image.new("RGB", (800, 450), "#FFFFFF").save(image_path)
            generated_job_id = "real-editable-20260704000000"
            with (
                patch.object(runner, "_generated_job_id", return_value=generated_job_id),
                patch.object(runner, "_run_pipeline", side_effect=RuntimeError("boom")),
            ):
                exit_code = runner.main(
                    [
                        "run",
                        "--input-images",
                        str(image_path),
                        "--output-dir",
                        str(root / "out"),
                    ]
                )
            report = json.loads((root / "out" / "report.json").read_text(encoding="utf-8"))

        self.assertEqual(exit_code, 1)
        self.assertEqual(
            report["artifact_root"],
            str((root / "out" / "artifacts" / generated_job_id).resolve()),
        )

    def test_subprocess_timeout_kills_process_tree(self):
        import scripts.run_real_generative_editable_pptx as runner

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            started = root / "grandchild-started.txt"
            marker = root / "grandchild-marker.txt"
            report_path = root / "report.json"
            code = (
                "import subprocess,sys,time;"
                "subprocess.Popen([sys.executable,'-c',"
                "\"import pathlib,sys,time; pathlib.Path(sys.argv[1]).write_text('started'); time.sleep(2); pathlib.Path(sys.argv[2]).write_text('leaked')\","
                "sys.argv[1],sys.argv[2]]);"
                "time.sleep(5)"
            )

            result = runner._run_subprocess_json(
                [sys.executable, "-c", code, str(started), str(marker)],
                timeout_seconds=1,
                report_path=report_path,
                timeout_payload={"status": "failed"},
            )
            time.sleep(2.5)
            started_exists = started.exists()
            marker_exists = marker.exists()

        self.assertEqual(result["error_type"], "TimeoutExpired")
        self.assertTrue(started_exists)
        self.assertFalse(marker_exists)

    def test_subprocess_timeout_returns_when_detached_descendant_holds_stdout_pipe(self):
        import scripts.run_real_generative_editable_pptx as runner

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            report_path = root / "report.json"
            code = (
                "import subprocess,sys,time;"
                "subprocess.Popen([sys.executable,'-c','import time; time.sleep(4)'],"
                "start_new_session=True);"
                "time.sleep(5)"
            )
            started_at = time.monotonic()

            result = runner._run_subprocess_json(
                [sys.executable, "-c", code],
                timeout_seconds=1,
                report_path=report_path,
                timeout_payload={"status": "failed"},
            )
            elapsed = time.monotonic() - started_at

        self.assertEqual(result["error_type"], "TimeoutExpired")
        self.assertLess(elapsed, 3.0)

    def test_subprocess_timeout_preserves_existing_child_report(self):
        import scripts.run_real_generative_editable_pptx as runner

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            report_path = root / "report.json"
            code = (
                "import json,pathlib,subprocess,sys,time;"
                "pathlib.Path(sys.argv[1]).write_text(json.dumps({'status':'passed','child':'done'}));"
                "subprocess.Popen([sys.executable,'-c','import time; time.sleep(4)'],"
                "start_new_session=True);"
                "time.sleep(5)"
            )

            result = runner._run_subprocess_json(
                [sys.executable, "-c", code, str(report_path)],
                timeout_seconds=1,
                report_path=report_path,
                timeout_payload={"status": "failed", "stage": "test"},
            )
            stored = json.loads(report_path.read_text(encoding="utf-8"))

        self.assertEqual(result["status"], "passed")
        self.assertEqual(result["child"], "done")
        self.assertEqual(result["subprocess_cleanup"]["error_type"], "TimeoutExpired")
        self.assertEqual(stored["status"], "passed")
        self.assertEqual(stored["child"], "done")

    def test_isolated_pages_mode_writes_page_reports(self):
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            input_dir = root / "inputs"
            input_dir.mkdir()
            image_path = input_dir / "slide_1.png"
            Image.new("RGB", (800, 450), "#FFFFFF").save(image_path)

            stdout = io.StringIO()
            with redirect_stdout(stdout):
                exit_code = main(
                    [
                        "run",
                        "--use-fake",
                        "--isolate-pages",
                        "--input-images",
                        str(image_path),
                        "--output-dir",
                        str(root / "out"),
                        "--job-id",
                        "isolated",
                    ]
                )

            result = json.loads(stdout.getvalue().splitlines()[-1])
            report = json.loads(Path(result["report_path"]).read_text(encoding="utf-8"))
            page_report_exists = Path(report["page_reports"][0]["report_path"]).exists()

        self.assertEqual(exit_code, 0)
        self.assertEqual(report["status"], "passed")
        self.assertEqual(len(report["page_reports"]), 1)
        self.assertEqual(report["page_reports"][0]["page"], 1)
        self.assertTrue(page_report_exists)

    def test_isolated_pages_top_level_report_includes_child_diagnostics(self):
        import scripts.run_real_generative_editable_pptx as runner

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            image_path = root / "slide_1.png"
            Image.new("RGB", (800, 450), "#FFFFFF").save(image_path)
            args = argparse.Namespace(
                input_images=[str(image_path)],
                input_glob=[],
                slides=0,
                job_id="isolated",
                aspect_ratio="16:9",
                provider_timeout=5,
                fallback_policy="fail",
                use_fake=False,
                page_wall_timeout=10,
            )

            child_result = {
                "status": "failed",
                "output_path": str(root / "out" / "page-01" / "isolated-page-01.pptx"),
                "artifact_root": str(root / "out" / "page-01" / "artifacts" / "isolated-page-01"),
                "object_stats": {"totals": {"PICTURE": 1, "TEXT_BOX": 3}},
                "reconstruction_issues": [{"code": "source_raster_guardrail_degraded"}],
                "validation_status": "degraded",
                "validation_issues": ["source_raster_guardrail_degraded"],
                "active_stage": "base_clean_background",
                "last_stage_event": {
                    "stage": "base_clean_background",
                    "status": "started",
                },
            }
            with patch.object(runner, "_run_subprocess_json", return_value=child_result):
                result = runner._run_isolated_pages(args, root / "out")

            report = json.loads(Path(result["report_path"]).read_text(encoding="utf-8"))
            page_report = report["page_reports"][0]

        self.assertEqual(page_report["object_stats"]["totals"]["TEXT_BOX"], 3)
        self.assertEqual(page_report["reconstruction_issues"][0]["code"], "source_raster_guardrail_degraded")
        self.assertEqual(page_report["validation_status"], "degraded")
        self.assertEqual(page_report["validation_issues"], ["source_raster_guardrail_degraded"])
        self.assertEqual(page_report["active_stage"], "base_clean_background")
        self.assertEqual(page_report["last_stage_event"]["status"], "started")

    def test_isolated_pages_forwards_selected_pipeline_mode_to_child(self):
        import scripts.run_real_generative_editable_pptx as runner

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            image_path = root / "slide_1.png"
            Image.new("RGB", (800, 450), "#FFFFFF").save(image_path)
            args = argparse.Namespace(
                input_images=[str(image_path)],
                input_glob=[],
                slides=0,
                job_id="isolated",
                aspect_ratio="16:9",
                provider_timeout=5,
                fallback_policy="fail",
                use_fake=True,
                page_wall_timeout=10,
                mode="vlm_first",
            )
            captured_commands = []

            def fake_subprocess(command, **kwargs):
                captured_commands.append(command)
                return {"status": "passed"}

            with patch.object(runner, "_run_subprocess_json", side_effect=fake_subprocess):
                runner._run_isolated_pages(args, root / "out")

        self.assertIn("--mode", captured_commands[0])
        self.assertEqual(
            captured_commands[0][captured_commands[0].index("--mode") + 1],
            "vlm_first",
        )

    def test_isolated_pages_forwards_pipeline_page_timeout_to_child(self):
        import scripts.run_real_generative_editable_pptx as runner

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            image_path = root / "slide_1.png"
            Image.new("RGB", (800, 450), "#FFFFFF").save(image_path)
            args = argparse.Namespace(
                input_images=[str(image_path)],
                input_glob=[],
                slides=0,
                job_id="isolated",
                aspect_ratio="16:9",
                provider_timeout=5,
                provider_max_attempts=3,
                provider_retry_backoff=1.0,
                pipeline_page_timeout=240.0,
                fallback_policy="fail",
                use_fake=True,
                page_wall_timeout=10,
                mode="vlm_first",
            )
            captured_commands = []

            def fake_subprocess(command, **kwargs):
                captured_commands.append(command)
                return {"status": "passed"}

            with patch.object(runner, "_run_subprocess_json", side_effect=fake_subprocess):
                runner._run_isolated_pages(args, root / "out")

        command = captured_commands[0]
        self.assertIn("--pipeline-page-timeout", command)
        self.assertEqual(command[command.index("--pipeline-page-timeout") + 1], "240.0")


def _solid_png(path: Path) -> Path:
    Image.new("RGB", (16, 16), "#2563EB").save(path)
    return path


if __name__ == "__main__":
    unittest.main()
